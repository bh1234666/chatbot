from app.llm.tools.workspace import inspect_file


def test_inspect_file_pptx_includes_table_text(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    pptx_path = tmp_path / "slides.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "数据页"
    table = slide.shapes.add_table(
        3, 2, Inches(0.5), Inches(1.5), Inches(8), Inches(1.5)
    ).table
    table.cell(0, 0).text = "指标"
    table.cell(0, 1).text = "数值"
    table.cell(1, 0).text = "X"
    table.cell(1, 1).text = "12"
    table.cell(2, 0).text = "Y"
    table.cell(2, 1).text = "29"
    prs.save(pptx_path)

    result = inspect_file(str(tmp_path), "slides.pptx")

    assert result["ok"] is True
    metadata = result["metadata"]
    assert metadata["slide_count"] == 1
    assert metadata["slide_titles"] == ["数据页"]
    assert "X" in metadata["text_preview"]
    assert "12" in metadata["text_preview"]
    assert "Y" in metadata["slide_texts"][0]
    assert "29" in metadata["slide_texts"][0]
