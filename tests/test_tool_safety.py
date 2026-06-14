import asyncio
import inspect
import json
import sys
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import patch

import pytest
from PIL import Image

sys.path.insert(0, str(Path(__file__).parent.parent))


def _assert_broad_code_guard_fact(facts, *, task_id: str | None = None):
    assert isinstance(facts, list)
    broad = [
        fact for fact in facts
        if isinstance(fact, dict) and fact.get("issue") == "broad_code_task_before_spawn"
    ]
    assert broad
    fact = broad[0]
    assert fact["kind"] == "guard_observation"
    assert fact["needs_attention"] is True
    assert fact.get("signals") or fact.get("expected_outputs")
    assert "suggested_action" not in fact
    assert "suggestion" not in fact
    if task_id is not None:
        assert fact["task_id"] == task_id
    return fact


def _assert_sanitized_task_has_broad_fact(result, *, task_id: str):
    assert not isinstance(result, str)
    task = next(t for t in result if t.get("task_id") == task_id)
    return _assert_broad_code_guard_fact(task.get("guard_observations"), task_id=task_id)


def test_todo_write_schema_warns_parallel_helpers_keep_single_in_progress():
    from app.llm.tools.registry import TODO_WRITE_SCHEMA

    description = TODO_WRITE_SCHEMA["function"]["description"]

    assert "At most one todo may be `in_progress`" in description
    assert "multiple helpers run in parallel" in description
    assert "keep other parallel branches `pending`" in description
    assert "not to carry long prose" in description
    assert "不承载长正文" in description


def test_edit_helper_prompt_requests_resources_instead_of_docx_placeholders():
    from app.llm.tools.delegate import _select_helper_system

    prompt = _select_helper_system("edit")
    assert "When a resource is missing, request or report the needed resource" in prompt
    assert "Deliverables should contain confirmed content, not placeholders for absent resources" in prompt
    assert "not template tokens such as TODO, TKTK, INSERT, PLACEHOLDER_*" in prompt
    assert "do not insert template placeholder tokens" in prompt
    assert "Rewrite internal-only source notes into user-facing conclusions" in prompt
    assert "缺资源时请求主线程" in prompt


def test_helper_report_verification_recommendation_does_not_invite_rechecking_when_passed():
    from app.llm.tools.delegate import _select_helper_system

    prompt = _select_helper_system("edit")
    assert "Use `recommend: no` when you already performed the requested checks" in prompt
    assert "Use `recommend: yes` only when a specific unverified risk" in prompt


def test_text_mojibake_quality_helpers_repair_common_gbk_case():
    from app.llm.tools.delegate_quality import (
        repair_common_mojibake_text,
        text_mojibake_warnings,
    )

    original = "\u4e8c\u4e00\u4e8c\u4e00"
    mojibake = original.encode("utf-8").decode("gbk")

    assert text_mojibake_warnings(original) == []
    assert text_mojibake_warnings(mojibake)
    repaired, info = repair_common_mojibake_text(mojibake)
    assert repaired == original
    assert info and info["repaired_score"] < info["original_score"]


async def test_chat_stream_cancels_abort_waiter_after_stream_opens(monkeypatch):
    from app.llm import client as llm_client

    cancelled = []
    real_ensure_future = asyncio.ensure_future

    def tracking_ensure_future(awaitable, *args, **kwargs):
        task = real_ensure_future(awaitable, *args, **kwargs)
        if getattr(awaitable, "__name__", "") == "wait":
            cancelled.append(task)
        return task

    class FakeDelta:
        content = "ok"

    class FakeChoice:
        delta = FakeDelta()

    class FakeChunk:
        choices = [FakeChoice()]

    class FakeStream:
        def __aiter__(self):
            return self

        async def __anext__(self):
            if getattr(self, "done", False):
                raise StopAsyncIteration
            self.done = True
            return FakeChunk()

    async def fake_retry(fn, label, **kwargs):
        return FakeStream()

    monkeypatch.setattr(llm_client, "_retry", fake_retry)
    monkeypatch.setattr(asyncio, "ensure_future", tracking_ensure_future)

    abort_event = asyncio.Event()
    tokens = [tok async for tok in llm_client.chat_stream([{"role": "user", "content": "hi"}], abort_event=abort_event)]

    assert tokens == ["ok"]
    assert cancelled
    assert all(task.cancelled() for task in cancelled)


async def test_chat_stream_cancels_internal_tasks_when_cancelled_before_open(monkeypatch):
    from app.llm import client as llm_client

    created = []
    real_ensure_future = asyncio.ensure_future

    def tracking_ensure_future(awaitable, *args, **kwargs):
        task = real_ensure_future(awaitable, *args, **kwargs)
        created.append(task)
        return task

    async def fake_retry(fn, label, **kwargs):
        await asyncio.Event().wait()

    monkeypatch.setattr(llm_client, "_retry", fake_retry)
    monkeypatch.setattr(asyncio, "ensure_future", tracking_ensure_future)

    abort_event = asyncio.Event()
    agen = llm_client.chat_stream([{"role": "user", "content": "hi"}], abort_event=abort_event)
    task = asyncio.create_task(agen.__anext__())
    await asyncio.sleep(0)
    task.cancel()

    with pytest.raises(asyncio.CancelledError):
        await task

    assert created
    assert all(t.cancelled() for t in created)


async def test_stream_iteration_error_is_recoverable_idle_timeout(monkeypatch):
    from types import SimpleNamespace

    from app.llm import client as llm_client

    class BrokenStream:
        def __aiter__(self):
            return self

        async def __anext__(self):
            raise RuntimeError("peer closed connection without sending complete message body")

        async def close(self):
            return None

    async def fake_retry(fn, **kwargs):
        return BrokenStream()

    monkeypatch.setattr(llm_client, "_retry", fake_retry)

    _, collector, reason = await llm_client._call_llm_streaming_with_idle(
        the_client=SimpleNamespace(chat=SimpleNamespace(completions=SimpleNamespace(create=lambda **kw: None))),
        model="fake",
        provider=None,
        msgs=[{"role": "user", "content": "hi"}],
        tools=[],
        extra_body={},
        abort_event=None,
        iter_no=1,
        task_id=None,
        idle_timeout=1,
        first_chunk_timeout=1,
    )

    assert reason == "idle_timeout"
    assert "peer closed connection" in collector.last_error


async def test_dispatch_rejects_main_thread_implementation_aliases():
    from app.llm.tools import registry

    for name in ("write", "write_file", "run", "run_command", "shell", "python", "edit_file", "progress_note"):
        result = await registry.dispatch(
            name,
            {"command": "echo hi", "path": "x.txt", "content": "hi"},
            archive_id="archive",
            group_id="group",
            user_id="user",
            workspace_dir="",
        )
        data = json.loads(result)
        assert data["ok"] is False
        assert "not available in the main thread" in data["error"]
        assert "instead of retrying blocked aliases" in data["error"]


async def test_workspace_missing_action_returns_fact_without_execution(tmp_path):
    from app.llm.tools import registry

    result = await registry.dispatch(
        "workspace",
        {"path": "accidental.txt", "content": "should not be written"},
        archive_id="archive",
        group_id="group",
        user_id="user",
        workspace_dir=str(tmp_path),
    )

    data = json.loads(result)
    assert data["ok"] is False
    assert data["error"] == "workspace_action_required"
    assert data["candidate_actions_from_fields"] == ["write"]
    assert not (tmp_path / "accidental.txt").exists()



async def test_translate_shared_cd_python_collapses_duplicate_helpers_shared_prefix():
    from app.llm.tools.workspace import _translate_windows_command

    translated = _translate_windows_command(
        'cd _helpers_shared/hw_redraw_charts && python draw_q5_2.py',
        r'F:\chatbot\ws',
    )

    assert translated == 'python _helpers_shared/hw_redraw_charts/draw_q5_2.py'


async def test_todo_write_demotes_extra_in_progress_items(tmp_path):
    from app.llm.tools.workspace import handle_todo_write

    result = await handle_todo_write(str(tmp_path), [
        {"id": "1", "content": "并行分支 A", "status": "in_progress"},
        {"id": "2", "content": "并行分支 B", "status": "in_progress"},
        {"id": "3", "content": "最终汇总", "status": "pending"},
    ])

    assert result["ok"] is True
    assert result["counts"]["in_progress"] == 1
    assert result["counts"]["pending"] == 2
    assert "normalization_warning" in result
    assert result["demoted_in_progress"][0]["id"] == "2"
    from app.llm.tools.delegate import _auto_correct_obvious_helper_kind

    kind, reason = _auto_correct_obvious_helper_kind(
        "code",
        "使用 matplotlib 读取 bench_runner_assemble_bench_out.csv，生成 PNG 图表 chart.png。",
    )

    assert kind == "code"
    assert reason is None


async def test_read_file_rejects_office_and_images(tmp_path):
    from app.llm.tools.workspace import handle_read_file

    docx = tmp_path / "report.docx"
    docx.write_bytes(b"PK\x03\x04fake-docx")
    image = tmp_path / "img.jpg"
    image.write_bytes(b"\xff\xd8\xfffake-jpeg")

    docx_result = await handle_read_file(str(tmp_path), "report.docx")
    image_result = await handle_read_file(str(tmp_path), "img.jpg")

    assert docx_result["ok"] is False
    assert docx_result["error"] == "binary_or_structured_file_not_readable_as_text"
    assert "office(action='read')" in docx_result["suggested_tools"]
    assert docx_result["next_call_fact"] == {
        "tool": "office",
        "action": "read",
        "path": "report.docx",
    }
    assert image_result["ok"] is False
    assert image_result["error"] == "binary_or_structured_file_not_readable_as_text"
    assert "ocr" in image_result["suggested_tools"]


async def test_registry_read_file_accepts_file_path_alias(tmp_path):
    from app.llm.tools import registry

    (tmp_path / "notes.txt").write_text("hello\nworld\n", encoding="utf-8")

    raw = await registry._handle_read_file(str(tmp_path), {"file_path": "notes.txt"})
    data = json.loads(raw)

    assert data["ok"] is True
    assert data["path"] == "notes.txt"
    assert "hello" in data["content"]


async def test_read_and_search_file_reject_directory_with_recovery_hint(tmp_path):
    from app.llm.tools.workspace import handle_read_file, handle_search_in_file

    (tmp_path / "src").mkdir()
    (tmp_path / "src" / "main.py").write_text("print('hello')\n", encoding="utf-8")

    read_result = await handle_read_file(str(tmp_path), "src")
    search_result = await handle_search_in_file(str(tmp_path), "src", "hello")

    for result in (read_result, search_result):
        assert result["ok"] is False
        assert result["error"] == "path_is_directory"
        assert "search_across_files" in result["suggested_tools"]
        assert "concrete file" in result["_next_action_instruction"]


async def test_read_file_structured_rejection_survives_binary_detection(tmp_path):
    from app.llm.tools.workspace import _check_file_readable

    image = tmp_path / "img.jpg"
    image.write_bytes(b"\x00\xff\xd8\xffbinary-jpeg")

    size, err = _check_file_readable(str(image))

    assert size == image.stat().st_size
    assert err is not None
    assert err["ok"] is False
    assert err["error"] == "binary_or_structured_file_not_readable_as_text"
    assert err["file_category"] == "image"
    assert err["file_type"] == "JPEG image"
    assert "ocr" in err["suggested_tools"]


async def test_read_file_unknown_binary_returns_structured_rejection(tmp_path):
    from app.llm.tools.workspace import _check_file_readable

    blob = tmp_path / "payload.binx"
    blob.write_bytes(b"\x00\x01\x02binary")

    size, err = _check_file_readable(str(blob))

    assert size == blob.stat().st_size
    assert err is not None
    assert err["ok"] is False
    assert err["error"] == "binary_or_structured_file_not_readable_as_text"
    assert err["file_category"] == "unknown"
    assert err["file_type"] == "unknown extension (.binx)"
    assert err["path"] == "payload.binx"
    assert err["suggested_tools"] == ["inspect_file"]


async def test_search_files_includes_workspace_matches(monkeypatch, tmp_path):
    from app.llm.tools import registry

    (tmp_path / "gen_docx_通信原理作业_第五章.docx").write_bytes(b"PK\x03\x04fake-docx")

    async def fake_search_files(*args, **kwargs):
        return []

    monkeypatch.setattr(registry.kb_mem, "search_files", fake_search_files)

    result = json.loads(await registry._handle_search_files(
        "archive", "group", str(tmp_path), {"query": "gen_docx", "limit": 10},
    ))

    assert result["count"] == 0
    assert result["workspace_count"] == 1
    assert result["workspace_matches"][0]["workspace_path"] == "gen_docx_通信原理作业_第五章.docx"


async def test_search_files_falls_back_to_workspace_when_kb_pool_unavailable(monkeypatch, tmp_path):
    from app.llm.tools import registry

    (tmp_path / "kb_fallback_docx_001.docx").write_bytes(b"PK\x03\x04fake-docx")

    async def fake_search_files(*args, **kwargs):
        raise RuntimeError("DB pool not initialized; call init_pool() first")

    logged: list[tuple[str, str]] = []
    monkeypatch.setattr(registry.kb_mem, "search_files", fake_search_files)
    monkeypatch.setattr(registry.debug, "log", lambda category, msg='', payload=None: logged.append((category, str(msg))))

    result = json.loads(await registry._handle_search_files(
        "archive", "group", str(tmp_path), {"query": "kb_fallback_docx_001", "limit": 10},
    ))

    assert result["kb_unavailable"] is True
    assert "DB pool not initialized" in result["kb_error"]
    assert result["workspace_count"] == 1
    assert result["workspace_matches"][0]["workspace_path"] == "kb_fallback_docx_001.docx"
    assert not any(category == "search_files.kb_unavailable" for category, _ in logged)


async def test_search_files_splits_multi_filename_workspace_queries(monkeypatch, tmp_path):
    from app.llm.tools import registry

    (tmp_path / "common_v2.h").write_text("// header", encoding="utf-8")
    (tmp_path / "bench_tree_ctx.c").write_text("// bench", encoding="utf-8")

    async def fake_search_files(*args, **kwargs):
        return []

    monkeypatch.setattr(registry.kb_mem, "search_files", fake_search_files)

    result = json.loads(await registry._handle_search_files(
        "archive", "group", str(tmp_path), {"query": "common_v2.h bench_tree_ctx", "limit": 10},
    ))

    paths = {m["workspace_path"] for m in result["workspace_matches"]}
    assert "common_v2.h" in paths
    assert "bench_tree_ctx.c" in paths


async def test_search_files_preserves_summary_and_current_user_attribution(monkeypatch, tmp_path):
    from app.llm.tools import registry

    async def fake_search_files(*args, **kwargs):
        return [{
            "id": "kb_old_same_user",
            "headline": "older same-user homework file",
            "content": "Summary: contains chapter 6 homework questions and waveform requirements.",
            "filename": "task.docx",
            "workspace_path": "group_files/old_task.docx",
            "archive_id": "archive",
            "group_id": "group",
            "uploader_uin": "u1",
            "uploader_name": "Alice",
            "upload_time": 2_000_000_000,
            "download_status": "done",
        }]

    monkeypatch.setattr(registry.kb_mem, "search_files", fake_search_files)

    result = json.loads(await registry._handle_search_files(
        "archive",
        "group",
        str(tmp_path),
        {"query": "homework task.docx", "limit": 10},
        current_user_id="u1",
        current_user_name="Alice",
    ))

    item = result["items"][0]
    assert item["kb_node_id"] == "kb_old_same_user"
    assert item["current_user_relation"] == "same_speaker_upload"
    assert item["content"].startswith("Summary: contains chapter 6")
    assert item["source_attribution"]["scope"] == "shared_group_file_index"
    assert item["source_attribution"]["kb_node_id"] == "kb_old_same_user"
    assert item["source_attribution"]["current_user_match"] is True
    assert item["source_attribution"]["current_user_relation"] == "same_speaker_upload"
    guidance = result["selection_guidance"]
    guidance_text = json.dumps(guidance, ensure_ascii=False)
    assert "content summary" in guidance_text
    assert "current_user_relation" in guidance_text
    assert "Older same-speaker files remain valid" in guidance_text
    assert "Other-user uploads can be current shared context" in guidance_text
    assert "not hard prohibitions" in guidance_text


async def test_search_files_keeps_same_name_different_uploaders_distinct(monkeypatch, tmp_path):
    from app.llm.tools import registry

    async def fake_search_files(*args, **kwargs):
        return [
            {
                "id": "kb_same_user_task",
                "headline": "Alice task document",
                "content": "Summary: Alice's older task file with chapter 6 questions.",
                "filename": "task.docx",
                "workspace_path": "group_files/alice_task.docx",
                "archive_id": "archive",
                "group_id": "group",
                "uploader_uin": "u1",
                "uploader_name": "Alice",
                "upload_time": 2_000_000_000,
                "download_status": "done",
            },
            {
                "id": "kb_other_user_task",
                "headline": "Bob task document",
                "content": "Summary: Bob's newer task file with a different worksheet.",
                "filename": "task.docx",
                "workspace_path": "group_files/bob_task.docx",
                "archive_id": "archive",
                "group_id": "group",
                "uploader_uin": "u2",
                "uploader_name": "Bob",
                "upload_time": 2_000_000_060,
                "download_status": "done",
            },
        ]

    monkeypatch.setattr(registry.kb_mem, "search_files", fake_search_files)

    result = json.loads(await registry._handle_search_files(
        "archive",
        "group",
        str(tmp_path),
        {"query": "task.docx", "limit": 10},
        current_user_id="u1",
        current_user_name="Alice",
    ))

    assert result["count"] == 2
    by_id = {item["kb_node_id"]: item for item in result["items"]}
    assert set(by_id) == {"kb_same_user_task", "kb_other_user_task"}
    assert by_id["kb_same_user_task"]["current_user_relation"] == "same_speaker_upload"
    assert by_id["kb_other_user_task"]["current_user_relation"] == "other_user_upload"
    assert by_id["kb_same_user_task"]["content"].startswith("Summary: Alice")
    assert by_id["kb_other_user_task"]["content"].startswith("Summary: Bob")
    assert by_id["kb_same_user_task"]["source_attribution"]["kb_node_id"] == "kb_same_user_task"
    assert by_id["kb_other_user_task"]["source_attribution"]["uploader_id"] == "u2"


def test_ocr_bridge_builds_local_mineru_env_from_hf_snapshots(tmp_path, monkeypatch):
    from app.llm.tools import ocr_bridge

    mineru_dir = tmp_path / "mineru"
    hf_dir = mineru_dir / "hf"
    pipeline = hf_dir / "models--opendatalab--PDF-Extract-Kit-1.0" / "snapshots" / "pipe123"
    vlm = hf_dir / "models--opendatalab--MinerU2.5-Pro-2604-1.2B" / "snapshots" / "vlm123"
    pipeline.mkdir(parents=True)
    vlm.mkdir(parents=True)

    monkeypatch.setattr(ocr_bridge, "_MINERU_DIR", mineru_dir)
    monkeypatch.setattr(ocr_bridge, "_MINERU_HF_DIR", hf_dir)
    monkeypatch.delenv("MINERU_MODEL_SOURCE", raising=False)
    monkeypatch.delenv("MINERU_TOOLS_CONFIG_JSON", raising=False)
    monkeypatch.delenv("HF_HOME", raising=False)
    monkeypatch.delenv("HUGGINGFACE_HUB_CACHE", raising=False)

    env = ocr_bridge._mineru_env()
    cfg_path = mineru_dir / "mineru.local.json"
    cfg = json.loads(cfg_path.read_text(encoding="utf-8"))

    assert env["MINERU_MODEL_SOURCE"] == "local"
    assert env["MINERU_TOOLS_CONFIG_JSON"] == str(cfg_path)
    assert env["HF_HOME"] == str(hf_dir)
    assert env["HUGGINGFACE_HUB_CACHE"] == str(hf_dir)
    assert cfg["models-dir"]["pipeline"] == str(pipeline)
    assert cfg["models-dir"]["vlm"] == str(vlm)
    assert cfg_path.read_bytes()[:3] != b"\xef\xbb\xbf"


def test_ocr_bridge_prefers_mineru_then_falls_back_to_legacy(tmp_path, monkeypatch):
    from app.llm.tools import ocr_bridge

    image_path = tmp_path / "sample.png"
    Image.new("RGB", (80, 30), "white").save(image_path)

    monkeypatch.setattr(ocr_bridge, "_ocr_with_mineru_file", lambda path, timeout=120: ocr_bridge.OcrResult(ok=False, error="mineru failed"))
    monkeypatch.setattr(ocr_bridge, "_ocr_with_legacy_file", lambda path, timeout=120: ocr_bridge.OcrResult(ok=True, text="回退成功", score=0.9))

    result = ocr_bridge.ocr_file(image_path)

    assert result.ok is True
    assert result.text == "回退成功"


def test_ocr_bridge_base64_falls_back_when_mineru_temp_path_fails(monkeypatch):
    from app.llm.tools import ocr_bridge

    png_b64 = "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAusB9WnS8b0AAAAASUVORK5CYII="

    calls = {"n": 0}

    def fake_tempfile(data, timeout):
        calls["n"] += 1
        return ocr_bridge.OcrResult(ok=False, error="mineru temp failed")

    monkeypatch.setattr(ocr_bridge, "_mineru_available", lambda: True)
    monkeypatch.setattr(ocr_bridge, "_ocr_via_tempfile", fake_tempfile)
    monkeypatch.setattr(ocr_bridge, "_call_legacy_headless", lambda args, timeout=120: ocr_bridge.OcrResult(ok=True, text="legacy text", score=0.8))

    result = ocr_bridge.ocr_base64(png_b64)

    assert calls["n"] == 1
    assert result.ok is True
    assert result.text == "legacy text"


def test_ocr_bridge_mineru_requires_hot_service(tmp_path, monkeypatch):
    from app.llm.tools import ocr_bridge

    image_path = tmp_path / "sample.png"
    Image.new("RGB", (80, 30), "white").save(image_path)

    calls = {"client": 0}

    monkeypatch.setattr(ocr_bridge, "_mineru_available", lambda: True)
    monkeypatch.setattr(ocr_bridge, "_wait_mineru_bg_api_url", lambda config, timeout: None)
    monkeypatch.setattr(ocr_bridge, "_run_mineru_client", lambda *args, **kwargs: calls.__setitem__("client", calls["client"] + 1))

    result = ocr_bridge._ocr_with_mineru_file(image_path, timeout=1)

    assert result.ok is False
    assert "hot service" in result.error
    assert calls["client"] == 0


def test_mineru_waits_for_active_tts_instead_of_killing(tmp_path, monkeypatch):
    from app.llm.tools import ocr_bridge

    image_path = tmp_path / "sample.png"
    Image.new("RGB", (80, 30), "white").save(image_path)

    monkeypatch.setattr(ocr_bridge, "_mineru_available", lambda: True)
    monkeypatch.setattr(ocr_bridge, "_tts_headless_running", lambda: True)
    monkeypatch.setattr(ocr_bridge, "_umi_worker_active", lambda: False)
    monkeypatch.setattr(
        ocr_bridge,
        "_terminate_idle_umi_processes",
        lambda: (_ for _ in ()).throw(AssertionError("must not kill while tts active")),
    )
    monkeypatch.setattr(
        ocr_bridge,
        "_run_mineru_client",
        lambda *args, **kwargs: (_ for _ in ()).throw(AssertionError("must not start mineru client")),
    )

    result = ocr_bridge._ocr_with_mineru_file(image_path, timeout=1)

    assert result.ok is False
    assert result.engine == "mineru"
    assert "GPU resource busy" in result.error


def test_terminate_umi_processes_skips_when_worker_active(monkeypatch):
    from app.llm.tools import ocr_bridge

    monkeypatch.setattr(ocr_bridge, "_umi_worker_active", lambda: True)
    monkeypatch.setattr(
        ocr_bridge.subprocess,
        "run",
        lambda *args, **kwargs: (_ for _ in ()).throw(AssertionError("must not taskkill active UmiOCR")),
    )

    ocr_bridge._terminate_umi_processes()


def test_terminate_idle_umi_processes_clears_stale_worker_handle(monkeypatch):
    from app.llm.tools import ocr_bridge

    class DoneProc:
        def poll(self):
            return 0

    monkeypatch.setattr(ocr_bridge, "_UMI_WORKER_PROC", DoneProc())
    monkeypatch.setattr(ocr_bridge, "_umi_worker_active", lambda: False)
    monkeypatch.setattr(ocr_bridge.sys, "platform", "linux")

    ocr_bridge._terminate_idle_umi_processes()

    assert ocr_bridge._UMI_WORKER_PROC is None


def test_ocr_bridge_converts_unsupported_image_suffix_for_legacy(tmp_path):
    from app.llm.tools import ocr_bridge

    src = tmp_path / "frame.heic"
    Image.new("RGB", (32, 16), "white").save(src, format="PNG")

    prepared, cleanup = ocr_bridge._ensure_supported_image_path(src)
    try:
        assert prepared != src
        assert prepared.suffix.lower() == ".png"
        assert prepared.is_file()
    finally:
        if cleanup is not None:
            cleanup()


def test_ocr_bridge_pdf_legacy_fallback_references_renderer():
    from app.llm.tools import ocr_bridge
    import inspect

    assert callable(ocr_bridge._render_pdf_to_images)
    assert "_render_pdf_to_images" in inspect.getsource(ocr_bridge._ocr_pdf_with_legacy)


def test_ocr_bridge_pdf_legacy_fallback_uses_renderer(monkeypatch, tmp_path):
    from app.llm.tools import ocr_bridge

    page = tmp_path / "page.png"
    page.write_bytes(b"fake")
    cleanup_called = False

    def fake_cleanup():
        nonlocal cleanup_called
        cleanup_called = True

    monkeypatch.setattr(ocr_bridge, "_legacy_available", lambda: True)
    monkeypatch.setattr(ocr_bridge, "_render_pdf_to_images", lambda path: ([page], fake_cleanup))
    monkeypatch.setattr(
        ocr_bridge,
        "_ocr_with_legacy_file",
        lambda path, timeout=120: ocr_bridge.OcrResult(ok=True, text="pdf page text", score=1.0),
    )

    result = ocr_bridge._ocr_pdf_with_legacy(tmp_path / "source.pdf")

    assert result.ok is True
    assert result.text == "pdf page text"
    assert cleanup_called is True


def test_windows_python_c_translation_preserves_stderr_redirect_suffix(tmp_path):
    from app.llm.tools.workspace import _translate_windows_command

    command = 'python -c "print(\"hello\")" 2>nul'
    translated = _translate_windows_command(command, str(tmp_path))

    assert translated.startswith("python _py_cmd_")
    assert translated.endswith(".py 2>nul")
    script_name = translated.split(" ", 2)[1]
    assert (tmp_path / script_name).read_text(encoding="utf-8") == 'print("hello")'


def test_windows_translate_unix_null_device_output_targets(tmp_path):
    from app.llm.tools.workspace import _translate_windows_command

    translated = _translate_windows_command(
        'curl -s -o /dev/null -w "%{http_code}" http://127.0.0.1:1234/',
        str(tmp_path),
    )
    assert "-o NUL" in translated
    assert "/dev/null" not in translated
    assert "http://127.0.0.1:1234/" in translated

    stderr_redirect = _translate_windows_command("python check.py 2>/dev/null", str(tmp_path))
    assert stderr_redirect == "python check.py 2>NUL"


def test_windows_cmd_wrapped_multiline_python_c_translation_writes_temp_script(tmp_path):
    from app.llm.tools.workspace import _translate_windows_command

    command = 'cmd /c python -c "\nprint(\"hello\")\nprint(\"done\")" 2>nul'
    translated = _translate_windows_command(command, str(tmp_path))

    assert translated.startswith("cmd /c python _py_cmd_")
    assert translated.endswith(".py 2>nul")
    script_name = translated.split("python ", 1)[1].rsplit(" ", 1)[0]
    assert (tmp_path / script_name).read_text(encoding="utf-8") == '\nprint("hello")\nprint("done")'


async def test_read_file_redundant_warning_gives_next_action(tmp_path):
    from app.llm.tools.workspace import handle_read_file

    path = tmp_path / "notes.txt"
    path.write_text("line1\nline2\nline3", encoding="utf-8")

    first = await handle_read_file(str(tmp_path), "notes.txt")
    second = await handle_read_file(str(tmp_path), "notes.txt")

    assert first["ok"] is True
    assert second["ok"] is True
    assert second["_already_read_full"] is True
    assert "复用已有上下文" in second["_redundant_read_warning"]
    assert "事实：全文已在本轮返回过" in second["_next_action_instruction"]
    assert "read_minimal_fragment" in second["_available_next_actions"]


async def test_read_file_fragment_paging_does_not_trigger_outline_fallback(tmp_path):
    from app.llm.tools.workspace import handle_read_file

    path = tmp_path / "manifest.json"
    path.write_text("\n".join(f'{{"item": {i}}}' for i in range(1, 21)), encoding="utf-8")

    results = []
    for start in (1, 5, 9, 13, 17):
        results.append(
            await handle_read_file(
                str(tmp_path),
                "manifest.json",
                start_line=start,
                end_line=start + 3,
            )
        )

    assert all(result["ok"] is True for result in results)
    assert not any(result.get("_outline_mode") for result in results)
    assert results[-1]["shown_range"] == [17, 20]
    assert results[-1].get("_fragment_read_count", 0) >= 4


async def test_read_file_duplicate_fragment_omits_repeated_content(tmp_path):
    from app.llm.tools.workspace import handle_read_file

    path = tmp_path / "notes.txt"
    path.write_text("\n".join(f"line {i}" for i in range(1, 11)), encoding="utf-8")

    first = await handle_read_file(str(tmp_path), "notes.txt", start_line=3, end_line=6)
    second = await handle_read_file(str(tmp_path), "notes.txt", start_line=3, end_line=6)
    forced = await handle_read_file(str(tmp_path), "notes.txt", start_line=3, end_line=6, force=True)

    assert first["ok"] is True
    assert "3:" in first["content"]
    assert second["ok"] is True
    assert second["content"] == ""
    assert second["content_omitted"] is True
    assert second["content_omitted_reason"] == "duplicate_read_range_already_returned"
    assert second["coverage_reason"] == "covered_by_prior_fragment_read"
    assert "force=true" in second["note"]
    assert "3:" in forced["content"]


async def test_read_file_subfragment_after_full_read_omits_repeated_content(tmp_path):
    from app.llm.tools.workspace import handle_read_file

    path = tmp_path / "brief.md"
    path.write_text("\n".join(f"section line {i}" for i in range(1, 9)), encoding="utf-8")

    full = await handle_read_file(str(tmp_path), "brief.md")
    sub = await handle_read_file(str(tmp_path), "brief.md", start_line=4, end_line=5)

    assert full["ok"] is True
    assert full["shown_range"] == [1, 8]
    assert sub["ok"] is True
    assert sub["content"] == ""
    assert sub["content_omitted"] is True
    assert sub["coverage_reason"] == "covered_by_prior_full_read"
    assert sub["covered_by_ranges"] == [[1, 8]]


async def test_todo_write_overuse_returns_direct_execution_instruction(tmp_path):
    from app.llm.tools.workspace import handle_todo_write

    todos = [{"id": "1", "content": "执行任务", "status": "in_progress"}]
    result = None
    for _ in range(30):
        result = await handle_todo_write(str(tmp_path), todos)

    assert result is not None
    assert result["todo_write_count"] == 30
    assert result["throttle_warning"]["issue"] == "frequent_todo_write_calls"
    assert result["throttle_warning"]["count"] == 30
    assert "does not execute work" in result["throttle_warning"]["fact"]
    assert "No workspace artifact is created by todo_write itself." in result["next_action_facts"]
    assert "next_action_instruction" not in result


async def test_todo_write_long_content_returns_fact_warning_without_rejecting(tmp_path):
    from app.llm.tools.workspace import handle_todo_write

    long_body = "Write section. " + ("This is document body text. " * 40)
    result = await handle_todo_write(str(tmp_path), [
        {"id": "1", "content": long_body, "status": "in_progress"},
    ])

    assert result["ok"] is True
    assert result["counts"]["total"] == 1
    warning = result["content_length_warning"]
    assert warning["issue"] == "todo_content_long_for_planning_state"
    assert warning["long_items"][0]["id"] == "1"
    assert "submitted todos were still recorded" in warning["fact"]
    assert "workspace" in warning["fact"]
    assert "office" in warning["fact"]
    assert result["next_action_facts"]


def test_normalize_tool_call_args_recovers_double_encoded_json_string():
    from app.llm.client import _normalize_tool_call_args_for_dispatch

    raw = json.dumps('{"action":"append","path":"report.docx","blocks":[{"type":"paragraph","text":"ok"}]}')

    args, err, repaired = _normalize_tool_call_args_for_dispatch(raw)

    assert err is None
    assert repaired is True
    assert args["action"] == "append"
    assert args["path"] == "report.docx"


def test_normalize_tool_call_args_recovers_large_office_blocks_with_inner_quotes_and_newlines():
    from app.llm.client import _normalize_tool_call_args_for_dispatch

    raw = '{"action":"append","path":"gen_docx_通信原理作业_第五章.docx","blocks":[{"type":"paragraph","text":"第一段包含未转义引号"眼睛"并且跨行\n第二行继续说明"}]}'

    args, err, repaired = _normalize_tool_call_args_for_dispatch(raw)

    assert err is None
    assert repaired is True
    assert args["action"] == "append"
    assert args["path"] == "gen_docx_通信原理作业_第五章.docx"
    assert args["blocks"][0]["text"] == '第一段包含未转义引号"眼睛"并且跨行\n第二行继续说明'


async def test_office_append_empty_args_suggests_complete_args(tmp_path):
    from docx import Document
    from app.llm.tools.office import handle_office

    Document().save(tmp_path / "report.docx")
    result = json.loads(await handle_office(str(tmp_path), {"action": "append", "path": "report.docx"}))

    assert result["ok"] is False
    assert result["error"] == "invalid_or_empty_args"
    assert result["recovery"]["action"] == "retry_with_complete_args"
    assert "blocks" in result["recovery"]["required_fields"]


async def test_office_missing_path_recovery_is_fact_not_fixed_batch_rule(tmp_path):
    from app.llm.tools.office import handle_office

    result = json.loads(await handle_office(str(tmp_path), {"action": "append"}))

    assert result["ok"] is False
    assert result["error"] == "invalid_or_empty_args"
    assert "path" in result["recovery"]["required_fields"]
    assert "same blocks" not in result["hint"].lower()
    assert "intended blocks" in result["next_call_fact"]


async def test_office_docx_read_supports_block_paging(tmp_path):
    from docx import Document
    from app.llm.tools.office import handle_office

    doc = Document()
    for idx in range(8):
        doc.add_paragraph(f"paragraph {idx}")
    doc.save(tmp_path / "paged.docx")

    result = json.loads(await handle_office(str(tmp_path), {
        "action": "read",
        "path": "paged.docx",
        "start_block": 2,
        "max_blocks": 3,
    }))

    assert result["ok"] is True
    assert result["returned_block_start"] == 2
    assert result["returned_block_end"] == 4
    assert result["returned_block_count"] == 3
    assert result["has_more_blocks"] is True
    assert result["next_start_block"] == 5
    assert [p["text"] for p in result["paragraphs"]] == [
        "paragraph 2",
        "paragraph 3",
        "paragraph 4",
    ]


async def test_office_docx_read_save_to_writes_segment_text(tmp_path):
    from docx import Document
    from app.llm.tools.office import handle_office

    doc = Document()
    doc.add_heading("Report", level=1)
    doc.add_paragraph("first paragraph")
    doc.add_paragraph("second paragraph")
    doc.save(tmp_path / "report.docx")

    result = json.loads(await handle_office(str(tmp_path), {
        "action": "read",
        "path": "report.docx",
        "save_to": "full_text.txt",
    }))

    assert result["ok"] is True
    assert result["saved_to"] == "full_text.txt"
    saved = (tmp_path / "full_text.txt").read_text(encoding="utf-8")
    assert "Report" in saved
    assert "first paragraph" in saved
    assert "second paragraph" in saved


async def test_office_verify_numbers_skips_context_metadata(tmp_path):
    from docx import Document
    from app.llm.tools.office import handle_office

    (tmp_path / "bench.csv").write_text(
        "distribution,op_type,mean_time_ns\n"
        "zipfian,lookup,38.5\n"
        "random,delete,203.9\n",
        encoding="utf-8",
    )
    doc = Document()
    doc.add_paragraph(
        "Benchmarks ran on a 13th Gen Intel Core i9-13980HX with 48 GB RAM, "
        "Windows 11, GCC 13.2, and Section 3.1 reports 38.5 ns for zipfian lookup."
    )
    doc.save(tmp_path / "paper.docx")

    result = json.loads(await handle_office(str(tmp_path), {
        "action": "verify_numbers",
        "path": "paper.docx",
        "csv_paths": ["bench.csv"],
        "tolerance": 0.05,
    }))

    assert result["ok"] is True
    assert result["mismatches_count"] == 0
    skipped = result["skipped_context_numbers_sample"]
    assert any(item["number_str"] == "13980" for item in skipped)
    assert any(item["number_str"] == "13.2" for item in skipped)
    assert "candidate CSV-backed numeric claims" in result["hint"]
    assert "强烈建议" not in result["hint"]


async def test_office_docx_append_keeps_backup_and_valid_zip(tmp_path):
    from docx import Document
    from app.llm.tools.office import handle_office

    doc = Document()
    doc.add_paragraph("before")
    doc.save(tmp_path / "report.docx")

    result = json.loads(await handle_office(str(tmp_path), {
        "action": "append",
        "path": "report.docx",
        "blocks": [{"type": "paragraph", "text": "after"}],
    }))

    assert result["ok"] is True
    assert (tmp_path / "report.docx.bak").exists()
    reread = json.loads(await handle_office(str(tmp_path), {
        "action": "read",
        "path": "report.docx",
    }))
    assert reread["ok"] is True
    texts = [p["text"] for p in reread["paragraphs"]]
    assert "before" in texts
    assert "after" in texts


async def test_office_docx_fill_empty_headings_batches_heading_repair(tmp_path):
    from docx import Document
    from app.llm.tools.office import handle_office

    doc = Document()
    doc.add_heading("Compression Report", level=0)
    doc.add_heading("", level=1)
    doc.add_heading("", level=2)
    doc.save(tmp_path / "report.docx")

    result = json.loads(await handle_office(str(tmp_path), {
        "action": "fill_empty_headings",
        "path": "report.docx",
        "headings": [
            {
                "level": 1,
                "text": "引言与背景",
                "after_blocks": [{"type": "paragraph", "text": "正文第一段。"}],
            },
            {
                "level": 2,
                "text": "实验设计",
                "after_blocks": [{"type": "paragraph", "text": "正文第二段。"}],
            },
        ],
    }))

    assert result["ok"] is True
    assert result["filled_count"] == 2
    assert result["remaining_empty_headings"] == 0

    reread = json.loads(await handle_office(str(tmp_path), {
        "action": "read",
        "path": "report.docx",
    }))
    texts = [p["text"] for p in reread["paragraphs"]]
    assert "引言与背景" in texts
    assert "实验设计" in texts
    assert "正文第一段。" in texts
    assert "正文第二段。" in texts


def test_office_embedded_image_ocr_reports_batch_offsets(tmp_path, monkeypatch):
    from PIL import Image
    from docx import Document
    from app.llm.tools import ocr_bridge
    from app.llm.tools.ocr_bridge import OcrResult

    image_paths = []
    colors = ["white", "red", "blue"]
    for idx in range(3):
        path = tmp_path / f"img{idx}.png"
        Image.new("RGB", (20 + idx, 20 + idx), colors[idx]).save(path)
        image_paths.append(path)

    doc = Document()
    for path in image_paths:
        doc.add_picture(str(path))
    doc.save(tmp_path / "images.docx")

    def fake_ocr_file(path, timeout=30):
        return OcrResult(ok=True, text=f"text from {Path(path).name}", score=1.0, tier="fast", engine="fake")

    monkeypatch.setattr(ocr_bridge, "ocr_file", fake_ocr_file)

    result = ocr_bridge.ocr_office_images(
        tmp_path / "images.docx",
        max_images=1,
        image_offset=1,
        max_workers=1,
    )

    assert result["ok"] is True
    assert result["total_images"] == 3
    assert result["image_offset"] == 1
    assert result["processed_images"] == 1
    assert result["has_more_images"] is True
    assert result["next_image_offset"] == 2
    assert result["ocr_count"] == 1


def test_delegate_does_not_auto_correct_document_output_to_edit():
    from app.llm.tools.delegate import _auto_correct_obvious_helper_kind

    kind, reason = _auto_correct_obvious_helper_kind(
        "code",
        "整理结果并生成 report.docx，包含表格和结论。",
    )

    assert kind == "code"
    assert reason is None


def test_delegate_does_not_auto_correct_mixed_code_and_chart_task():
    from app.llm.tools.delegate import _auto_correct_obvious_helper_kind

    kind, reason = _auto_correct_obvious_helper_kind(
        "code",
        "实现 .c 算法并编译 benchmark，最后用 matplotlib 生成 PNG 性能图。",
    )

    assert kind == "code"
    assert reason is None


def test_delegate_does_not_auto_correct_general_document_output_to_edit():
    from app.llm.tools.delegate import _auto_correct_obvious_helper_kind

    kind, reason = _auto_correct_obvious_helper_kind(
        "general",
        "请整理答案并生成 通信原理作业_第五章.docx。",
    )

    assert kind == "general"
    assert reason is None


def test_delegate_does_not_auto_correct_document_output_from_expected_outputs():
    from app.llm.tools.delegate import _auto_correct_obvious_helper_kind

    kind, reason = _auto_correct_obvious_helper_kind(
        "general",
        "请生成完整课程作业。",
        ["通信原理作业_第五章.docx"],
    )

    assert kind == "general"
    assert reason is None


def test_delegate_keeps_data_analysis_report_task_as_code_for_split():
    from app.llm.tools.delegate import _auto_correct_obvious_helper_kind

    kind, reason = _auto_correct_obvious_helper_kind(
        "code",
        "读取 sales.csv 做数据分析和回归分析，生成 analysis_results.csv，最终输出 final_report.docx。",
        ["analysis_results.csv", "final_report.docx"],
    )

    assert kind == "code"
    assert reason is None


def test_delegate_does_not_rewrite_edit_data_analysis_report_task():
    from app.llm.tools.delegate import _auto_correct_obvious_helper_kind

    kind, reason = _auto_correct_obvious_helper_kind(
        "edit",
        "读取 metrics.csv，做数据分析、统计汇总和可视化中间数据，最终生成 final_report.docx。",
        ["analysis_results.csv", "final_report.docx"],
    )

    assert kind == "edit"
    assert reason is None


def test_delegate_does_not_rewrite_general_directory_statistics_to_code():
    from app.llm.tools.delegate import _auto_correct_obvious_helper_kind

    kind, reason = _auto_correct_obvious_helper_kind(
        "general",
        "统计当前目录 .py 文件数量、总行数、总字符数，并输出最大的 5 个 Python 文件。需要运行脚本遍历目录。",
    )

    assert kind == "general"
    assert reason is None


def test_delegate_does_not_rewrite_general_markdown_report_kind():
    from app.llm.tools.delegate import _auto_correct_obvious_helper_kind

    kind, reason = _auto_correct_obvious_helper_kind(
        "general",
        "Create docs/algorithm_report.md from existing benchmark results and README notes.",
        ["_env/docs/algorithm_report.md"],
    )

    assert kind == "general"
    assert reason is None


def test_office_output_detection_uses_expected_outputs_when_present():
    from app.llm.tools.workspace_utils import _has_office_document_output

    assert _has_office_document_output(
        "Create a framework contract for a final Word/docx paper workflow.",
        ["paper_contract.json"],
    ) is False
    assert _has_office_document_output(
        "Create the final Word/docx paper.",
        ["paper.docx"],
    ) is True


def test_delegate_does_not_rewrite_mixed_general_code_and_markdown_report():
    from app.llm.tools.delegate import _auto_correct_obvious_helper_kind

    kind, reason = _auto_correct_obvious_helper_kind(
        "general",
        "Implement A* and Floyd-Warshall, add benchmark coverage, then create docs/algorithm_report.md.",
        ["_env/src/algolab/graph.py", "_env/docs/algorithm_report.md"],
    )

    assert kind == "general"
    assert reason is None


def test_delegate_does_not_rewrite_code_markdown_only_documentation_task():
    from app.llm.tools.delegate import _auto_correct_obvious_helper_kind

    kind, reason = _auto_correct_obvious_helper_kind(
        "code",
        (
            "Write documentation for the graph algorithm project. Read existing "
            "_env/src/algolab/graph.py and _env/src/algolab/benchmark.py first, "
            "then update README.md and docs/algorithm_report.md."
        ),
        ["_env/README.md", "_env/docs/algorithm_report.md"],
    )

    assert kind == "code"
    assert reason is None


def test_delegate_does_not_rewrite_read_text_evidence_task():
    from app.llm.tools.delegate import _auto_correct_obvious_helper_kind

    kind, reason = _auto_correct_obvious_helper_kind(
        "read",
        (
            "Write a Python OCR script for images and PDFs, extract all visible text, "
            "and save visual text evidence to _helpers_shared/ocr_results.txt."
        ),
        ["_helpers_shared/ocr_results.txt"],
    )

    assert kind == "read"
    assert reason is None


def test_delegate_does_not_rewrite_code_ocr_task_to_text_report_edit():
    from app.llm.tools.delegate import _auto_correct_obvious_helper_kind

    kind, reason = _auto_correct_obvious_helper_kind(
        "code",
        (
            "Write a Python OCR script for images and PDFs, extract all visible text, "
            "and save visual text evidence to _helpers_shared/ocr_results.txt."
        ),
        ["_helpers_shared/ocr_results.txt"],
    )

    assert kind == "code"
    assert reason is None


async def test_delegate_sanitize_corrects_general_source_stats_to_code(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    result = await _sanitize_and_validate_tasks(
        {
            "tasks": [{
                "task_id": "project_stats",
                "kind": "general",
                "prompt": "Walk the directory, count Python files, total line count, character count, and largest files.",
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    assert isinstance(result, str)
    data = json.loads(result)
    assert data["ok"] is False
    assert data["error_kind"] == "unsupported_helper_kind"
    assert data["task_id"] == "project_stats"
    assert "general helper" in data["hint"]


def test_delegate_does_not_auto_correct_image_clarity_kind():
    from app.llm.tools.delegate import _auto_correct_obvious_helper_kind

    kind, reason = _auto_correct_obvious_helper_kind(
        "draw",
        "查看图片 _downloaded_media/img_real_test.png，判断图片中文字/内容是否清晰可辨。只需报告清晰或不清晰。",
    )

    assert kind == "draw"
    assert reason is None


def test_read_helper_prompt_requires_upgrade_for_clarity_checks():
    from app.llm.tools import delegate
    from app.llm.tools import helper_prompt_catalog

    prompt = delegate._HELPER_SYSTEM_READ
    source = inspect.getsource(helper_prompt_catalog)
    assert "clarity/readability judgments" in prompt
    assert "清晰度或可辨性判断" in source
    assert "suitable max_tier" in source
    assert "编号/数值/标签读取" in source
    assert "精确视觉证据使用 allow_upgrade" in source
    assert "no_stronger_tier is true" in prompt
    assert "engine_config.cache_hit=true" in prompt
    assert "keep each file/tier attempt purposeful" in prompt
    assert "Use additional search only for named gaps" in prompt


async def test_ocr_rejects_path_in_image_base64(tmp_path):
    from app.llm.tools.registry import _handle_ocr

    result = json.loads(await _handle_ocr(str(tmp_path), {"image_base64": "uploads/card.png"}))

    assert result["ok"] is False
    assert "改用 image_path" in result["error"]
    assert result["suggested_args"]["image_path"] == "uploads/card.png"


async def test_ocr_adds_purpose_quality_warning_for_missing_target_fields(tmp_path, monkeypatch):
    from app.llm.tools import registry
    from app.llm.tools.ocr_bridge import OcrResult

    image = tmp_path / "card.png"
    image.write_bytes(b"fake image bytes")

    async def fake_run_gpu_ocr(fn, *args, **kwargs):
        return OcrResult(ok=True, text="## E2E OCR CARD 9137", score=1.0, tier="accurate", engine="mineru")

    monkeypatch.setattr(registry, "run_gpu_ocr", fake_run_gpu_ocr)

    result = json.loads(await registry._handle_ocr(str(tmp_path), {
        "image_path": "card.png",
        "allow_upgrade": True,
        "max_tier": "accurate",
        "purpose": "读取卡片编号、三个数值和 PASS/FAIL 结果",
    }))

    assert result["ok"] is True
    warning = result["purpose_quality_warning"]
    assert warning["signal"] == "ocr_purpose_incomplete"
    assert "numbers" in warning["missing"]
    assert "labels_or_result" in warning["missing"]


async def test_ocr_marks_result_sufficient_when_purpose_fields_are_present(tmp_path, monkeypatch):
    from app.llm.tools import registry
    from app.llm.tools.ocr_bridge import OcrResult

    image = tmp_path / "card.png"
    image.write_bytes(b"fake image bytes")

    async def fake_run_gpu_ocr(fn, *args, **kwargs):
        return OcrResult(
            ok=True,
            text="E2E OCR CARD 9137\nALPHA 42.5\nBETA 18.0\nGAMMA 39.5\nRESULT PASS",
            score=1.0,
            tier="balanced",
            engine="mineru",
            next_tier="accurate",
        )

    monkeypatch.setattr(registry, "run_gpu_ocr", fake_run_gpu_ocr)

    result = json.loads(await registry._handle_ocr(str(tmp_path), {
        "image_path": "card.png",
        "allow_upgrade": True,
        "max_tier": "accurate",
        "purpose": "读取卡片编号、三个数值和 PASS/FAIL 结果",
    }))

    assert result["ok"] is True
    assert result["sufficient_for_purpose"] is True
    assert "purpose_quality_warning" not in result


async def test_ocr_save_to_writes_text_for_segmented_reading(tmp_path, monkeypatch):
    from app.llm.tools import registry
    from app.llm.tools.ocr_bridge import OcrResult

    image = tmp_path / "card.png"
    image.write_bytes(b"fake image bytes")

    async def fake_run_gpu_ocr(fn, *args, **kwargs):
        return OcrResult(ok=True, text="line1\nline2\nline3", score=1.0, tier="accurate", engine="mineru")

    monkeypatch.setattr(registry, "run_gpu_ocr", fake_run_gpu_ocr)

    result = json.loads(await registry._handle_ocr(str(tmp_path), {
        "image_path": "card.png",
        "allow_upgrade": True,
        "max_tier": "accurate",
        "save_to": "ocr/card.txt",
    }))

    assert result["ok"] is True
    assert result["saved_to"] == "ocr/card.txt"
    assert result["saved_text_lines"] == 3
    assert "text" not in result
    assert (tmp_path / "ocr" / "card.txt").read_text(encoding="utf-8") == "line1\nline2\nline3"


async def test_ocr_raw_text_path_is_exposed_as_workspace_file(tmp_path, monkeypatch):
    from app.llm.tools import registry
    from app.llm.tools.ocr_bridge import OcrResult

    image = tmp_path / "card.png"
    image.write_bytes(b"fake image bytes")
    raw_dir = tmp_path / "outside_raw"
    raw_dir.mkdir()
    raw_file = raw_dir / "raw.txt"
    raw_file.write_text("raw line 1\nraw line 2", encoding="utf-8")

    async def fake_run_gpu_ocr(fn, *args, **kwargs):
        return OcrResult(
            ok=True,
            text="folded text",
            score=1.0,
            tier="balanced",
            engine="mineru",
            raw_text_path=str(raw_file),
            raw_text_len=21,
        )

    monkeypatch.setattr(registry, "run_gpu_ocr", fake_run_gpu_ocr)

    result = json.loads(await registry._handle_ocr(str(tmp_path), {
        "image_path": "card.png",
        "allow_upgrade": True,
    }))

    assert result["ok"] is True
    assert "raw_text_path" not in result
    rel = result["raw_text_workspace_path"]
    assert rel.startswith("ocr_raw/")
    assert (tmp_path / rel).read_text(encoding="utf-8") == "raw line 1\nraw line 2"
    assert "read_file" in result["raw_text_hint"]


def test_delegate_does_not_auto_correct_pdf_document_output_to_edit():
    from app.llm.tools.delegate import _auto_correct_obvious_helper_kind

    kind, reason = _auto_correct_obvious_helper_kind(
        "general",
        "请整理材料并输出 final_report.pdf。",
    )

    assert kind == "general"
    assert reason is None


async def test_delegate_sanitize_preserves_explicit_helper_kind_for_docx_output(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    result = await _sanitize_and_validate_tasks(
        {
            "tasks": [{
                "task_id": "hw_514_v2",
                "kind": "general",
                "mode": "hard",
                "prompt": "请生成一份完整的 Word 文档作业解答。",
                "expected_outputs": ["通信原理作业_第五章.docx"],
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    assert isinstance(result, str)
    data = json.loads(result)
    assert data["ok"] is False
    assert data["error_kind"] == "unsupported_helper_kind"
    assert data["task_id"] == "hw_514_v2"
    assert "general helper" in data["hint"]


async def test_dispatch_helper_aliases_route_to_workspace_tools(tmp_path):
    from app.llm.tools import registry

    result = await registry.dispatch(
        "write",
        {"path": "note.txt", "content": "hello"},
        archive_id="archive",
        group_id="group",
        user_id="user",
        workspace_dir=str(tmp_path),
        caller="helper",
    )
    data = json.loads(result)
    assert data["ok"] is True
    assert data["action"] == "write"
    assert data.get("_alias_note")

    result1b = await registry.dispatch(
        "write_workspace_file",
        {"path": "note2.txt", "content": "hello again"},
        archive_id="archive",
        group_id="group",
        user_id="user",
        workspace_dir=str(tmp_path),
        caller="helper",
    )
    data1b = json.loads(result1b)
    assert data1b["ok"] is True
    assert data1b["action"] == "write"
    assert data1b.get("_alias_note")

    result2 = await registry.dispatch(
        "run",
        {"command": "python -c \"print('ok')\"", "timeout_sec": 5},
        archive_id="archive",
        group_id="group",
        user_id="user",
        workspace_dir=str(tmp_path),
        caller="helper",
    )
    data2 = json.loads(result2)
    assert data2["ok"] is True
    assert data2["action"] == "run"
    assert data2.get("_alias_note")


async def test_dispatch_helper_caller_sets_owner_for_bash_and_resource_tools(tmp_path):
    from app.core.core_processes import current_owner
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import registry

    (tmp_path / "_env").mkdir()
    target = tmp_path / "_env" / "app.js"
    target.write_text("old\n", encoding="utf-8")
    env = EnvironmentContext(
        root_dir=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
        project_key="project",
    )

    with runtime_context("environment", env):
        result = await registry.dispatch(
            "bash",
            {
                "command": (
                    "python -c \"from pathlib import Path; "
                    "Path('_env/app.js').write_text('new', encoding='utf-8')\""
                ),
                "timeout_sec": 10,
            },
            archive_id="archive",
            group_id="group",
            user_id="user",
            workspace_dir=str(tmp_path),
            caller="helper",
        )

    data = json.loads(result)
    assert data["ok"] is True
    assert target.read_text(encoding="utf-8") == "new"
    assert current_owner() == "main:unknown"

    resource = await registry.dispatch(
        "request_resource",
        {"kind": "read", "reason": "need exact source evidence", "needed_outputs": ["evidence.txt"]},
        archive_id="archive",
        group_id="group",
        user_id="user",
        workspace_dir=str(tmp_path),
        caller="helper",
    )
    resource_data = json.loads(resource)
    assert resource_data["ok"] is False
    assert resource_data["action"] == "request_resource"
    assert resource_data["suggested_helper_kind"] == "read"


async def test_environment_edit_file_missing_project_path_points_to_staged_edit_path(tmp_path):
    import json

    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import registry

    staged = tmp_path / "_env" / "service"
    staged.mkdir(parents=True)
    (staged / "render.py").write_text("def render():\n    return 'old'\n", encoding="utf-8")
    env = EnvironmentContext(
        root_dir=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
        project_key="project",
    )

    with runtime_context("environment", env):
        result = await registry.dispatch(
            "edit_file",
            {
                "path": "service/render.py",
                "old_str": "return 'old'",
                "new_str": "return 'new'",
            },
            archive_id="archive",
            group_id="group",
            user_id="user",
            workspace_dir=str(tmp_path),
            caller="helper",
        )

    data = json.loads(result)
    assert data["ok"] is False
    assert data["staged_candidate"] == "_env/service/render.py"
    assert data["staged_candidate_exists"] is True
    assert "workspace editing" in data["error"]
    assert "Retry read_file" not in data["error"]


async def test_dispatch_helper_caller_large_read_does_not_hit_main_guard(tmp_path):
    from app.llm.tools import registry

    (tmp_path / "large.txt").write_text("x" * 130_000, encoding="utf-8")

    result = await registry.dispatch(
        "read_file",
        {"path": "large.txt"},
        archive_id="archive",
        group_id="group",
        user_id="user",
        workspace_dir=str(tmp_path),
        caller="helper",
    )
    data = json.loads(result)

    assert data["ok"] is True
    assert data.get("error_kind") != "main_thread_large_read_should_delegate_or_target"


def test_main_thread_tools_do_not_expose_implementation_tool_names():
    from app.llm.tools.registry import ROUND2_TOOLS

    names = {tool["function"]["name"] for tool in ROUND2_TOOLS}

    assert not {
        "python", "bash", "edit_file", "multi_edit", "insert_in_file",
        "office", "progress_note", "ocr", "tts",
    } & names


def test_main_thread_tool_metadata_side_effects_are_consistent():
    from app.llm.tools import registry

    by_name = {meta.name: meta for meta in registry.MAIN_THREAD_TOOL_METAS}

    for name in ("expand_warm", "expand_cold", "expand_kb", "search_files"):
        assert by_name[name].read_only is True
        assert by_name[name].side_effect == "none"

    for name in ("workspace", "commit_to_main", "delegate", "processes"):
        assert by_name[name].read_only is False
        assert by_name[name].side_effect in {"workspace", "external"}
        assert by_name[name].requires_permission == "generate_file"

    assert by_name["mark_avoid_mention"].side_effect == "memory"
    assert by_name["ask_user_question"].side_effect == "external"
    assert by_name["ocr"].main_thread_allowed is False
    assert by_name["tts"].main_thread_allowed is False


def test_tool_meta_helpers_filter_only_main_thread_tools():
    from app.llm.tools.tool_meta import ToolMeta, meta_by_name, schemas_for_main_thread

    allowed_schema = {"type": "function", "function": {"name": "allowed"}}
    blocked_schema = {"type": "function", "function": {"name": "blocked"}}
    metas = [
        ToolMeta("allowed", allowed_schema, read_only=True, side_effect="none", requires_permission="chat"),
        ToolMeta("blocked", blocked_schema, read_only=False, side_effect="workspace", requires_permission="generate_file", main_thread_allowed=False),
    ]

    assert schemas_for_main_thread(metas) == [allowed_schema]
    assert set(meta_by_name(metas)) == {"allowed", "blocked"}


def test_helper_tool_filter_keeps_edit_helper_off_code_execution_tools():
    from app.llm.tools.delegate import _HELPER_TOOLS, _filter_tools_for_kind

    names = {tool["function"]["name"] for tool in _filter_tools_for_kind("edit", _HELPER_TOOLS)}

    assert "python" not in names
    assert "bash" not in names
    assert "office" in names
    assert "edit_file" in names


def test_helper_kind_schema_includes_resource_helpers():
    from app.llm.tools.delegate import SPAWN_HELPER_TOOL_SCHEMA, VALID_HELPER_KINDS
    from app.llm.tools.registry import DELEGATE_TOOL_SCHEMA

    delegate_kind = DELEGATE_TOOL_SCHEMA["function"]["parameters"]["properties"]["tasks"]["items"]["properties"]["kind"]

    assert "tts" in VALID_HELPER_KINDS
    assert "read" in VALID_HELPER_KINDS
    assert "project_map" in VALID_HELPER_KINDS
    assert "file_summary" in VALID_HELPER_KINDS
    assert "impact_review" in VALID_HELPER_KINDS
    assert "inventory" in VALID_HELPER_KINDS
    assert "summarize" in VALID_HELPER_KINDS
    assert "tts" in SPAWN_HELPER_TOOL_SCHEMA["function"]["parameters"]["properties"]["kind"]["enum"]
    assert "read" in SPAWN_HELPER_TOOL_SCHEMA["function"]["parameters"]["properties"]["kind"]["enum"]
    assert "project_map" in SPAWN_HELPER_TOOL_SCHEMA["function"]["parameters"]["properties"]["kind"]["enum"]
    assert "file_summary" in delegate_kind["enum"]
    assert "impact_review" in delegate_kind["enum"]
    assert "inventory" in delegate_kind["enum"]
    assert "summarize" not in delegate_kind["enum"]
    assert "tts" in delegate_kind["enum"]
    assert "read" in delegate_kind["enum"]


def test_task_quality_guard_keeps_summarize_environment_only():
    from app.llm.tools.delegate import _task_quality_guard_environment_helper_text

    helper_line, suggested_line, principle = _task_quality_guard_environment_helper_text()
    assert "summarize" not in helper_line
    assert "summarize" not in suggested_line
    assert "summarize" not in principle
    assert "project_map" in principle


def test_task_quality_guard_keeps_selected_source_files_as_file_summary():
    from pathlib import Path

    prompt = Path("app/llm/aux_prompts.py").read_text(encoding="utf-8")
    assert "summarizing selected source/config files in a code project -> file_summary" in prompt
    assert "source/config files" in prompt


def test_delegate_schema_clarifies_kind_routing_for_draw_code_verify():
    from app.llm.tools.registry import DELEGATE_TOOL_SCHEMA

    text = DELEGATE_TOOL_SCHEMA["function"]["description"]
    kind_text = (
        DELEGATE_TOOL_SCHEMA["function"]["parameters"]["properties"]["tasks"]["items"]
        ["properties"]["kind"]["description"]
    )

    assert "implementation" in text
    assert "verification" in text
    assert "charts" in text
    assert "programs, HTML, scripts" in kind_text
    assert "directory/source statistics" in kind_text
    assert "project_map" in kind_text
    assert "`draw` for final image/chart files" in kind_text
    assert "`verify` for checking existing code/images/documents" in kind_text
    assert "prepared archive contents" in kind_text
    assert "executable file-preparation steps" in kind_text


def test_helper_tool_filter_keeps_main_thread_alias_boundary():
    from app.llm.tools.delegate import _HELPER_TOOLS, _filter_tools_for_kind

    code_names = {tool["function"]["name"] for tool in _filter_tools_for_kind("code", _HELPER_TOOLS)}
    edit_names = {tool["function"]["name"] for tool in _filter_tools_for_kind("edit", _HELPER_TOOLS)}
    verify_names = {tool["function"]["name"] for tool in _filter_tools_for_kind("verify", _HELPER_TOOLS)}

    assert "python" in code_names
    assert "bash" in code_names
    assert "python" not in edit_names
    assert "bash" not in edit_names
    assert "office" in edit_names
    assert "workspace" in edit_names
    assert "edit_file" in edit_names
    assert "multi_edit" in edit_names
    assert "read_file" in edit_names
    assert "code_index" not in edit_names
    assert "search_across_files" not in edit_names
    assert "read_function" not in edit_names
    assert "processes" not in edit_names
    assert "ocr" not in code_names
    assert "tts" not in code_names
    assert "ocr" not in edit_names
    assert "tts" not in edit_names
    for names in (verify_names,):
        assert "office" not in names
        assert "edit_file" not in names
        assert "multi_edit" not in names
        assert "insert_in_file" not in names
    for names in ():
        assert "code_index" not in names
        assert "read_function" not in names
        assert "search_across_files" not in names
    assert "bash" not in verify_names
    assert "python" in verify_names
    assert "workspace" in verify_names

    verify_tools = _filter_tools_for_kind("verify", _HELPER_TOOLS)
    verify_workspace = next(tool for tool in verify_tools if tool["function"]["name"] == "workspace")
    verify_action = verify_workspace["function"]["parameters"]["properties"]["action"]
    assert verify_action["enum"] == ["run", "locate"]


async def test_edit_helper_workspace_runtime_guard_allows_existing_check_run(tmp_path):
    from app.core.core_processes import reset_current_helper_kind, set_current_helper_kind
    from app.llm.tools import registry

    (tmp_path / "verify_report.py").write_text("print('PASS')\n", encoding="utf-8")
    token = set_current_helper_kind("edit")
    try:
        run_data = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "run", "command": "python verify_report.py", "timeout_sec": 5},
        ))
    finally:
        reset_current_helper_kind(token)

    assert run_data["ok"] is True
    assert run_data["returncode"] == 0
    assert "PASS" in run_data["stdout"]


async def test_edit_helper_workspace_runtime_guard_rejects_non_check_run(tmp_path):
    from app.core.core_processes import reset_current_helper_kind, set_current_helper_kind
    from app.llm.tools import registry

    token = set_current_helper_kind("edit")
    try:
        run_data = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "run", "command": "python calc.py", "timeout_sec": 5},
        ))
    finally:
        reset_current_helper_kind(token)

    assert run_data["ok"] is False
    assert run_data["blocked_reason"] == "edit_helper_workspace_run_limited_to_existing_checks"
    assert run_data["suggested_helper_kind"] == "code"


def test_general_helper_prompt_is_removed_from_visible_catalog():
    import app.llm.tools.delegate as delegate

    assert not hasattr(delegate, "_HELPER_SYSTEM_GENERAL")
    assert "general" not in delegate.VALID_HELPER_KINDS


def test_all_selected_helper_prompts_include_consistency_contract():
    from app.llm.tools.delegate import _select_helper_system
    from app.llm.tools.helper_kinds import VALID_HELPER_KINDS

    for kind in VALID_HELPER_KINDS:
        prompt = _select_helper_system(kind)
        assert "## Helper Consistency Contract" in prompt
        assert "Stay inside your assigned helper kind" in prompt
        assert "Use measured evidence for exact claims" in prompt
        assert "terminal_reason=resource_required" in prompt
        assert "所有 helper 共用" in prompt


def test_selected_helper_prompts_include_actual_tool_boundary():
    from app.llm.tools.delegate import _HELPER_TOOLS, _filter_tools_for_kind, _select_helper_system
    from app.llm.tools.helper_kinds import VALID_HELPER_KINDS, _normalize_helper_kind_mode

    for kind in VALID_HELPER_KINDS:
        prompt = _select_helper_system(kind)
        tool_names = {
            tool["function"]["name"]
            for tool in _filter_tools_for_kind(kind, _HELPER_TOOLS)
            if isinstance(tool, dict)
        }
        normalized_kind, _ = _normalize_helper_kind_mode(kind)
        assert "## Actual Tool Boundary" in prompt
        assert f"Helper kind: `{normalized_kind}`." in prompt
        assert "Available tools: (none listed)" not in prompt
        available_line = prompt.split("Available tools:", 1)[1].split("\n", 1)[0]
        for tool_name in tool_names:
            assert tool_name in available_line
        for required in ("bash", "python", "workspace", "office", "ocr", "tts"):
            if required not in tool_names:
                assert required in prompt.split("Unavailable capability names in this helper:", 1)[1].split("\n", 1)[0]
        if kind == "draw":
            assert "bash" not in tool_names
            assert "python" in tool_names


def test_helper_user_prompt_adds_ordered_acceptance_facts_only_when_needed():
    from app.llm.tools.delegate_runner import _build_helper_user_prompt

    ordered_prompt = _build_helper_user_prompt(
        prompt="Reproduce the current failure before editing the parser, then verify the fix.",
        dynamic_prompt_prefix_parts=[],
        kind="code",
        task_contract_context='{"acceptance_checks":["baseline before change"]}',
    )
    plain_prompt = _build_helper_user_prompt(
        prompt="Fix the parser and verify it.",
        dynamic_prompt_prefix_parts=[],
        kind="code",
        task_contract_context='{"acceptance_checks":["tests pass"]}',
    )

    assert "## Ordered Acceptance Facts" in ordered_prompt
    assert "pre-change evidence before editing" in ordered_prompt
    assert "## Ordered Acceptance Facts" not in plain_prompt


def test_helper_user_prompt_turns_main_apply_tools_into_boundary_facts():
    from app.llm.tools.delegate_runner import _build_helper_user_prompt

    prompt = _build_helper_user_prompt(
        prompt=(
            "### Deliverable 1: report.md (create with env_apply_create)\n"
            "Write the report from supplied evidence."
        ),
        dynamic_prompt_prefix_parts=[],
        kind="edit",
        task_contract_context='{"expected_outputs":["report.md"]}',
    )

    assert "## Helper Tool Boundary Facts" in prompt
    assert "main-process routing, apply, or helper-result consumption steps" in prompt
    assert "not executable helper tools" in prompt
    assert "Create text/Markdown/Office artifacts with your visible edit" in prompt
    assert "### Deliverable 1: report.md (create with env_apply_create)" in prompt


def test_helper_source_field_provenance_marks_structured_input_summaries():
    from app.llm.tools.delegate import _mark_unverified_helper_source_interpretations

    prompt, observations = _mark_unverified_helper_source_interpretations(
        "Create artifact.md. Source summary says service cost is $220 for 3 nights.",
        input_files=["source.json", "profile.yaml"],
    )

    assert "## Source Field Provenance" in prompt
    assert "read the raw field names/values/notes from input_files" in prompt
    assert "A total, package, included" in prompt
    assert observations == [{"kind": "source_field_provenance_added", "input_files": 2}]


def test_helper_source_field_provenance_skips_plain_input_prompts():
    from app.llm.tools.delegate import _mark_unverified_helper_source_interpretations

    original = "Create artifact.md from the listed files and preserve the requested headings."
    prompt, observations = _mark_unverified_helper_source_interpretations(
        original,
        input_files=["outline.md"],
    )

    assert prompt == original
    assert observations == []


def test_code_and_verify_helper_prompts_preserve_statistic_units():
    from app.llm.tools.delegate import _HELPER_SYSTEM_CODE, _HELPER_SYSTEM_VERIFY

    for text in (_HELPER_SYSTEM_CODE, _HELPER_SYSTEM_VERIFY):
        assert "Characters, bytes, file size, line count, and file count are" in text
        assert "not interchangeable" in text or "different metrics" in text
    assert "统计结果要区分字符、字节、行数" in _HELPER_SYSTEM_CODE
    assert "核对统计单位" in _HELPER_SYSTEM_VERIFY


def test_code_helper_prompt_uses_search_evidence_for_symbol_migrations():
    from app.llm.tools.delegate import _HELPER_SYSTEM_CODE

    assert "identifier, API, schema, field, import, or contract migrations" in _HELPER_SYSTEM_CODE
    assert "use search/index evidence to discover impacted references" in _HELPER_SYSTEM_CODE
    assert "verify remaining old/new references" in _HELPER_SYSTEM_CODE
    assert "迁移/契约变更" in _HELPER_SYSTEM_CODE


def test_project_analysis_helper_filters_are_read_only_and_light():
    from app.llm.tools.delegate import _HELPER_TOOLS, _filter_tools_for_kind

    blocked = {
        "workspace", "office", "edit_file", "multi_edit", "insert_in_file",
        "bash", "python", "processes", "delegate", "spawn_helper", "wait_helper",
        "ocr", "tts",
    }
    required = {"read_file", "search_files", "search_in_file", "code_index", "progress_note"}
    for kind in ("project_map", "file_summary", "impact_review"):
        names = {
            tool["function"]["name"]
            for tool in _filter_tools_for_kind(kind, _HELPER_TOOLS)
            if isinstance(tool, dict)
        }
        assert not (blocked & names)
        assert required <= names


def test_inventory_helper_filter_allows_readonly_statistics_without_artifact_tools():
    from app.llm.tools.delegate import _HELPER_TOOLS, _filter_tools_for_kind

    tools = _filter_tools_for_kind("inventory", _HELPER_TOOLS)
    names = {
        tool["function"]["name"]
        for tool in tools
        if isinstance(tool, dict)
    }

    assert {"read_file", "search_files", "search_in_file", "code_index", "progress_note"} <= names
    assert {"workspace", "python"} <= names
    assert "bash" not in names
    assert not ({"office", "edit_file", "multi_edit", "insert_in_file", "ocr", "tts", "delegate"} & names)
    workspace = next(tool for tool in tools if tool["function"]["name"] == "workspace")
    action = workspace["function"]["parameters"]["properties"]["action"]
    assert action["enum"] == ["run", "locate", "write", "mkdir"]


def test_inventory_helper_prompt_declares_inventory_workflow():
    from app.llm.tools.delegate import _select_helper_system

    prompt = _select_helper_system("inventory")
    assert "environment project inventory helper" in prompt
    assert "first-pass inventory" in prompt
    assert "README" in prompt
    assert "file counts by suffix/category" in prompt
    assert "temporary scripts or notes under `_scratch/`" in prompt
    assert "Write only scratch inventory notes under allowed scratch paths" in prompt
    assert "inventory helper 只在项目模式做工程摸底" in prompt


def test_legacy_summarize_helper_alias_uses_inventory_workflow():
    from app.llm.tools.delegate import _HELPER_TOOLS, _filter_tools_for_kind, _select_helper_system
    from app.llm.tools.helper_kinds import _normalize_helper_kind_mode

    assert _normalize_helper_kind_mode("summarize") == ("inventory", "easy")
    assert _filter_tools_for_kind("summarize", _HELPER_TOOLS) == _filter_tools_for_kind("inventory", _HELPER_TOOLS)
    prompt = _select_helper_system("summarize")
    assert "environment project inventory helper" in prompt


async def test_inventory_helper_workspace_runtime_guard_allows_only_temporary_writes(tmp_path):
    from app.core.core_processes import reset_current_helper_kind, set_current_helper_kind
    from app.llm.tools import registry

    token = set_current_helper_kind("inventory")
    try:
        scratch_py = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "write", "path": "_scratch/inventory_probe.py", "content": "print('ok')"},
        ))
        assert scratch_py["ok"] is True
        assert (tmp_path / "_scratch" / "inventory_probe.py").read_text(encoding="utf-8") == "print('ok')"

        scratch_note = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "write", "path": "scratch/notes.md", "content": "# notes"},
        ))
        assert scratch_note["ok"] is True

        env_write = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "write", "path": "_env/foo.py", "content": "x"},
        ))
        assert env_write["ok"] is False
        assert env_write["blocked_reason"] == "inventory_helper_workspace_path_forbidden"

        deliverable_write = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "write", "path": "_scratch/report.docx", "content": "x"},
        ))
        assert deliverable_write["ok"] is False
        assert deliverable_write["blocked_reason"] == "inventory_helper_workspace_path_forbidden"

        locate_data = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "locate", "pattern": "*.txt"},
        ))
        assert locate_data["ok"] is True
    finally:
        reset_current_helper_kind(token)


async def test_environment_workspace_locate_reports_chat_workspace_scope_fact(tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import registry

    project = tmp_path / "project"
    chat_ws = tmp_path / "chat"
    project.mkdir()
    chat_ws.mkdir()
    (project / "app.py").write_text("print('project')\n", encoding="utf-8")

    env = EnvironmentContext(
        root_dir=str(project),
        archive_id="arch_test",
        group_id="env_user_test",
        user_id="user",
        project_key="project",
    )
    with runtime_context("environment", env):
        result = json.loads(await registry._handle_workspace(
            str(chat_ws),
            {"action": "locate", "pattern": "*.py"},
        ))

    fact = result["environment_workspace_scope_fact"]
    assert result["ok"] is True
    assert result["searched_scope"] == "chat_workspace_only"
    assert result["project_search_performed"] is False
    assert result["project_search_tools"] == ["env_list_tree", "env_search", "env_inventory", "env_read"]
    assert "not the real environment project root" in fact
    assert "transient chat workspace" in fact
    assert "env_list_tree" in fact
    assert "not evidence that a project file is absent" in fact
    assert "helper" not in fact.lower()


async def test_environment_search_files_reports_chat_scope_fact(tmp_path, monkeypatch):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import registry

    project = tmp_path / "project"
    chat_ws = tmp_path / "chat"
    project.mkdir()
    chat_ws.mkdir()
    (project / "users.db").write_bytes(b"sqlite")
    (chat_ws / "note.txt").write_text("chat workspace only\n", encoding="utf-8")

    async def fake_search_files(*args, **kwargs):
        return []

    monkeypatch.setattr(registry.kb_mem, "search_files", fake_search_files)

    env = EnvironmentContext(
        root_dir=str(project),
        archive_id="arch_test",
        group_id="env_user_test",
        user_id="user",
        project_key="project",
    )
    with runtime_context("environment", env):
        result = json.loads(await registry._handle_search_files(
            "arch_test",
            "env_user_test",
            str(chat_ws),
            {"query": "users.db", "limit": 5},
        ))

    fact = result["environment_search_scope_fact"]
    assert result["searched_scope"] == "chat_history_and_chat_workspace_only"
    assert result["project_search_performed"] is False
    assert result["project_search_tools"] == ["env_list_tree", "env_search", "env_inventory", "env_read"]
    assert "does not search the real environment project root" in fact
    assert "transient chat workspaces" in fact
    assert "not evidence that a project file is absent" in fact
    assert "helper" not in fact.lower()


def test_resource_helper_tool_filters_isolate_tts_and_read():
    from app.llm.tools.delegate import _HELPER_TOOLS, _filter_tools_for_kind

    tts_tools = _filter_tools_for_kind("tts", _HELPER_TOOLS)
    read_tools = _filter_tools_for_kind("read", _HELPER_TOOLS)
    legacy_ocr_tools = _filter_tools_for_kind("ocr", _HELPER_TOOLS)
    tts_names = {tool["function"]["name"] for tool in tts_tools}
    read_names = {tool["function"]["name"] for tool in read_tools}
    legacy_ocr_names = {tool["function"]["name"] for tool in legacy_ocr_tools}

    assert "tts" in tts_names
    assert "ocr" not in tts_names
    assert "bash" not in tts_names
    assert "office" not in tts_names
    assert "todo_write" not in tts_names
    assert "todo_read" not in tts_names
    assert "search_files" not in tts_names
    assert "recall_thread" not in tts_names
    assert "progress_note" not in tts_names
    assert "read_skill" not in tts_names

    assert "ocr" in read_names
    assert "office" in read_names
    assert "read_file" in read_names
    assert "inspect_file" in read_names
    assert "tts" not in read_names
    assert "bash" not in read_names
    assert "python" not in read_names
    assert "edit_file" not in read_names
    assert legacy_ocr_names == read_names
    read_workspace = next(tool for tool in read_tools if tool["function"]["name"] == "workspace")
    action = read_workspace["function"]["parameters"]["properties"]["action"]
    assert action["enum"] == ["write"]
    assert set(read_workspace["function"]["parameters"]["required"]) == {"action", "path", "content"}
    tts_workspace = next(tool for tool in tts_tools if tool["function"]["name"] == "workspace")
    tts_action = tts_workspace["function"]["parameters"]["properties"]["action"]
    assert tts_action["enum"] == ["write"]
    assert set(tts_workspace["function"]["parameters"]["required"]) == {"action", "path", "content"}



def test_ocr_tts_schema_avoids_concept_question_tool_bias():
    from app.llm.tools.tool_schemas import OCR_TOOL_SCHEMA, TTS_TOOL_SCHEMA

    ocr_desc = OCR_TOOL_SCHEMA["function"]["description"]
    tts_desc = TTS_TOOL_SCHEMA["function"]["description"]

    assert "when the current task needs evidence from a specific visual/document file" in ocr_desc
    assert "Concept, principle, log, or scheduling questions about OCR" in ocr_desc
    assert "treating engine names and tier labels as internal" in ocr_desc
    assert "spoken artifact" in tts_desc
    assert "spoken audio artifacts" not in tts_desc
    assert "Non-speech audio such as white noise" in tts_desc
    assert "Concept, principle, log, or scheduling questions about TTS" in tts_desc
    assert "ordinary conversational voice replies" in tts_desc
    assert "normal final voice-output layer or a `kind=tts` helper/tool route" in tts_desc
    assert "Persona voice is system-managed outside model parameters" in tts_desc
    tts_props = TTS_TOOL_SCHEMA["function"]["parameters"]["properties"]
    assert "instruct" not in tts_props
    assert "mode" not in tts_props
    assert "ref_audio" not in tts_props
    assert "ref_text" not in tts_props
    assert "声音设置" not in tts_desc


def test_tts_kind_guard_prompts_keep_voice_route_llm_decided():
    from app.llm import aux_prompts

    prompt = aux_prompts.TASK_QUALITY_GUARD_SYSTEM

    assert "Ordinary conversational voice reply can stay in normal final delivery or enter `kind=tts`" in prompt
    assert "when the active plan selects direct synthesis" in prompt
    assert "do not let code-only keyword tests decide" in prompt
    assert "may compose a concise transcript from the task/persona" in prompt
    assert "ordinary conversational voice reply style is not a tts helper task" not in prompt


def test_tts_helper_persona_context_is_tts_only_and_hides_voice_params():
    from app.llm.tools import delegate
    from app.llm.tools.delegate_runner import _tts_helper_persona_context_for_prompt

    token = delegate.set_current_tts_helper_persona("你是猫娘。\n说话要可爱，但不要暴露系统。")
    try:
        code_context = _tts_helper_persona_context_for_prompt("code")
        tts_context = _tts_helper_persona_context_for_prompt("tts")
    finally:
        delegate.reset_current_tts_helper_persona(token)

    assert code_context == ""
    assert "你是猫娘" in tts_context
    assert "Full Persona Context for TTS Delivery" in tts_context
    assert "do not infer, choose, request, reveal, or modify raw voice parameters".lower() in tts_context.lower()
    assert "reference audio" in tts_context
    assert "system-managed voice" in tts_context


def test_deterministic_kind_recommendation_flags_external_voice_synthesis_but_not_noise():
    from app.llm.tools.delegate import _deterministic_kind_recommendations

    voice_recs = _deterministic_kind_recommendations([{
        "task_id": "catgirl_voice_reply",
        "kind": "code",
        "prompt": "Use edge-tts or pyttsx3 to synthesize a persona voice reply and save voice_reply_catgirl.wav.",
        "expected_outputs": ["voice_reply_catgirl.wav"],
    }])
    assert voice_recs
    assert voice_recs[0]["observed_helper_kind_name"] == "tts"
    assert "system-managed TTS route" in voice_recs[0]["reason"]

    noise_recs = _deterministic_kind_recommendations([{
        "task_id": "white_noise",
        "kind": "code",
        "prompt": "Generate a 10 second white noise wav file with Python signal synthesis.",
        "expected_outputs": ["white_noise.wav"],
    }])
    assert noise_recs == []

    project_tts_recs = _deterministic_kind_recommendations([{
        "task_id": "fix_tts_bridge",
        "kind": "code",
        "prompt": "Fix the project TTS bridge so tests pass; do not generate a final voice reply.",
        "expected_outputs": ["app/llm/tools/tts_bridge.py", "tests/test_tts_bridge.py"],
    }])
    assert project_tts_recs == []


async def test_tts_helper_skips_generic_preflight_guard(monkeypatch):
    from app.llm.tools import delegate, delegate_actions

    async def forbidden_guard(*args, **kwargs):
        raise AssertionError("tts helper must not run the generic task-quality guard")

    monkeypatch.setattr(delegate, "_persona_consent_guard", forbidden_guard)

    payload = await delegate_actions._run_delegate_preflight_guard(
        {},
        [{
            "task_id": "voice_reply",
            "kind": "tts",
            "mode": "easy",
            "prompt": "Generate the already-authorized voice reply from the supplied text.",
            "expected_outputs": ["reply.wav"],
        }],
        "trace_tts_preflight_skip",
    )

    assert payload is None


async def test_mixed_tts_batch_preflight_guard_filters_tts_only(monkeypatch):
    from app.llm.tools import delegate, delegate_actions

    seen = []

    async def fake_guard(persona, user_message, tasks):
        seen.append(tasks)
        return True, "", [], []

    monkeypatch.setattr(delegate, "_persona_consent_guard", fake_guard)

    payload = await delegate_actions._run_delegate_preflight_guard(
        {},
        [
            {
                "task_id": "voice_reply",
                "kind": "tts",
                "mode": "easy",
                "prompt": "Generate the already-authorized voice reply.",
                "expected_outputs": ["reply.wav"],
            },
            {
                "task_id": "read_source",
                "kind": "read",
                "mode": "easy",
                "prompt": "Read the supplied source material.",
                "expected_outputs": ["read_evidence.txt"],
            },
        ],
        "trace_tts_preflight_filter",
    )

    assert payload is None
    assert len(seen) == 1
    assert [task["task_id"] for task in seen[0]] == ["read_source"]


async def test_external_tts_runtime_guard_warns_once_and_allows_retry(tmp_path, monkeypatch):
    from app.core.core_processes import reset_current_helper_kind, set_current_helper_kind
    from app.llm.tools import registry

    registry._reset_external_tts_runtime_guard_for_tests()
    script = tmp_path / "synth.py"
    script.write_text(
        "import edge_tts\n"
        "text = 'hello voice reply'\n"
        "voice_reply_path = 'voice_reply.wav'\n",
        encoding="utf-8",
    )
    calls = []

    async def fake_run(workspace_dir, command, timeout_sec=None, abort_event=None):
        calls.append((workspace_dir, command, timeout_sec))
        return {"ok": True, "command": command}

    monkeypatch.setattr(registry.ws_tool, "handle_run", fake_run)
    token = set_current_helper_kind("code")
    try:
        first = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "run", "command": "python synth.py", "timeout_sec": 5},
            caller_kind="helper",
        ))
        second = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "run", "command": "python synth.py", "timeout_sec": 5},
            caller_kind="helper",
        ))
    finally:
        reset_current_helper_kind(token)
        registry._reset_external_tts_runtime_guard_for_tests()

    assert first["ok"] is False
    assert first["blocked_reason"] == "external_tts_engine_requires_system_tts_tool"
    assert first["retry_allowed_on_next_attempt"] is True
    assert first["suggested_tool"] == "tts"
    assert "speech/voice output" in first["fact"]
    assert "Non-speech audio such as white noise" in first["fact"]
    assert "voice/audio" not in first["fact"]
    assert second["ok"] is True
    assert len(calls) == 1


async def test_external_tts_runtime_guard_does_not_block_non_speech_audio(tmp_path, monkeypatch):
    from app.core.core_processes import reset_current_helper_kind, set_current_helper_kind
    from app.llm.tools import registry

    registry._reset_external_tts_runtime_guard_for_tests()
    script = tmp_path / "noise.py"
    script.write_text(
        "import wave\n"
        "output = 'white_noise.wav'\n",
        encoding="utf-8",
    )
    calls = []

    async def fake_run(workspace_dir, command, timeout_sec=None, abort_event=None):
        calls.append(command)
        return {"ok": True, "command": command}

    monkeypatch.setattr(registry.ws_tool, "handle_run", fake_run)
    token = set_current_helper_kind("code")
    try:
        data = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "run", "command": "python noise.py", "timeout_sec": 5},
            caller_kind="helper",
        ))
    finally:
        reset_current_helper_kind(token)
        registry._reset_external_tts_runtime_guard_for_tests()

    assert data["ok"] is True
    assert calls == ["python noise.py"]


def test_ocr_workspace_schema_frames_txt_as_internal_evidence():
    from app.llm.tools.tool_schemas import OCR_WORKSPACE_TOOL_SCHEMA

    fn = OCR_WORKSPACE_TOOL_SCHEMA["function"]
    assert "visual text evidence" in fn["description"]
    assert "internal source material" in fn["description"]
    assert "not user-facing copy" in fn["description"]

    props = fn["parameters"]["properties"]
    assert "Fixed value" in props["action"]["description"]
    assert "固定为写入文本材料" in props["action"]["description"]
    assert "Confirmed text" in props["content"]["description"]
    assert "可确认内容" in props["content"]["description"]


def test_read_helper_tools_exclude_shell_and_mutating_edit_access():
    from app.llm.tools.delegate import _HELPER_TOOLS
    from app.llm.tools.helper_kinds import _filter_tools_for_kind

    names = {
        tool["function"]["name"]
        for tool in _filter_tools_for_kind("read", _HELPER_TOOLS)
        if isinstance(tool, dict)
    }

    assert {"ocr", "office", "inspect_file", "read_file", "search_in_file", "workspace", "request_resource"} <= names
    assert "bash" not in names
    assert "python" not in names
    assert "edit_file" not in names
    assert "multi_edit" not in names


def test_helpers_cannot_spawn_or_wait_for_helpers():
    from app.llm.tools.delegate import _HELPER_TOOLS, _filter_tools_for_kind

    for kind in ("code", "edit", "verify", "draw", "tts", "read", "project_map", "file_summary", "impact_review", "inventory", "summarize"):
        names = {
            tool["function"]["name"]
            for tool in _filter_tools_for_kind(kind, _HELPER_TOOLS)
            if isinstance(tool, dict)
        }
        assert "delegate" not in names
        assert "spawn_helper" not in names
        assert "wait_helper" not in names
        assert "processes" not in names


def test_helpers_do_not_mutate_long_term_preferences_directly():
    from app.llm.tools.delegate import _HELPER_TOOLS, _filter_tools_for_kind

    for kind in ("code", "edit", "verify", "draw", "tts", "read", "project_map", "file_summary", "impact_review", "inventory", "summarize"):
        names = {
            tool["function"]["name"]
            for tool in _filter_tools_for_kind(kind, _HELPER_TOOLS)
            if isinstance(tool, dict)
        }
        assert "mark_avoid_mention" not in names


def test_delegate_schema_says_helpers_do_not_spawn_helpers():
    from app.llm.tools.tool_schemas import DELEGATE_TOOL_SCHEMA

    description = DELEGATE_TOOL_SCHEMA["function"]["description"]
    assert "main process remains the coordinator" in description
    assert "manages dependencies" in description
    assert "waits or resumes based on heartbeats" in description
    assert "helper 现在可以**自己再 spawn_helper**" not in description


async def test_legacy_spawn_wait_helper_runtime_paths_are_disabled(tmp_path):
    from app.core.core_processes import reset_current_owner, set_current_owner
    from app.llm.tools.delegate import handle_spawn_helper, handle_wait_helper

    token = set_current_owner("helper:trace:edit_task")
    try:
        spawn_data = json.loads(await handle_spawn_helper(
            {"task_id": "draw_resource", "prompt": "draw chart", "kind": "draw"},
            archive_id="a",
            group_id="g",
            user_id="u",
            helper_workspace=str(tmp_path),
        ))
        wait_data = json.loads(await handle_wait_helper({"task_id": "draw_resource"}))
    finally:
        reset_current_owner(token)

    assert spawn_data["ok"] is False
    assert spawn_data["error"] == "helper_spawn_disabled"
    assert "coordinated by the main process through delegate" in spawn_data["blocked_reason"]
    assert spawn_data["recovery_facts"]["coordinator"] == "main_process"
    assert "spawn_resource_helper" in spawn_data["recovery_facts"]["available_shapes"]
    assert "Only the main process can create helpers" in spawn_data["dispatch_boundary_facts"]
    assert wait_data["ok"] is False
    assert wait_data["error"] == "helper_wait_disabled"
    assert "coordinated by the main process" in wait_data["blocked_reason"]
    assert wait_data["recovery_facts"]["coordinator"] == "main_process"
    assert "inspect_helper_state" in wait_data["recovery_facts"]["available_shapes"]
    assert "Only the main process can inspect global helper state" in wait_data["dispatch_boundary_facts"]


async def test_read_helper_workspace_runtime_guard_only_allows_text_evidence(tmp_path):
    from app.core.core_processes import reset_current_helper_kind, set_current_helper_kind
    from app.llm.tools import registry

    token = set_current_helper_kind("read")
    try:
        run_data = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "run", "command": "python crop.py", "timeout_sec": 5},
        ))
        assert run_data["ok"] is False
        assert run_data["blocked_reason"] == "read_helper_workspace_action_forbidden"

        script_data = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "write", "path": "crop.py", "content": "print('crop')"},
        ))
        assert script_data["ok"] is False
        assert script_data["blocked_reason"] == "read_helper_workspace_path_forbidden"

        text_data = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "write", "path": "ocr_result.txt", "content": "# OCR\ntext"},
        ))
        assert text_data["ok"] is True
        assert (tmp_path / "ocr_result.txt").read_text(encoding="utf-8") == "# OCR\ntext"
    finally:
        reset_current_helper_kind(token)


async def test_tts_helper_workspace_runtime_guard_only_allows_text_notes(tmp_path):
    from app.core.core_processes import reset_current_helper_kind, set_current_helper_kind
    from app.llm.tools import registry

    token = set_current_helper_kind("tts")
    try:
        run_data = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "run", "command": "python synth.py", "timeout_sec": 5},
        ))
        assert run_data["ok"] is False
        assert run_data["blocked_reason"] == "tts_helper_workspace_action_forbidden"

        wav_data = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "write", "path": "voice.wav", "content": "not audio"},
        ))
        assert wav_data["ok"] is False
        assert wav_data["blocked_reason"] == "tts_helper_workspace_path_forbidden"

        text_data = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "write", "path": "tts_manifest.txt", "content": "voice file generated"},
        ))
        assert text_data["ok"] is True
        assert (tmp_path / "tts_manifest.txt").read_text(encoding="utf-8") == "voice file generated"
    finally:
        reset_current_helper_kind(token)


async def test_tts_helper_runtime_uses_filtered_tool_list(monkeypatch, tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    helper_ws.mkdir()
    main_ws.mkdir()
    captured = {}

    async def fake_loop(msgs, tools, **kwargs):
        captured["tool_names"] = [
            tool.get("function", {}).get("name", "")
            for tool in tools
            if isinstance(tool, dict)
        ]
        captured["helper_kind"] = kwargs.get("helper_kind")
        (helper_ws / "greeting.wav").write_bytes(b"RIFFxxxxWAVE")
        return '```json\n{"files": ["greeting.wav"]}\n```', msgs

    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))

    result = await delegate._run_one_helper(
        task_id="tts_greeting_runtime_tools",
        prompt="Generate greeting.wav with system TTS.",
        main_workspace=str(main_ws),
        helper_workspace=str(helper_ws),
        archive_id="archive",
        group_id="group",
        user_id="user",
        resume=False,
        local_abort=asyncio.Event(),
        wait_for_register=asyncio.Event(),
        user_lang="en",
        kind="tts",
        mode="easy",
        helper_think=False,
        expected_outputs=["greeting.wav"],
    )

    assert result["ok"] is True
    assert captured["helper_kind"] == "tts"
    assert captured["tool_names"] == [
        "workspace",
        "fetch_to_temp",
        "request_resource",
        "inspect_file",
        "read_file",
        "search_in_file",
        "fetch_indexed_file",
        "tts",
    ]
    assert "fetch_group_file" not in captured["tool_names"]
    assert "todo_write" not in captured["tool_names"]
    assert "todo_read" not in captured["tool_names"]
    assert "progress_note" not in captured["tool_names"]
    assert "recall_thread" not in captured["tool_names"]
    assert "read_skill" not in captured["tool_names"]
    assert "search_files" not in captured["tool_names"]
    assert "bash" not in captured["tool_names"]
    assert "python" not in captured["tool_names"]


async def test_general_helper_cannot_use_workspace_even_if_tool_called(tmp_path):
    from app.core.core_processes import reset_current_helper_kind, set_current_helper_kind
    from app.llm.tools import registry

    token = set_current_helper_kind("general")
    try:
        data = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "write", "path": "general.txt", "content": "internal"},
        ))
        assert data["ok"] is False
        assert data["blocked_reason"] == "general_helper_workspace_forbidden"
        assert data["suggested_tool"] == "request_resource"
        assert not (tmp_path / "general.txt").exists()
    finally:
        reset_current_helper_kind(token)


async def test_verify_helper_workspace_runtime_guard_allows_only_validation(tmp_path):
    from app.core.core_processes import reset_current_helper_kind, set_current_helper_kind
    from app.llm.tools import registry

    token = set_current_helper_kind("verify")
    try:
        write_data = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "write", "path": "verify.txt", "content": "internal"},
        ))
        assert write_data["ok"] is False
        assert write_data["blocked_reason"] == "verify_helper_workspace_action_forbidden"
        assert not (tmp_path / "verify.txt").exists()

        locate_data = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "locate", "pattern": "*.txt"},
        ))
        assert locate_data["ok"] is True
    finally:
        reset_current_helper_kind(token)


async def test_helper_env_write_contract_blocks_undeclared_project_copy(tmp_path):
    from app.core.core_processes import (
        reset_current_helper_expected_outputs,
        reset_current_helper_kind,
        set_current_helper_expected_outputs,
        set_current_helper_kind,
    )
    from app.llm.tools import registry

    (tmp_path / "_env" / "src" / "algolab").mkdir(parents=True)
    (tmp_path / "_env" / "src" / "algolab" / "graph.py").write_text("old\n", encoding="utf-8")
    (tmp_path / "_env" / "tests").mkdir(parents=True)
    (tmp_path / "_env" / "tests" / "test_graph.py").write_text("old\n", encoding="utf-8")

    kind_token = set_current_helper_kind("code")
    outputs_token = set_current_helper_expected_outputs(["_env/tests/test_graph.py"])
    try:
        blocked = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "write", "path": "_env/src/algolab/graph.py", "content": "new\n"},
        ))
        assert blocked["ok"] is False
        assert blocked["blocked_reason"] == "helper_env_write_outside_expected_outputs"
        assert blocked["matching_helper_kind"] == "code"
        assert blocked["observed_recovery_tool"] == "request_resource"
        assert blocked["observed_recovery_shape"]["needed_outputs"] == ["_env/src/algolab/graph.py"]

        allowed = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "write", "path": "_env/tests/test_graph.py", "content": "new\n"},
        ))
        assert allowed["ok"] is True
        assert (tmp_path / "_env" / "tests" / "test_graph.py").read_text(encoding="utf-8") == "new\n"

        scratch = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "write", "path": "scratch_notes.txt", "content": "notes\n"},
        ))
        assert scratch["ok"] is True
    finally:
        reset_current_helper_expected_outputs(outputs_token)
        reset_current_helper_kind(kind_token)


async def test_helper_env_write_scope_allows_component_directory_without_changing_acceptance(tmp_path):
    from app.core.core_processes import (
        reset_current_helper_expected_outputs,
        reset_current_helper_kind,
        reset_current_helper_write_scopes,
        set_current_helper_expected_outputs,
        set_current_helper_kind,
        set_current_helper_write_scopes,
    )
    from app.llm.tools import registry

    (tmp_path / "_env" / "src" / "agent").mkdir(parents=True)
    (tmp_path / "_env" / "ui").mkdir(parents=True)

    kind_token = set_current_helper_kind("code")
    outputs_token = set_current_helper_expected_outputs(["_env/src/agent/core.py"])
    scopes_token = set_current_helper_write_scopes(["_env/src/agent"])
    try:
        sibling = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "write", "path": "_env/src/agent/tools.py", "content": "TOOLS = {}\n"},
        ))
        assert sibling["ok"] is True

        outside = json.loads(await registry._handle_workspace(
            str(tmp_path),
            {"action": "write", "path": "_env/ui/app.js", "content": "console.log('x');\n"},
        ))
        assert outside["ok"] is False
        assert outside["blocked_reason"] == "helper_env_write_outside_expected_outputs"
        assert outside["matching_helper_kind"] == "code"
        assert outside["observed_recovery_shape"]["needed_outputs"] == ["_env/ui/app.js"]
    finally:
        reset_current_helper_write_scopes(scopes_token)
        reset_current_helper_expected_outputs(outputs_token)
        reset_current_helper_kind(kind_token)


def test_environment_expected_outputs_expand_to_staged_paths():
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools.delegate import _expand_environment_expected_outputs

    env = EnvironmentContext(
        root_dir="F:/project",
        archive_id="arch",
        group_id="group",
        user_id="user",
        project_key="proj",
    )
    prompt = "Edit `_env/src/taskboard/models.py` and `_env/src/taskboard/storage.py`."
    with runtime_context("environment", env):
        expanded = _expand_environment_expected_outputs(prompt, ["models.py", "storage.py"])
    assert expanded == ["_env/src/taskboard/models.py", "_env/src/taskboard/storage.py"]


def test_environment_expected_outputs_expand_paths_with_spaces():
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools.delegate import _expand_environment_expected_outputs

    env = EnvironmentContext(
        root_dir="F:/project",
        archive_id="arch",
        group_id="group",
        user_id="user",
        project_key="proj",
    )
    name = "包涵 - 2026.5-8月 口语话题更新 (2026.5.15修改).docx"
    prompt = f"Update _env/docs/{name} after reading all evidence."
    with runtime_context("environment", env):
        expanded = _expand_environment_expected_outputs(prompt, [name])
    assert expanded == [f"_env/docs/{name}"]


async def test_workspace_write_allows_existing_environment_project_copy_with_fact(tmp_path):
    from app.llm.tools import registry

    target = tmp_path / "_env" / "app" / "llm" / "tools" / "environment.py"
    target.parent.mkdir(parents=True)
    target.write_text("def existing():\n    return True\n", encoding="utf-8")

    staged = json.loads(await registry._handle_workspace(
        str(tmp_path),
        {
            "action": "write",
            "path": "_env/app/llm/tools/environment.py",
            "content": "short",
        },
    ))
    assert staged["ok"] is True
    assert staged["staged_project_copy"] is True
    assert staged["staged_project_path"] == "_env/app/llm/tools/environment.py"
    assert "env_apply_replace" in staged["suggested_next_tools"]
    assert target.read_text(encoding="utf-8") == "short"

    created = json.loads(await registry._handle_workspace(
        str(tmp_path),
        {
            "action": "write",
            "path": "_env/app/llm/tools/new_environment_tests.py",
            "content": "def test_new():\n    assert True\n",
        },
    ))
    assert created["ok"] is True


async def test_main_thread_blocks_editing_existing_environment_project_copy(tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools.workspace import handle_edit_file, handle_insert_in_file, handle_multi_edit

    target = tmp_path / "_env" / "app.js"
    target.parent.mkdir(parents=True)
    target.write_text('const form = document.getElementById("contact-formm");\n', encoding="utf-8")
    env = EnvironmentContext(
        root_dir="F:/project",
        archive_id="arch",
        group_id="group",
        user_id="user",
        project_key="proj",
    )

    with runtime_context("environment", env):
        edit = await handle_edit_file(
            str(tmp_path),
            "_env/app.js",
            'document.getElementById("contact-formm")',
            'document.getElementById("contact-form")',
        )
        multi = await handle_multi_edit(
            str(tmp_path),
            "_env/app.js",
            [{
                "old_str": 'document.getElementById("contact-formm")',
                "new_str": 'document.getElementById("contact-form")',
            }],
        )
        insert = await handle_insert_in_file(str(tmp_path), "_env/app.js", -1, "// done\n")

    assert edit["ok"] is True
    assert edit["staged_project_copy"] is True
    assert "real project file" in edit["pending_project_apply_fact"]
    assert multi["ok"] is False
    assert "old_str not found" in json.dumps(multi["failures"], ensure_ascii=False)
    assert insert["ok"] is True
    assert insert["staged_project_copy"] is True
    assert target.read_text(encoding="utf-8") == 'const form = document.getElementById("contact-form");\n\n// done\n'


async def test_chat_mode_allows_editing_env_named_workspace_file(tmp_path):
    from app.llm.tools.workspace import handle_edit_file

    target = tmp_path / "_env" / "notes.txt"
    target.parent.mkdir(parents=True)
    target.write_text("alpha beta\n", encoding="utf-8")

    result = await handle_edit_file(str(tmp_path), "_env/notes.txt", "alpha", "bravo")

    assert result["ok"] is True
    assert target.read_text(encoding="utf-8") == "bravo beta\n"


async def test_main_thread_allows_compact_utility_python_script(tmp_path):
    from app.llm.tools import registry

    result = json.loads(await registry._handle_workspace(
        str(tmp_path),
        {
            "action": "write",
            "path": "_append_screenshots.py",
            "content": (
                "import re\n"
                "main = open('_env/screenshots_content.txt', 'r', encoding='utf-8').read()\n"
                "patch = open('ocr_raw/append.txt', 'r', encoding='utf-8').read()\n"
                "updated = re.sub('PARTIAL', 'PASS', main) + '\\n' + patch\n"
                "open('_env/screenshots_content.txt', 'w', encoding='utf-8').write(updated)\n"
                "print(len(updated.splitlines()))\n"
            ),
        },
    ))

    assert result["ok"] is True
    assert (tmp_path / "_append_screenshots.py").exists()


async def test_main_thread_project_source_python_requests_delegate(tmp_path):
    from app.llm.tools import registry

    result = json.loads(await registry._handle_workspace(
        str(tmp_path),
        {
            "action": "write",
            "path": "db_index_project/rbtree.py",
            "content": "class RedBlackTree:\n    def insert(self, key):\n        pass\n",
        },
    ))

    assert result["ok"] is False
    assert result["error_kind"] == "main_thread_project_artifact_should_delegate"
    assert result["recovery_facts"]["matching_helper_kind"] in {"code", "edit"}
    assert not (tmp_path / "db_index_project" / "rbtree.py").exists()


async def test_main_thread_allows_small_temp_verification_script(tmp_path):
    from app.llm.tools import registry

    result = json.loads(await registry._handle_workspace(
        str(tmp_path),
        {
            "action": "write",
            "path": ".temp/_verify_docx.py",
            "content": "from docx import Document\nprint(len(Document('db_index_paper.docx').paragraphs))\n",
        },
    ))

    assert result["ok"] is True
    assert (tmp_path / ".temp" / "_verify_docx.py").is_file()


def test_env_apply_create_rejects_large_main_authored_project_text(tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import environment

    project = tmp_path / "project"
    workspace = tmp_path / "workspace"
    project.mkdir()
    workspace.mkdir()
    env = EnvironmentContext(root_dir=str(project), archive_id="arch", group_id="group", user_id="user", project_key="proj")

    content = "\n".join(f"- contract line {i}: module boundary and acceptance rule" for i in range(60))
    with runtime_context("environment", env):
        result = environment._handle_apply_create(
            str(workspace),
            {"path": "db_index_project/CONTRACT.md", "content": content},
        )

    assert result["ok"] is False
    assert result["error_kind"] == "main_thread_project_artifact_create_should_delegate"
    assert result["recovery_facts"]["matching_helper_kind"] == "edit"
    assert result["candidate_preserved"] is True
    candidate_path = workspace / result["candidate_workspace_path"]
    assert candidate_path.read_text(encoding="utf-8") == content
    assert result["recovery_facts"]["input_files"] == [result["candidate_workspace_path"]]
    assert "candidate_apply_arguments" not in result["recovery_facts"]
    assert "candidate_apply_fact" not in result
    assert "not a clean producer-owned output by itself" in result["candidate_handoff_fact"]
    assert "Avoid pasting" in result["candidate_preservation_fact"]
    assert not (project / "db_index_project" / "CONTRACT.md").exists()


def test_env_apply_create_rejects_directory_targets(tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import environment

    project = tmp_path / "project"
    project.mkdir()
    env = EnvironmentContext(root_dir=str(project), archive_id="arch", group_id="group", user_id="user", project_key="proj")

    with runtime_context("environment", env):
        result = environment._handle_apply_create(str(tmp_path / "workspace"), {"path": "contracts"})

    assert result["ok"] is False
    assert result["error_kind"] == "env_apply_create_requires_file_target"
    assert not (project / "contracts").exists()


def test_env_apply_create_existing_file_suggests_replace_not_delete(tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import environment

    project = tmp_path / "project"
    project.mkdir()
    (project / "result.csv").write_text("old\n", encoding="utf-8")
    env = EnvironmentContext(
        root_dir=str(project),
        archive_id="arch_existing_create",
        group_id="env_user_test",
        user_id="user",
        project_key="project",
    )

    with runtime_context("environment", env):
        result = environment._handle_apply_create(
            str(tmp_path / "workspace"),
            {"path": "result.csv", "content": "new\n"},
        )

    assert result["ok"] is False
    assert result["error_kind"] == "env_apply_create_target_exists"
    assert "env_apply_replace" in result["existing_file_fact"]
    assert "Deleting the file first is not required" in result["existing_file_fact"]
    assert result["suggested_tools"] == ["env_diff", "env_apply_replace", "env_run"]


def test_env_apply_create_refuses_failed_helper_staged_file(tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import environment

    project = tmp_path / "project"
    workspace = tmp_path / "workspace"
    staged = workspace / "_env" / "db_index_project" / "src" / "base.py"
    staged.parent.mkdir(parents=True)
    staged.write_text("class Index:\n    pass\n", encoding="utf-8")
    project.mkdir()
    env = EnvironmentContext(root_dir=str(project), archive_id="arch", group_id="group", user_id="user", project_key="proj")
    environment.record_env_helper_outputs(
        str(workspace),
        task_id="scaffold_core",
        files=["_env/db_index_project/src/base.py"],
        ok=False,
        terminal_reason="outputs_missing",
        outputs_complete=False,
        kind="code",
        mode="easy",
    )

    with runtime_context("environment", env):
        result = environment._handle_apply_create(
            str(workspace),
            {
                "path": "db_index_project/src/base.py",
                "workspace_path": "_env/db_index_project/src/base.py",
            },
        )

    assert result["ok"] is False
    assert result["error_kind"] == "staged_environment_file_not_ready"
    assert result["source_task_id"] == "scaffold_core"
    assert not (project / "db_index_project" / "src" / "base.py").exists()


def test_env_apply_create_refuses_unknown_staged_new_project_file(tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import environment

    project = tmp_path / "project"
    workspace = tmp_path / "workspace"
    staged = workspace / "_env" / "db_index_project" / "src" / "base.py"
    staged.parent.mkdir(parents=True)
    staged.write_text("class Index:\n    pass\n", encoding="utf-8")
    project.mkdir()
    env = EnvironmentContext(root_dir=str(project), archive_id="arch", group_id="group", user_id="user", project_key="proj")

    with runtime_context("environment", env):
        result = environment._handle_apply_create(
            str(workspace),
            {
                "path": "db_index_project/src/base.py",
                "workspace_path": "_env/db_index_project/src/base.py",
            },
        )

    assert result["ok"] is False
    assert result["error_kind"] == "staged_environment_file_without_ready_provenance"
    assert not (project / "db_index_project" / "src" / "base.py").exists()


def test_env_apply_create_allows_ready_helper_staged_file(tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import environment

    project = tmp_path / "project"
    workspace = tmp_path / "workspace"
    staged = workspace / "_env" / "db_index_project" / "src" / "base.py"
    staged.parent.mkdir(parents=True)
    staged.write_text("class Index:\n    pass\n", encoding="utf-8")
    project.mkdir()
    env = EnvironmentContext(root_dir=str(project), archive_id="arch", group_id="group", user_id="user", project_key="proj")
    environment.record_env_helper_outputs(
        str(workspace),
        task_id="scaffold_core",
        files=["_env/db_index_project/src/base.py"],
        ok=True,
        terminal_reason="completed",
        outputs_complete=True,
        kind="code",
        mode="easy",
    )

    with runtime_context("environment", env):
        result = environment._handle_apply_create(
            str(workspace),
            {
                "path": "db_index_project/src/base.py",
                "workspace_path": "_env/db_index_project/src/base.py",
            },
        )

    assert result["ok"] is True
    assert (project / "db_index_project" / "src" / "base.py").read_text(encoding="utf-8") == "class Index:\n    pass\n"
    assert result["acceptance_fact"]["helper_owned"] is True
    assert "producer boundary" in result["acceptance_fact"]["fact"]
    assert "final intended project state" in result["acceptance_fact"]["fact"]


def test_env_apply_replace_reports_helper_owned_final_state_fact(tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import environment

    project = tmp_path / "project"
    workspace = tmp_path / "workspace"
    staged = workspace / "_env" / "src" / "settings.py"
    staged.parent.mkdir(parents=True)
    staged.write_text("VALUE = 2\n", encoding="utf-8")
    project_file = project / "src" / "settings.py"
    project_file.parent.mkdir(parents=True)
    project_file.write_text("VALUE = 1\n", encoding="utf-8")
    project.mkdir(exist_ok=True)
    env = EnvironmentContext(root_dir=str(project), archive_id="arch", group_id="group", user_id="user", project_key="proj")
    environment.record_env_helper_outputs(
        str(workspace),
        task_id="settings_fix",
        files=["_env/src/settings.py"],
        ok=True,
        terminal_reason="completed",
        outputs_complete=True,
        kind="code",
        mode="easy",
    )

    with runtime_context("environment", env):
        result = environment._handle_apply_replace(
            str(workspace),
            {
                "path": "src/settings.py",
                "workspace_path": "_env/src/settings.py",
                "expected_hash": environment._sha256(project_file),
            },
        )

    assert result["ok"] is True
    assert project_file.read_text(encoding="utf-8") == "VALUE = 2\n"
    assert result["acceptance_fact"]["helper_owned"] is True
    assert "content quality remains at the producer boundary" in result["acceptance_fact"]["fact"]
    assert "A check that ran before a later project apply covers the earlier state" in result["acceptance_fact"]["fact"]


def test_env_diff_missing_project_target_reports_create_path_facts(tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import environment

    project = tmp_path / "project"
    workspace = tmp_path / "workspace"
    staged = workspace / "_env" / "eu_active_2026.csv"
    staged.parent.mkdir(parents=True)
    staged.write_text("email,region\nx@example.com,EU\n", encoding="utf-8")
    project.mkdir()
    env = EnvironmentContext(root_dir=str(project), archive_id="arch", group_id="group", user_id="user", project_key="proj")

    with runtime_context("environment", env):
        result = environment._handle_diff(
            str(workspace),
            {
                "path": "eu_active_2026.csv",
                "workspace_path": "_env/eu_active_2026.csv",
            },
        )

    assert result["ok"] is True
    assert result["diff_available"] is False
    assert result["fact_kind"] == "env_diff_project_target_missing"
    assert result["project_file_exists"] is False
    assert result["workspace_path_exists"] is True
    assert result["recovery_facts"] == {
        "matching_tool_shape": "env_apply_create",
        "tool": "env_apply_create",
        "arguments": {
            "path": "eu_active_2026.csv",
            "workspace_path": "_env/eu_active_2026.csv",
        },
    }
    assert not (project / "eu_active_2026.csv").exists()


async def test_main_thread_blocks_large_new_env_project_artifact_write(tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import registry

    project = tmp_path / "project"
    project.mkdir()
    env = EnvironmentContext(root_dir=str(project), archive_id="arch", group_id="group", user_id="user", project_key="proj")
    content = "class RedBlackTree:\n" + "\n".join(f"    def method_{i}(self): return {i}" for i in range(80))

    with runtime_context("environment", env):
        result = json.loads(await registry._handle_workspace(
            str(tmp_path / "workspace"),
            {
                "action": "write",
                "path": "_env/db_index_project/src/rbtree.py",
                "content": content,
            },
        ))

    assert result["ok"] is False
    assert result["error_kind"] == "main_thread_project_artifact_should_delegate"
    assert result["recovery_facts"]["matching_helper_kind"] == "code"
    assert not (tmp_path / "workspace" / "_env" / "db_index_project" / "src" / "rbtree.py").exists()


async def test_request_resource_is_helper_only_and_structured(tmp_path):
    from app.core.core_processes import (
        reset_current_helper_kind,
        reset_current_helper_proc_id,
        set_current_helper_kind,
        set_current_helper_proc_id,
    )
    from app.llm.tools import registry

    main_data = json.loads(await registry._handle_request_resource(
        str(tmp_path),
        {"kind": "draw", "reason": "need chart"},
    ))
    assert main_data["ok"] is False
    assert "helper" in main_data["error"]

    token = set_current_helper_proc_id("proc123")
    try:
        data = json.loads(await registry._handle_request_resource(
            str(tmp_path),
            {
                "kind": "draw",
                "reason": "缺少图表 PNG",
                "needed_outputs": ["chart.png"],
                "resume_instruction": "嵌入图表后继续写最终版",
            },
        ))
    finally:
        reset_current_helper_proc_id(token)

    assert data["ok"] is False
    assert data["action"] == "request_resource"
    assert data["requires_main_resource"] is True
    assert data["matching_helper_kind"] == "draw"
    assert "resource" in data["resource_resolution_facts"]
    assert data["suggested_helper_kind"] == "draw"
    assert data["needed_outputs"] == ["chart.png"]

    kind_token = set_current_helper_kind("edit")
    try:
        data_without_proc = json.loads(await registry._handle_request_resource(
            str(tmp_path),
            {"kind": "draw", "reason": "缺少图表 PNG"},
        ))
    finally:
        reset_current_helper_kind(kind_token)
    assert data_without_proc["requires_main_resource"] is True


def test_ocr_cache_lock_uses_shared_project_cache_dir():
    from app.llm.tools import ocr_bridge

    cache_key = "abc123:accurate/table"
    with ocr_bridge._ocr_cache_file_lock(cache_key):
        expected = ocr_bridge._OCR_CACHE_LOCK_DIR / "abc123_accurate_table.lock"
        assert expected.exists()
        assert ocr_bridge._PROJECT_ROOT == Path(__file__).parent.parent
        assert ocr_bridge._OCR_CACHE_DIR == ocr_bridge._PROJECT_ROOT / "output" / "ocr_cache"


def test_ocr_cache_higher_tier_satisfies_lower_request(monkeypatch, tmp_path):
    from app.llm.tools import ocr_bridge

    image = tmp_path / "sample.jpg"
    image.write_bytes(b"sparse mineru complement image bytes")
    file_hash = ocr_bridge._file_sha256(image)
    accurate = ocr_bridge.OcrResult(
        ok=True,
        text="accurate cached text",
        score=0.99,
        engine="mineru",
        tier="accurate",
        engine_config={"source": "test"},
    )
    ocr_bridge._store_cached_tier(file_hash, "accurate", accurate, image)

    def fail_run_tier(*args, **kwargs):
        raise AssertionError("higher-tier cache should avoid running OCR")

    monkeypatch.setattr(ocr_bridge, "_run_tier", fail_run_tier)

    result = ocr_bridge.ocr_file_tiered(image, tier="fast", allow_upgrade=False, max_tier="accurate")

    assert result.ok is True
    assert result.text == "accurate cached text"
    assert result.tier == "accurate"
    assert result.engine_config["cache_hit"] is True
    assert result.engine_config["cache_tier"] == "accurate"
    assert result.engine_config["cache_satisfies_requested_tier"] == "fast"


def test_ocr_sparse_mineru_result_uses_legacy_complement(monkeypatch, tmp_path):
    from app.llm.tools import ocr_bridge

    image = tmp_path / "sample.jpg"
    image.write_bytes(b"sparse mineru complement image bytes")

    calls = []

    def fake_run_tier(path, cfg, timeout):
        calls.append(cfg.tier)
        return ocr_bridge.OcrResult(
            ok=True,
            text="## E2E OCR CARD 9137",
            score=1.0,
            engine="mineru",
            tier=cfg.tier,
        )

    def fake_legacy(path, timeout=120):
        return ocr_bridge.OcrResult(
            ok=True,
            text="E2E OCR CARD 9137\nALPHA 42.5\nBETA 18.0\nGAMMA 39.5\nRESULT PASS",
            score=0.9,
            engine="umi",
        )

    monkeypatch.setattr(ocr_bridge, "_run_tier", fake_run_tier)
    monkeypatch.setattr(ocr_bridge, "_ocr_with_legacy_file", fake_legacy)
    monkeypatch.setattr(ocr_bridge, "_load_cached_at_or_above", lambda *args, **kwargs: None)
    monkeypatch.setattr(ocr_bridge, "_store_cached_tier", lambda *args, **kwargs: None)

    result = ocr_bridge.ocr_file_tiered(image, tier="fast", allow_upgrade=True, max_tier="accurate")

    assert calls == ["fast", "balanced", "accurate"]
    assert "ALPHA 42.5" in result.text
    assert result.engine_config["legacy_complement_engine"] == "umi"
    assert any(f.get("signal") == "ocr_legacy_complement_used" for f in result.quality_flags)


def test_ocr_cache_rejects_legacy_engine_for_mineru_tier(monkeypatch, tmp_path):
    from app.llm.tools import ocr_bridge

    image = tmp_path / "sample.jpg"
    image.write_bytes(b"fake image bytes")
    file_hash = ocr_bridge._file_sha256(image)
    stale_accurate = ocr_bridge.OcrResult(
        ok=True,
        text="stale umi text",
        score=0.9,
        engine="umi",
        tier="accurate",
    )
    cache_path = ocr_bridge._ocr_cache_path(file_hash, "accurate")
    cache_path.parent.mkdir(parents=True, exist_ok=True)
    cache_path.write_text(
        json.dumps({
            "cache_version": ocr_bridge._OCR_CACHE_VERSION,
            "created_at": 1,
            "source": {"sha256": file_hash},
            "result": ocr_bridge._ocr_result_to_dict(stale_accurate),
        }),
        encoding="utf-8",
    )

    calls = []

    def fake_run_tier(path, cfg, timeout):
        calls.append(cfg.tier)
        return ocr_bridge.OcrResult(
            ok=True,
            text="fresh mineru text",
            score=1.0,
            engine="mineru",
            tier=cfg.tier,
        )

    monkeypatch.setattr(ocr_bridge, "_run_tier", fake_run_tier)

    result = ocr_bridge.ocr_file_tiered(image, tier="accurate", allow_upgrade=False, max_tier="accurate")

    assert calls == ["accurate"]
    assert result.ok is True
    assert result.text == "fresh mineru text"
    assert result.engine == "mineru"
    assert result.engine_config.get("cache_hit") is not True


async def test_gpu_resource_total_gate_serializes_when_capacity_one():
    from app.llm.tools.gpu_resources import GpuResourceLimiter

    limiter = GpuResourceLimiter(total=1, ocr=2, tts=2, mineru=1, umiocr=1, budget_mb=8000)
    active = 0
    peak = 0

    async def worker(kind: str):
        nonlocal active, peak
        async with limiter.async_scope(kind):
            active += 1
            peak = max(peak, active)
            await asyncio.sleep(0.02)
            active -= 1

    await asyncio.gather(worker("ocr"), worker("tts"))

    assert peak == 1


async def test_gpu_resource_allows_different_kinds_when_total_allows():
    from app.llm.tools.gpu_resources import GpuResourceLimiter

    limiter = GpuResourceLimiter(total=2, ocr=1, tts=1, mineru=1, umiocr=1, budget_mb=8000)
    active = 0
    peak = 0
    both_entered = asyncio.Event()
    release = asyncio.Event()

    async def worker(kind: str):
        nonlocal active, peak
        async with limiter.async_scope(kind):
            active += 1
            peak = max(peak, active)
            if active == 2:
                both_entered.set()
            await release.wait()
            active -= 1

    tasks = [asyncio.create_task(worker("ocr")), asyncio.create_task(worker("tts"))]
    try:
        await asyncio.wait_for(both_entered.wait(), timeout=1.0)
    finally:
        release.set()
        await asyncio.gather(*tasks)

    assert peak == 2


async def test_gpu_resource_mineru_stays_single_task_by_default():
    from app.llm.tools.gpu_resources import GpuResourceLimiter

    limiter = GpuResourceLimiter(total=2, ocr=2, tts=2, mineru=1, umiocr=1, budget_mb=8000)
    active = 0
    peak = 0

    async def worker():
        nonlocal active, peak
        async with limiter.async_scope("mineru"):
            active += 1
            peak = max(peak, active)
            await asyncio.sleep(0.02)
            active -= 1

    await asyncio.gather(worker(), worker())

    assert peak == 1


async def test_gpu_resource_reentrant_across_to_thread():
    from app.llm.tools.gpu_resources import GpuResourceLimiter

    limiter = GpuResourceLimiter(total=1, ocr=1, tts=1, mineru=1, umiocr=1, budget_mb=8000)

    def nested_sync() -> str:
        with limiter.scope("ocr"):
            return "ok"

    async with limiter.async_scope("ocr"):
        result = await asyncio.to_thread(nested_sync)

    assert result == "ok"


async def test_gpu_resource_budget_serializes_when_weights_exceed_budget():
    from app.llm.tools.gpu_resources import GpuResourceLimiter

    limiter = GpuResourceLimiter(
        total=3,
        ocr=3,
        tts=3,
        mineru=3,
        umiocr=3,
        budget_mb=5000,
        costs_mb={"mineru": 3000, "tts": 3000, "ocr": 0, "umiocr": 1000},
    )
    active = 0
    peak = 0

    async def worker(kind: str):
        nonlocal active, peak
        async with limiter.async_scope(kind):
            active += 1
            peak = max(peak, active)
            await asyncio.sleep(0.02)
            active -= 1

    await asyncio.gather(worker("mineru"), worker("tts"))

    assert peak == 1


async def test_gpu_resource_budget_allows_zero_cost_grouping_under_total_cap():
    from app.llm.tools.gpu_resources import GpuResourceLimiter

    limiter = GpuResourceLimiter(
        total=3,
        ocr=3,
        tts=1,
        mineru=1,
        umiocr=1,
        budget_mb=3000,
        costs_mb={"ocr": 0, "mineru": 3000, "tts": 3000, "umiocr": 1500},
    )
    active = 0
    peak = 0

    async def worker():
        nonlocal active, peak
        async with limiter.async_scope("ocr"):
            active += 1
            peak = max(peak, active)
            await asyncio.sleep(0.02)
            active -= 1

    await asyncio.gather(worker(), worker(), worker())

    assert peak == 3


def test_mineru_bg_api_url_recovers_default_port_when_state_files_missing(monkeypatch, tmp_path):
    from app.llm.tools import ocr_bridge

    port_file = tmp_path / ".mineru_api_port"
    config_file = tmp_path / ".mineru_api_config.json"
    monkeypatch.setattr(ocr_bridge, "_PORT_FILE", port_file)
    monkeypatch.setattr(ocr_bridge, "_CONFIG_FILE", config_file)
    monkeypatch.setattr(ocr_bridge, "_DEFAULT_MINERU_API_PORT", 51111)
    monkeypatch.setattr(ocr_bridge, "_mineru_http_alive", lambda port, timeout=2: port == 51111)

    cfg = ocr_bridge._default_mineru_config(
        None,
        backend=ocr_bridge._mineru_backend_default(),
        formula=True,
        table=True,
        image_analysis=False,
    )

    assert ocr_bridge._mineru_bg_api_url(cfg) == "http://127.0.0.1:51111"
    assert port_file.read_text(encoding="ascii") == "51111"
    assert json.loads(config_file.read_text(encoding="utf-8")) == ocr_bridge._mineru_service_config(cfg)


def test_mineru_service_config_uses_settings_concurrency(monkeypatch):
    from app.llm.tools import ocr_bridge

    monkeypatch.delenv("MINERU_API_MAX_CONCURRENT_REQUESTS", raising=False)
    monkeypatch.setattr(ocr_bridge.settings, "mineru_concurrency", 2)
    cfg = ocr_bridge._default_mineru_config(
        None,
        backend=ocr_bridge._mineru_backend_default(),
        formula=True,
        table=True,
        image_analysis=False,
    )

    assert ocr_bridge._mineru_service_config(cfg)["max_concurrent"] == "2"
    assert ocr_bridge._mineru_env(cfg)["MINERU_API_MAX_CONCURRENT_REQUESTS"] == "2"


def test_mineru_file_lock_only_required_for_single_concurrency(monkeypatch):
    from app.llm.tools import ocr_bridge

    monkeypatch.setattr(ocr_bridge.settings, "mineru_concurrency", 1)
    assert ocr_bridge._mineru_exclusive_process_lock_needed() is True

    monkeypatch.setattr(ocr_bridge.settings, "mineru_concurrency", 2)
    assert ocr_bridge._mineru_exclusive_process_lock_needed() is False


async def test_dispatch_requires_resource_helpers_for_main_thread_ocr_tts(tmp_path):
    from app.llm.tools import registry

    for tool_name, helper_kind in (("ocr", "read"), ("tts", "tts")):
        result = await registry.dispatch(
            tool_name,
            {"image_path": "x.png", "text": "hello"},
            archive_id="archive",
            group_id="group",
            user_id="user",
            workspace_dir=str(tmp_path),
            caller="main",
        )
        data = json.loads(result)
        assert data["ok"] is False
        assert data["error"] == "main_thread_resource_must_delegate"
        assert data["suggested_helper_kind"] == helper_kind


async def test_tts_handler_resolves_persona_voice_instruct_from_archive(monkeypatch, tmp_path):
    from app.llm.tools import registry
    from app.llm.tools.tts_bridge import TtsResult

    captured = {}

    monkeypatch.setattr("app.llm.tools.tts_bridge.is_available", lambda: True)

    async def fake_run_gpu_tts(func, text, **kwargs):
        captured["func"] = func.__name__
        captured["text"] = text
        captured["kwargs"] = kwargs
        output = kwargs["output"]
        Path(output).write_bytes(b"RIFFxxxxWAVE")
        return TtsResult(ok=True, paths=[output], durations=[0.5])

    async def fake_get_persona(_archive_id):
        from app.memory import persona_files
        pf = persona_files.load_persona("助手")
        assert pf is not None
        return pf.content

    assert not hasattr(registry, "tts_persona_guard")
    monkeypatch.setattr(registry, "run_gpu_tts", fake_run_gpu_tts)
    monkeypatch.setattr("app.memory.archive.get_persona", fake_get_persona)

    token = registry.set_current_voice_instruct("")
    ref_token = registry.set_current_voice_ref_audio("")
    try:
        raw = await registry._handle_tts(
            str(tmp_path),
            {
                "text": "你好",
                "language": "Chinese",
                "output_filename": "hello.wav",
                "instruct": "female, child, very high pitch",
                "mode": "voice_clone",
                "ref_audio": "fake.wav",
                "ref_text": "fake",
            },
            archive_id="arch_test",
        )
    finally:
        registry.reset_current_voice_instruct(token)
        registry.reset_current_voice_ref_audio(ref_token)

    data = json.loads(raw)
    assert data["ok"] is True
    assert captured["func"] == "tts_design"
    assert captured["kwargs"]["instruct"] == "male, young adult, moderate pitch"
    assert captured["kwargs"]["cwd"] == str(tmp_path)
    assert captured["kwargs"]["output"] == str(tmp_path / "hello.wav")
    assert data["paths"] == ["hello.wav"]
    assert data["voice_reply_file_candidate"] == "hello.wav"
    assert data["deliverable_candidate"] == "hello.wav"
    assert "delivery_guidance" in data
    assert "mode" not in data
    assert "voice_locked" not in data
    assert "voice_clone_ref" not in data
    assert "push_note" not in data
    assert "push_supported" not in data


async def test_tts_handler_resolves_short_identity_persona_voice(monkeypatch, tmp_path):
    from app.llm.tools import registry
    from app.llm.tools.tts_bridge import TtsResult

    captured = {}

    monkeypatch.setattr("app.llm.tools.tts_bridge.is_available", lambda: True)

    async def fake_run_gpu_tts(func, text, **kwargs):
        captured["func"] = func.__name__
        captured["kwargs"] = kwargs
        output = kwargs["output"]
        Path(output).write_bytes(b"RIFFxxxxWAVE")
        return TtsResult(ok=True, paths=[output], durations=[0.5])

    async def fake_get_persona(_archive_id):
        return (
            "\u4f60\u662f\u4e00\u4e2a\u52a9\u624b\u3002\n\n"
            "## \u6838\u5fc3\u539f\u5219\n"
            "- \u4e25\u683c\u9075\u5faa\u7528\u6237\u7684\u6bcf\u4e00\u6761\u547d\u4ee4\u3002"
        )

    assert not hasattr(registry, "tts_persona_guard")
    monkeypatch.setattr(registry, "run_gpu_tts", fake_run_gpu_tts)
    monkeypatch.setattr("app.memory.archive.get_persona", fake_get_persona)

    token = registry.set_current_voice_instruct("")
    ref_token = registry.set_current_voice_ref_audio("")
    try:
        raw = await registry._handle_tts(
            str(tmp_path),
            {"text": "hello", "language": "Chinese", "output_filename": "hello.wav"},
            archive_id="arch_test",
        )
    finally:
        registry.reset_current_voice_instruct(token)
        registry.reset_current_voice_ref_audio(ref_token)

    data = json.loads(raw)
    assert data["ok"] is True
    assert captured["func"] == "tts_design"
    assert captured["kwargs"]["instruct"] == "male, young adult, moderate pitch"


async def test_tts_push_flag_returns_candidate_facts_without_user_facing_limitation(monkeypatch, tmp_path):
    from app.llm.tools import registry
    from app.llm.tools.tts_bridge import TtsResult

    monkeypatch.setattr("app.llm.tools.tts_bridge.is_available", lambda: True)

    async def fake_run_gpu_tts(func, text, **kwargs):
        Path(kwargs["output"]).write_bytes(b"RIFFxxxxWAVE")
        return TtsResult(ok=True, paths=[kwargs["output"]], durations=[0.5])

    monkeypatch.setattr(registry, "run_gpu_tts", fake_run_gpu_tts)

    token = registry.set_current_voice_instruct("male, young adult, moderate pitch")
    ref_token = registry.set_current_voice_ref_audio("")
    try:
        raw = await registry._handle_tts(
            str(tmp_path),
            {
                "text": "你好",
                "language": "Chinese",
                "output_filename": "reply.wav",
                "push": True,
            },
            archive_id="",
        )
    finally:
        registry.reset_current_voice_instruct(token)
        registry.reset_current_voice_ref_audio(ref_token)

    data = json.loads(raw)
    dumped = json.dumps(data, ensure_ascii=False)
    assert data["ok"] is True
    assert data["paths"] == ["reply.wav"]
    assert data["voice_reply_file_candidate"] == "reply.wav"
    assert data["deliverable_candidate"] == "reply.wav"
    assert data["push_ignored"] is True
    assert "push_note" not in data
    assert "push_supported" not in data
    assert "不会自动发出" not in dumped
    assert "不能发" not in dumped


def test_tts_design_requires_system_voice_profile():
    from app.llm.tools.tts_bridge import tts_design

    result = tts_design("hello", instruct="")
    assert result.ok is False
    assert "voice profile" in (result.error or "")


def test_helper_tool_call_bloat_thresholds_stay_bounded():
    from pathlib import Path
    import re

    src = Path("app/llm/client.py").read_text(encoding="utf-8")
    warn = re.search(r"helper_tool_arg_bloat_warn_at\s*=\s*([0-9_]+)", src)
    close = re.search(r"helper_tool_arg_bloat_close_at\s*=\s*([0-9_]+)", src)

    assert warn and close
    assert int(warn.group(1).replace("_", "")) == 14_000
    assert int(close.group(1).replace("_", "")) == 24_000
    assert "helper_tool_call_bloat" in src


def test_helper_tool_call_bloat_allows_larger_write_like_payloads():
    from app.llm.client import _helper_tool_arg_bloat_thresholds

    warn, close, kind = _helper_tool_arg_bloat_thresholds(
        "workspace",
        '{"action":"write","path":"report.md","content":"' + ("x" * 100),
    )

    assert warn == 14_000
    assert close == 48_000
    assert kind == "write_like"


def test_helper_tool_call_bloat_closes_source_writes_earlier():
    from app.llm.client import _helper_tool_arg_bloat_thresholds

    warn, close, kind = _helper_tool_arg_bloat_thresholds(
        "workspace",
        '{"action":"write","path":"_env/src/rb_tree.c","content":"' + ("x" * 100),
    )

    assert warn == 10_000
    assert close == 18_000
    assert kind == "source_write"


def test_helper_tool_call_bloat_treats_docx_builder_script_as_write_like():
    from app.llm.client import _helper_tool_arg_bloat_thresholds

    args = (
        '{"action":"write","path":"build_paper.py","content":"'
        'from docx import Document\\n'
        'doc = Document()\\n'
        'doc.add_paragraph(\\"body\\")\\n'
        'doc.save(\\"db_index_paper.docx\\")'
    )

    warn, close, kind = _helper_tool_arg_bloat_thresholds("workspace", args)

    assert warn == 14_000
    assert close == 48_000
    assert kind == "write_like"


def test_helper_tool_call_bloat_keeps_general_payloads_bounded():
    from app.llm.client import _helper_tool_arg_bloat_thresholds

    warn, close, kind = _helper_tool_arg_bloat_thresholds(
        "workspace",
        '{"action":"locate","pattern":"' + ("x" * 100),
    )

    assert warn == 14_000
    assert close == 24_000
    assert kind == "general"


def test_preflight_guard_preserves_full_helper_envelope_fields():
    from pathlib import Path

    src = Path("app/llm/tools/delegate_actions.py").read_text(encoding="utf-8")
    assert '"input_files": s.get("input_files", [])' in src
    assert '"acceptance_checks": s.get("acceptance_checks", [])' in src
    assert "input_files=spec.get(\"input_files\") or []" in src


async def test_preflight_kind_retry_preserves_input_files_and_acceptance(monkeypatch):
    from app.llm.tools import delegate
    from app.llm.tools import delegate_actions

    async def fake_guard(*_args, **_kwargs):
        return True, "", [], [], {}

    monkeypatch.setattr(delegate, "_persona_consent_guard", fake_guard)

    helper_specs = [
        {
            "task_id": "paper_framework",
            "kind": "code",
            "mode": "easy",
            "prompt": "Create a document-style paper outline and save it as framework.md.",
            "framework": "shared contract",
            "input_files": ["_env/project_inventory.md"],
            "expected_outputs": ["_helpers_shared/framework.md"],
            "acceptance_checks": ["outline has sections", "output matrix exists"],
        }
    ]

    payload = await delegate_actions._run_delegate_preflight_guard(
        {"_guard_task_specs": helper_specs},
        helper_specs,
        "trace_preflight_preserve",
    )

    # P133: prose-vs-code is no longer a deterministic mismatch — guard LLM and
    # main thread judge based on context. A code/easy markdown framework task
    # passes preflight cleanly. Pin: physical hard constraints still apply
    # (covered separately in office/exec output tests).
    assert payload is None


def test_guard_attention_reports_existing_workspace_candidate_input(tmp_path):
    from app.llm.tools import delegate_actions

    candidate = tmp_path / "_env" / ".blocked_creates" / "triage_report.txt"
    candidate.parent.mkdir(parents=True)
    candidate.write_text("draft report\n", encoding="utf-8")

    helper_specs = [{
        "task_id": "triage-report",
        "kind": "edit",
        "mode": "easy",
        "prompt": "Revise the preserved candidate report if needed.",
        "input_files": ["_env/.blocked_creates/triage_report.txt"],
        "expected_outputs": ["_env/triage_report.txt"],
    }]

    attached = delegate_actions._attach_guard_attention_facts(
        helper_specs,
        main_workspace=str(tmp_path),
    )
    facts = attached[0].get("guard_observations") or []
    availability = [
        fact for fact in facts
        if fact.get("source") == "workspace_input_file_availability"
    ]

    assert availability
    fact = availability[0]
    assert fact["issue"] == "explicit_input_files_exist_in_main_workspace"
    assert fact["workspace_input_count"] == 1
    assert fact["workspace_input_files"][0]["path"] == "_env/.blocked_creates/triage_report.txt"
    assert fact["workspace_input_files"][0]["size_bytes"] == candidate.stat().st_size


def test_external_voice_synthesis_fact_reaches_preflight_guard_attention():
    from app.llm.tools import delegate_actions

    helper_specs = [{
        "task_id": "catgirl_voice_reply",
        "kind": "code",
        "mode": "easy",
        "prompt": (
            "Generate a Chinese TTS voice reply as voice_reply_catgirl.wav. "
            "Use Python TTS libraries such as gTTS/edge-tts/pyttsx3; if one fails, try another."
        ),
        "expected_outputs": ["voice_reply_catgirl.wav"],
        "acceptance_checks": ["real spoken TTS voice, not mathematical noise"],
    }]

    attached = delegate_actions._attach_guard_attention_facts(helper_specs)
    observations = attached[0].get("guard_observations") or []
    kind_facts = [
        fact for fact in observations
        if fact.get("source") == "deterministic_kind_check"
    ]

    assert kind_facts
    assert kind_facts[0]["observed_helper_kind_name"] == "tts"
    assert "system-managed TTS route" in kind_facts[0]["reason"]
    assert "Non-speech audio" in kind_facts[0]["reason"]


async def test_tts_handler_fails_without_system_voice_profile(monkeypatch, tmp_path):
    from app.llm.tools import registry

    monkeypatch.setattr("app.llm.tools.tts_bridge.is_available", lambda: True)

    async def fake_get_persona(_archive_id):
        return "无匹配人设正文"

    monkeypatch.setattr("app.memory.archive.get_persona", fake_get_persona)

    token = registry.set_current_voice_instruct("")
    ref_token = registry.set_current_voice_ref_audio("")
    try:
        raw = await registry._handle_tts(
            str(tmp_path),
            {"text": "你好", "language": "Chinese"},
            archive_id="arch_test",
        )
    finally:
        registry.reset_current_voice_instruct(token)
        registry.reset_current_voice_ref_audio(ref_token)

    data = json.loads(raw)
    assert data["ok"] is False
    assert "voice profile" in data["error"]


def test_windows_python_c_translation_writes_temp_script(tmp_path):
    from app.llm.tools.workspace import _translate_windows_command

    translated = _translate_windows_command('python -c "import docx; print(docx.__version__)"', str(tmp_path))

    assert translated.startswith("python _py_cmd_")
    assert translated.endswith(".py")
    script_name = translated.split(" ", 1)[1]
    assert (tmp_path / script_name).read_text(encoding="utf-8") == "import docx; print(docx.__version__)"


def test_windows_python_cmd_shim_c_translation_writes_temp_script(tmp_path):
    from app.llm.tools.workspace import _translate_windows_command

    translated = _translate_windows_command('python3.cmd -c "\\nprint(123)\\n"', str(tmp_path))

    assert translated.startswith("python3.cmd _py_cmd_")
    assert translated.endswith(".py")
    script_name = translated.split(" ", 1)[1]
    assert (tmp_path / script_name).read_text(encoding="utf-8") == "\nprint(123)\n"


def test_windows_python_runner_forwarding_c_translation_writes_temp_script(tmp_path):
    from app.llm.tools.workspace import _translate_windows_command

    translated = _translate_windows_command('python3 .\\python3.cmd -c "\\nprint(123)\\n"', str(tmp_path))

    assert translated.startswith("python3.cmd _py_cmd_")
    assert translated.endswith(".py")
    script_name = translated.split(" ", 1)[1]
    assert (tmp_path / script_name).read_text(encoding="utf-8") == "\nprint(123)\n"


def test_edit_file_old_str_not_found_points_to_local_fragment_instead_of_full_reread(tmp_path):
    from app.llm.tools.workspace import handle_edit_file

    path = tmp_path / "sample.py"
    path.write_text("def demo():\n    return 1\n", encoding="utf-8")

    result = asyncio.run(handle_edit_file(str(tmp_path), "sample.py", "return 2", "return 3"))

    assert result["ok"] is False
    assert "old_str not found in sample.py" in result["error"]
    assert "read_file(start_line=..., end_line=...)" in result["error"]
    assert "whole-file reread is usually only useful" in result["next_action_instruction"]
    assert "read a relevant local fragment" in result["recovery_facts"]["available_recovery_shapes"]


def test_multi_edit_old_str_not_found_points_to_local_fragment(tmp_path):
    from app.llm.tools.workspace import handle_multi_edit

    path = tmp_path / "sample.py"
    path.write_text("def demo():\n    return 1\n", encoding="utf-8")

    result = asyncio.run(handle_multi_edit(str(tmp_path), "sample.py", [{
        "old_str": "return 2",
        "new_str": "return 3",
        "expected_count": 1,
    }]))

    assert result["ok"] is False
    assert "read only the local fragment" in result["failures"][0]


def test_windows_python_c_translation_handles_nested_quotes(tmp_path):
    from app.llm.tools.workspace import _translate_windows_command

    command = 'python -c "from PIL import Image; print(f\'size={Image}\')"'
    translated = _translate_windows_command(command, str(tmp_path))

    script_name = translated.split(" ", 1)[1]
    assert "print(f'size={Image}')" in (tmp_path / script_name).read_text(encoding="utf-8")


def test_windows_python_c_translation_decodes_escaped_newlines(tmp_path):
    from app.llm.tools.workspace import _translate_windows_command

    command = 'python -c "\\nprint(123)\\nprint(456)\\n"'
    translated = _translate_windows_command(command, str(tmp_path))

    script_name = translated.split(" ", 1)[1]
    assert (tmp_path / script_name).read_text(encoding="utf-8") == "\nprint(123)\nprint(456)\n"


def test_windows_python_c_translation_preserves_string_literal_newline_escapes(tmp_path):
    from app.llm.tools.workspace import _translate_windows_command

    command = (
        'python -c "\\n'
        "name = 'users'\\n"
        'print(f\\\'\\n--- {name} columns ---\\\')\\n'
        '"'
    )
    translated = _translate_windows_command(command, str(tmp_path))

    script_name = translated.split(" ", 1)[1]
    script = (tmp_path / script_name).read_text(encoding="utf-8")
    assert script.startswith("\nname = 'users'\n")
    assert "print(f'\\n--- {name} columns ---')" in script
    assert "print(f'\n---" not in script


def test_windows_python_c_translation_unescapes_shell_quoted_fstring(tmp_path):
    from app.llm.tools.workspace import _translate_windows_command

    command = 'python -c "print(f\'size={sorted(df[\\\"distribution\\\"].unique())}\')"'
    translated = _translate_windows_command(command, str(tmp_path))

    script_name = translated.split(" ", 1)[1]
    script = (tmp_path / script_name).read_text(encoding="utf-8")
    assert 'print(f\'size={sorted(df["distribution"].unique())}\')' in script


def test_windows_cmd_wrapped_python_c_translation_unescapes_fstring(tmp_path):
    from app.llm.tools.workspace import _translate_windows_command

    command = 'cmd /c python -c "print(f\'size={sorted(df[\\\"size\\\"].unique())}\')"'
    translated = _translate_windows_command(command, str(tmp_path))

    assert translated.startswith("cmd /c python _py_cmd_")
    script_name = translated.split("python ", 1)[1]
    script = (tmp_path / script_name).read_text(encoding="utf-8")
    assert 'print(f\'size={sorted(df["size"].unique())}\')' in script


def test_windows_set_env_python_c_translation_writes_temp_script(tmp_path):
    from app.llm.tools.workspace import _translate_windows_command

    command = 'set PYTHONPATH=_env/src && python -c "print(f\'ok={1}\')"'
    translated = _translate_windows_command(command, str(tmp_path))

    assert translated.startswith('set PYTHONPATH=_env/src&& python _py_cmd_')
    script_name = translated.split("python ", 1)[1]
    script = (tmp_path / script_name).read_text(encoding="utf-8")
    assert "print(f'ok={1}')" in script


def test_windows_set_env_translation_quotes_non_python_c_command(tmp_path):
    from app.llm.tools.workspace import _translate_windows_command

    command = "set PYTHONPATH=_env/src && python -m pytest tests -q"
    translated = _translate_windows_command(command, str(tmp_path))

    assert translated == 'set PYTHONPATH=_env/src&& python -m pytest tests -q'


def test_windows_python_c_translation_preserves_interpreter_flags(tmp_path):
    from app.llm.tools.workspace import _translate_windows_command

    command = 'python -B -c "import sys; sys.path.insert(0,\'src\'); exec(open(\'smoke_test.py\').read())"'
    translated = _translate_windows_command(command, str(tmp_path))

    assert translated.startswith("python -B _py_cmd_")
    script_name = translated.split("python -B ", 1)[1]
    script = (tmp_path / script_name).read_text(encoding="utf-8")
    assert "sys.path.insert(0,'src')" in script
    assert "exec(open('smoke_test.py').read())" in script


def test_windows_cd_and_python_c_translation_writes_temp_script(tmp_path):
    from app.llm.tools.workspace import _translate_windows_command

    (tmp_path / "_env").mkdir()
    command = 'cd _env && python -c "from snake_arcade.render import render_text; print(\'render_text OK\')"'
    translated = _translate_windows_command(command, str(tmp_path))

    assert translated.startswith("cd /d _env && python _py_cmd_")
    script_name = translated.split("python ", 1)[1]
    script = (tmp_path / "_env" / script_name).read_text(encoding="utf-8")
    assert "from snake_arcade.render import render_text" in script
    assert "print('render_text OK')" in script


def test_windows_unix_env_prefix_python_c_translation_writes_temp_script(tmp_path):
    from app.llm.tools.workspace import _translate_windows_command

    command = 'PYTHONPATH=_env/src PYTHONUTF8=1 python -c "print(\\"ok\\")"'
    translated = _translate_windows_command(command, str(tmp_path))

    assert translated.startswith('set PYTHONPATH=_env/src&& set PYTHONUTF8=1&& python _py_cmd_')
    script_name = translated.split("python ", 1)[1]
    script = (tmp_path / script_name).read_text(encoding="utf-8")
    assert 'print("ok")' in script


def test_windows_simple_cp_translation_writes_workspace_copy_script(tmp_path):
    from app.llm.tools.workspace import _translate_windows_command

    (tmp_path / "source.txt").write_text("hello", encoding="utf-8")
    translated = _translate_windows_command("cp source.txt copied.txt", str(tmp_path))

    assert translated.startswith("python _py_cmd_")
    script_name = translated.split("python ", 1)[1]
    script = tmp_path / script_name
    assert "shutil.copy2" in script.read_text(encoding="utf-8")


def test_env_run_normalization_reuses_workspace_python_c_translation(tmp_path):
    from app.llm.tools.environment import _normalize_env_command

    command = 'python -c "print(f\'size={sorted(df[\\\"distribution\\\"].unique())}\')"'
    translated, script_name, script_path = _normalize_env_command(command, tmp_path)

    assert "python " in translated
    assert script_name.endswith(".py")
    assert script_path is not None
    script = script_path.read_text(encoding="utf-8")
    assert 'print(f\'size={sorted(df["distribution"].unique())}\')' in script


def test_env_run_keeps_complex_shell_python_c_commands_untranslated(tmp_path):
    from app.llm.tools.environment import _normalize_env_command

    command = (
        'python -c "import importlib.util; print(importlib.util.find_spec(\'PIL\'))" '
        '2>&1 | findstr /V "Traceback"'
    )
    translated, script_name, script_path = _normalize_env_command(command, tmp_path)

    assert translated == command
    assert script_name == ""
    assert script_path is None


def test_workspace_python_error_hints_are_actionable():
    from app.llm.tools.workspace import _diagnose_build_failure

    syntax_hint = _diagnose_build_failure(
        'cmd /c python -c "print(\"',
        'SyntaxError: unterminated string literal (detected at line 1)',
        '',
        returncode=1,
    )
    assert "workspace(action='write')" in syntax_hint
    assert "python script.py" in syntax_hint

    key_hint = _diagnose_build_failure(
        "python script.py",
        "KeyError: 'algo'",
        '',
        returncode=1,
    )
    assert "columns" in key_hint
    assert "unique values" in key_hint


def test_p35_framework_contract_future_output_names_are_not_dependencies():
    from app.llm.tools.delegate_runner import _is_framework_future_output_reference

    prompt = (
        "Create only a compact framework contract with an output matrix and file naming convention. "
        "Future helpers will produce ds_rb_tree.md, ds_skiplist.md, experiment_results.csv, "
        "and db_index_comparison_paper.docx. Do not read those files in this helper."
    )

    for filename in [
        "ds_rb_tree.md",
        "ds_skiplist.md",
        "experiment_results.csv",
        "db_index_comparison_paper.docx",
    ]:
        assert _is_framework_future_output_reference(
            task_id="framework_contract",
            kind="code",
            prompt=prompt,
            filename=filename,
            expected_outputs=["_env/framework_contract.md"],
        )


def test_p35_does_not_hide_explicit_existing_dependency_reads():
    from app.llm.tools.delegate_runner import _is_framework_future_output_reference

    prompt = (
        "Create a framework note after reading the file ds_rb_tree.md from previous helper output. "
        "Use existing evidence before writing the contract."
    )

    assert not _is_framework_future_output_reference(
        task_id="framework_contract",
        kind="code",
        prompt=prompt,
        filename="ds_rb_tree.md",
        expected_outputs=["_env/framework_contract.md"],
    )


def test_extract_declared_files_ignores_natural_language_json_values():
    from app.llm.tools.delegate import _extract_declared_files

    report = '''## 不可行
```json
{"files": ["建议: no, 理由: 依赖 CSV 不存在，本任务无法启动。应先派 bench helper 生成 bench_runner_assemble_bench_out.csv。"]}
```
'''

    assert _extract_declared_files(report) == set()


def test_extract_declared_files_accepts_plain_filenames_from_json():
    from app.llm.tools.delegate import _extract_declared_files

    report = '```json\n{"files": ["chart1.png", "notes.txt"]}\n```'

    assert _extract_declared_files(report) == {"chart1.png", "notes.txt"}


def test_workspace_copy_keeps_single_helper_prefixed_data_file(tmp_path):
    from app.llm.tools.workspace import copy_workspace_contents

    src = tmp_path / "main"
    dst = tmp_path / "helper"
    src.mkdir()
    (src / "bench_runner_assemble_bench_out.csv").write_text("algo,n,time_ms\n", encoding="utf-8")

    copied = copy_workspace_contents(str(src), str(dst))

    assert copied == 1
    assert (dst / "bench_runner_assemble_bench_out.csv").exists()


def test_workspace_copy_still_filters_repeated_helper_artifact(tmp_path):
    from app.llm.tools.workspace import copy_workspace_contents

    src = tmp_path / "main"
    dst = tmp_path / "helper"
    src.mkdir()
    (src / "gen_charts_gen_charts_chart1.png").write_bytes(b"png")

    copied = copy_workspace_contents(str(src), str(dst))

    assert copied == 0
    assert not (dst / "gen_charts_gen_charts_chart1.png").exists()


def test_workspace_copy_filters_helper_internal_reports_and_ledgers(tmp_path):
    from app.llm.tools.workspace import copy_workspace_contents

    src = tmp_path / "main"
    dst = tmp_path / "helper"
    src.mkdir()
    (src / ".helper_old_task_full_report.txt").write_text("old report", encoding="utf-8")
    (src / ".helper_completions.json").write_text("{}", encoding="utf-8")
    (src / "answer.txt").write_text("ok", encoding="utf-8")

    copied = copy_workspace_contents(str(src), str(dst))

    assert copied == 1
    assert (dst / "answer.txt").exists()
    assert not (dst / ".helper_old_task_full_report.txt").exists()
    assert not (dst / ".helper_completions.json").exists()


def test_workspace_copy_downloaded_media_is_opt_in_for_ocr_helpers(tmp_path):
    from app.llm.tools.workspace import copy_workspace_contents

    src = tmp_path / "main"
    normal_dst = tmp_path / "normal_helper"
    ocr_dst = tmp_path / "ocr_helper"
    media = src / "_downloaded_media"
    media.mkdir(parents=True)
    (media / "img_test.jpg").write_bytes(b"jpg")
    (src / "note.txt").write_text("ok", encoding="utf-8")

    normal_count = copy_workspace_contents(str(src), str(normal_dst))
    ocr_count = copy_workspace_contents(str(src), str(ocr_dst), include_downloaded_media=True)

    assert normal_count == 1
    assert not (normal_dst / "_downloaded_media" / "img_test.jpg").exists()
    assert ocr_count == 2
    assert (ocr_dst / "_downloaded_media" / "img_test.jpg").exists()


def test_workspace_copy_environment_files_is_opt_in(tmp_path):
    from app.llm.tools.workspace import copy_workspace_contents

    src = tmp_path / "main"
    normal_dst = tmp_path / "normal_helper"
    env_dst = tmp_path / "env_helper"
    env_dir = src / "_env"
    env_dir.mkdir(parents=True)
    (env_dir / "snake.js").write_text("const score = 0;\n", encoding="utf-8")
    (src / "note.txt").write_text("ok", encoding="utf-8")

    normal_count = copy_workspace_contents(str(src), str(normal_dst))
    env_count = copy_workspace_contents(str(src), str(env_dst), include_environment_files=True)

    assert normal_count == 1
    assert not (normal_dst / "_env" / "snake.js").exists()
    assert env_count == 2
    assert (env_dst / "_env" / "snake.js").read_text(encoding="utf-8") == "const score = 0;\n"


def test_delegate_copyback_merges_modified_environment_files_only(tmp_path):
    from app.llm.tools.delegate_copyback import _copy_results_to_main
    from app.llm.tools.workspace_utils import take_workspace_snapshot

    main_ws = tmp_path / "main"
    helper_ws = tmp_path / "helper"
    helper_env = helper_ws / "_env"
    helper_env.mkdir(parents=True)
    main_ws.mkdir()
    (helper_env / "index.html").write_text("<canvas></canvas>\n", encoding="utf-8")
    (helper_env / "snake.js").write_text("const speed = 1;\n", encoding="utf-8")
    main_env = main_ws / "_env"
    main_env.mkdir()
    (main_env / "index.html").write_text("<canvas></canvas>\n", encoding="utf-8")
    (main_env / "snake.js").write_text("const speed = 1;\n", encoding="utf-8")
    snapshot = take_workspace_snapshot(str(helper_ws))

    (helper_env / "snake.js").write_text("const speed = 2;\n", encoding="utf-8")
    copied, stats, file_map = _copy_results_to_main(
        str(helper_ws),
        str(main_ws),
        "snake_js",
        fork_snapshot=snapshot,
    )

    assert copied == ["_env/snake.js"]
    assert file_map == [{
        "helper_name": "snake.js",
        "main_name": "_env/snake.js",
        "shared_name": None,
    }]
    assert stats["env_copied_count"] == 1
    assert stats["env_skipped_unchanged"] == 1
    assert stats["env_skipped_unchanged_files"] == ["_env/index.html"]
    assert (main_ws / "_env" / "snake.js").read_text(encoding="utf-8") == "const speed = 2;\n"
    assert (main_ws / "_env" / "index.html").read_text(encoding="utf-8") == "<canvas></canvas>\n"


def test_fetch_to_temp_uses_displayed_name_remap_alias(tmp_path):
    from app.llm.tools.workspace import fetch_to_temp, update_displayed_name_remap

    main_ws = tmp_path / "main"
    helper_ws = main_ws / ".temp" / "_delegate_user_chart"
    main_ws.mkdir()
    helper_ws.mkdir(parents=True)
    (main_ws / "bench_runner_assemble_bench_out.csv").write_text("algo,n,time_ms\n", encoding="utf-8")
    update_displayed_name_remap(
        str(main_ws),
        {"bench_runner_assemble_bench_out.csv": "bench_out.csv"},
    )

    copied, skipped = fetch_to_temp(
        str(main_ws),
        str(helper_ws),
        paths=["bench_out.csv"],
        source="main",
    )

    assert copied == ["bench_out.csv"]
    assert skipped == []
    assert (helper_ws / "bench_out.csv").read_text(encoding="utf-8") == "algo,n,time_ms\n"


def test_fetch_to_temp_rejects_workspace_root_copy(tmp_path):
    from app.llm.tools.workspace import fetch_to_temp

    main_ws = tmp_path / "main"
    helper_ws = main_ws / "_delegate_user_helper"
    helper_ws.mkdir(parents=True)
    (main_ws / "data.txt").write_text("ok\n", encoding="utf-8")

    copied, skipped = fetch_to_temp(
        str(main_ws),
        str(helper_ws),
        paths=["."],
        source="main",
    )

    assert copied == []
    assert skipped == ["."]
    assert not (helper_ws / "_delegate_user_helper").exists()


async def test_helper_fetch_to_temp_copies_into_helper_workspace(tmp_path):
    from app.llm.tools.registry import _handle_fetch_to_temp

    main_ws = tmp_path / "main"
    temp_ws = main_ws / ".temp"
    helper_ws = temp_ws / "_delegate_user_chart"
    helper_ws.mkdir(parents=True)
    (main_ws / "bench_runner_assemble_bench_out.csv").write_text("algo,n,time_ms\n", encoding="utf-8")

    result = json.loads(await _handle_fetch_to_temp(
        str(helper_ws),
        {"source": "main", "paths": ["bench_runner_assemble_bench_out.csv"]},
    ))

    assert result["ok"] is True
    assert result["copied"] == ["bench_runner_assemble_bench_out.csv"]
    assert (helper_ws / "bench_runner_assemble_bench_out.csv").exists()
    assert not (temp_ws / "bench_runner_assemble_bench_out.csv").exists()


async def test_recall_thread_returns_persisted_helper_task_contract(tmp_path):
    from app.core.core_processes import (
        ThreadContext,
        reset_current_thread_context,
        set_current_thread_context,
    )
    from app.llm.tools.workspace_transfer_tools import handle_recall_thread

    helper_ws = tmp_path / "main" / ".temp" / "_delegate_user_writer"
    helper_ws.mkdir(parents=True)
    (helper_ws / ".helper_task_contract.json").write_text(
        json.dumps(
            {
                "task_id": "writer",
                "helper_kind": "edit",
                "helper_mode": "normal",
                "resume": False,
                "goal_excerpt": "Write the final report from verified evidence.",
                "expected_outputs": ["final.docx"],
                "write_scopes": ["final.docx"],
                "acceptance_checks": ["docx opens", "claims match evidence"],
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )

    token = set_current_thread_context(
        ThreadContext(
            user_message="完成最终报告",
            plan_intent="helper task: writer (edit/normal)",
            plan_key_points=["helper task contract snapshot is available in .helper_task_contract.json"],
            plan_deliverables=["final.docx"],
            role_label="helper:writer",
        )
    )
    try:
        result = json.loads(await handle_recall_thread(str(helper_ws), {}))
    finally:
        reset_current_thread_context(token)

    contract = result["helper_task_contract"]
    assert contract["task_id"] == "writer"
    assert contract["helper_kind"] == "edit"
    assert contract["expected_outputs"] == ["final.docx"]
    assert contract["acceptance_checks"] == ["docx opens", "claims match evidence"]
    assert "automatic success/failure decision" in contract["truth_scope"]
    assert result["plan"]["deliverables"] == ["final.docx"]


async def test_helper_fetch_to_temp_stages_underscore_project_path_in_environment(tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools.registry import _handle_fetch_to_temp

    project = tmp_path / "project"
    main_ws = tmp_path / "main"
    helper_ws = main_ws / ".temp" / "_delegate_user_read"
    source = project / "_extracted" / "group" / "小组报告.docx"
    source.parent.mkdir(parents=True)
    source.write_bytes(b"PK\x03\x04docx")
    helper_ws.mkdir(parents=True)

    env = EnvironmentContext(
        root_dir=str(project),
        archive_id="arch_test",
        group_id="env_user_u",
        user_id="u",
        project_key="p",
    )
    with runtime_context("environment", env):
        result = json.loads(await _handle_fetch_to_temp(
            str(helper_ws),
            {"source": "main", "paths": ["_extracted/group/小组报告.docx"]},
        ))

    assert result["ok"] is True
    assert result["copied"] == ["_env/_extracted/group/小组报告.docx"]
    assert result["env_copied"] == ["_env/_extracted/group/小组报告.docx"]
    assert result["skipped"] == []
    assert result["normalized_paths"] is None
    assert (helper_ws / "_env" / "_extracted" / "group" / "小组报告.docx").is_file()


async def test_python_tool_modulenotfound_points_to_workspace_run():
    from app.llm.tools import registry

    result = await registry.dispatch(
        "python",
        {"code": "import definitely_missing_package_for_tool_test"},
        archive_id="archive",
        group_id="group",
        user_id="user",
        workspace_dir="",
        caller="helper",
    )
    data = json.loads(result)
    assert data["ok"] is False
    assert "ModuleNotFoundError" in data["error"]
    assert data["recovery_action"] == "switch_tool"
    assert data["retry_same_tool"] is False
    assert "workspace(action='run')" in data["recommended_tools"]
    assert "workspace(action='run'" in data["fix_hint"]


async def test_python_tool_open_security_error_points_to_workspace_write():
    from app.llm.tools import registry

    result = await registry.dispatch(
        "python",
        {"code": "open('report.md', 'w').write('x')"},
        archive_id="archive",
        group_id="group",
        user_id="user",
        workspace_dir="",
        caller="helper",
    )
    data = json.loads(result)
    assert data["ok"] is False
    assert "SecurityError" in data["error"]
    assert data["recovery_action"] == "switch_tool"
    assert data["retry_same_tool"] is False
    assert "workspace(action='write'" in data["fix_hint"]
    assert "read_file" in data["fix_hint"]


@pytest.mark.parametrize("primitive", ["compile", "exec", "eval"])
async def test_python_tool_dynamic_execution_security_error_points_to_workspace_run(primitive):
    from app.llm.tools import registry

    result = await registry.dispatch(
        "python",
        {"code": f"{primitive}('print(1)')"},
        archive_id="archive",
        group_id="group",
        user_id="user",
        workspace_dir="",
        caller="helper",
    )
    data = json.loads(result)
    assert data["ok"] is False
    assert "SecurityError" in data["error"]
    assert data["recovery_action"] == "switch_tool"
    assert data["retry_same_tool"] is False
    assert "workspace(action='run')" in data["fix_hint"]



def test_helper_missing_file_response_points_to_fetch_to_temp():
    from app.llm.tools.workspace import _file_not_found_response

    result = _file_not_found_response(
        "C:/repo/workspace/.temp/_delegate_user_chart",
        "bench_runner_assemble_bench_out.csv",
    )

    assert result["ok"] is False
    assert "fetch_to_temp(source='main'" in result["error"]
    assert "只使用清单或结果明确暴露的路径" in result["error"]


def test_helper_missing_project_path_points_to_existing_env_copy(tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools.workspace import _file_not_found_response

    helper_ws = tmp_path / "main" / ".temp" / "_delegate_user_read"
    staged = helper_ws / "_env" / "_extracted" / "group" / "小组报告.docx"
    staged.parent.mkdir(parents=True)
    staged.write_bytes(b"PK\x03\x04docx")
    env = EnvironmentContext(
        root_dir=str(tmp_path / "project"),
        archive_id="arch_test",
        group_id="env_user_u",
        user_id="u",
        project_key="p",
    )

    with runtime_context("environment", env):
        result = _file_not_found_response(
            str(helper_ws),
            "_extracted/group/小组报告.docx",
        )

    assert "`_env/_extracted/group/小组报告.docx`" in result["error"]
    assert "请直接改用它" in result["error"]


def test_helper_missing_environment_file_points_to_main_fetch():
    from app.llm.tools.workspace import _file_not_found_response

    result = _file_not_found_response(
        "C:/repo/workspace/.temp/_delegate_user_env_tests",
        "_env/app/tests/test_environment.py",
    )

    assert result["ok"] is False
    assert "stages them under `_env/...`" in result["error"]
    assert "locate/search in the sandbox" in result["error"]
    assert "Stay with staged relative paths" in result["error"]
    assert "exact project_path" in result["error"]
    assert "fetch_to_temp(source='main'" in result["error"]
    assert "request_resource" in result["error"]


def test_all_helper_kinds_advertise_hard_as_same_kind_mode():
    from app.llm.tools.helper_kinds import HELPER_CONFIGS, VALID_HELPER_KINDS

    unsupported = [
        kind for kind in VALID_HELPER_KINDS
        if not HELPER_CONFIGS[kind].get("supports_hard_mode")
    ]

    assert unsupported == []


def test_hard_helper_prompt_sets_same_kind_standards_for_non_code_helpers():
    from app.llm.tools.delegate import _select_helper_system

    for kind, expected in [
        ("read", "For read hard mode"),
        ("edit", "For edit hard mode"),
        ("draw", "For draw hard mode"),
        ("tts", "For TTS hard mode"),
        ("verify", "For verify hard mode"),
        ("project_map", "For project-analysis hard mode"),
    ]:
        prompt = _select_helper_system(kind, "hard")
        assert "Hard mode is a richer same-kind workflow" in prompt
        assert expected in prompt
        assert "does not turn read into edit" in prompt
        assert "同类增强" in prompt


def test_helper_python_filenotfound_points_to_fetch_to_temp():
    from app.llm.tools.workspace import _diagnose_build_failure

    hint = _diagnose_build_failure(
        "python chart.py",
        "\nFileNotFoundError: [Errno 2] No such file or directory: 'bench_runner_assemble_bench_out.csv'",
        "",
        returncode=1,
        ws_dir="C:/repo/workspace/.temp/_delegate_user_chart",
    )

    assert "Locate the file" in hint
    assert "fetch_to_temp(source='main'" in hint
    assert "Use `_helpers_shared/...` only when that exact shared path is exposed" in hint


def test_helper_system_prompt_describes_environment_workspace_copies():
    from app.llm.tools.delegate import _select_helper_system

    prompt = _select_helper_system("code", "easy")

    assert "project/environment work" in prompt
    assert "sparse `_env/...` workspace copies" in prompt
    assert "helpers work through local staged copies" in prompt
    assert "keep absolute project paths for the main process" in prompt
    assert "缺依赖" in prompt
    assert "same-batch producer outputs" in prompt
    assert "guessed paths marked as unavailable" in prompt
    assert "`old_str` is the literal file text" in prompt
    assert "JSON logs escape quotes and backslashes" in prompt


def test_read_helper_prompt_requests_staged_project_files():
    from app.llm.tools.delegate import _select_helper_system

    prompt = _select_helper_system("read", "easy")

    assert "workspace-relative files" in prompt
    assert "staged `_env/...` copies" in prompt
    assert "request_resource" in prompt
    assert "allow_upgrade=true" in prompt


def test_ocr_missing_file_feedback_is_generic_and_actionable(tmp_path):
    import asyncio

    from app.llm.tools.registry import _handle_ocr

    raw = asyncio.run(_handle_ocr(str(tmp_path), {"image_path": "_env/image/missing.png"}))

    assert '"ok": false' in raw
    assert "env_fetch" in raw
    assert "workspace_relative" in raw
    assert "QQ" not in raw
    assert "NapCat" not in raw
    assert "桥接" not in raw


def test_ocr_rejects_plain_text_and_source_files_with_read_hints(tmp_path):
    import asyncio
    import json

    from app.llm.tools.registry import _handle_ocr

    (tmp_path / "module.py").write_text("print('hello')\n", encoding="utf-8")

    raw = asyncio.run(_handle_ocr(str(tmp_path), {"image_path": "module.py"}))
    data = json.loads(raw)

    assert data["ok"] is False
    assert data["wrong_tool"] == "ocr_for_plain_text"
    assert "read_file" in data["suggested_tools"]
    assert "code_index" in data["suggested_tools"]


def test_main_visible_prompts_do_not_recommend_final_kind():
    import re
    from app.core import context, orchestrator
    from app.llm import client as llm_client
    from app.llm.tools import registry
    from app.llm.tools.delegate import BUNDLED_SKILLS, SPAWN_HELPER_TOOL_SCHEMA
    obsolete_pairing_flag = "auto" + "_" + "final"

    visible_text = "\n".join([
        context.ROUND2_SYSTEM_TEMPLATE,
        BUNDLED_SKILLS["parallel-helpers-coordination"],
        SPAWN_HELPER_TOOL_SCHEMA["function"]["parameters"]["properties"]["kind"]["description"],
        registry.DELEGATE_TOOL_SCHEMA["function"]["description"],
        registry.DELEGATE_TOOL_SCHEMA["function"]["parameters"]["properties"]["tasks"]["items"]["properties"]["kind"]["description"],
    ])

    forbidden_recommendations = [
        r"升级\s*`?kind=['\"]?final",
        r"建议[^\n]{0,80}kind=['\"]?final",
        r"派[^\n]{0,80}kind=['\"]?final",
        r"重派[^\n]{0,80}kind=['\"]?final",
        r"用\s*`?kind=['\"]?final",
        r"直接\s*kind=['\"]?final",
        r"kind=['\"]?final[^\n]{0,40}重试",
        r"final twin",
        r"final helper",
    ]
    for pattern in forbidden_recommendations:
        assert not re.search(pattern, visible_text)
    assert obsolete_pairing_flag not in visible_text

    spawn_kind_enum = SPAWN_HELPER_TOOL_SCHEMA["function"]["parameters"]["properties"]["kind"]["enum"]
    delegate_kind_enum = registry.DELEGATE_TOOL_SCHEMA["function"]["parameters"]["properties"]["tasks"]["items"]["properties"]["kind"]["enum"]
    assert "final" not in spawn_kind_enum
    assert "final" not in delegate_kind_enum
    assert obsolete_pairing_flag not in registry.DELEGATE_TOOL_SCHEMA["function"]["parameters"]["properties"]

    source_text = "\n".join([
        orchestrator.__file__ and Path(orchestrator.__file__).read_text(encoding="utf-8"),
        llm_client.__file__ and Path(llm_client.__file__).read_text(encoding="utf-8"),
    ])
    for pattern in forbidden_recommendations:
        assert not re.search(pattern, source_text)


def test_document_round2_prompt_requires_edit_and_accepts_verified_existing_docx():
    from app.core import orchestrator

    prompts = orchestrator._build_round2_system_prompts(
        is_coding=False,
        is_document=True,
        needs_recall=False,
        parallelizable=False,
    )
    visible_text = "\n".join(p["content"] for p in prompts)

    assert "full Office documents belong to edit helpers" in visible_text
    assert "add it to deliverables instead of regenerating it" in visible_text
    assert "delegate `kind='edit'` for Office document production" in visible_text
    assert "let edit freeze with `request_resource`" in visible_text
    assert "producer self-verification report" in visible_text
    assert "consume the helper's structural facts" in visible_text
    assert "立即 delegate(kind:'edit' 纯文档 / 'general' 文档+图表)" not in visible_text
    assert "交给 edit/general helper" not in visible_text


async def test_delegate_guard_does_not_split_single_edit_office_output(monkeypatch):
    from app.llm import client as llm_client
    from app.llm.tools import delegate

    async def fake_chat_json(*args, **kwargs):
        return {
            "should_act": True,
            "reason": "五道题可拆",
            "split_recommendations": [{
                "task_id": "hw_docx",
                "should_split": True,
                "split_into": ["5.1", "5.2", "5.17", "5.18", "5.21"],
                "reason": "多题独立",
            }],
            "kind_recommendations": [],
        }

    monkeypatch.setattr(llm_client, "chat_json", fake_chat_json)

    should_act, reason, split_recs, kind_recs = await delegate._persona_consent_guard(
        persona="",
        user_message="把作业写成一个 Word 文档",
        tasks=[{
            "task_id": "hw_docx",
            "kind": "edit",
            "mode": "hard",
            "prompt": "请生成一份完整的 Word 文档，文件名：通信原理作业_第五章.docx。包含5.1、5.2、5.17、5.18、5.21。",
            "expected_outputs": ["通信原理作业_第五章.docx"],
        }],
    )

    assert should_act is True
    assert reason == "五道题可拆"
    assert split_recs == []
    assert kind_recs == []


async def test_delegate_guard_brief_keeps_long_prompt_head_and_tail(monkeypatch):
    from app.llm import client as llm_client
    from app.llm.tools import delegate

    captured = {}

    async def fake_chat_json(msgs, *args, **kwargs):
        captured["msgs"] = msgs
        return {"should_act": True, "reason": "allow"}

    monkeypatch.setattr(llm_client, "chat_json", fake_chat_json)

    prompt = (
        "HEAD_FACT: already read bounded source facts.\n"
        + "\n".join(f"middle detail {i}: supporting fact" for i in range(120))
        + "\nTAIL_OUTPUTS: create final_report.txt and reply_draft.txt.\n"
        + "TAIL_CHECKS: preserve all listed acceptance facts."
    )

    should_act, reason, split_recs, kind_recs = await delegate._persona_consent_guard(
        persona="",
        user_message="Create the final materials from the verified facts.",
        tasks=[{
            "task_id": "final_materials",
            "kind": "edit",
            "mode": "easy",
            "prompt": prompt,
            "expected_outputs": ["final_report.txt", "reply_draft.txt"],
        }],
    )

    assert should_act is True
    assert reason == "allow"
    assert split_recs == []
    assert kind_recs == []

    visible = "\n".join(str(m.get("content", "")) for m in captured["msgs"])
    assert "HEAD_FACT" in visible
    assert "TAIL_OUTPUTS" in visible
    assert "TAIL_CHECKS" in visible
    assert "middle omitted for guard brief" in visible
    assert "prompt total" in visible


async def test_spawn_async_entity_detection_prefers_section_headings_over_constraint_bullets(tmp_path):
    from app.llm.tools import delegate

    prompt = """绘制通信原理作业的五张 matplotlib 图。

### 图1: 5.2题 — AMI码接收机原理方框图
### 图2: 5.17题 — AMI码和HDB3码波形图
### 图3: 5.18题 — 双相码和CMI码波形图
### 图4: 5.21题 — 升余弦波形 T0=T 眼图
### 图5: 5.21题 — 升余弦波形 T0=2T 眼图

## 关键约束：
- 所有文字用中文
- 保存到 _helpers_shared/hw_redraw_charts/
- 完成后报告5张图的路径和尺寸
"""

    result = json.loads(await delegate.handle_delegate(
        str(tmp_path),
        {
            "action": "spawn_async",
            "tasks": [{
                "task_id": "hw_redraw_charts",
                "kind": "draw",
                "mode": "easy",
                "prompt": prompt,
                "expected_outputs": [
                    "_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png",
                    "_helpers_shared/hw_redraw_charts/q5_17_waveforms.png",
                    "_helpers_shared/hw_redraw_charts/q5_18_biphase.png",
                    "_helpers_shared/hw_redraw_charts/q5_21_eye_T0eqT.png",
                    "_helpers_shared/hw_redraw_charts/q5_21_eye_T0eq2T.png",
                ],
            }],
        },
        archive_id="archive",
        group_id="group",
        user_id="user",
    ))

    warnings = result.get("granularity_warnings") or []
    broad = [w for w in warnings if w.get("issue") == "overly_broad_entity_list"]
    assert broad
    assert broad[0]["entity_items"][:3] == [
        "图1: 5.2题 — AMI码接收机原理方框图",
        "图2: 5.17题 — AMI码和HDB3码波形图",
        "图3: 5.18题 — 双相码和CMI码波形图",
    ]
    assert "所有文字用中文" not in broad[0]["entity_items"]


async def test_spawn_async_granularity_warnings_are_model_visible_facts(monkeypatch, tmp_path):
    from app.llm.tools import delegate

    async def fake_guard(*args, **kwargs):
        return True, "", [], []

    monkeypatch.setattr(delegate, "_persona_consent_guard", fake_guard)

    prompt = """Implement five independent modules and produce one result per module.

1. alpha module
2. beta module
3. gamma module
4. delta module
5. epsilon module
"""

    result = json.loads(await delegate.handle_delegate(
        str(tmp_path),
        {
            "action": "spawn_async",
            "tasks": [{
                "task_id": "wide_impl",
                "kind": "code",
                "mode": "easy",
                "prompt": prompt,
                "expected_outputs": ["results_alpha.csv", "results_beta.csv", "results_gamma.csv", "results_delta.csv"],
            }],
        },
        archive_id="archive",
        group_id="group",
        user_id="user",
    ))

    warnings = result.get("granularity_warnings") or []
    assert warnings
    joined = json.dumps(warnings, ensure_ascii=False)
    assert "suggested_action" not in joined
    assert "suggestion" not in joined
    assert "suggested_kind" not in joined
    assert "suggested_mode" not in joined
    assert "observed_parallel_boundary_fact" in joined
    assert "details" in joined
    assert "structured granularity_warnings field" in result["hint"]
    assert "按结构化建议" not in result["hint"]


def test_tool_out_summary_uses_lowercase_error_for_expected_tool_failure():
    from app.llm.tools.registry import _tool_out_summary

    result = json.dumps({
        "ok": False,
        "error": "Traceback (most recent call last):\nModuleNotFoundError: No module named 'x'",
    })

    summary = _tool_out_summary("python", result)

    assert summary.startswith("error: ")
    assert not summary.startswith("ERROR:")


def test_tool_out_summary_marks_delegate_incomplete_as_blocked():
    from app.llm.tools.registry import _tool_out_summary

    result = json.dumps({
        "ok": True,
        "task_ok": False,
        "helpers_completed": 1,
        "helpers_still_running": 0,
        "incomplete_count": 1,
        "resource_required_count": 1,
        "failed_count": 0,
    })

    summary = _tool_out_summary("delegate", result)

    assert summary == "blocked: returned_results=1/1, success=0 (incomplete=1, resource_required=1)"


def test_tool_out_summary_marks_request_resource_as_frozen():
    from app.llm.tools.registry import _tool_out_summary

    result = json.dumps({
        "ok": False,
        "action": "request_resource",
        "requires_main_resource": True,
        "resource_kind": "code",
        "needed_outputs": ["_env/src/app.py"],
    })

    summary = _tool_out_summary("request_resource", result)

    assert summary == "frozen: requested code resource needed=_env/src/app.py"


async def test_agent_state_stage_update_reuses_existing_contract_goal():
    from app.core import agent_state, debug
    from app.llm.tools import registry

    trace_id = "trace-stage-update"
    debug.set_trace_id(trace_id)
    agent_state.reset_trace(trace_id)

    created = json.loads(await registry._handle_agent_state({
        "action": "upsert_contract",
        "task_id": "dijkstra",
        "goal": "Implement Dijkstra and verify tests.",
        "current_stage": "delegate",
    }))
    assert created["ok"] is True

    updated = json.loads(await registry._handle_agent_state({
        "action": "upsert_contract",
        "task_id": "dijkstra",
        "current_stage": "verify",
    }))

    assert updated["ok"] is True
    assert updated["contract"]["goal"] == "Implement Dijkstra and verify tests."
    assert updated["contract"]["current_stage"] == "verify"


async def test_agent_state_resource_update_missing_fields_returns_request_facts():
    from app.core import agent_state, debug
    from app.llm.tools import registry

    trace_id = "trace-resource-facts"
    debug.set_trace_id(trace_id)
    agent_state.reset_trace(trace_id)
    record = agent_state.register_resource_request(
        trace_id=trace_id,
        blocked_task_id="assemble_docx",
        blocked_kind="edit",
        request={
            "resource_kind": "draw",
            "needed_outputs": ["_helpers_shared/chart.png"],
            "reason": "chart missing",
        },
    )

    missing = json.loads(await registry._handle_agent_state({
        "action": "update_resource_request",
    }))

    assert missing["ok"] is False
    assert missing["resource_requests"][0]["request_id"] == record["request_id"]
    assert missing["resource_requests"][0]["needed_outputs"] == ["_helpers_shared/chart.png"]


async def test_agent_state_resource_update_accepts_model_visible_ready_alias():
    from app.core import agent_state, debug
    from app.llm.tools import registry

    trace_id = "test_resource_ready_alias"
    debug.set_trace_id(trace_id)
    agent_state.reset_trace(trace_id)
    try:
        record = agent_state.register_resource_request(
            trace_id=trace_id,
            blocked_task_id="doc",
            blocked_kind="edit",
            request={"resource_kind": "code", "needed_outputs": ["verification_report.txt"]},
        )
        updated = json.loads(await registry._handle_agent_state({
            "action": "update_resource_request",
            "request_id": record["request_id"],
            "status": "ready",
            "satisfied_by": ["verification_report.txt"],
        }))
    finally:
        debug.set_trace_id("")

    assert updated["ok"] is True
    assert updated["resource_request"]["state"] == agent_state.RESOURCE_READY
    assert updated["status"]["blocked_work"] == []
    assert updated["status"]["ready_to_resume_work"][0]["satisfied_by"] == ["verification_report.txt"]


async def test_edit_helper_python_docx_script_rejection_points_to_office_blocks(tmp_path):
    from app.core.core_processes import reset_current_helper_kind, set_current_helper_kind
    from app.llm.tools import registry

    token = set_current_helper_kind("edit")
    helper_ws = tmp_path / "_delegate_user_doc"
    helper_ws.mkdir()
    content = """
from docx import Document
doc = Document()
doc.add_heading('Report', 1)
doc.add_paragraph('Body')
doc.save('report.docx')
""" + "\n# filler\n" * 3000
    try:
        raw = await registry._handle_workspace(
            str(helper_ws),
            {"action": "write", "path": "build_docx.py", "content": content},
        )
    finally:
        reset_current_helper_kind(token)

    result = json.loads(raw)
    assert result["ok"] is False
    assert result["blocked_kind"] == "edit"
    assert result["matching_helper_kind"] == "edit"
    assert result["observed_recovery_tool"] == "request_resource"
    assert result["observed_recovery_shape"]["resource_kind"] == "edit"
    assert "office(action='write'" in result["observed_recovery_shape"]["shape"]
    assert result["suggested_helper_kind"] == "edit"
    assert "office(action='write'" in result["suggested_request"]
    assert "office blocks" in result["error"]


async def test_spawn_async_dependency_detection_accepts_helpers_shared_expected_outputs(tmp_path):
    from app.llm.tools import delegate

    prompt = """生成最终文档，并引用 q5_2_block_diagram.png。"""

    result = json.loads(await delegate.handle_delegate(
        str(tmp_path),
        {
            "action": "spawn_async",
            "tasks": [
                {
                    "task_id": "hw_redraw_charts",
                    "kind": "draw",
                    "resume": False,
                    "fork_from": "",
                    "prompt": "生成图并保存到 _helpers_shared/hw_redraw_charts/q5_2_block_diagram.png",
                    "expected_outputs": ["_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png"],
                },
                {
                    "task_id": "hw_docx_final",
                    "kind": "edit",
                    "resume": False,
                    "fork_from": "",
                    "prompt": prompt,
                    "expected_outputs": ["通信原理作业_第五章.docx"],
                },
            ],
        },
        archive_id="archive",
        group_id="group",
        user_id="user",
    ))

    warnings = result.get("granularity_warnings") or []
    dependency = [w for w in warnings if w.get("issue") == "unsatisfied_dependency" and w.get("task_id") == "hw_docx_final"]
    assert dependency == []


async def test_spawn_async_dependency_detection_accepts_completed_declared_files(tmp_path):
    from app.llm.tools import delegate
    from collections import deque

    trace_id = "trace-declared-files-ok"
    delegate._completion_ledger[trace_id] = deque([{
        "task_id": "hw_redraw_charts",
        "kind": "draw",
        "outputs_complete": True,
        "delivered_files": [],
        "declared_files": ["_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png"],
    }], maxlen=50)

    result = json.loads(await delegate._handle_delegate_spawn_async(
        str(tmp_path),
        {"action": "spawn_async", "tasks": []},
        archive_id="archive",
        group_id="group",
        user_id="user",
        cleaned_tasks=[{
            "task_id": "hw_docx_final",
            "kind": "edit",
            "mode": "easy",
            "resume": False,
            "fork_from": "",
            "prompt": "生成最终文档，并引用 q5_2_block_diagram.png。",
            "expected_outputs": ["通信原理作业_第五章.docx"],
        }],
        main_owner="main:test",
        user_lang_now="zh",
        trace_id=trace_id,
    ))

    warnings = result.get("granularity_warnings") or []
    dependency = [w for w in warnings if w.get("issue") == "unsatisfied_dependency" and w.get("task_id") == "hw_docx_final"]
    assert dependency == []


async def test_repeated_task_attempts_warn_without_hard_reject(monkeypatch, tmp_path):
    from app.core import debug
    from app.core import core_processes
    from app.core.core_processes import ProcessRegistry
    from app.llm.tools import delegate
    from collections import deque

    trace_id = "trace-repeat-warn"
    debug.set_trace_id(trace_id)
    delegate._completion_ledger[trace_id] = deque([
        {
            "task_id": "repair_engine",
            "ok": False,
            "outputs_complete": False,
            "terminal_reason": "failed",
        }
        for _ in range(4)
    ], maxlen=50)

    fake_registry = ProcessRegistry()
    monkeypatch.setattr(core_processes, "_registry", fake_registry)

    async def fake_run_one_helper(**kwargs):
        return {
            "task_id": kwargs["task_id"],
            "ok": True,
            "report": "continued after warning",
            "files": [],
            "terminal_reason": "completed",
        }

    monkeypatch.setattr(delegate, "_run_one_helper", fake_run_one_helper)

    result = json.loads(await delegate.handle_delegate(
        str(tmp_path),
        {
            "action": "spawn_async",
            "tasks": [{
                "task_id": "repair_engine",
                "kind": "code",
                "mode": "hard",
                "resume": True,
                "prompt": "Continue from the previous failure report and verify the fix.",
                "expected_outputs": ["engine.py"],
            }],
        },
        archive_id="archive",
        group_id="group",
        user_id="user",
    ))

    assert result["ok"] is True
    assert result["action"] == "spawn_async"
    assert result["spawned"][0]["task_id"] == "repair_engine"
    warnings = result.get("granularity_warnings") or []
    repeat = [w for w in warnings if w.get("issue") == "resume_attempt_cap"]
    assert repeat
    assert repeat[0]["previous_attempts"] == 4
    assert "resume_attempt_hard_cap" not in json.dumps(result, ensure_ascii=False)


async def test_handle_delegate_spawn_empty_args_reports_json_hint(tmp_path):
    from app.llm.tools import delegate

    result = json.loads(await delegate.handle_delegate(
        str(tmp_path),
        {"action": "spawn"},
        archive_id="archive",
        group_id="group",
        user_id="user",
    ))

    assert result["ok"] is False
    assert result["error"] == "tasks must be a non-empty array"
    assert "An empty task list does not resume existing helpers" in result["hint"]


async def test_handle_delegate_wraps_top_level_single_task_args(tmp_path, monkeypatch):
    from app.llm.tools import delegate
    from app.llm.tools import delegate_actions

    async def fake_spawn_async(main_workspace, args, **kwargs):
        return json.dumps({
            "ok": True,
            "action": kwargs.get("cleaned_tasks", [{}])[0].get("task_id"),
            "tasks": args.get("tasks"),
            "schema_repair": args.get("_schema_repair_fact"),
        })

    monkeypatch.setattr(delegate_actions, "_handle_delegate_spawn_async", fake_spawn_async)

    result = json.loads(await delegate.handle_delegate(
        str(tmp_path),
        {
            "action": "spawn_async",
            "task_id": "browser-evidence",
            "kind": "code",
            "mode": "easy",
            "prompt": "Collect browser evidence.",
            "input_files": ["docs/index.html"],
            "expected_outputs": ["browser_evidence.txt"],
            "acceptance_checks": ["browser evidence exists"],
            "dispatch_reason": "Top-level fields should be schema-shape repaired.",
        },
        archive_id="archive",
        group_id="group",
        user_id="user",
    ))

    assert result["ok"] is True
    assert result["action"] == "browser-evidence"
    assert result["tasks"][0]["task_id"] == "browser-evidence"
    assert result["tasks"][0]["expected_outputs"] == ["browser_evidence.txt"]
    assert "wrapped them as tasks=[...]" in result["schema_repair"]


def test_delegate_top_level_single_task_normalizer_preserves_fields():
    from app.llm.tools.delegate_actions import _normalize_top_level_delegate_task_args

    args = {
        "action": "spawn",
        "task_id": "patch-client",
        "kind": "code",
        "mode": "easy",
        "prompt": "Patch the client.",
        "input_files": ["report_client.py"],
        "expected_outputs": ["_env/report_client.py", "_env/api_notes.md"],
        "acceptance_checks": ["tests pass"],
        "dispatch_reason": "Single task was supplied at top level.",
        "wait_window_sec": 120,
    }

    assert _normalize_top_level_delegate_task_args(args) is True
    assert args["_normalized_top_level_task"] is True
    assert args["wait_window_sec"] == 120
    assert args["tasks"] == [{
        "task_id": "patch-client",
        "prompt": "Patch the client.",
        "kind": "code",
        "mode": "easy",
        "dispatch_reason": "Single task was supplied at top level.",
        "input_files": ["report_client.py"],
        "expected_outputs": ["_env/report_client.py", "_env/api_notes.md"],
        "acceptance_checks": ["tests pass"],
    }]


async def test_delegate_allows_16_tasks_but_rejects_17(tmp_path):
    from app.llm.tools import delegate

    sixteen = [
        {"task_id": f"t{i}", "kind": "edit", "prompt": "整理一句话"}
        for i in range(16)
    ]
    ok_result = await delegate._sanitize_and_validate_tasks(
        {"action": "spawn", "tasks": sixteen},
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )
    assert isinstance(ok_result, list)
    assert len(ok_result) == 16

    seventeen = [
        {"task_id": f"t{i}", "kind": "edit", "prompt": "整理一句话"}
        for i in range(17)
    ]
    rejected = json.loads(await delegate._sanitize_and_validate_tasks(
        {"action": "spawn", "tasks": seventeen},
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    ))
    assert rejected["ok"] is False
    assert "At most 16 helper tasks can be spawned" in rejected["error"]


async def test_handle_delegate_sanitize_empty_dict_reports_json_hint(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    result = await _sanitize_and_validate_tasks(
        {},
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    assert isinstance(result, str)
    data = json.loads(result)
    assert data["ok"] is False
    assert data["error"] == "tool_call_args_json_broken"
    assert "empty dict" in data["delegate_error"]
    assert "valid JSON arguments" in data["delegate_hint"]
    assert data["raw_args_excerpt"] == "{}"


async def test_delegate_sanitize_recovers_prom_prompt_alias(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    result = await _sanitize_and_validate_tasks(
        {
            "action": "spawn",
            "tasks": [{
                "task_id": "assemble_docx",
                "kind": "edit",
                "prom": "Assemble db_index_paper.docx from staged evidence.",
                "expected_outputs": ["db_index_paper.docx"],
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    assert isinstance(result, list)
    assert result[0]["prompt"] == "Assemble db_index_paper.docx from staged evidence."
    assert result[0]["expected_outputs"] == ["db_index_paper.docx"]


async def test_delegate_marks_main_generated_command_recipes_as_unverified(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    result = await _sanitize_and_validate_tasks(
        {
            "action": "spawn",
            "tasks": [{
                "task_id": "query_data",
                "kind": "code",
                "mode": "easy",
                "prompt": (
                    "Goal: query users.db and write result.csv.\n\n"
                    "Steps:\n"
                    "1. Run `python3 -c \"import sqlite3; print('probe')\"`.\n"
                    "2. Write result.csv.\n\n"
                    "Use `python3` command (which maps to the shim python3.cmd).\n"
                ),
                "input_files": ["users.db", "verify_results.py"],
                "expected_outputs": ["result.csv"],
                "acceptance_checks": ["python3 verify_results.py prints PASS"],
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    assert not isinstance(result, str)
    prompt = result[0]["prompt"]
    assert "## Command Recipe Provenance" in prompt
    assert "non-authoritative examples" in prompt
    assert "freshly verified in this helper workspace" in prompt
    assert "Launcher fact: the main-thread prompt mentioned a `python3`/shim launcher assumption" in prompt
    assert "which maps to the shim python3.cmd" not in prompt
    assert result[0]["acceptance_checks"] == ["python3 verify_results.py prints PASS"]


async def test_delegate_does_not_add_command_recipe_block_without_command_recipe(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    result = await _sanitize_and_validate_tasks(
        {
            "action": "spawn",
            "tasks": [{
                "task_id": "small_patch",
                "kind": "code",
                "mode": "easy",
                "prompt": "Fix the parser edge case and report the verification evidence.",
                "input_files": ["parser.py"],
                "expected_outputs": ["parser.py"],
                "acceptance_checks": ["parser handles empty input"],
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    assert not isinstance(result, str)
    assert "## Command Recipe Provenance" not in result[0]["prompt"]


async def test_delegate_marks_exhaustive_scope_expansion_as_provenance(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    result = await _sanitize_and_validate_tasks(
        {
            "action": "spawn",
            "tasks": [{
                "task_id": "schema_audit",
                "kind": "code",
                "mode": "easy",
                "prompt": (
                    "Audit the database. The user now says: \"If anything in the schema is weird, double-check before assuming.\"\n\n"
                    "Do the following:\n"
                    "1. Dump all data from every table.\n"
                    "2. Report EVERYTHING in a complete evidence report.\n"
                ),
                "input_files": ["users.db", "result.csv"],
                "expected_outputs": ["_helpers_shared/schema_audit.txt"],
                "acceptance_checks": ["All row counts shown", "CSV cross-checked against DB data"],
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    assert not isinstance(result, str)
    prompt = result[0]["prompt"]
    assert "## Scope Provenance" in prompt
    assert "exhaustive audit or reporting expansion written by the main process" in prompt
    assert "smallest sufficient structured probe/report" in prompt
    assert "Dump all data from every table" in prompt
    assert result[0]["acceptance_checks"] == ["All row counts shown", "CSV cross-checked against DB data"]


async def test_handle_delegate_spawn_empty_tasks_keeps_resume_all_hint(tmp_path):
    from app.llm.tools import delegate

    result = json.loads(await delegate.handle_delegate(
        str(tmp_path),
        {"action": "spawn", "tasks": []},
        archive_id="archive",
        group_id="group",
        user_id="user",
    ))

    assert result["ok"] is False
    assert result["error"] == "tasks must be a non-empty array"
    assert "An empty task list does not resume existing helpers" in result["hint"]
    assert "without calling delegate again" in result["hint"]


async def test_delegate_missing_kind_returns_preserved_retry_envelope(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    result = await _sanitize_and_validate_tasks(
        {
            "action": "spawn",
            "tasks": [{
                "task_id": "tables_bench",
                "prompt": "Build benchmark comparison tables from existing algorithm evidence.",
                "framework": "Shared paper framework contract.",
                "input_files": ["_helpers_shared/paper_framework.md"],
                "expected_outputs": ["tables.md"],
                "acceptance_checks": ["table covers all algorithms"],
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    data = json.loads(result)
    task = data["next_delegate_shape"]["tasks"][0]
    assert data["error"] == "unsupported_helper_kind"
    assert task["task_id"] == "tables_bench"
    assert task["framework"] == "Shared paper framework contract."
    assert task["input_files"] == ["_helpers_shared/paper_framework.md"]
    assert task["expected_outputs"] == ["tables.md"]
    assert task["acceptance_checks"] == ["table covers all algorithms"]


async def test_guard_ignores_legacy_kind_recommendation_when_allowed():
    from app.llm.tools.delegate_wait import _build_guard_intervention

    helper_specs = [{
        "task_id": "paper_tables",
        "kind": "code",
        "mode": "hard",
        "prompt": "Write the theoretical comparison table section.",
        "framework": "Paper contract with output matrix and acceptance checks.",
        "dispatch_reason": "Keep this as one helper because it is a single table section from verified evidence.",
        "input_files": ["algo_evidence.md"],
        "expected_outputs": ["paper_tables.md"],
        "acceptance_checks": ["no benchmark fabrication"],
    }]
    payload = await _build_guard_intervention(
        (
            True,
            "",
            [],
            [{
                "task_id": "paper_tables",
                "current_kind": "code",
                "suggested_kind": "edit",
                "suggested_mode": "easy",
                "reason": "Document-style markdown artifacts belong to edit.",
            }],
            {},
        ),
        trace_id="test_kind_retry_shape",
        cancel_helpers=False,
        helper_specs=helper_specs,
    )

    assert payload is None


async def test_delegate_dispatch_reason_reaches_guard_and_helper_envelope(monkeypatch, tmp_path):
    from app.llm.tools import delegate
    from app.llm.tools.delegate_framework import format_helper_request_envelope

    seen_guard_tasks: list[list[dict]] = []

    async def fake_guard(*args, **kwargs):
        seen_guard_tasks.append(args[2])
        return True, "allow", [], [], {}

    monkeypatch.setattr(delegate, "_persona_consent_guard", fake_guard)
    monkeypatch.setattr(delegate, "_detect_missing_unified_framework", lambda cleaned: None)

    cleaned = await delegate._sanitize_and_validate_tasks(
        {
            "tasks": [{
                "task_id": "single_report",
                "kind": "edit",
                "mode": "easy",
                "prompt": "Write the single report from the provided evidence.",
                "dispatch_reason": "One output file, bounded evidence, and a single acceptance path.",
                "input_files": ["evidence.txt"],
                "expected_outputs": ["report.md"],
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    assert not isinstance(cleaned, str)
    assert cleaned[0]["dispatch_reason"] == "One output file, bounded evidence, and a single acceptance path."

    guard_result = await delegate._persona_consent_guard("", "", cleaned)
    assert guard_result[0] is True
    assert seen_guard_tasks[0][0]["dispatch_reason"] == cleaned[0]["dispatch_reason"]

    envelope = format_helper_request_envelope(cleaned[0])
    assert "### Main Dispatch Reason" in envelope
    assert "One output file, bounded evidence" in envelope


async def test_guard_ignores_legacy_framework_block_when_allowed():
    from app.llm.tools import delegate
    from app.llm.tools.delegate_wait import _build_guard_intervention

    helper_specs = [{
        "task_id": "algo_rbt",
        "kind": "code",
        "mode": "easy",
        "prompt": "Implement only the red-black tree slice.",
        "input_files": ["_helpers_shared/paper_framework.md"],
        "expected_outputs": ["rbtree.py", "results_rbt.csv"],
        "acceptance_checks": ["uses shared csv schema"],
    }]
    key = delegate._framework_block_counter_key("test_framework_retry_shape", ["algo_rbt"])
    delegate._guard_framework_block_trace_total.pop(key, None)
    payload = await _build_guard_intervention(
        (
            True,
            "",
            [],
            [],
            {
                "block": True,
                "task_ids": ["algo_rbt"],
                "reason": "Need a shared benchmark framework.",
            },
        ),
        trace_id="test_framework_retry_shape",
        cancel_helpers=False,
        helper_specs=helper_specs,
    )

    assert payload is None


def test_tool_call_args_normalize_completes_truncated_multiline_delegate_args():
    from app.llm.client import _normalize_tool_call_args_for_dispatch

    raw = '{"action":"spawn","tasks":[{"task_id":"draw","prompt":"第一行\n第二行'
    args, err, repaired = _normalize_tool_call_args_for_dispatch(raw)

    assert err is None
    assert repaired is True
    assert args == {
        "action": "spawn",
        "tasks": [{"task_id": "draw", "prompt": "第一行\n第二行"}],
    }


async def test_chat_with_tools_loop_repairs_truncated_multiline_delegate_json_and_dispatches(monkeypatch):
    from app.llm import client as llm_client
    from app.llm.client import chat_with_tools_loop

    dispatched = []

    async def fake_streaming_call(**kwargs):
        tool_call = SimpleNamespace(
            id="call_delegate_repaired",
            function=SimpleNamespace(
                name="delegate",
                arguments='{"action":"spawn","tasks":[{"task_id":"draw","prompt":"第一行\n第二行',
            ),
        )
        message = SimpleNamespace(
            role="assistant",
            content=None,
            tool_calls=[tool_call],
            reasoning_content=None,
        )
        response = SimpleNamespace(choices=[SimpleNamespace(message=message)])
        return response, SimpleNamespace(), "ok"

    async def fake_dispatcher(name, args):
        dispatched.append((name, args))
        return json.dumps({"ok": True, "results": [{"task_id": "draw", "report": "done"}]})

    monkeypatch.setattr(llm_client, "_call_llm_streaming_with_idle", fake_streaming_call)

    await chat_with_tools_loop(
        [{"role": "user", "content": "run delegate"}],
        tools=[{"type": "function", "function": {"name": "delegate", "parameters": {}}}],
        dispatcher=fake_dispatcher,
        max_iter=1,
    )

    assert dispatched == [
        ("delegate", {"action": "spawn", "tasks": [{"task_id": "draw", "prompt": "第一行\n第二行"}]}),
    ]


async def test_chat_with_tools_loop_blocks_malformed_delegate_json_before_dispatch(monkeypatch):
    from app.llm import client as llm_client
    from app.llm.client import chat_with_tools_loop

    dispatched = []

    calls = 0

    async def fake_streaming_call(**kwargs):
        nonlocal calls
        calls += 1
        if calls == 1:
            tool_call = SimpleNamespace(
                id="call_delegate_broken",
                function=SimpleNamespace(
                    name="delegate",
                    arguments='{"action":"spawn","tasks":[{"task_id":"draw"}]]',
                ),
            )
            message = SimpleNamespace(
                role="assistant",
                content=None,
                tool_calls=[tool_call],
                reasoning_content=None,
            )
        else:
            message = SimpleNamespace(
                role="assistant",
                content=json.dumps({"intent": "参数错误", "key_points": []}, ensure_ascii=False),
                tool_calls=None,
                reasoning_content=None,
            )
        response = SimpleNamespace(choices=[SimpleNamespace(message=message)])
        return response, SimpleNamespace(), "ok"

    async def fake_dispatcher(name, args):
        dispatched.append((name, args))
        return json.dumps({"ok": True})

    monkeypatch.setattr(llm_client, "_call_llm_streaming_with_idle", fake_streaming_call)

    content, msgs = await chat_with_tools_loop(
        [{"role": "user", "content": "run delegate"}],
        tools=[{"type": "function", "function": {"name": "delegate", "parameters": {}}}],
        dispatcher=fake_dispatcher,
        max_iter=1,
    )

    assert dispatched == []
    tool_msgs = [m for m in msgs if m.get("role") == "tool"]
    assert tool_msgs
    result = json.loads(tool_msgs[-1]["content"])
    assert result["ok"] is False
    assert result["error"] == "tool_call_args_json_broken"
    assert result["tool_name"] == "delegate"
    assert "was not dispatched to the real tool" in result["hint"]
    assert "工具参数 JSON 解析失败" in result["hint"]


def test_main_thread_source_write_abort_has_dedicated_non_error_path():
    import inspect
    from app.llm import client as llm_client
    from app.llm import client_tools_loop

    stream_src = inspect.getsource(llm_client._call_llm_streaming_with_idle)
    loop_src = inspect.getsource(client_tools_loop.chat_with_tools_loop)

    assert 'return collector.to_response(model), collector, "main_thread_source_write"' in stream_src
    assert 'stream_exit_reason == "main_thread_source_write"' in loop_src
    assert "LLM stream controlled abort" in loop_src
    assert '"idle_timeout", "main_thread_source_write", "helper_tool_call_bloat"' in loop_src
    assert "SOURCE_WRITE_DELEGATION_HINT" in loop_src
    from app.llm.tools.runtime_hints import SOURCE_WRITE_DELEGATION_HINT
    assert "[SYSTEM_HINT/source_write_delegation]" in SOURCE_WRITE_DELEGATION_HINT
    assert "The recoverable facts are the target path, content size" in SOURCE_WRITE_DELEGATION_HINT
    assert "does not add evidence" in SOURCE_WRITE_DELEGATION_HINT
    assert "thinking_disabled retry" in loop_src


async def test_office_append_after_adaptive_shrink_returns_arg_size_warning(tmp_path):
    from docx import Document
    from app.llm.tools.office import handle_office, report_office_failure, set_office_adaptive_key

    set_office_adaptive_key("test_after_adaptive_shrink_warning")
    report_office_failure("test_after_adaptive_shrink_warning", "test floor sizing warning")
    report_office_failure("test_after_adaptive_shrink_warning", "test floor sizing warning")
    report_office_failure("test_after_adaptive_shrink_warning", "test floor sizing warning")
    Document().save(tmp_path / "report.docx")
    blocks = [{"type": "paragraph", "text": "x" * 1300}]
    result = json.loads(await handle_office(str(tmp_path), {
        "action": "append",
        "path": "report.docx",
        "blocks": blocks,
    }))

    assert result["ok"] is True
    assert result["blocks_appended"] == 1
    assert result["arg_size_warning"]["blocks_count"] == 1
    assert result["arg_size_warning"]["current_blocks_limit"] == 12
    assert result["arg_size_warning"]["adaptive_floor_reached"] is True


async def test_office_append_warns_when_total_text_is_too_large_after_adaptive_shrink(tmp_path):
    from docx import Document
    from app.llm.tools.office import handle_office, report_office_failure, set_office_adaptive_key

    set_office_adaptive_key("test_total_text_warning")
    report_office_failure("test_total_text_warning", "test floor sizing warning")
    report_office_failure("test_total_text_warning", "test floor sizing warning")
    report_office_failure("test_total_text_warning", "test floor sizing warning")
    Document().save(tmp_path / "report.docx")
    blocks = [
        {"type": "paragraph", "text": "a" * 1100},
        {"type": "paragraph", "text": "b" * 1100},
        {"type": "paragraph", "text": "c" * 1100},
        {"type": "paragraph", "text": "d" * 1100},
    ]
    result = json.loads(await handle_office(str(tmp_path), {
        "action": "append",
        "path": "report.docx",
        "blocks": blocks,
    }))

    assert result["ok"] is True
    assert result["blocks_appended"] == 4
    assert result["arg_size_warning"]["blocks_count"] == 4
    assert result["arg_size_warning"]["total_text_chars"] == 4400
    assert result["arg_size_warning"]["total_text_too_large"] is True
    assert result["arg_size_warning"]["too_many_blocks"] is False
    assert "blocks≤12" in result["arg_size_warning"]["next_action_instruction"]


async def test_office_append_long_text_block_returns_arg_size_warning(tmp_path):
    from docx import Document
    from app.llm.tools.office import handle_office, report_office_failure, set_office_adaptive_key

    set_office_adaptive_key("test_long_text_block_warning")
    report_office_failure("test_long_text_block_warning", "test floor sizing warning")
    report_office_failure("test_long_text_block_warning", "test floor sizing warning")
    report_office_failure("test_long_text_block_warning", "test floor sizing warning")
    Document().save(tmp_path / "report.docx")
    result = json.loads(await handle_office(str(tmp_path), {
        "action": "append",
        "path": "report.docx",
        "blocks": [{"type": "paragraph", "text": "x" * 1300}],
    }))

    assert result["ok"] is True
    assert result["arg_size_warning"]["oversize_text_blocks"][0]["index"] == 0
    assert result["arg_size_warning"]["oversize_text_blocks"][0]["chars"] == 1300


async def test_office_append_rejects_blocks_above_hard_limit(tmp_path):
    from docx import Document
    from app.llm.tools.office import handle_office, report_office_failure, set_office_adaptive_key

    set_office_adaptive_key("test_blocks_above_hard_limit")
    report_office_failure("test_blocks_above_hard_limit", "test floor sizing warning")
    report_office_failure("test_blocks_above_hard_limit", "test floor sizing warning")
    report_office_failure("test_blocks_above_hard_limit", "test floor sizing warning")
    Document().save(tmp_path / "report.docx")
    blocks = [{"type": "paragraph", "text": f"段落 {i}"} for i in range(13)]
    result = json.loads(await handle_office(str(tmp_path), {
        "action": "append",
        "path": "report.docx",
        "blocks": blocks,
    }))

    assert result["ok"] is False
    assert result["error"] == "office_args_too_large"
    assert result["arg_size_warning"]["too_many_blocks_hard"] is True
    assert result["arg_size_warning"]["blocks_count"] == 13
    assert result["recovery"]["action"] == "retry_with_smaller_blocks"


async def test_office_append_at_hard_limit_still_succeeds_without_warning(tmp_path):
    from docx import Document
    from app.llm.tools.office import handle_office, report_office_failure, set_office_adaptive_key

    set_office_adaptive_key("test_at_hard_limit_no_warning")
    report_office_failure("test_at_hard_limit_no_warning", "test floor sizing warning")
    report_office_failure("test_at_hard_limit_no_warning", "test floor sizing warning")
    report_office_failure("test_at_hard_limit_no_warning", "test floor sizing warning")
    Document().save(tmp_path / "report.docx")
    blocks = [{"type": "paragraph", "text": f"段落 {i}"} for i in range(12)]
    result = json.loads(await handle_office(str(tmp_path), {
        "action": "append",
        "path": "report.docx",
        "blocks": blocks,
    }))

    assert result["ok"] is True
    assert result["blocks_appended"] == 12
    assert "arg_size_warning" not in result


async def test_office_docx_table_cells_accept_text_style_objects(tmp_path):
    from docx import Document
    from app.llm.tools.office import handle_office

    Document().save(tmp_path / "report.docx")
    result = json.loads(await handle_office(str(tmp_path), {
        "action": "append",
        "path": "report.docx",
        "blocks": [{
            "type": "table",
            "header": True,
            "rows": [
                [{"text": "算法", "style": "strong"}, {"text": "插入", "style": "strong"}],
                [{"text": "红黑树"}, {"text": "O(log n)"}],
            ],
        }],
    }))

    assert result["ok"] is True
    doc = Document(tmp_path / "report.docx")
    assert doc.tables[0].cell(0, 0).text == "算法"
    assert "{'text'" not in doc.tables[0].cell(0, 0).text
    assert doc.tables[0].cell(0, 0).paragraphs[0].runs[0].bold is True


async def test_office_docx_table_cells_accept_serialized_text_style_objects(tmp_path):
    from docx import Document
    from app.llm.tools.office import handle_office

    Document().save(tmp_path / "report.docx")
    result = json.loads(await handle_office(str(tmp_path), {
        "action": "append",
        "path": "report.docx",
        "blocks": [{
            "type": "table",
            "header": True,
            "rows": [
                ["{'text': '算法', 'style': 'strong'}", "{'text': '插入', 'style': 'strong'}"],
                ["{'text': '红黑树'}", "{'text': 'O(log n)'}"],
            ],
        }],
    }))

    assert result["ok"] is True
    doc = Document(tmp_path / "report.docx")
    assert doc.tables[0].cell(0, 0).text == "算法"
    assert doc.tables[0].cell(1, 0).text == "红黑树"
    assert "{'text'" not in doc.tables[0].cell(0, 0).text
    assert doc.tables[0].cell(0, 0).paragraphs[0].runs[0].bold is True


async def test_office_docx_accepts_subtitle_as_paragraph_alias(tmp_path):
    from app.llm.tools.office import handle_office
    from docx import Document

    result = json.loads(await handle_office(str(tmp_path), {
        "action": "write",
        "path": "report.docx",
        "title": "Report",
        "blocks": [
            {"type": "subtitle", "text": "Author, Affiliation"},
            {"type": "paragraph", "text": "Body"},
        ],
    }))

    assert result["ok"] is True
    assert result["input_normalizations"]["subtitle_blocks_as_paragraph"] == 1
    doc = Document(tmp_path / "report.docx")
    assert "Author, Affiliation" in [p.text for p in doc.paragraphs]


async def test_office_docx_table_rows_accept_cells_object_alias(tmp_path):
    from app.llm.tools.office import handle_office
    from docx import Document

    Document().save(tmp_path / "report.docx")
    result = json.loads(await handle_office(str(tmp_path), {
        "action": "append",
        "path": "report.docx",
        "blocks": [{
            "type": "table",
            "rows": [
                {"cells": ["Operation", "Worst Case"]},
                {"cells": ["Search", "O(log n)"]},
            ],
        }],
    }))

    assert result["ok"] is True
    assert result["input_normalizations"]["table_rows_cells_objects_as_arrays"] == 2
    doc = Document(tmp_path / "report.docx")
    assert doc.tables[0].cell(0, 0).text == "Operation"
    assert doc.tables[0].cell(1, 1).text == "O(log n)"


async def test_office_docx_table_empty_row_error_reports_row_index(tmp_path):
    from docx import Document
    from app.llm.tools.office import handle_office

    Document().save(tmp_path / "report.docx")
    result = json.loads(await handle_office(str(tmp_path), {
        "action": "append",
        "path": "report.docx",
        "blocks": [{
            "type": "table",
            "rows": [
                ["算法", "插入"],
                ["", {"text": ""}],
            ],
        }],
    }))

    assert result["ok"] is False
    assert "row_index" in result["error"]
    assert "row_has_no_non_empty_cells" in result["error"]
    assert "paragraph/list" in result["error"]


async def test_office_docx_accepts_numbered_list_alias(tmp_path):
    from docx import Document
    from app.llm.tools.office import handle_office

    Document().save(tmp_path / "report.docx")
    result = json.loads(await handle_office(str(tmp_path), {
        "action": "append",
        "path": "report.docx",
        "blocks": [{
            "type": "numbered_list",
            "items": ["first", "second"],
        }],
    }))

    assert result["ok"] is True
    doc = Document(tmp_path / "report.docx")
    assert [p.text for p in doc.paragraphs if p.text] == ["first", "second"]
    assert all(p.style.name == "List Number" for p in doc.paragraphs if p.text)


async def test_office_docx_accepts_heading_level_alias(tmp_path):
    from docx import Document
    from app.llm.tools.office import handle_office

    Document().save(tmp_path / "report.docx")
    result = json.loads(await handle_office(str(tmp_path), {
        "action": "append",
        "path": "report.docx",
        "blocks": [{
            "type": "heading1",
            "text": "Introduction",
        }],
    }))

    assert result["ok"] is True
    doc = Document(tmp_path / "report.docx")
    para = next(p for p in doc.paragraphs if p.text == "Introduction")
    assert para.style.name == "Heading 1"


async def test_p35_existing_workspace_basename_is_not_missing_dependency(tmp_path):
    from app.llm.tools.delegate import _workspace_has_basename

    (tmp_path / "gen_521_analysis_answers_521.txt").write_text("answers", encoding="utf-8")

    existing = _workspace_has_basename(str(tmp_path))

    assert "gen_521_analysis_answers_521.txt" in existing
    assert "answers_521.txt" not in existing


async def test_spawn_async_dependency_detection_accepts_prefixed_expected_output_names():
    completed_files = {"q518_q51_517_answers_51_517.txt", "q52_q52_answers_52.txt"}
    will_produce_files = {"q518_answers_518.txt", "answers_518.txt"}
    referenced = {
        "q518_answers_518.txt",
        "q51_517_answers_51_517.txt",
        "q52_answers_52.txt",
    }

    unsatisfied = []
    for ref in sorted(referenced):
        ref_l = ref.lower()
        found = (
            ref_l in completed_files
            or ref_l in will_produce_files
            or any(f.endswith("_" + ref_l) for f in completed_files)
            or any(f.endswith("_" + ref_l) for f in will_produce_files)
        )
        if not found:
            unsatisfied.append(ref)

    assert unsatisfied == []


async def test_spawn_async_dependency_detection_accepts_helpers_shared_prefixed_outputs(tmp_path):
    from app.llm.tools import delegate
    from collections import deque

    trace_id = "trace-shared-prefixed-files-ok"
    delegate._completion_ledger[trace_id] = deque([{
        "task_id": "q51_517",
        "kind": "code",
        "outputs_complete": True,
        "delivered_files": ["q518_q51_517_answers_51_517.txt"],
        "declared_files": ["_helpers_shared/q518/q51_517_answers_51_517.txt"],
    }, {
        "task_id": "q52",
        "kind": "code",
        "outputs_complete": True,
        "delivered_files": ["q52_q52_answers_52.txt"],
        "declared_files": ["_helpers_shared/q52/q52_answers_52.txt"],
    }], maxlen=50)

    result = json.loads(await delegate._handle_delegate_spawn_async(
        str(tmp_path),
        {"action": "spawn_async", "tasks": []},
        archive_id="archive",
        group_id="group",
        user_id="user",
        cleaned_tasks=[{
            "task_id": "hw_docx_final",
            "kind": "edit",
            "mode": "easy",
            "resume": False,
            "fork_from": "",
            "prompt": "生成最终文档，并引用 q518_answers_518.txt、q51_517_answers_51_517.txt、q52_answers_52.txt。",
            "expected_outputs": ["通信原理作业_第五章.docx"],
        }, {
            "task_id": "q518",
            "kind": "code",
            "mode": "easy",
            "resume": False,
            "fork_from": "",
            "prompt": "生成 q518_answers_518.txt 并写入 _helpers_shared/q518/q518_answers_518.txt。",
            "expected_outputs": ["_helpers_shared/q518/q518_answers_518.txt"],
        }],
        main_owner="main:test",
        user_lang_now="zh",
        trace_id=trace_id,
    ))

    warnings = result.get("granularity_warnings") or []
    dependency = [w for w in warnings if w.get("issue") == "unsatisfied_dependency" and w.get("task_id") == "hw_docx_final"]
    assert dependency == []


async def test_workspace_has_basename_includes_helpers_shared_relative_paths(tmp_path):
    from app.llm.tools.delegate import _workspace_has_basename

    nested = tmp_path / "_helpers_shared" / "hw_redraw_charts"
    nested.mkdir(parents=True)
    (nested / "q5_2_block_diagram.png").write_bytes(b"png")

    existing = _workspace_has_basename(str(tmp_path))

    assert "q5_2_block_diagram.png" in existing
    assert "_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png" in existing
    assert "hw_redraw_charts/q5_2_block_diagram.png" in existing


async def test_removed_general_helper_returns_explicit_replan_error(monkeypatch, tmp_path):
    from app.llm.tools import delegate

    result = json.loads(await delegate.handle_delegate(
        str(tmp_path),
        {
            "action": "spawn",
            "wait_window_sec": 0,
            "tasks": [{
                "task_id": "hw_legacy_docx_prompt",
                "kind": "general",
                "prompt": "请生成完整课程作业并保存为 homework_report.docx。",
            }],
        },
        archive_id="archive",
        group_id="group",
        user_id="user",
    ))

    assert result["ok"] is False
    assert result["error"] == "unsupported_helper_kind"
    assert result["allowed_kinds"]
    assert "edit" in result["allowed_kinds"]
    assert "general" not in result["allowed_kinds"]
    assert "homework_report.docx" in result["original_prompt"]


async def test_sanitize_unknown_legacy_pairing_flag_does_not_generate_hard_twin(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    cleaned = await _sanitize_and_validate_tasks(
        {
            "legacy_pairing_flag": True,
            "tasks": [{
                "task_id": "impl",
                "kind": "code",
                "mode": "easy",
                "prompt": "生成 impl.c 并编译测试。",
                "expected_outputs": ["impl.c"],
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    assert not isinstance(cleaned, str)
    by_id = {task["task_id"]: task for task in cleaned}
    assert set(by_id) == {"impl"}
    assert "paired_with" not in by_id["impl"]


async def test_sanitize_unknown_legacy_pairing_flag_does_not_log_generated_pair(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    logged: list[tuple[str, str]] = []

    with patch("app.llm.tools.delegate.debug.log", lambda category, msg='', payload=None: logged.append((category, msg))):
        cleaned = await _sanitize_and_validate_tasks(
            {
                "legacy_pairing_flag": True,
                "tasks": [{
                    "task_id": "impl",
                    "kind": "code",
                    "mode": "easy",
                    "prompt": "生成 impl.c 并编译测试。",
                    "expected_outputs": ["impl.c"],
                }],
            },
            main_workspace=str(tmp_path),
            archive_id="archive",
            group_id="group",
            user_id="user",
        )

    assert not isinstance(cleaned, str)
    assert not any(category == "delegate.explicit_hard_backup_pair" for category, _ in logged)
    assert not any(task["task_id"] == "impl_hard" for task in cleaned)


async def test_sanitize_unknown_legacy_pairing_flag_does_not_pair_hard_suffix_alone(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    cleaned = await _sanitize_and_validate_tasks(
        {
            "legacy_pairing_flag": True,
            "tasks": [{
                "task_id": "impl_hard",
                "kind": "code",
                "mode": "hard",
                "prompt": "生成 impl.c 并编译测试。",
                "expected_outputs": ["impl.c"],
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    assert not isinstance(cleaned, str)
    assert [task["task_id"] for task in cleaned] == ["impl_hard"]


async def test_sanitize_easy_code_task_does_not_auto_pair_hard_backup(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    cleaned = await _sanitize_and_validate_tasks(
        {
            "tasks": [{
                "task_id": "impl",
                "kind": "code",
                "mode": "easy",
                "prompt": "生成 impl.c 并编译测试。",
                "expected_outputs": ["impl.c"],
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    assert not isinstance(cleaned, str)
    by_id = {task["task_id"]: task for task in cleaned}
    assert set(by_id) == {"impl"}
    assert "paired_with" not in by_id["impl"]


async def test_environment_broad_hard_pairing_work_is_rejected_before_spawn(monkeypatch, tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import delegate

    async def fake_guard(*args, **kwargs):
        return True, "", [], [], {}

    monkeypatch.setattr(delegate, "_persona_consent_guard", fake_guard)
    monkeypatch.setattr(
        delegate,
        "_detect_missing_unified_framework",
        lambda cleaned: None,
    )

    tasks = [
        {
            "task_id": f"env_task_{i}",
            "kind": "code",
            "mode": "easy",
            "prompt": (
                "Refactor this multi-file _env/ project slice: update _env/src/pkg/module.py, "
                "_env/src/pkg/service.py, _env/src/pkg/store.py, _env/tests/test_module.py, "
                "and _env/README.md; run compile and smoke validation."
            ),
            "expected_outputs": [
                f"_env/src/pkg/module_{i}.py",
                f"_env/src/pkg/service_{i}.py",
                f"_env/src/pkg/store_{i}.py",
                f"_env/tests/test_module_{i}.py",
                f"_env/docs/notes_{i}.md",
            ],
        }
        for i in range(11)
    ]

    env = EnvironmentContext(
        root_dir=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
        project_key="project",
    )
    with runtime_context("environment", env):
        result = await delegate._sanitize_and_validate_tasks(
            {"tasks": tasks},
            main_workspace=str(tmp_path),
            archive_id="archive",
            group_id="group",
            user_id="user",
        )

    assert not isinstance(result, str)
    primary_tasks = [
        task for task in result
        if str(task.get("task_id") or "").startswith("env_task_")
        and not str(task.get("task_id") or "").endswith("_hard")
    ]
    assert primary_tasks
    assert all(task.get("guard_observations") for task in primary_tasks)
    assert {
        observation.get("issue")
        for task in primary_tasks
        for observation in task.get("guard_observations", [])
    } == {"broad_code_task_before_spawn"}


async def test_environment_small_code_task_does_not_auto_pair_hard(monkeypatch, tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import delegate

    async def fake_guard(*args, **kwargs):
        return True, "", [], [], {}

    monkeypatch.setattr(delegate, "_persona_consent_guard", fake_guard)
    monkeypatch.setattr(delegate, "_detect_missing_unified_framework", lambda cleaned: None)

    env = EnvironmentContext(
        root_dir=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
        project_key="project",
    )
    with runtime_context("environment", env):
        cleaned = await delegate._sanitize_and_validate_tasks(
            {
                "tasks": [{
                    "task_id": "dijkstra_impl",
                    "kind": "code",
                    "mode": "easy",
                    "prompt": "Edit _env/src/algolab/graph.py and verify _env/tests/test_graph.py.",
                    "expected_outputs": ["_env/src/algolab/graph.py", "_env/tests/test_graph.py"],
                }]
            },
            main_workspace=str(tmp_path),
            archive_id="archive",
            group_id="group",
            user_id="user",
        )

    assert not isinstance(cleaned, str)
    assert [task["task_id"] for task in cleaned] == ["dijkstra_impl"]
    assert cleaned[0]["write_scopes"] == ["_env/src/algolab", "_env/tests"]


async def test_environment_single_algorithm_with_tests_readme_does_not_auto_pair_hard(monkeypatch, tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import delegate

    async def fake_guard(*args, **kwargs):
        return True, "", [], [], {}

    monkeypatch.setattr(delegate, "_persona_consent_guard", fake_guard)
    monkeypatch.setattr(delegate, "_detect_missing_unified_framework", lambda cleaned: None)

    prompt = (
        "You are implementing weighted-graph Dijkstra shortest paths in an existing "
        "src-layout _env/ Python project. Update _env/src/algolab/graph.py, "
        "_env/tests/test_graph.py, and _env/README.md. Include project context, "
        "a complete algorithm specification, six pytest cases, README usage notes, "
        "and a verification command. Existing README context mentions graph algorithms "
        "and benchmark reports, but this request implements one algorithm only. "
        "Do not create benchmarks or a long report. "
        + "Acceptance detail. " * 80
    )
    env = EnvironmentContext(
        root_dir=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
        project_key="project",
    )
    with runtime_context("environment", env):
        cleaned = await delegate._sanitize_and_validate_tasks(
            {
                "tasks": [{
                    "task_id": "dijkstra_impl",
                    "kind": "code",
                    "mode": "easy",
                    "prompt": prompt,
                    "expected_outputs": [
                        "_env/src/algolab/graph.py",
                        "_env/tests/test_graph.py",
                        "_env/README.md",
                    ],
                }]
            },
            main_workspace=str(tmp_path),
            archive_id="archive",
            group_id="group",
            user_id="user",
        )

    assert not isinstance(cleaned, str)
    assert [task["task_id"] for task in cleaned] == ["dijkstra_impl"]


async def test_environment_project_code_task_too_broad_rejected_before_hard_pairing(monkeypatch, tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import delegate

    async def fake_guard(*args, **kwargs):
        return True, "", [], [], {}

    monkeypatch.setattr(delegate, "_persona_consent_guard", fake_guard)
    monkeypatch.setattr(delegate, "_detect_missing_unified_framework", lambda cleaned: None)

    env = EnvironmentContext(
        root_dir=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
        project_key="project",
    )
    with runtime_context("environment", env):
        result = await delegate._sanitize_and_validate_tasks(
            {
                "tasks": [{
                    "task_id": "project_refactor",
                    "kind": "code",
                    "mode": "easy",
                    "prompt": (
                        "Refactor this multi-file _env/ project. Update _env/src/app.py, "
                        "_env/src/core.py, _env/src/store.py, _env/tests/test_core.py, "
                        "and _env/README.md; run compile and smoke validation."
                    ),
                    "expected_outputs": [
                        "_env/src/app.py",
                        "_env/src/core.py",
                        "_env/src/store.py",
                        "_env/tests/test_core.py",
                        "_env/README.md",
                    ],
                }]
            },
            main_workspace=str(tmp_path),
            archive_id="archive",
            group_id="group",
            user_id="user",
        )

    _assert_sanitized_task_has_broad_fact(result, task_id="project_refactor")


def test_code_hard_pairing_no_longer_rejects_on_overflow_source():
    from pathlib import Path

    src = Path("app/llm/tools/delegate.py").read_text(encoding="utf-8")
    assert "code_hard_pairing_exceeds_delegate_task_limit" not in src
    assert "delegate.code_hard_explicit_paired" in src
    assert "_explicit_hard_by_base" in src


def test_delegate_hard_mode_suffix_has_single_english_first_definition():
    from pathlib import Path

    src = Path("app/llm/tools/helper_prompt_catalog.py").read_text(encoding="utf-8")
    assert src.count("_HARD_MODE_SUFFIX =") == 1
    assert "Hard mode is a richer same-kind workflow" in src
    assert "For code/coding hard mode" in src
    assert "For read hard mode" in src
    assert "Hard Mode Discipline" not in src


async def test_read_helper_env_analysis_outputs_become_internal_evidence(monkeypatch, tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import delegate

    async def fake_guard(*args, **kwargs):
        return True, "", [], [], {}

    monkeypatch.setattr(delegate, "_persona_consent_guard", fake_guard)
    monkeypatch.setattr(delegate, "_detect_missing_unified_framework", lambda cleaned: None)

    env = EnvironmentContext(
        root_dir=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
        project_key="project",
    )
    with runtime_context("environment", env):
        cleaned = await delegate._sanitize_and_validate_tasks(
            {
                "tasks": [{
                    "task_id": "read_filesystem",
                    "kind": "read",
                    "mode": "easy",
                    "prompt": (
                        "Read staged project source files and write "
                        "_env/filesystem_analysis.txt as internal evidence."
                    ),
                    "expected_outputs": ["_env/filesystem_analysis.txt"],
                    "acceptance_checks": ["_env/filesystem_analysis.txt exists"],
                }]
            },
            main_workspace=str(tmp_path),
            archive_id="archive",
            group_id="group",
            user_id="user",
        )

    assert not isinstance(cleaned, str)
    task = cleaned[0]
    assert task["expected_outputs"] == ["filesystem_analysis.txt"]
    assert "_env/filesystem_analysis.txt" not in task["prompt"]
    assert "filesystem_analysis.txt" in task["prompt"]
    assert task["acceptance_checks"] == ["filesystem_analysis.txt exists"]


async def test_delegate_reports_read_only_analysis_file_output_fact_to_guard(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    result = await _sanitize_and_validate_tasks(
        {
            "tasks": [{
                "task_id": "audit_context",
                "kind": "read",
                "mode": "easy",
                "prompt": (
                    "Read-only analysis. Inspect app/core/context.py and write "
                    "findings to _evidence_context_assembly.txt."
                ),
                "expected_outputs": ["_evidence_context_assembly.txt"],
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    assert not isinstance(result, str)
    task = result[0]
    observations = task.get("guard_observations") or []
    assert task["task_id"] == "audit_context"
    assert "_evidence_context_assembly.txt" in task["expected_outputs"]
    fact = next(
        item for item in observations
        if item.get("issue") == "read_only_analysis_output_conflict"
    )
    assert fact["task_id"] == "audit_context"
    assert "_evidence_context_assembly.txt" in fact["expected_outputs"]
    assert "guard should decide" in fact["details"]


async def test_delegate_code_helper_can_own_env_outputs_when_tests_are_read_only(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    result = await _sanitize_and_validate_tasks(
        {
            "tasks": [{
                "task_id": "migrate_contract",
                "kind": "code",
                "mode": "easy",
                "prompt": (
                    "Update project source from customer_name to account_name. "
                    "Do not modify tests. 不要改动测试文件。"
                ),
                "input_files": [
                    "_env/contracts/customer_event.py",
                    "_env/service/render.py",
                ],
                "expected_outputs": [
                    "_env/contracts/customer_event.py",
                    "_env/service/render.py",
                ],
                "acceptance_checks": ["pytest contracts/tests service/tests"],
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    assert not isinstance(result, str)
    assert result[0]["kind"] == "code"
    assert result[0]["expected_outputs"] == [
        "_env/contracts/customer_event.py",
        "_env/service/render.py",
    ]


async def test_delegate_reports_each_read_only_analysis_file_output_fact(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    result = await _sanitize_and_validate_tasks(
        {
            "tasks": [
                {
                    "task_id": "read_context",
                    "kind": "read",
                    "mode": "easy",
                    "prompt": "Read-only analysis. Inspect context files and write _evidence_context.txt.",
                    "expected_outputs": ["_evidence_context.txt"],
                },
                {
                    "task_id": "read_tools",
                    "kind": "read",
                    "mode": "easy",
                    "prompt": "Analysis only; inspect tool files and save _evidence_tools.txt.",
                    "expected_outputs": ["_evidence_tools.txt"],
                },
            ],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    assert not isinstance(result, str)
    by_id = {task["task_id"]: task for task in result}
    assert set(by_id) == {"read_context", "read_tools"}
    for task_id, output in (
        ("read_context", "_evidence_context.txt"),
        ("read_tools", "_evidence_tools.txt"),
    ):
        observations = by_id[task_id].get("guard_observations") or []
        fact = next(
            item for item in observations
            if item.get("issue") == "read_only_analysis_output_conflict"
        )
        assert fact["task_id"] == task_id
        assert fact["expected_outputs"] == [output]


async def test_delegate_prompt_output_candidates_do_not_become_expected_outputs(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    result = await _sanitize_and_validate_tasks(
        {
            "tasks": [{
                "task_id": "write_report",
                "kind": "edit",
                "mode": "easy",
                "prompt": "Summarize the evidence and generate final_report.md for the user.",
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    assert not isinstance(result, str)
    task = result[0]
    assert task["expected_outputs"] == []
    assert "Candidate Output Facts" in task["prompt"]
    assert "final_report.md" in task["prompt"]
    assert task["guard_observations"][0]["issue"] == "undeclared_output_candidates"
    assert task["guard_observations"][0]["candidate_outputs"] == ["final_report.md"]


async def test_delegate_reports_helper_produced_read_input_fact_to_guard(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    result = await _sanitize_and_validate_tasks(
        {
            "tasks": [{
                "task_id": "read_prior_report",
                "kind": "read",
                "mode": "easy",
                "prompt": "Read helper_summary.md and extract missing evidence.",
                "input_files": ["helper_summary.md"],
                "expected_outputs": ["prior_report_evidence.txt"],
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    assert not isinstance(result, str)
    task = result[0]
    observations = task.get("guard_observations") or []
    fact = next(
        item for item in observations
        if item.get("issue") == "read_helper_targets_helper_produced_artifacts"
    )
    assert fact["task_id"] == "read_prior_report"
    assert fact["inputs"] == ["helper_summary.md"]
    assert "guard should decide" in fact["details"]


async def test_delegate_unknown_helper_kind_returns_fact_not_code_fallback(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    result = await _sanitize_and_validate_tasks(
        {
            "tasks": [{
                "task_id": "browser_check",
                "kind": "browser",
                "mode": "easy",
                "prompt": "Open the app and verify the form behavior.",
                "expected_outputs": ["browser_report.md"],
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    assert isinstance(result, str)
    data = json.loads(result)
    assert data["ok"] is False
    assert data["error"] == "unsupported_helper_kind"
    assert data["requested_kind"] == "browser"
    assert "code" in data["allowed_kinds"]
    assert "browser" not in data["allowed_kinds"]
    assert data["next_delegate_shape"]["tasks"][0]["expected_outputs"] == ["browser_report.md"]


async def test_delegate_ocr_kind_returns_tool_capability_fact(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    result = await _sanitize_and_validate_tasks(
        {
            "tasks": [{
                "task_id": "ocr_scan",
                "kind": "ocr",
                "mode": "easy",
                "prompt": "OCR the scanned PDF and save evidence.",
                "expected_outputs": ["ocr_evidence.txt"],
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    assert isinstance(result, str)
    data = json.loads(result)
    assert data["ok"] is False
    assert data["error"] == "unsupported_helper_kind"
    assert data["requested_kind"] == "ocr"
    assert "read" in data["allowed_kinds"]
    assert "ocr" not in data["allowed_kinds"]
    assert "helper tools/capabilities" in data["fact"]


async def test_readonly_project_analysis_kind_with_outputs_returns_fact(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    result = await _sanitize_and_validate_tasks(
        {
            "tasks": [{
                "task_id": "project_map",
                "kind": "project_map",
                "mode": "easy",
                "prompt": "Map the project architecture and write _helpers_shared/project_map.txt.",
                "expected_outputs": ["_helpers_shared/project_map.txt"],
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    assert not isinstance(result, str)
    task = result[0]
    assert task["task_id"] == "project_map"
    assert task["kind"] == "project_map"
    assert task["expected_outputs"] == ["_helpers_shared/project_map.txt"]
    observations = task.get("guard_observations") or []
    fact = next(
        item for item in observations
        if item.get("issue") == "read_only_project_analysis_output_conflict"
    )
    assert fact["current_kind"] == "project_map"
    assert fact["expected_outputs"] == ["_helpers_shared/project_map.txt"]
    assert "guard should decide" in fact["details"]


async def test_broad_code_task_is_blocked_before_hard_pairing(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    result = await _sanitize_and_validate_tasks(
        {
            "tasks": [{
                "task_id": "compress_all",
                "kind": "code",
                "mode": "easy",
                "prompt": """
实现统一接口 compression.h、compression.c、Makefile，并一次性完成 6 种压缩算法：
1. Huffman
2. LZ77
3. LZW
4. BWT
5. Arithmetic coding
6. RLE
同时跑 5×5×3 benchmark，输出完整 CSV。
""",
                "expected_outputs": [
                    "compression.h",
                    "compression.c",
                    "huffman.c",
                    "lz77.c",
                    "lzw.c",
                    "results_all.csv",
                ],
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    _assert_sanitized_task_has_broad_fact(result, task_id="compress_all")


def test_delegate_wait_guard_feedback_is_english_first_and_recoverable():
    from pathlib import Path

    src = Path("app/llm/tools/delegate_wait.py").read_text(encoding="utf-8")
    assert "The task-quality guard blocked this delegation." in src
    assert "free-form guard feedback, not a structured program decision" in src
    assert "dispatch_reason" in src
    assert '"error": "guard_blocked"' in src
    assert '"error_kind": "guard_blocked"' in src
    assert "task_too_broad_should_split" not in src
    assert "task_kind_mismatch" not in src
    assert "recovery_plan" not in src
    assert "next_delegate_shape" not in src
    assert "same_base_kind_as_original_boundary" not in src
    assert "**立即拆分并重派**" not in src
    assert "**立即用正确 base kind 重派**" not in src
    assert "铁律" not in src


def test_delegate_spawn_uses_preflight_guard_without_post_spawn_blocking_guard():
    from pathlib import Path

    src = Path("app/llm/tools/delegate_actions.py").read_text(encoding="utf-8")
    assert "_preflight_guard_payload = await _run_delegate_preflight_guard" in src
    assert "delegate.post_spawn_guard.skipped" in src
    assert "persona_guard_{(debug.current_trace_id()" not in src
    assert "no post-spawn blocking guard started" in src


def test_delegate_start_logged_only_after_preflight_passes():
    from pathlib import Path

    sanitize_src = Path("app/llm/tools/delegate.py").read_text(encoding="utf-8")
    actions_src = Path("app/llm/tools/delegate_actions.py").read_text(encoding="utf-8")
    assert '"delegate.start"' not in sanitize_src
    assert actions_src.index("if _preflight_guard_payload is not None:") < actions_src.index('"delegate.start"')


def test_task_quality_guard_allows_benchmark_code_helpers():
    from pathlib import Path

    prompt = Path("app/llm/aux_prompts.py").read_text(encoding="utf-8")
    assert "Code helpers may run benchmarks, compile, debug, compute data" in prompt
    assert "browser-automation evidence that needs Playwright/Puppeteer/Selenium/Chromium-style runtime commands" in prompt
    assert "If the delegation should not run as-is, return should_act=false" in prompt
    assert "runtime will only read `should_act` and `reason`" in prompt


def test_browser_automation_guard_fact_marks_runtime_capability_without_hard_decision():
    from app.llm.tools.delegate_actions import _attach_guard_attention_facts

    tasks = [{
        "task_id": "browser_probe",
        "kind": "code",
        "prompt": (
            "Use Playwright Chromium to visit http://127.0.0.1:5555/, observe the page, "
            "save a screenshot, and report the visible API facts."
        ),
        "acceptance_checks": ["Browser automation evidence is reported before any project edit."],
    }]

    attached = _attach_guard_attention_facts(tasks)
    observations = attached[0].get("guard_observations") or []
    fact = next(
        item for item in observations
        if item.get("issue") == "browser_automation_evidence_capability"
    )
    assert fact["needs_attention"] is False
    assert fact["current_kind"] == "code"
    assert "kind='code' can be appropriate" in fact["details"]


def test_plain_http_fetch_does_not_add_browser_automation_guard_fact():
    from app.llm.tools.delegate_actions import _attach_guard_attention_facts

    tasks = [{
        "task_id": "http_fetch",
        "kind": "read",
        "prompt": "Fetch http://127.0.0.1:5555/ with requests and extract endpoint text.",
        "acceptance_checks": ["Save extracted text."],
    }]

    attached = _attach_guard_attention_facts(tasks)
    observations = attached[0].get("guard_observations") or []
    assert not any(
        item.get("issue") == "browser_automation_evidence_capability"
        for item in observations
    )


async def test_explicit_hard_sibling_pairs_without_extra_twin(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    cleaned = await _sanitize_and_validate_tasks(
        {
            "legacy_pairing_flag": True,
            "tasks": [
                {
                    "task_id": "impl",
                    "kind": "code",
                    "mode": "easy",
                    "prompt": "生成 impl.c 并编译测试。",
                    "expected_outputs": ["impl.c"],
                },
                {
                    "task_id": "impl_hard",
                    "kind": "code",
                    "mode": "hard",
                    "prompt": "生成 impl.c 并编译测试。",
                    "expected_outputs": ["impl.c"],
                },
            ],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    assert not isinstance(cleaned, str)
    by_id = {task["task_id"]: task for task in cleaned}
    assert set(by_id) == {"impl", "impl_hard"}
    assert by_id["impl"]["paired_with"] == "impl_hard"
    assert by_id["impl_hard"]["paired_with"] == "impl"


async def test_broad_code_detection_blocks_explicit_broad_hard_task(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    result = await _sanitize_and_validate_tasks(
        {
            "tasks": [{
                "task_id": "compress_all_hard",
                "kind": "code",
                "mode": "hard",
                "prompt": "实现 6 种压缩算法并输出 compression.h compression.c Makefile results_all.csv。",
                "expected_outputs": [
                    "compression.h",
                    "compression.c",
                    "huffman.c",
                    "lz77.c",
                    "lzw.c",
                    "results_all.csv",
                ],
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    _assert_sanitized_task_has_broad_fact(result, task_id="compress_all_hard")


async def test_database_index_paper_hard_helper_must_split_before_spawn(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    result = await _sanitize_and_validate_tasks(
        {
            "tasks": [{
                "task_id": "classic_analysis",
                "kind": "code",
                "mode": "hard",
                "prompt": (
                    "\u6bd4\u8f83\u7ea2\u9ed1\u6811\u3001\u8df3\u8868\u3001B\u6811\u3001B+\u6811\u7b49\u6570\u636e\u5e93\u7ef4\u62a4\u7b97\u6cd5\u7684\u6027\u80fd\uff0c"
                    "\u5e76\u4e3a\u8bba\u6587\u63d0\u4f9b\u56db\u79cd\u7ecf\u5178\u7b97\u6cd5\u7684\u6bd4\u8f83\u8868\u3001\u5b9e\u9a8c\u7d20\u6750\u548c\u6587\u732e\u5f0f\u5206\u6790\u3002"
                ),
                "expected_outputs": ["classic_algorithm_analysis.md"],
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    fact = _assert_sanitized_task_has_broad_fact(result, task_id="classic_analysis")
    assert fact["signals"]


async def test_hard_pair_preflight_guard_blocks_before_helper_start(monkeypatch, tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import delegate

    async def fake_guard(*args, **kwargs):
        return False, "split project first", [{
            "task_id": "project_refactor",
            "reason": "Too many independent project areas for one helper.",
            "split_into": ["project_inventory", "project_api", "project_tests"],
        }], [], {}

    started: list[str] = []

    async def fake_run_one_helper(*args, **kwargs):
        started.append(kwargs.get("task_id", "?"))
        return {"task_id": kwargs.get("task_id", "?"), "ok": True, "report": "unexpected"}

    monkeypatch.setattr(delegate, "_persona_consent_guard", fake_guard)
    monkeypatch.setattr(delegate, "_run_one_helper", fake_run_one_helper)
    monkeypatch.setattr(delegate, "_detect_missing_unified_framework", lambda cleaned: None)
    workflow_events: list[dict] = []
    monkeypatch.setattr(
        "app.core.environment_events.publish_workflow_event",
        lambda payload: workflow_events.append(payload),
    )

    env = EnvironmentContext(
        root_dir=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
        project_key="project",
    )
    with runtime_context("environment", env):
        result = json.loads(await delegate.handle_delegate(
            str(tmp_path),
            {
                "action": "spawn",
                "wait_window_sec": 0,
                "tasks": [{
                    "task_id": "project_refactor",
                    "kind": "code",
                    "mode": "easy",
                    "prompt": (
                        "Refactor the project under _env/src, update _env/tests, "
                        "refresh _env/README.md, and adjust _env/package.json. "
                        "This is a multi-file project maintenance task with several modules."
                    ),
                    "expected_outputs": [
                        "_env/src/a.py",
                        "_env/src/b.py",
                        "_env/tests/test_a.py",
                    ],
                }],
            },
            archive_id="archive",
            group_id="group",
            user_id="user",
        ))

    assert result["ok"] is False
    assert result["error"] == "guard_blocked"
    assert "split project first" in result["reason"]
    assert result["preflight_guard"] is True
    assert result["helpers_initially_spawned"] == 0
    assert started == []
    assert workflow_events
    assert workflow_events[-1]["kind"] == "helper_blocked"
    assert workflow_events[-1]["blocked_count"] >= 1
    assert any(
        task.get("task_id") == "project_refactor"
        for task in workflow_events[-1]["blocked_tasks"]
    )


async def test_preflight_guard_allows_when_guard_returns_allow_even_with_legacy_split_fields(monkeypatch, tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import delegate

    async def fake_guard(*args, **kwargs):
        return True, "split first", [{
            "task_id": "framework_contract",
            "reason": "The prompt mentions several later algorithms.",
            "split_into": ["contract", "implementations", "paper"],
        }], [], {}

    started: list[str] = []

    async def fake_run_one_helper(*args, **kwargs):
        started.append(kwargs.get("task_id", "?"))
        return {
            "task_id": kwargs.get("task_id", "?"),
            "ok": True,
            "report": "VERDICT: PASS\nDeclared files: contracts/algorithm_study_plan.md",
            "files": ["contracts/algorithm_study_plan.md"],
            "terminal_reason": "completed",
            "outputs_check": {"outputs_complete": True, "quality_warnings": []},
            "_post_helper_action": "output_json_directly",
        }

    monkeypatch.setattr(delegate, "_persona_consent_guard", fake_guard)
    monkeypatch.setattr(delegate, "_run_one_helper", fake_run_one_helper)
    monkeypatch.setattr(delegate, "_detect_missing_unified_framework", lambda cleaned: None)

    env = EnvironmentContext(
        root_dir=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
        project_key="project",
    )
    with runtime_context("environment", env):
        result = json.loads(await delegate.handle_delegate(
            str(tmp_path),
            {
                "action": "spawn",
                "wait_window_sec": 1,
                "tasks": [{
                    "task_id": "framework_contract",
                    "kind": "code",
                    "mode": "easy",
                    "prompt": (
                        "Create a compact shared framework contract for Red-Black Tree, Skip List, "
                        "B-Tree, B+Tree, and a new invented structure. Define schema, evidence map, "
                        "validation checks, and merge order only. Later implementation helpers will use it."
                    ),
                    "expected_outputs": ["contracts/algorithm_study_plan.md"],
                }],
            },
            archive_id="archive",
            group_id="group",
            user_id="user",
        ))

    assert result["ok"] is True
    assert result["helpers_initially_spawned"] == 1
    assert started == ["framework_contract"]


async def test_preflight_guard_allows_single_text_synthesis_when_source_split_is_advice(monkeypatch, tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import delegate

    async def fake_guard(*args, **kwargs):
        return True, "source reading can be separated", [{
            "task_id": "summary_assembly",
            "reason": "Separate source-material reading from final text synthesis.",
            "split_into": ["read_profile", "read_data", "summary_synthesis"],
        }], [], {}

    started: list[str] = []

    async def fake_run_one_helper(*args, **kwargs):
        started.append(kwargs.get("task_id", "?"))
        helper_workspace = Path(str(kwargs.get("helper_workspace") or tmp_path))
        helper_workspace.mkdir(parents=True, exist_ok=True)
        (helper_workspace / "summary.md").write_text("# Summary\n", encoding="utf-8")
        return {
            "task_id": kwargs.get("task_id", "?"),
            "ok": True,
            "report": "VERDICT: PASS\nDeclared files: summary.md",
            "files": ["summary.md"],
            "terminal_reason": "completed",
            "outputs_check": {"outputs_complete": True, "quality_warnings": []},
            "_post_helper_action": "output_json_directly",
        }

    monkeypatch.setattr(delegate, "_persona_consent_guard", fake_guard)
    monkeypatch.setattr(delegate, "_run_one_helper", fake_run_one_helper)
    monkeypatch.setattr(delegate, "_detect_missing_unified_framework", lambda cleaned: None)

    env = EnvironmentContext(
        root_dir=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
        project_key="project",
    )
    with runtime_context("environment", env):
        result = json.loads(await delegate.handle_delegate(
            str(tmp_path),
            {
                "action": "spawn",
                "wait_window_sec": 1,
                "tasks": [{
                    "task_id": "summary_assembly",
                    "kind": "code",
                    "mode": "easy",
                    "prompt": (
                        "Write summary.md from already confirmed profile.yaml and data.json facts. "
                        "The relevant budget, constraint, and source facts are listed inline here."
                    ),
                    "expected_outputs": ["summary.md"],
                }],
            },
            archive_id="archive",
            group_id="group",
            user_id="user",
        ))

    assert result.get("error") != "task_too_broad_should_split"
    assert started == ["summary_assembly"]


async def test_preflight_guard_allows_single_text_synthesis_with_bounded_input_files(monkeypatch, tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import delegate

    for name, content in {
        "profile.yaml": "budget: 100\n",
        "data.json": '{"items":[]}\n',
        "verify_required.py": "print('PASS')\n",
    }.items():
        (tmp_path / name).write_text(content, encoding="utf-8")

    async def fake_guard(*args, **kwargs):
        return True, "source reading can be separated", [{
            "task_id": "artifact_assembly",
            "reason": "Separate source-material reading from final text synthesis.",
            "split_into": ["read_sources", "artifact_synthesis"],
        }], [], {}

    started: list[str] = []

    async def fake_run_one_helper(*args, **kwargs):
        started.append(kwargs.get("task_id", "?"))
        helper_workspace = Path(str(kwargs.get("helper_workspace") or tmp_path))
        helper_workspace.mkdir(parents=True, exist_ok=True)
        (helper_workspace / "artifact.md").write_text("# Artifact\n", encoding="utf-8")
        return {
            "task_id": kwargs.get("task_id", "?"),
            "ok": True,
            "report": "VERDICT: PASS\nDeclared files: artifact.md",
            "files": ["artifact.md"],
            "terminal_reason": "completed",
            "outputs_check": {"outputs_complete": True, "quality_warnings": []},
            "_post_helper_action": "output_json_directly",
        }

    monkeypatch.setattr(delegate, "_persona_consent_guard", fake_guard)
    monkeypatch.setattr(delegate, "_run_one_helper", fake_run_one_helper)
    monkeypatch.setattr(delegate, "_detect_missing_unified_framework", lambda cleaned: None)

    env = EnvironmentContext(
        root_dir=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
        project_key="project",
    )
    with runtime_context("environment", env):
        result = json.loads(await delegate.handle_delegate(
            str(tmp_path),
            {
                "action": "spawn",
                "wait_window_sec": 1,
                "tasks": [{
                    "task_id": "artifact_assembly",
                    "kind": "code",
                    "mode": "easy",
                    "prompt": (
                        "Create artifact.md from the listed source files. Preserve raw field names, "
                        "constraints, and verifier facts from the inputs before computing the final text."
                    ),
                    "input_files": ["profile.yaml", "data.json", "verify_required.py"],
                    "expected_outputs": ["artifact.md"],
                }],
            },
            archive_id="archive",
            group_id="group",
            user_id="user",
        ))

    assert result.get("error") != "task_too_broad_should_split"
    assert started == ["artifact_assembly"]


def test_single_text_source_split_softens_light_embedded_source_facts():
    from app.llm.tools.delegate import _should_soften_source_read_split_for_single_text_output

    task = {
        "task_id": "artifact_assembly",
        "kind": "edit",
        "prompt": (
            "Create artifact.md from the evidence below.\n"
            "1. profile.yaml\n2. source.json\n3. verify_required.py\n"
            "4. constraints\n5. labels\n6. budget facts\n"
        ),
        "expected_outputs": ["artifact.md"],
    }
    rec = {
        "task_id": "artifact_assembly",
        "reason": "6 source material items/groups should be read first",
        "split_into": ["read_sources_batch_1", "read_sources_batch_2"],
    }

    assert _should_soften_source_read_split_for_single_text_output(task, rec) is True


def test_single_text_source_split_keeps_heavy_source_material_block():
    from app.llm.tools.delegate import _should_soften_source_read_split_for_single_text_output

    task = {
        "task_id": "artifact_assembly",
        "kind": "edit",
        "prompt": (
            "Create artifact.md after extracting from source materials:\n"
            "1. a.docx\n2. b.pdf\n3. c.pdf\n4. d.docx\n5. e.pdf\n6. f.docx\n"
        ),
        "expected_outputs": ["artifact.md"],
    }
    rec = {
        "task_id": "artifact_assembly",
        "reason": "6 source material items/groups should be read first",
        "split_into": ["read_sources_batch_1", "read_sources_batch_2"],
    }

    assert _should_soften_source_read_split_for_single_text_output(task, rec) is False


def test_verified_evidence_final_assembly_does_not_emit_source_split_fact():
    from app.llm.tools.delegate import _deterministic_source_read_split_recommendations

    prompt = (
        "Create 4 text files.\n\n"
        "## Classification Evidence (already verified - use these facts)\n"
        "1. **01.txt** - TIER: needs_me.\n"
        "2. **02.txt** - TIER: noise.\n"
        "3. **03.txt** - TIER: noise.\n"
        "4. **04.txt** - TIER: can_wait.\n"
        "5. **05.txt** - TIER: noise.\n"
        "6. **06.txt** - TIER: needs_me.\n"
        "7. **07.txt** - TIER: noise.\n"
        "8. **08.txt** - TIER: noise.\n\n"
        "## Files to create\n"
        "### File 1: `inbox_triage_report.md`\n"
        "### File 2: `draft_client_outage.txt`\n"
        "### File 3: `draft_legal_msa.txt`\n"
        "### File 4: `phishing_flagged.txt`\n"
    )

    recs = _deterministic_source_read_split_recommendations([{
        "task_id": "report_and_drafts",
        "kind": "edit",
        "mode": "easy",
        "prompt": prompt,
        "_source_count_hint": 8,
    }])

    assert recs == []


def test_already_read_material_final_assembly_does_not_emit_source_split_fact():
    from app.llm.tools.delegate import (
        _deterministic_compact_text_owner_observations,
        _deterministic_source_read_split_recommendations,
        _prompt_has_verified_or_embedded_evidence_for_final_assembly,
    )

    prompt = (
        "You have read all 8 source files already. Produce exactly three output files.\n\n"
        "## Source materials already read\n"
        "1. **01.txt** - sender and urgency facts.\n"
        "2. **02.txt** - newsletter facts.\n"
        "3. **03.txt** - suspicious link facts.\n"
        "4. **04.txt** - status update facts.\n"
        "5. **05.txt** - ambiguous request facts.\n"
        "6. **06.txt** - legal deadline facts.\n"
        "7. **07.txt** - meetup facts.\n"
        "8. **08.txt** - recruiter outreach facts.\n\n"
        "## Output file 1: classified_items.txt\n"
        "## Output file 2: summary_report.txt\n"
        "## Output file 3: urgent_reply_draft.txt\n"
    )

    task = {
        "task_id": "material_outputs",
        "kind": "edit",
        "mode": "easy",
        "prompt": prompt,
        "expected_outputs": [
            "classified_items.txt",
            "summary_report.txt",
            "urgent_reply_draft.txt",
        ],
        "_source_count_hint": 8,
    }

    assert _prompt_has_verified_or_embedded_evidence_for_final_assembly(prompt) is True
    assert _deterministic_source_read_split_recommendations([task]) == []
    owner_facts = _deterministic_compact_text_owner_observations([task])
    assert owner_facts and owner_facts[0]["issue"] == "compact_text_owner_shape"


def test_compact_text_owner_shape_fact_reaches_guard_attention():
    from app.llm.tools import delegate_actions
    from app.llm.tools.delegate import _deterministic_compact_text_owner_observations

    tasks = [{
        "task_id": "report_and_drafts",
        "kind": "edit",
        "mode": "easy",
        "prompt": (
            "Use verified evidence below to create final report.md and draft.txt.\n"
            "## Verified Evidence\n"
            "File 1: confirmed fact.\n"
        ),
        "expected_outputs": ["report.md", "draft.txt"],
    }]

    direct = _deterministic_compact_text_owner_observations(tasks)
    assert direct and direct[0]["issue"] == "compact_text_owner_shape"

    attached = delegate_actions._attach_guard_attention_facts([dict(tasks[0])])
    observations = attached[0].get("guard_observations") or []
    assert any(
        item.get("source") == "compact_text_owner_shape_check"
        and item.get("issue") == "compact_text_owner_shape"
        for item in observations
    )


def test_source_material_ref_count_ignores_slash_aliases():
    from app.llm.tools.delegate import _source_material_ref_count

    prompt = "\n".join([
        "- Service A / public label",
        "- Service B / alternate name",
        "- Option C / user-facing alias",
    ])

    assert _source_material_ref_count(prompt) == 0


def test_source_material_ref_count_keeps_explicit_batches():
    from app.llm.tools.delegate import _source_material_ref_count

    prompt = "\n".join([
        "- batch 1: first source group",
        "- folder 2: second source group",
        "- 第三组: 其他材料",
    ])

    assert _source_material_ref_count(prompt) == 3


async def test_broad_source_material_code_task_is_split_before_helper_start(monkeypatch, tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import delegate

    guard_calls: list[list[dict]] = []

    async def permissive_guard(*args, **kwargs):
        guard_calls.append(args[2])
        return True, "allow", [], [], {}

    started: list[str] = []

    async def fake_run_one_helper(*args, **kwargs):
        started.append(kwargs.get("task_id", "?"))
        return {"task_id": kwargs.get("task_id", "?"), "ok": True, "report": "unexpected"}

    monkeypatch.setattr(delegate, "_persona_consent_guard", permissive_guard)
    monkeypatch.setattr(delegate, "_run_one_helper", fake_run_one_helper)
    monkeypatch.setattr(delegate, "_detect_missing_unified_framework", lambda cleaned: None)

    files = "\n".join(f"- _extracted/group{i}/student{i}.docx" for i in range(1, 13))
    env = EnvironmentContext(
        root_dir=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
        project_key="project",
    )
    with runtime_context("environment", env):
        result = json.loads(await delegate.handle_delegate(
            str(tmp_path),
            {
                "action": "spawn",
                "wait_window_sec": 0,
                "tasks": [{
                    "task_id": "extract_reports",
                    "kind": "code",
                    "mode": "hard",
                    "prompt": (
                        "Run a python-docx script to read all source reports and extract author, title, "
                        "summary, and keywords from these Office files:\n"
                        f"{files}\n"
                        "Save internal evidence first, then later work may synthesize a final report."
                    ),
                    "expected_outputs": ["personal_reports_summary.txt"],
                }],
            },
            archive_id="archive",
            group_id="group",
            user_id="user",
        ))

    assert result["ok"] is True
    assert started == ["extract_reports"]
    assert guard_calls
    observations = guard_calls[0][0].get("guard_observations") or []
    assert any(
        isinstance(item, dict)
        and item.get("source") == "deterministic_split_check"
        and item.get("observed_split_boundary_names")
        for item in observations
    )


async def test_wait_window_running_helper_is_not_counted_as_failed(monkeypatch, tmp_path):
    from app.llm.tools import delegate

    started = asyncio.Event()

    async def fake_guard(*args, **kwargs):
        return True, "allow", [], [], {}

    async def fake_run_one_helper(*args, **kwargs):
        started.set()
        task_id = kwargs.get("task_id", "?")
        if task_id == "fast_read":
            return {"task_id": task_id, "ok": True, "report": "fast evidence"}
        await asyncio.sleep(60)
        return {"task_id": task_id, "ok": True, "report": "late"}

    monkeypatch.setattr(delegate, "_persona_consent_guard", fake_guard)
    monkeypatch.setattr(delegate, "_run_one_helper", fake_run_one_helper)
    monkeypatch.setattr(delegate, "_detect_missing_unified_framework", lambda cleaned: None)

    result = json.loads(await delegate.handle_delegate(
        str(tmp_path),
        {
            "action": "spawn",
            "wait_window_sec": 30,
            "tasks": [
                {
                    "task_id": "fast_read",
                    "kind": "read",
                    "mode": "easy",
                    "prompt": "Read one short source file and produce concise evidence.",
                    "expected_outputs": ["fast_read_evidence.txt"],
                },
                {
                    "task_id": "slow_read",
                    "kind": "read",
                    "mode": "easy",
                    "prompt": "Read one slow source file and produce concise evidence.",
                    "expected_outputs": ["slow_read_evidence.txt"],
                },
            ],
        },
        archive_id="archive",
        group_id="group",
        user_id="user",
    ))

    assert started.is_set()
    assert result["helpers_completed"] == 1
    assert result["helpers_still_running"] == 1
    assert result["failed_count"] == 0
    assert result["incomplete_count"] == 0
    assert "task_ok=false" in result["_ok_field_meaning"]
    assert [item["task_id"] for item in result["results"]] == ["fast_read"]


def test_source_material_reading_detector_treats_script_as_method():
    from app.llm.tools.delegate import (
        _deterministic_source_read_split_recommendations,
        _looks_like_source_material_reading,
        _source_material_ref_count,
    )

    prompt = (
        "Use python-docx to read the body text of all reports.\n"
        + "\n".join(f"_extracted/g/student_{i}.docx" for i in range(10))
    )

    assert _looks_like_source_material_reading(prompt)
    assert _source_material_ref_count(prompt) == 10
    recs = _deterministic_source_read_split_recommendations([
        {"task_id": "read_all", "kind": "code", "mode": "hard", "prompt": prompt}
    ])
    assert recs and recs[0]["task_id"] == "read_all"
    assert recs[0]["observed_split_boundary_names"][0].startswith("read_sources_batch_")


def test_source_material_reading_uses_manifest_count_hint():
    from app.llm.tools.delegate import _deterministic_source_read_split_recommendations

    recs = _deterministic_source_read_split_recommendations([
        {
            "task_id": "ielts_synthesis",
            "kind": "code",
            "mode": "easy",
            "prompt": "Read all files in the current directory and organize the source material by IELTS sections.",
            "_source_count_hint": 18,
        }
    ])

    assert recs and recs[0]["task_id"] == "ielts_synthesis"
    assert recs[0]["observed_split_boundary_names"][:2] == ["read_sources_batch_1", "read_sources_batch_2"]


def test_source_material_reading_split_ignores_framework_contract_without_inputs():
    from app.llm.tools.delegate import _deterministic_source_read_split_recommendations

    recs = _deterministic_source_read_split_recommendations([
        {
            "task_id": "framework_contract",
            "kind": "edit",
            "mode": "easy",
            "prompt": (
                "Create a paper framework contract with ten chapters, section ownership, "
                "file naming, acceptance checks, merge order, and final Word validation plan."
            ),
            "expected_outputs": ["framework_contract.md"],
        }
    ])

    assert recs == []


def test_compact_text_material_read_batches_attach_guard_fact():
    from app.llm.tools.delegate import _deterministic_compact_text_bundle_split_observations

    recs = _deterministic_compact_text_bundle_split_observations([
        {
            "task_id": "read_part_a",
            "kind": "read",
            "prompt": "Read these source files and write internal evidence.",
            "input_files": ["notes/a.txt", "notes/b.md", "data/c.json"],
            "expected_outputs": ["part_a_evidence.txt"],
        },
        {
            "task_id": "read_part_b",
            "kind": "read",
            "prompt": "Read these source files and write internal evidence.",
            "input_files": ["notes/d.txt", "notes/e.md"],
            "expected_outputs": ["part_b_evidence.txt"],
        },
    ])

    assert recs
    assert recs[0]["issue"] == "compact_text_material_bundle_split"
    assert recs[0]["observed_split_boundary_names"] == ["read_part_a", "read_part_b"]


def test_compact_text_material_read_batches_ignore_heavy_inputs():
    from app.llm.tools.delegate import _deterministic_compact_text_bundle_split_observations

    recs = _deterministic_compact_text_bundle_split_observations([
        {
            "task_id": "read_part_a",
            "kind": "read",
            "prompt": "Extract source materials.",
            "input_files": ["reports/a.pdf", "reports/b.docx"],
            "expected_outputs": ["part_a_evidence.txt"],
        },
        {
            "task_id": "read_part_b",
            "kind": "read",
            "prompt": "Extract source materials.",
            "input_files": ["reports/c.pdf", "reports/d.docx"],
            "expected_outputs": ["part_b_evidence.txt"],
        },
    ])

    assert recs == []


def test_same_batch_output_overlap_attaches_guard_fact():
    from app.llm.tools.delegate import _deterministic_same_batch_output_overlap_observations

    recs = _deterministic_same_batch_output_overlap_observations([
        {
            "task_id": "easy_impl",
            "kind": "code",
            "mode": "easy",
            "expected_outputs": ["src/result.py"],
        },
        {
            "task_id": "hard_impl",
            "kind": "code",
            "mode": "hard",
            "expected_outputs": ["_env/src/result.py"],
        },
    ])

    assert recs
    assert recs[0]["issue"] == "same_batch_expected_output_overlap"
    assert recs[0]["expected_outputs"] == ["src/result.py"]


async def test_recent_output_overlap_reaches_preflight_guard(monkeypatch, tmp_path):
    from app.core import debug
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import delegate
    from app.llm.tools.delegate_state import _add_to_completion_ledger

    trace_id = "trace_recent_overlap_test"
    previous_trace_id = debug.current_trace_id()
    debug.set_trace_id(trace_id)
    guard_calls: list[list[dict]] = []
    started: list[str] = []

    async def permissive_guard(*args, **kwargs):
        guard_calls.append(args[2])
        return True, "allow", [], [], {}

    async def fake_run_one_helper(*args, **kwargs):
        started.append(kwargs.get("task_id", "?"))
        return {
            "task_id": kwargs.get("task_id", "?"),
            "ok": True,
            "report": "done",
            "files": ["summary.md"],
            "terminal_reason": "completed",
            "outputs_check": {"outputs_complete": True, "quality_warnings": []},
        }

    monkeypatch.setattr(delegate, "_persona_consent_guard", permissive_guard)
    monkeypatch.setattr(delegate, "_run_one_helper", fake_run_one_helper)
    monkeypatch.setattr(delegate, "_detect_missing_unified_framework", lambda cleaned: None)
    await _add_to_completion_ledger(
        trace_id,
        "previous_summary",
        {
            "ok": True,
            "files": ["summary.md"],
            "terminal_reason": "completed",
            "outputs_check": {"outputs_complete": True},
        },
    )

    env = EnvironmentContext(
        root_dir=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
        project_key="project",
    )
    try:
        with runtime_context("environment", env):
            result = json.loads(await delegate.handle_delegate(
                str(tmp_path),
                {
                    "action": "spawn",
                    "wait_window_sec": 0,
                    "tasks": [{
                        "task_id": "new_summary",
                        "kind": "edit",
                        "mode": "easy",
                        "prompt": "Create the final summary from current evidence.",
                        "expected_outputs": ["_env/summary.md"],
                    }],
                },
                archive_id="archive",
                group_id="group",
                user_id="user",
            ))
    finally:
        debug.set_trace_id(previous_trace_id)

    assert result["ok"] is True
    assert started == ["new_summary"]
    assert guard_calls
    observations = guard_calls[0][0].get("guard_observations") or []
    assert any(
        isinstance(item, dict)
        and item.get("source") == "recent_output_overlap_check"
        and item.get("issue") == "recent_helper_expected_output_overlap"
        for item in observations
    )


async def test_ready_artifact_overlap_reaches_preflight_guard(monkeypatch, tmp_path):
    from app.core import agent_state, debug
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import delegate

    trace_id = "trace_ready_artifact_overlap_test"
    previous_trace_id = debug.current_trace_id()
    debug.set_trace_id(trace_id)
    agent_state.reset_trace(trace_id)
    guard_calls: list[list[dict]] = []
    started: list[str] = []

    async def permissive_guard(*args, **kwargs):
        guard_calls.append(args[2])
        return True, "allow", [], [], {}

    async def fake_run_one_helper(*args, **kwargs):
        started.append(kwargs.get("task_id", "?"))
        return {
            "task_id": kwargs.get("task_id", "?"),
            "ok": True,
            "report": "done",
            "files": ["summary.md"],
            "terminal_reason": "completed",
            "outputs_check": {"outputs_complete": True, "quality_warnings": []},
        }

    monkeypatch.setattr(delegate, "_persona_consent_guard", permissive_guard)
    monkeypatch.setattr(delegate, "_run_one_helper", fake_run_one_helper)
    monkeypatch.setattr(delegate, "_detect_missing_unified_framework", lambda cleaned: None)

    agent_state.register_artifact(
        trace_id=trace_id,
        path="summary.md",
        artifact_type="report",
        created_by="main",
        status=agent_state.ARTIFACT_READY,
        verified_by="pytest",
    )

    env = EnvironmentContext(
        root_dir=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
        project_key="project",
    )
    try:
        with runtime_context("environment", env):
            result = json.loads(await delegate.handle_delegate(
                str(tmp_path),
                {
                    "action": "spawn",
                    "wait_window_sec": 0,
                    "tasks": [{
                        "task_id": "ready_artifact_summary",
                        "kind": "edit",
                        "mode": "easy",
                        "prompt": "Revise the final summary if needed.",
                        "expected_outputs": ["_env/summary.md"],
                    }],
                },
                archive_id="archive",
                group_id="group",
                user_id="user",
            ))
    finally:
        agent_state.reset_trace(trace_id)
        debug.set_trace_id(previous_trace_id)

    assert result["ok"] is True
    assert started == ["ready_artifact_summary"]
    assert guard_calls
    observations = guard_calls[0][0].get("guard_observations") or []
    assert any(
        isinstance(item, dict)
        and item.get("source") == "ready_artifact_overlap_check"
        and item.get("issue") == "ready_artifact_expected_output_overlap"
        for item in observations
    )


def test_broad_code_guard_allows_single_framework_contract_file():
    from app.llm.tools.delegate import _detect_broad_code_task_warning_v2

    result = _detect_broad_code_task_warning_v2([{
        "task_id": "framework_contract",
        "kind": "code",
        "mode": "easy",
        "prompt": (
            "Create a shared framework contract and outline for a long report. "
            "Include the document structure, evidence map, validation checks, and merge order.\n"
            "1. title\n2. abstract\n3. background\n4. method\n5. comparison\n6. conclusion\n"
            "## Goal\n## Outline\n## Evidence Map\n## Validation"
        ),
        "expected_outputs": ["_shared/framework_contract.md"],
    }])

    assert result is None


def test_broad_code_guard_allows_framework_contract_task_with_neutral_output_name():
    from app.llm.tools.delegate import _detect_broad_code_task_warning_v2

    result = _detect_broad_code_task_warning_v2([{
        "task_id": "framework_contract",
        "kind": "code",
        "mode": "easy",
        "prompt": (
            "Create a compact shared framework contract for a database index algorithm paper. "
            "Define implementation slices, benchmark schema, evidence map, validation checks, and merge order. "
            "Mention Red-Black Tree, Skip List, B-Tree, B+Tree, and a new invented structure only as later slices."
        ),
        "expected_outputs": ["contracts/algorithm_study_plan.md"],
    }])

    assert result is None


def test_broad_code_guard_allows_framework_contract_with_manifest_output():
    from app.llm.tools.delegate import _detect_broad_code_task_warning_v2

    result = _detect_broad_code_task_warning_v2([{
        "task_id": "framework_contract",
        "kind": "code",
        "mode": "easy",
        "prompt": (
            "Create a compact shared framework contract file for a mixed-language local agent project. "
            "The contract must define canonical directory layout, interface contracts, file ownership, "
            "validation commands, and merge order. The project summary mentions Python backend, "
            "JavaScript browser UI, C native utility, tests, fixtures, docs, root Makefile, and "
            "scripts/check_project.py as later slices. Write the contract and also write a temporary "
            "inventory listing all 16 files with their owning slice. Keep it compact."
        ),
        "expected_outputs": [
            "_helpers_shared/framework_contract.md",
            "_helpers_shared/file_manifest.json",
        ],
        "acceptance_checks": [
            "framework_contract.md covers all sections",
            "file_manifest.json lists files with owning slices",
        ],
    }])

    assert result is None


def test_broad_code_guard_allows_framework_setup_with_global_goal_context():
    from app.llm.tools.delegate import _detect_broad_code_task_warning_v2

    result = _detect_broad_code_task_warning_v2([{
        "task_id": "framework_setup",
        "kind": "code",
        "mode": "easy",
        "prompt": (
            "Create only a compact shared framework contract for the database-index paper project. "
            "Define interfaces, evidence map, validation checks, output ownership, segment boundaries, and merge order. "
            "The overall user goal later includes implementing algorithms, running benchmarks, and generating final docx, "
            "but do not implement algorithms, run benchmarks, or write the final report in this helper. "
            "Later helpers will receive the compact contract through their framework field."
        ),
        "expected_outputs": ["_helpers_shared/framework_setup.md"],
    }])

    assert result is None


def test_broad_code_guard_allows_compact_framework_with_one_validation_script():
    from app.llm.tools.delegate import _detect_broad_code_task_warning_v2

    result = _detect_broad_code_task_warning_v2([{
        "task_id": "framework_contract",
        "kind": "code",
        "mode": "easy",
        "prompt": (
            "Create a compact shared framework contract for a database-index paper project. "
            "Define goal, comparison criteria, paper sections, downstream helper output matrix, "
            "file naming, validation checks, and merge order. Mention Red-Black Tree, Skip List, "
            "B-Tree, B+Tree, and NewStructure only as later bounded slices. "
            "Write one lightweight script that checks the contract file shape. "
            "Do not implement algorithms, run benchmarks, or write paper body."
        ),
        "expected_outputs": [
            "_helpers_shared/framework_contract.md",
            "scripts/check_framework.py",
        ],
        "acceptance_checks": [
            "framework_contract.md names downstream slices",
            "check_framework.py validates the contract shape",
        ],
    }])

    assert result is None


def test_broad_code_guard_blocks_framework_that_also_owns_implementation():
    from app.llm.tools.delegate import _detect_broad_code_task_warning_v2

    result = _detect_broad_code_task_warning_v2([{
        "task_id": "framework_contract",
        "kind": "code",
        "mode": "easy",
        "prompt": (
            "Create the framework contract, implement Red-Black Tree, Skip List, B-Tree, "
            "and B+Tree, run benchmarks, and assemble the final Word paper."
        ),
        "expected_outputs": [
            "_helpers_shared/framework_contract.md",
            "scripts/check_framework.py",
            "src/rbtree.py",
            "src/skiplist.py",
            "benchmark/results.json",
            "paper/DB_Index_Paper.docx",
        ],
    }])

    _assert_broad_code_guard_fact(result, task_id="framework_contract")


def test_prose_framework_code_hard_recommends_edit_easy():
    from app.llm.tools.delegate import _deterministic_kind_recommendations

    recs = _deterministic_kind_recommendations([{
        "task_id": "paper_framework_contract",
        "kind": "code",
        "mode": "hard",
        "prompt": (
            "Create a compact paper framework contract for a database-index algorithm paper. "
            "Include chapter outline, theoretical comparison standards, new-algorithm design notes, "
            "document acceptance checklist, and final Word assembly plan."
        ),
        "expected_outputs": ["_env/contract.md"],
    }])

    # P133: prose-vs-code is no longer a deterministic recommendation — only
    # physical hard constraints (office/binary outputs from non-edit/non-code
    # helpers) survive. A code/hard markdown framework prompt passes cleanly.
    assert recs == []


def test_single_compact_framework_guard_does_not_exempt_hard_or_long_prompt():
    from app.llm.tools.delegate import _is_single_compact_framework_contract_task_for_guard

    base = {
        "task_id": "framework_contract",
        "kind": "edit",
        "mode": "easy",
        "prompt": (
            "Create only a compact shared framework contract. Define outline, evidence map, "
            "validation checks, output ownership, segment boundaries, and merge order."
        ),
        "expected_outputs": ["_env/contract.md"],
    }
    assert _is_single_compact_framework_contract_task_for_guard(base) is True

    hard = dict(base, mode="hard")
    assert _is_single_compact_framework_contract_task_for_guard(hard) is False

    long_prompt = dict(base, prompt=base["prompt"] + "\n" + ("Detailed chapter body. " * 200))
    assert _is_single_compact_framework_contract_task_for_guard(long_prompt) is False


def test_helper_large_text_write_warning_allows_medium_report_write():
    from app.llm.tools.registry import _helper_large_text_write_warning

    content = "# Paper\n\n" + "\n\n".join(
        f"## Section {i}\n" + ("数据库索引算法比较与实验分析。" * 80)
        for i in range(12)
    )

    result = _helper_large_text_write_warning("edit", "analysis/red_black_tree.md", content)

    assert result is None


def test_helper_large_text_write_warning_blocks_near_stream_limit_report():
    from app.llm.tools.registry import _helper_large_text_write_warning

    content = "# Paper\n\n" + "\n\n".join(
        f"## Section {i}\n" + ("数据库索引算法比较与实验分析。" * 140)
        for i in range(24)
    )

    result = _helper_large_text_write_warning("edit", "analysis/red_black_tree.md", content)

    assert result is not None
    assert result["error"] == "helper_large_text_write_should_segment"
    assert result["recovery_action"] == "continue_same_task_segmented"
    assert result["recovery_facts"]["same_task"] is True
    assert "stream safety limit" in result["hint"]
    assert "section files" in result["hint"]
    assert "Output files" in result["hint"]


def test_helper_large_text_write_warning_allows_compact_skeleton():
    from app.llm.tools.registry import _helper_large_text_write_warning

    content = "# Paper Skeleton\n\n## Abstract\n\n## Background\n\n## Method\n\n"

    assert _helper_large_text_write_warning("edit", "analysis/skeleton.md", content) is None


def test_helper_large_text_write_warning_does_not_block_read_evidence():
    from app.llm.tools.registry import _helper_large_text_write_warning

    content = "\n".join(f"source line {i}: evidence" for i in range(500))

    assert _helper_large_text_write_warning("read", "evidence.txt", content) is None


def test_broad_code_guard_returns_split_templates_for_expected_outputs():
    from app.llm.tools.delegate import _detect_broad_code_task_warning_v2

    result = _detect_broad_code_task_warning_v2([{
        "task_id": "algorithm_benchmarks",
        "kind": "code",
        "mode": "hard",
        "prompt": (
            "Implement Red-Black Tree, Skip List, B-Tree, B+Tree, and SSB-Tree. "
            "Benchmark performance and produce a paper-ready comparison.\n"
            "## Red-Black Tree\n## Skip List\n## B-Tree\n## B+Tree\n## SSB-Tree\n"
        ),
        "expected_outputs": [
            "rb_tree.py",
            "skip_list.py",
            "btree.py",
            "bplus_tree.py",
            "benchmark_results.csv",
            "paper.docx",
        ],
    }])

    _assert_broad_code_guard_fact(result, task_id="algorithm_benchmarks")


async def test_greenfield_complete_project_helper_is_not_vertical_slice(tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools.delegate import _detect_broad_code_task_warning_v2

    env = EnvironmentContext(
        root_dir=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
        project_key="project",
    )
    with runtime_context("environment", env):
        result = _detect_broad_code_task_warning_v2([{
            "task_id": "impl_full",
            "kind": "code",
            "mode": "hard",
            "prompt": (
                "Build the complete greenfield mixed agent project from scratch. "
                "The _env/ directory is the project root. Create ALL files listed in the framework. "
                "Every file must work and be complete. After creating all files, run pytest and check_project.py."
            ),
            "framework": "shared project contract",
            "expected_outputs": [
                "_env/README.md",
                "_env/pyproject.toml",
                "_env/src/agent/core.py",
                "_env/src/agent/tools.py",
                "_env/src/agent/server.py",
                "_env/ui/index.html",
                "_env/ui/app.js",
                "_env/native/util.c",
                "_env/native/Makefile",
                "_env/tests/test_core.py",
                "_env/fixtures/sample_input.txt",
                "_env/docs/usage.md",
                "_env/scripts/check_project.py",
            ],
            "acceptance_checks": ["all files exist", "pytest passes", "check_project passes"],
        }])

    _assert_broad_code_guard_fact(result, task_id="impl_full")


async def test_greenfield_small_first_vertical_slice_still_allowed(tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools.delegate import _detect_broad_code_task_warning_v2

    env = EnvironmentContext(
        root_dir=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
        project_key="project",
    )
    with runtime_context("environment", env):
        result = _detect_broad_code_task_warning_v2([{
            "task_id": "first_vertical_slice",
            "kind": "code",
            "mode": "easy",
            "prompt": (
                "Create the first vertical slice for a greenfield project with a tiny core module, "
                "one smoke test, and one check script. Verify with pytest."
            ),
            "framework": "shared project contract",
            "expected_outputs": [
                "_env/src/agent/core.py",
                "_env/tests/test_core.py",
                "_env/scripts/check_project.py",
                "_env/README.md",
            ],
            "acceptance_checks": ["pytest smoke passes", "check_project.py runs"],
        }])

    assert result is None


def test_broad_code_guard_splits_outputs_after_framework_exists():
    from app.llm.tools.delegate import _detect_broad_code_task_warning_v2

    result = _detect_broad_code_task_warning_v2([{
        "task_id": "algorithm_benchmarks",
        "kind": "code",
        "mode": "hard",
        "framework": "Shared interfaces, benchmark schema, output ownership, validation checks, and merge order are fixed.",
        "prompt": (
            "Implement Red-Black Tree, Skip List, B-Tree, B+Tree, and SSB-Tree. "
            "Benchmark performance and produce a paper-ready comparison.\n"
            "## Red-Black Tree\n## Skip List\n## B-Tree\n## B+Tree\n## SSB-Tree\n"
        ),
        "expected_outputs": [
            "rb_tree.py",
            "skip_list.py",
            "btree.py",
            "bplus_tree.py",
            "benchmark_results.csv",
            "paper.docx",
        ],
    }])

    _assert_broad_code_guard_fact(result, task_id="algorithm_benchmarks")


def test_framework_prompts_require_exact_output_matrix_and_non_shared_final_artifacts():
    from app.core.environment_prompt import ENVIRONMENT_PROMPT_ADDON
    from app.llm.tools.delegate_framework import broad_framework_guard_warnings
    from app.llm.tools.helper_prompt_catalog import _HELPER_CONSISTENCY_CONTRACT
    from app.llm.tools.tool_schemas import DELEGATE_TOOL_SCHEMA

    schema_text = json.dumps(DELEGATE_TOOL_SCHEMA, ensure_ascii=False)
    assert "exact output matrix" in schema_text
    assert "`_helpers_shared/...` files are handoff evidence" in schema_text
    assert "user-facing artifacts" in schema_text
    assert "clean non-shared workspace files" in schema_text
    assert "The first framework helper owns only this structural contract" in schema_text
    assert "evidence-backed analysis, research claims, citations" in schema_text
    assert "final numeric values, citations, conclusions" in schema_text
    # ENVIRONMENT_PROMPT_ADDON was compressed (2026-06): same contract, tighter
    # wording — output matrix + bounded slots + producer ownership must remain.
    assert "output matrix" in ENVIRONMENT_PROMPT_ADDON
    assert "merge order" in ENVIRONMENT_PROMPT_ADDON
    assert "It defines bounded slots, not final content" in ENVIRONMENT_PROMPT_ADDON
    assert "belong to producer helpers" in ENVIRONMENT_PROMPT_ADDON
    assert "Define slots, dependencies, ownership, and acceptance rather than filling those slots" in _HELPER_CONSISTENCY_CONTRACT

    warnings = broad_framework_guard_warnings([{
        "task_id": "framework_contract",
        "kind": "code",
        "mode": "hard",
        "prompt": "Build a framework and implement several algorithms, benchmarks, and final document.",
        "expected_outputs": ["framework.md", "analysis.md", "paper.docx"],
    }])

    joined = json.dumps(warnings, ensure_ascii=False)
    assert "downstream output matrix" in joined
    assert "_helpers_shared/...` is handoff evidence" in joined
    assert "slots, dependencies, acceptance, and downstream output matrix" in joined
    assert "ownership boundaries" in joined
    assert "implementation, experiment, evidence, chart" in joined
    assert "final-document work" in joined


def test_broad_code_guard_allows_split_greenfield_framework_fanout():
    from app.llm.tools.delegate import _detect_broad_code_task_warning_v2

    framework = (
        "Project contract: small local agent. Components: core, ui, native, tests, docs, scripts. "
        "Each helper owns only its listed _env paths and must run its local acceptance checks."
    )
    tasks = [
        {
            "task_id": "core_backend",
            "kind": "code",
            "mode": "easy",
            "framework": framework,
            "prompt": (
                "Implement the Python backend/core module only. Use stdlib, clean APIs, "
                "and verify imports plus demo execution."
            ),
            "expected_outputs": [
                "_env/core/__init__.py",
                "_env/core/agent.py",
                "_env/core/tools.py",
                "_env/core/llm.py",
                "_env/core/memory.py",
                "_env/run.py",
                "_env/requirements.txt",
            ],
            "acceptance_checks": [
                "All files exist",
                "python import check passes",
                "python _env/run.py prints output",
            ],
        },
        {
            "task_id": "ui_frontend",
            "kind": "code",
            "mode": "easy",
            "framework": framework,
            "prompt": "Create only the vanilla browser UI and mock fallback.",
            "expected_outputs": ["_env/ui/index.html", "_env/ui/app.js", "_env/ui/style.css"],
            "acceptance_checks": ["HTML references JS/CSS", "JS has fetch and mock fallback"],
        },
        {
            "task_id": "native_util",
            "kind": "code",
            "mode": "easy",
            "framework": framework,
            "prompt": "Create only the C tokenizer utility.",
            "expected_outputs": ["_env/native/tokenizer.c", "_env/native/tokenizer.h", "_env/native/Makefile"],
            "acceptance_checks": ["C files exist", "gcc smoke compile command is named"],
        },
        {
            "task_id": "tests_fixtures",
            "kind": "code",
            "mode": "easy",
            "framework": framework,
            "prompt": "Create only pytest tests and JSON fixture for the existing API contract.",
            "expected_outputs": [
                "_env/tests/__init__.py",
                "_env/tests/test_agent.py",
                "_env/tests/test_tools.py",
                "_env/tests/test_memory.py",
                "_env/fixtures/sample_chat.json",
            ],
            "acceptance_checks": ["pytest files exist", "fixture is valid JSON"],
        },
        {
            "task_id": "docs_readme",
            "kind": "code",
            "mode": "easy",
            "framework": framework,
            "prompt": "Create only the project README.",
            "expected_outputs": ["_env/docs/README.md"],
            "acceptance_checks": ["README has quickstart and validation commands"],
        },
        {
            "task_id": "check_script",
            "kind": "code",
            "mode": "easy",
            "framework": framework,
            "prompt": "Create only the project validation script.",
            "expected_outputs": ["_env/scripts/check_project.py"],
            "acceptance_checks": ["script exists", "script validates expected project paths"],
        },
    ]

    assert _detect_broad_code_task_warning_v2(tasks) is None


def test_broad_code_guard_suggests_contract_first_for_framework_pipeline():
    from app.llm.tools.delegate import _detect_broad_code_task_warning_v2

    result = _detect_broad_code_task_warning_v2([{
        "task_id": "framework_init",
        "kind": "code",
        "mode": "easy",
        "prompt": (
            "Create a framework contract, index interface, benchmark schema, and paper outline for "
            "Red-Black Tree, Skip List, B-Tree, B+Tree, and a new invented database index algorithm."
        ),
        "expected_outputs": [
            "contracts/framework.md",
            "src/index_interface.py",
            "benchmarks/schema.json",
            "paper/outline.md",
        ],
    }])

    _assert_broad_code_guard_fact(result, task_id="framework_init")


def test_broad_code_guard_suggests_contract_first_for_scaffold_pipeline():
    from app.llm.tools.delegate import _detect_broad_code_task_warning_v2

    result = _detect_broad_code_task_warning_v2([{
        "task_id": "scaffold",
        "kind": "code",
        "mode": "easy",
        "framework": "",
        "prompt": (
            "Create the scaffold for a database-index research project. It will compare red-black tree, skip list, "
            "B-tree, and B+ tree implementations, run benchmarks, and later support a rigorous paper and docx report."
        ),
        "expected_outputs": [
            "_env/db_index_project/src/rbtree.py",
            "_env/db_index_project/src/skiplist.py",
            "_env/db_index_project/src/btree.py",
            "_env/db_index_project/src/bptree.py",
            "_env/db_index_project/README.md",
        ],
    }])

    _assert_broad_code_guard_fact(result, task_id="scaffold")


def test_broad_code_guard_maps_shared_outputs_to_helper_shared():
    from app.llm.tools.delegate import _detect_broad_code_task_warning_v2

    result = _detect_broad_code_task_warning_v2([{
        "task_id": "shared_batch",
        "kind": "code",
        "mode": "easy",
        "prompt": "Implement A, B, C, D and benchmark.\n## A\n## B\n## C\n## D\n",
        "expected_outputs": [
            "_shared/framework_contract.json",
            "a.py",
            "b.py",
            "c.py",
            "benchmark.csv",
        ],
    }])

    _assert_broad_code_guard_fact(result, task_id="shared_batch")


def test_removed_general_framework_kind_does_not_get_legacy_recommendation():
    from app.llm.tools.delegate import _deterministic_kind_recommendations

    recs = _deterministic_kind_recommendations([{
        "task_id": "framework_builder",
        "kind": "general",
        "mode": "hard",
        "prompt": "Create the shared framework contract, schema, outline, evidence map, and validation checks for later peer helpers.",
        "expected_outputs": ["_helpers_shared/paper_framework.txt"],
    }])

    assert recs == []


async def test_read_fanout_framework_block_suggests_delegate_contract_helper(monkeypatch, tmp_path):
    from app.llm.tools import delegate

    async def fake_guard(*args, **kwargs):
        return False, "needs framework", [], [], {
            "block": True,
            "task_ids": ["read_pptx", "read_xlsx", "read_reports"],
            "reason": "Broad multi-part reading needs a shared evidence schema.",
        }

    started: list[str] = []

    async def fake_run_one_helper(*args, **kwargs):
        started.append(kwargs.get("task_id", "?"))
        return {"task_id": kwargs.get("task_id", "?"), "ok": True, "report": "unexpected"}

    monkeypatch.setattr(delegate, "_persona_consent_guard", fake_guard)
    monkeypatch.setattr(delegate, "_run_one_helper", fake_run_one_helper)

    monkeypatch.setattr(delegate, "_detect_missing_unified_framework", lambda cleaned: None)

    result = json.loads(await delegate.handle_delegate(
        str(tmp_path),
        {
            "action": "spawn",
            "wait_window_sec": 0,
            "tasks": [
                {
                    "task_id": "read_pptx",
                    "kind": "read",
                    "prompt": "Read three pptx files and write _helpers_shared/pptx_evidence.txt.",
                    "expected_outputs": ["_helpers_shared/pptx_evidence.txt"],
                },
                {
                    "task_id": "read_xlsx",
                    "kind": "read",
                    "prompt": "Read xlsx files and write _helpers_shared/xlsx_evidence.txt.",
                    "expected_outputs": ["_helpers_shared/xlsx_evidence.txt"],
                },
                {
                    "task_id": "read_reports",
                    "kind": "read",
                    "prompt": "Read docx reports and write _helpers_shared/report_evidence.txt.",
                    "expected_outputs": ["_helpers_shared/report_evidence.txt"],
                },
            ],
        },
        archive_id="archive",
        group_id="group",
        user_id="user",
    ))

    assert result["ok"] is False
    assert result["error_kind"] == "guard_blocked"
    assert result["helpers_initially_spawned"] == 0
    assert "needs framework" in result["reason"]
    assert started == []


async def test_markdown_report_sanitize_corrects_before_guard(tmp_path):
    from app.llm.tools import delegate

    result = await delegate._sanitize_and_validate_tasks(
        {
            "tasks": [{
                "task_id": "project_report",
                "kind": "code",
                "mode": "easy",
                "prompt": (
                    "Create a final project report under _env/docs/report.md from _env/src, "
                    "_env/tests, _env/README.md, and _env/package.json. This is a multi-file project task."
                ),
                "expected_outputs": [
                    "_env/docs/report.md",
                    "_env/docs/summary.md",
                    "_env/docs/changelog.md",
                    "_env/docs/api.md",
                    "_env/docs/testing.md",
                ],
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    _assert_sanitized_task_has_broad_fact(result, task_id="project_report")


async def test_explicit_hard_helper_preflight_guard_blocks_before_helper_start(monkeypatch, tmp_path):
    from app.llm.tools import delegate

    async def fake_guard(*args, **kwargs):
        return False, "split first", [{
            "task_id": "graph_algos",
            "reason": "Implementation, benchmark, tests, README, data, and report are separable work areas.",
            "split_into": ["graph_core", "graph_benchmark", "graph_tests", "graph_docs"],
        }], [], {}

    started: list[str] = []

    async def fake_run_one_helper(*args, **kwargs):
        started.append(kwargs.get("task_id", "?"))
        return {"task_id": kwargs.get("task_id", "?"), "ok": True, "report": "unexpected"}

    monkeypatch.setattr(delegate, "_persona_consent_guard", fake_guard)
    monkeypatch.setattr(delegate, "_run_one_helper", fake_run_one_helper)
    monkeypatch.setattr(delegate, "_detect_missing_unified_framework", lambda cleaned: None)

    result = json.loads(await delegate.handle_delegate(
        str(tmp_path),
        {
            "action": "spawn",
            "wait_window_sec": 0,
            "tasks": [{
                "task_id": "graph_algos",
                "kind": "code",
                "mode": "hard",
                "prompt": (
                    "Implement Dijkstra, A*, Floyd-Warshall, benchmark.py, tests, README, "
                    "data/edges.csv, and docs/algorithm_report.md; run verification."
                ),
                "expected_outputs": [
                    "src/algolab/graph.py",
                    "src/algolab/benchmark.py",
                    "tests/test_graph.py",
                    "README.md",
                    "data/edges.csv",
                    "docs/algorithm_report.md",
                ],
            }],
        },
        archive_id="archive",
        group_id="group",
        user_id="user",
    ))

    assert result["ok"] is False
    assert result["error"] == "guard_blocked"
    assert "split first" in result["reason"]
    assert started == []


async def test_guard_intervention_ignores_extra_guard_tuple_fields_when_allowed(monkeypatch):
    from app.llm.tools import delegate_wait

    killed: list[str] = []

    class FakeRegistry:
        async def cancel_all_helpers_in_trace(self, trace_id: str) -> int:
            killed.append(trace_id)
            return 0

    monkeypatch.setattr(delegate_wait, "_sync_delegate_action_globals", lambda: None)
    monkeypatch.setattr(delegate_wait, "_sync_delegate_globals", lambda: None)
    monkeypatch.setattr(delegate_wait, "proc_registry", lambda: FakeRegistry())

    payload = await delegate_wait._build_guard_intervention(
        (
            True,
            "split first",
            [{
                "task_id": "large_task",
                "reason": "Needs smaller helpers.",
                "split_into": ["part_a", "part_b"],
            }],
            [],
            {},
            {"future_extension": True},
        ),
        trace_id="trace_extra_fields",
        cancel_helpers=True,
        helper_specs=[],
    )

    assert payload is None
    assert killed == []


async def test_framework_guard_allows_second_attempt_in_same_trace(monkeypatch, tmp_path):
    from app.llm.tools import delegate

    async def fake_guard(*args, **kwargs):
        return True, "needs framework", [], [], {
            "block": True,
            "task_ids": ["impl_rbtree", "impl_skiplist", "impl_btree"],
            "reason": "Need shared benchmark framework.",
        }

    started: list[str] = []

    async def fake_run_one_helper(*args, **kwargs):
        started.append(kwargs.get("task_id", "?"))
        return {
            "task_id": kwargs.get("task_id", "?"),
            "ok": True,
            "report": "done",
            "terminal_reason": "completed",
            "outputs_check": {"outputs_complete": True},
        }

    monkeypatch.setattr(delegate, "_persona_consent_guard", fake_guard)
    monkeypatch.setattr(delegate, "_run_one_helper", fake_run_one_helper)
    monkeypatch.setattr(delegate.debug, "current_trace_id", lambda: "trace_framework_loop_release")
    delegate._guard_framework_block_trace_total["trace_framework_loop_release"] = 1

    try:
        result = json.loads(await delegate.handle_delegate(
            str(tmp_path),
            {
                "action": "spawn",
                "wait_window_sec": 0,
                "tasks": [
                    {
                        "task_id": "impl_rbtree",
                        "kind": "code",
                        "mode": "hard",
                        "prompt": "写一个自包含单文件程序，内嵌统一CSV基准测试。",
                        "expected_outputs": ["benchmark_rbtree.csv"],
                    },
                    {
                        "task_id": "impl_skiplist",
                        "kind": "code",
                        "mode": "hard",
                        "prompt": "写一个自包含单文件程序，内嵌统一CSV基准测试。",
                        "expected_outputs": ["benchmark_skiplist.csv"],
                    },
                    {
                        "task_id": "impl_btree",
                        "kind": "code",
                        "mode": "hard",
                        "prompt": "写一个自包含单文件程序，内嵌统一CSV基准测试。",
                        "expected_outputs": ["benchmark_btree.csv"],
                    },
                ],
            },
            archive_id="archive",
            group_id="group",
            user_id="user",
        ))
    finally:
        delegate._guard_framework_block_trace_total.pop("trace_framework_loop_release", None)

    assert result["ok"] is True
    assert result["helpers_initially_spawned"] == 3
    assert set(started) == {"impl_rbtree", "impl_skiplist", "impl_btree"}


async def test_framework_guard_allows_read_peer_batch_with_explicit_framework(monkeypatch, tmp_path):
    from app.llm.tools import delegate

    async def fake_guard(*args, **kwargs):
        return True, "needs framework", [], [], {
            "block": True,
            "task_ids": ["rb_analysis", "skip_analysis", "btree_analysis"],
            "reason": "Comparable analysis needs a shared framework.",
        }

    started: list[str] = []

    async def fake_run_one_helper(*args, **kwargs):
        started.append(kwargs.get("task_id", "?"))
        return {
            "task_id": kwargs.get("task_id", "?"),
            "ok": True,
            "report": "VERDICT: PASS",
            "terminal_reason": "completed",
            "outputs_check": {"outputs_complete": True},
        }

    framework = (
        "共享分析框架：目标是比较红黑树、跳表、B树并形成论文证据。"
        "每个分析必须覆盖数据结构定义、核心操作伪代码、时间复杂度、空间复杂度、"
        "缓存效率、并发特性、范围查询、磁盘友好性、实际应用场景等比较维度。"
        "输出为 UTF-8 中文 evidence txt 文件。验收包括维度完整、术语一致、可合并。"
    )

    monkeypatch.setattr(delegate, "_persona_consent_guard", fake_guard)
    monkeypatch.setattr(delegate, "_run_one_helper", fake_run_one_helper)
    monkeypatch.setattr(delegate.debug, "current_trace_id", lambda: "trace_read_framework_present")

    result = json.loads(await delegate.handle_delegate(
        str(tmp_path),
        {
            "action": "spawn",
            "wait_window_sec": 0,
            "tasks": [
                {
                    "task_id": "rb_analysis",
                    "kind": "read",
                    "mode": "hard",
                    "framework": framework,
                    "prompt": "按共享框架分析红黑树，保存 rb_analysis.txt。",
                    "expected_outputs": ["rb_analysis.txt"],
                },
                {
                    "task_id": "skip_analysis",
                    "kind": "read",
                    "mode": "hard",
                    "framework": framework,
                    "prompt": "按共享框架分析跳表，保存 skip_analysis.txt。",
                    "expected_outputs": ["skip_analysis.txt"],
                },
                {
                    "task_id": "btree_analysis",
                    "kind": "read",
                    "mode": "hard",
                    "framework": framework,
                    "prompt": "按共享框架分析B树，保存 btree_analysis.txt。",
                    "expected_outputs": ["btree_analysis.txt"],
                },
            ],
        },
        archive_id="archive",
        group_id="group",
        user_id="user",
    ))

    assert result["ok"] is True
    assert result["helpers_initially_spawned"] == 3
    assert set(started) == {"rb_analysis", "skip_analysis", "btree_analysis"}


async def test_broad_code_guard_allows_single_algorithm_with_tests_and_readme(tmp_path):
    from app.llm.tools.delegate import _sanitize_and_validate_tasks

    prompt = (
        "Implement weighted-graph Dijkstra shortest paths.\n\n"
        "## Read first\n"
        "Inspect src/algolab/graph.py, tests/test_graph.py, and README.md.\n\n"
        "## Implementation\n"
        "Add one Dijkstra implementation with path reconstruction and negative-weight rejection.\n\n"
        "## Tests\n"
        "Add 3-5 pytest cases covering a normal path, unreachable nodes, and negative weights.\n\n"
        "## Documentation\n"
        "Update the README usage snippet.\n\n"
        "## Verification\n"
        "Run compile, pytest, and one smoke command before replying.\n"
    )

    cleaned = await _sanitize_and_validate_tasks(
        {
            "tasks": [{
                "task_id": "dijkstra_impl",
                "kind": "code",
                "mode": "easy",
                "prompt": prompt,
                "expected_outputs": [
                    "src/algolab/graph.py",
                    "tests/test_graph.py",
                    "README.md",
                ],
            }],
        },
        main_workspace=str(tmp_path),
        archive_id="archive",
        group_id="group",
        user_id="user",
    )

    assert not isinstance(cleaned, str)
    assert cleaned[0]["task_id"] == "dijkstra_impl"
    assert "paired_with" not in cleaned[0]


def test_delegate_cleanup_preserves_underscore_task_ids(tmp_path):
    from app.core.delegate_cleanup import (
        _delegate_dir_tag,
        _delegate_task_id,
        cleanup_inactive_delegate_dirs,
    )

    extra_dir = tmp_path / "_delegate_user_extra_task"
    active_dir = tmp_path / "_delegate_user_compress_all"
    old_dir = tmp_path / "_delegate_user_old_task"
    stale_dir = tmp_path / "_delegate_user_stale_task"
    for path in (extra_dir, active_dir, old_dir, stale_dir):
        path.mkdir()

    assert _delegate_dir_tag(active_dir.name) == "user_compress"
    assert _delegate_task_id(active_dir.name, current_tag="user") == "compress_all"

    cleaned = cleanup_inactive_delegate_dirs(
        str(tmp_path),
        "user",
        active_task_ids=["compress_all"],
        keep_resume_task_ids=["old_task"],
        max_keep=1,
    )

    assert cleaned == 1
    assert active_dir.exists()
    assert old_dir.exists()
    assert stale_dir.exists()
    assert not extra_dir.exists()


def test_legacy_paired_hard_task_detection_accepts_numbered_suffix():
    from app.llm.tools.delegate import _is_legacy_paired_hard_task

    assert _is_legacy_paired_hard_task({"task_id": "impl_hard_2", "mode": "hard"}) is True
    assert _is_legacy_paired_hard_task({"task_id": "impl_hard_2", "mode": "easy"}) is False
    assert _is_legacy_paired_hard_task({"task_id": "impl_final_2", "mode": "hard"}) is True


async def test_legacy_spawn_pairing_flag_blocks_completed_primary_and_hard_pair(monkeypatch, tmp_path):
    from app.core import core_processes
    from app.core.core_processes import ProcessRegistry
    from app.llm.tools import delegate

    fake_registry = ProcessRegistry()
    await fake_registry.mark_recently_completed("impl")
    await fake_registry.mark_recently_completed("impl_hard")
    monkeypatch.setattr(core_processes, "_registry", fake_registry)

    called = False

    async def fake_run_one_helper(**kwargs):
        nonlocal called
        called = True
        return {"task_id": kwargs["task_id"], "ok": True, "files": []}

    monkeypatch.setattr(delegate, "_run_one_helper", fake_run_one_helper)

    result = json.loads(await delegate.handle_delegate(
        str(tmp_path),
        {
            "action": "spawn",
            "legacy_pairing_flag": True,
            "tasks": [{
                "task_id": "impl",
                "kind": "code",
                "mode": "easy",
                "prompt": "生成 impl.c。",
                "expected_outputs": ["impl.c"],
            }],
        },
        archive_id="archive",
        group_id="group",
        user_id="user",
    ))

    assert result["ok"] is True
    assert result["already_completed"] is True
    assert set(result["duplicate_task_ids"]) == {"impl"}
    assert called is False


async def test_legacy_spawn_pairing_flag_runs_explicit_hard_task(monkeypatch, tmp_path):
    from app.llm.tools import delegate
    from app.llm.tools import delegate_actions

    captured = {}

    async def fake_run_one_helper(**kwargs):
        captured["kwargs"] = kwargs
        return {
            "task_id": kwargs["task_id"],
            "ok": True,
            "report": "done",
            "files": [],
            "terminal_reason": "completed",
        }

    monkeypatch.setattr(delegate, "_run_one_helper", fake_run_one_helper)

    async def fake_guard(*args, **kwargs):
        return True, "test guard pass", [], []

    monkeypatch.setattr(delegate, "_persona_consent_guard", fake_guard)
    monkeypatch.setattr(delegate_actions, "_persona_consent_guard", fake_guard, raising=False)

    result = json.loads(await delegate.handle_delegate(
        str(tmp_path),
        {
            "action": "spawn",
            "wait_window_sec": 0,
            "legacy_pairing_flag": True,
            "tasks": [{
                "task_id": "impl_hard",
                "kind": "code",
                "mode": "hard",
                "prompt": "生成 impl.c。",
                "expected_outputs": ["impl.c"],
            }],
        },
        archive_id="archive",
        group_id="group",
        user_id="user",
    ))

    assert result["ok"] is True
    assert captured["kwargs"]["task_id"] == "impl_hard"
    assert captured["kwargs"]["mode"] == "hard"


def test_copy_results_maps_declared_docx_with_helper_prefixed_name(tmp_path):
    from app.llm.tools.delegate import _copy_results_to_main

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    helper_ws.mkdir()
    main_ws.mkdir()
    (helper_ws / "gen_docx_通信原理作业_第五章.docx").write_bytes(b"PK\x03\x04fake-docx")

    copied, stats, file_map = _copy_results_to_main(
        str(helper_ws),
        str(main_ws),
        task_id="supplement_521",
        declared_files={"通信原理作业_第五章.docx"},
    )

    assert copied == ["supplement_521_gen_docx_通信原理作业_第五章.docx"]
    assert stats["capped"] is False
    assert file_map[0]["helper_name"] == "gen_docx_通信原理作业_第五章.docx"


def test_copy_results_treats_declared_env_outputs_as_satisfied_and_skips_test_cache(tmp_path):
    from app.llm.tools.delegate import _copy_results_to_main

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    source_dir = helper_ws / "_env" / "contracts"
    cache_dir = helper_ws / "_env" / ".pytest_cache" / "v" / "cache"
    source_dir.mkdir(parents=True)
    cache_dir.mkdir(parents=True)
    main_ws.mkdir()
    (source_dir / "customer_event.py").write_text(
        "def validate_event(payload):\n    return payload['account_name']\n",
        encoding="utf-8",
    )
    (cache_dir / "nodeids").write_text("contracts/tests/test_schema.py::test_ok\n", encoding="utf-8")

    copied, stats, file_map = _copy_results_to_main(
        str(helper_ws),
        str(main_ws),
        task_id="migrate-field",
        declared_files={"_env/contracts/customer_event.py"},
        expected_outputs=["_env/contracts/customer_event.py"],
        helper_kind="code",
    )

    assert copied == ["_env/contracts/customer_event.py"]
    assert stats["declared_satisfied_by_existing_copy"] == ["_env/contracts/customer_event.py"]
    assert stats["env_skipped_unexpected_new"] == []
    assert not (main_ws / "_env" / ".pytest_cache").exists()
    assert file_map == [{
        "helper_name": "contracts/customer_event.py",
        "main_name": "_env/contracts/customer_event.py",
        "shared_name": None,
    }]


def test_delegate_report_uses_main_workspace_name_for_prefixed_copy():
    visible_copied_back = ["assemble_docx_db_index_paper.docx"]
    file_map = [{
        "helper_name": "db_index_paper.docx",
        "main_name": "assemble_docx_db_index_paper.docx",
        "shared_name": None,
    }]

    helper_alias_by_main = {
        str(m.get("main_name")): str(m.get("helper_name"))
        for m in file_map
        if isinstance(m, dict)
        and m.get("main_name")
        and m.get("helper_name")
        and str(m.get("main_name")) != str(m.get("helper_name"))
    }

    def displayed_file_list(paths):
        displayed = []
        seen = set()
        for path in paths:
            name = str(path or "").replace("\\", "/").strip()
            if name in seen:
                continue
            seen.add(name)
            displayed.append(name)
        return displayed

    display = displayed_file_list(visible_copied_back)
    aliases = {
        main_name: helper_name
        for main_name, helper_name in helper_alias_by_main.items()
        if main_name in set(visible_copied_back)
    }

    assert display == ["assemble_docx_db_index_paper.docx"]
    assert aliases == {"assemble_docx_db_index_paper.docx": "db_index_paper.docx"}


def test_copy_results_merges_nested_helpers_shared_outputs(tmp_path):
    from app.llm.tools.delegate import _copy_results_to_main

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    nested = helper_ws / "_helpers_shared" / "hw_redraw_charts"
    nested.mkdir(parents=True)
    main_ws.mkdir()
    (nested / "q5_2_block_diagram.png").write_bytes(b"\x89PNG\r\n\x1a\n" + b"x" * 3000)

    copied, stats, file_map = _copy_results_to_main(
        str(helper_ws),
        str(main_ws),
        task_id="hw_redraw_charts",
        declared_files={"q5_2_block_diagram.png"},
    )

    assert copied == ["_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png"]
    assert (main_ws / "_helpers_shared" / "hw_redraw_charts" / "q5_2_block_diagram.png").is_file()
    assert stats["capped"] is False
    assert file_map == []


def test_extract_declared_files_preserves_helpers_shared_paths():
    from app.llm.tools.delegate import _extract_declared_files

    report = '```json\n{"files": ["_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png"]}\n```'

    assert _extract_declared_files(report) == {"_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png"}


async def test_helper_result_summary_treats_shared_declared_outputs_as_delivered(tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    nested = helper_ws / "_helpers_shared" / "hw_redraw_charts"
    nested.mkdir(parents=True)
    main_ws.mkdir()
    (nested / "q5_2_block_diagram.png").write_bytes(b"\x89PNG\r\n\x1a\n" + b"x" * 3000)
    (helper_ws / ".helper_hw_redraw_charts_full_report.txt").write_text(
        '```json\n{"files": ["_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png"]}\n```',
        encoding="utf-8",
    )

    async def fake_loop(*args, **kwargs):
        nested.mkdir(parents=True, exist_ok=True)
        (nested / "q5_2_block_diagram.png").write_bytes(b"\x89PNG\r\n\x1a\n" + b"y" * 3000)
        return '```json\n{"files": ["_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png"]}\n```', []

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))
    try:
        result = await delegate._run_one_helper(
            task_id="hw_redraw_charts",
            prompt="生成图",
            main_workspace=str(main_ws),
            helper_workspace=str(helper_ws),
            archive_id="archive",
            group_id="group",
            user_id="user",
            resume=False,
            local_abort=asyncio.Event(),
            wait_for_register=asyncio.Event(),
            user_lang="zh",
            kind="draw",
            mode="easy",
            helper_think=False,
            expected_outputs=["q5_2_block_diagram.png"],
        )
    finally:
        monkeypatch.undo()

    assert result["files"] == ["_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png"]
    assert result["declared_files"] == ["_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png"]
    assert result["declared_but_missing"] == []
    assert result["delivered_but_not_declared"] == []
    assert result["outputs_check"]["outputs_missing"] == []
    assert result["outputs_check"]["outputs_complete"] is True


async def test_helper_accepts_machine_readable_output_files_report_without_repair(tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    helper_ws.mkdir(parents=True)
    main_ws.mkdir()
    calls: list[str] = []

    async def fake_loop(msgs, *args, **kwargs):
        calls.append("\n".join(str(m.get("content") or "") for m in msgs if isinstance(m, dict)))
        (helper_ws / "_env").mkdir(exist_ok=True)
        (helper_ws / "_env" / "analysis.md").write_text("verified analysis", encoding="utf-8")
        return 'Output files: {"files": ["_env/analysis.md"]}', msgs

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))
    try:
        result = await delegate._run_one_helper(
            task_id="analysis",
            prompt="Write analysis.md.",
            main_workspace=str(main_ws),
            helper_workspace=str(helper_ws),
            archive_id="archive",
            group_id="group",
            user_id="user",
            resume=False,
            local_abort=asyncio.Event(),
            wait_for_register=asyncio.Event(),
            user_lang="en",
            kind="code",
            mode="easy",
            helper_think=False,
            expected_outputs=["_env/analysis.md"],
        )
    finally:
        monkeypatch.undo()

    assert len(calls) == 1
    assert result["ok"] is True
    assert result["terminal_reason"] == "completed"
    assert result["outputs_check"]["outputs_complete"] is True
    assert "_env/analysis.md" in result["main_available_files"]
    assert "_env/analysis.md" in result["staged_project_files"]
    assert "Staged project files available" in result["report"]
    assert "project tests run before apply validate the old project state" in result["report"]
    assert "project validation commands run before env_apply_* validate the old real project state" in result["pending_project_apply_fact"]
    assert "env_diff is the compact first inspection" in result["post_helper_usage_hint"]
    assert "do not use read_file to re-check helper-owned content" in result["post_helper_usage_hint"]
    assert "post-apply test validates the staged edit" in result["post_helper_usage_hint"]


def test_code_repair_expected_outputs_do_not_keep_renamed_source_inputs():
    from app.llm.tools.delegate import _augment_code_repair_expected_outputs

    augmented = _augment_code_repair_expected_outputs(
        kind="code",
        prompt=(
            "Update the contract implementation. Rename `contracts/customer_event.py` "
            "to `contracts/account_event.py`, update imports, and edit service/render.py."
        ),
        input_files=[
            "contracts/customer_event.py",
            "contracts/tests/test_schema.py",
            "service/render.py",
        ],
        expected_outputs=[
            "contracts/account_event.py",
            "contracts/tests/test_schema.py",
            "service/render.py",
        ],
    )

    assert "_env/contracts/customer_event.py" not in augmented
    assert augmented == [
        "contracts/account_event.py",
        "contracts/tests/test_schema.py",
        "service/render.py",
    ]


def test_code_repair_expected_outputs_keep_edited_source_inputs():
    from app.llm.tools.delegate import _augment_code_repair_expected_outputs

    augmented = _augment_code_repair_expected_outputs(
        kind="code",
        prompt="Patch service/render.py so it uses account_name and then run tests.",
        input_files=["service/render.py", "service/tests/test_client.py"],
        expected_outputs=[],
    )

    assert augmented == ["_env/service/render.py"]


async def test_helper_report_format_repair_skips_resource_required_state(monkeypatch, tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    helper_ws.mkdir(parents=True)
    main_ws.mkdir()
    calls = 0

    async def fake_loop(msgs, *args, dispatcher=None, **kwargs):
        nonlocal calls
        calls += 1
        await dispatcher("request_resource", {
            "kind": "draw",
            "reason": "need chart",
            "needed_outputs": ["chart.png"],
            "resume_instruction": "embed chart.png",
        })
        return 'Output files: {"files": ["paper.docx"]}', msgs

    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))

    result = await delegate._run_one_helper(
        task_id="paper_edit",
        prompt="Write paper.docx with chart.png.",
        main_workspace=str(main_ws),
        helper_workspace=str(helper_ws),
        archive_id="archive",
        group_id="group",
        user_id="user",
        resume=False,
        local_abort=asyncio.Event(),
        wait_for_register=asyncio.Event(),
        user_lang="en",
        kind="edit",
        mode="easy",
        helper_think=False,
        expected_outputs=["paper.docx"],
    )

    assert calls == 1
    assert result["ok"] is False
    assert result["terminal_reason"] == "resource_required"
    assert result["resource_required"]["matching_helper_kind"] == "draw"
    assert result["resource_required"]["suggested_helper_kind"] == "draw"
    assert result["declared_but_missing"] == []
    assert result["pending_declared_outputs"] == ["paper.docx"]
    assert "Pending until requested resource is provided" in result["report"]
    assert "Declared but not produced" not in result["report"]


async def test_helpers_shared_outputs_are_handoff_not_user_visible(monkeypatch, tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    helper_ws.mkdir(parents=True)
    main_ws.mkdir()

    async def fake_loop(msgs, *args, **kwargs):
        shared = helper_ws / "_helpers_shared"
        shared.mkdir(exist_ok=True)
        (shared / "algo_rbtree.md").write_text("analysis evidence", encoding="utf-8")
        return '## Output files\n```json\n{"files": ["_helpers_shared/algo_rbtree.md"]}\n```', msgs

    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))

    result = await delegate._run_one_helper(
        task_id="algo_rbtree",
        prompt="Write shared analysis evidence.",
        main_workspace=str(main_ws),
        helper_workspace=str(helper_ws),
        archive_id="archive",
        group_id="group",
        user_id="user",
        resume=False,
        local_abort=asyncio.Event(),
        wait_for_register=asyncio.Event(),
        user_lang="en",
        kind="code",
        mode="easy",
        helper_think=False,
        expected_outputs=["_helpers_shared/algo_rbtree.md"],
    )

    assert result["ok"] is True
    assert result["outputs_check"]["outputs_complete"] is True
    assert result["workspace_files"] == ["_helpers_shared/algo_rbtree.md"]
    assert result["user_visible_files"] == []
    assert "_post_helper_action" not in result
    assert "not user-facing deliverables" in result["post_helper_usage_hint"]


async def test_helper_outputs_check_accepts_readonly_shared_expected_via_writable_shared(tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    shared = helper_ws / "_helpers_shared"
    shared.mkdir(parents=True)
    main_ws.mkdir()
    astar_code = (
        "from heapq import heappop, heappush\n\n"
        "def a_star(start, goal, neighbors, heuristic):\n"
        "    open_set = [(0, start)]\n"
        "    came_from = {}\n"
        "    g_score = {start: 0}\n"
        "    while open_set:\n"
        "        _, node = heappop(open_set)\n"
        "        if node == goal:\n"
        "            path = [node]\n"
        "            while node in came_from:\n"
        "                node = came_from[node]\n"
        "                path.append(node)\n"
        "            return list(reversed(path))\n"
        "        for nxt, cost in neighbors(node):\n"
        "            cand = g_score[node] + cost\n"
        "            if cand < g_score.get(nxt, float('inf')):\n"
        "                came_from[nxt] = node\n"
        "                g_score[nxt] = cand\n"
        "                heappush(open_set, (cand + heuristic(nxt, goal), nxt))\n"
        "    return []\n"
    )
    (shared / "astar.py").write_text(astar_code, encoding="utf-8")
    (helper_ws / ".helper_astar_impl_full_report.txt").write_text(
        '```json\n{"files": ["_helpers_shared/astar.py"]}\n```',
        encoding="utf-8",
    )

    async def fake_loop(*args, **kwargs):
        shared.mkdir(parents=True, exist_ok=True)
        (shared / "astar.py").write_text(astar_code, encoding="utf-8")
        return '```json\n{"files": ["_helpers_shared/astar.py"]}\n```', []

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))
    try:
        result = await delegate._run_one_helper(
            task_id="astar_impl",
            prompt="Implement A* using the shared graph scaffold.",
            main_workspace=str(main_ws),
            helper_workspace=str(helper_ws),
            archive_id="archive",
            group_id="group",
            user_id="user",
            resume=False,
            local_abort=asyncio.Event(),
            wait_for_register=asyncio.Event(),
            user_lang="zh",
            kind="code",
            mode="easy",
            helper_think=False,
            expected_outputs=["_shared/astar.py"],
        )
    finally:
        monkeypatch.undo()

    assert result["ok"] is True
    assert result["terminal_reason"] == "completed"
    assert result["files"] == []
    assert result["workspace_files"] == []
    assert result["user_visible_files"] == []
    assert (main_ws / "_helpers_shared" / "astar.py").is_file()
    assert result["outputs_check"]["outputs_missing"] == []
    assert result["outputs_check"]["outputs_complete"] is True
    assert result["outputs_check"]["shared_protocol_matches"] == [{
        "expected_readonly": "_shared/astar.py",
        "delivered_writable_shared": "_helpers_shared/astar.py",
    }]
    assert "helper 对只读" in result["post_helper_usage_hint"]


async def test_helper_outputs_check_accepts_declared_shared_output_directory(monkeypatch, tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    helper_ws.mkdir(parents=True)
    main_ws.mkdir()

    async def fake_loop(*args, **kwargs):
        data_dir = helper_ws / "_helpers_shared" / "test_data"
        data_dir.mkdir(parents=True, exist_ok=True)
        (data_dir / "text_english_10kb.dat").write_bytes(b"a" * 1024)
        (data_dir / "random_bytes_10kb.dat").write_bytes(b"b" * 1024)
        return '```json\n{"files": ["_helpers_shared/test_data/"]}\n```', []

    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))

    result = await delegate._run_one_helper(
        task_id="framework",
        prompt="Generate shared test data files.",
        main_workspace=str(main_ws),
        helper_workspace=str(helper_ws),
        archive_id="archive",
        group_id="group",
        user_id="user",
        resume=False,
        local_abort=asyncio.Event(),
        wait_for_register=asyncio.Event(),
        user_lang="en",
        kind="code",
        mode="easy",
        helper_think=False,
        expected_outputs=["_helpers_shared/test_data/"],
    )

    assert result["ok"] is True
    assert result["terminal_reason"] == "completed"
    assert result["declared_but_missing"] == []
    assert result["outputs_check"]["outputs_missing"] == []
    assert result["outputs_check"]["outputs_complete"] is True
    assert result["outputs_check"]["delivered_count"] >= 1
    assert (main_ws / "_helpers_shared" / "test_data" / "text_english_10kb.dat").is_file()


def test_declared_mapping_accepts_helper_prefixed_docx_delivery():
    from app.llm.tools.delegate import _matches_declared_output_via_mapping

    file_map = [{
        "helper_name": "gen_docx_通信原理作业_第五章.docx",
        "main_name": "supplement_521_gen_docx_通信原理作业_第五章.docx",
        "shared_name": None,
    }]

    assert _matches_declared_output_via_mapping(
        "supplement_521_gen_docx_通信原理作业_第五章.docx",
        {"gen_docx_通信原理作业_第五章.docx"},
        file_map,
    ) is True


async def test_helper_result_summary_treats_helper_prefixed_docx_as_declared(tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    helper_ws.mkdir()
    main_ws.mkdir()
    (helper_ws / "gen_docx_通信原理作业_第五章.docx").write_bytes(b"PK\x03\x04" + b"x" * 6000)
    (helper_ws / ".helper_supplement_521_full_report.txt").write_text(
        '```json\n{"files": ["gen_docx_通信原理作业_第五章.docx"]}\n```',
        encoding="utf-8",
    )

    async def fake_loop(*args, **kwargs):
        (helper_ws / "gen_docx_通信原理作业_第五章.docx").write_bytes(b"PK\x03\x04" + b"y" * 6000)
        return '```json\n{"files": ["gen_docx_通信原理作业_第五章.docx"]}\n```', []

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))
    try:
        result = await delegate._run_one_helper(
            task_id="supplement_521",
            prompt="生成文档",
            main_workspace=str(main_ws),
            helper_workspace=str(helper_ws),
            archive_id="archive",
            group_id="group",
            user_id="user",
            resume=False,
            local_abort=asyncio.Event(),
            wait_for_register=asyncio.Event(),
            user_lang="zh",
            kind="edit",
            mode="easy",
            helper_think=False,
            expected_outputs=["gen_docx_通信原理作业_第五章.docx"],
        )
    finally:
        monkeypatch.undo()

    assert result["files"] == ["gen_docx_通信原理作业_第五章.docx"]
    assert result["declared_files"] == ["gen_docx_通信原理作业_第五章.docx"]
    assert result["declared_but_missing"] == []
    assert result["delivered_but_not_declared"] == []
    assert result["outputs_check"]["outputs_missing"] == []
    assert result["outputs_check"]["outputs_complete"] is True
    assert result["outputs_check"]["producer_self_verified"] is True
    assert result["_post_helper_action"] == "output_json_directly"

    from app.core.filesystem import FileRegistry

    registry = FileRegistry.load(
        scope_id=f"workspace:{main_ws.resolve()}",
        workspace_root=main_ws,
    )
    record = registry.find_by_workspace_path("gen_docx_通信原理作业_第五章.docx")
    assert record is not None
    assert record.verified is True
    assert record.metadata["outputs_complete"] is True
    assert record.metadata["producer_self_verified"] is True


async def test_environment_helper_env_outputs_are_validated_without_user_delivery(tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    env_src = helper_ws / "_env" / "src" / "taskboard"
    env_tests = helper_ws / "_env" / "tests"
    env_src.mkdir(parents=True)
    env_tests.mkdir(parents=True)
    main_ws.mkdir()

    async def fake_loop(*args, **kwargs):
        env_src.mkdir(parents=True, exist_ok=True)
        env_tests.mkdir(parents=True, exist_ok=True)
        (env_src / "cli.py").write_text(
            "import argparse\n\n"
            "def build_parser():\n"
            "    parser = argparse.ArgumentParser(prog='taskboard')\n"
            "    parser.add_argument('--version', action='store_true')\n"
            "    return parser\n\n"
            "def main(argv=None):\n"
            "    args = build_parser().parse_args(argv)\n"
            "    return 'taskboard 1.0' if args.version else 'taskboard ready'\n",
            encoding="utf-8",
        )
        (env_tests / "test_cli.py").write_text(
            "from taskboard.cli import main\n\n"
            "def test_main_default():\n"
            "    assert main([]) == 'taskboard ready'\n\n"
            "def test_main_version():\n"
            "    assert main(['--version']) == 'taskboard 1.0'\n",
            encoding="utf-8",
        )
        return "Edited the environment workspace copies and verified tests.", []

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))
    try:
        result = await delegate._run_one_helper(
            task_id="taskboard_cli_pytest_fix",
            prompt="Edit _env/src/taskboard/cli.py and _env/tests/test_cli.py.",
            main_workspace=str(main_ws),
            helper_workspace=str(helper_ws),
            archive_id="archive",
            group_id="group",
            user_id="user",
            resume=False,
            local_abort=asyncio.Event(),
            wait_for_register=None,
            user_lang="zh",
            kind="code",
            mode="easy",
            helper_think=False,
            expected_outputs=["_env/src/taskboard/cli.py", "_env/tests/test_cli.py"],
        )
    finally:
        monkeypatch.undo()

    assert result["ok"] is True
    assert result["terminal_reason"] == "completed"
    assert result["files"] == ["_env/src/taskboard/cli.py", "_env/tests/test_cli.py"]
    assert result["workspace_files"] == ["_env/src/taskboard/cli.py", "_env/tests/test_cli.py"]
    assert result["user_visible_files"] == []
    assert result["outputs_check"]["outputs_missing"] == []
    assert result["outputs_check"]["outputs_complete"] is True
    assert result["outputs_check"]["delivered_count"] == 2
    assert not any(
        warning.get("issue") == "suspicious_short_completion"
        for warning in result["outputs_check"]["quality_warnings"]
    )
    assert (main_ws / "_env" / "src" / "taskboard" / "cli.py").is_file()
    assert (main_ws / "_env" / "tests" / "test_cli.py").is_file()


async def test_environment_helper_matches_project_relative_expected_outputs(tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    env_src = helper_ws / "_env" / "src" / "algolab"
    env_tests = helper_ws / "_env" / "tests"
    env_src.mkdir(parents=True)
    env_tests.mkdir(parents=True)
    main_ws.mkdir()

    async def fake_loop(*args, **kwargs):
        env_src.mkdir(parents=True, exist_ok=True)
        env_tests.mkdir(parents=True, exist_ok=True)
        (env_src / "graph.py").write_text(
            "from heapq import heappop, heappush\n\n"
            "def shortest_path(graph, start, target):\n"
            "    queue = [(0, start, [])]\n"
            "    seen = set()\n"
            "    while queue:\n"
            "        cost, node, path = heappop(queue)\n"
            "        if node in seen:\n"
            "            continue\n"
            "        path = path + [node]\n"
            "        if node == target:\n"
            "            return cost, path\n"
            "        seen.add(node)\n"
            "        for nxt, weight in graph.get(node, []):\n"
            "            heappush(queue, (cost + weight, nxt, path))\n"
            "    return float('inf'), []\n",
            encoding="utf-8",
        )
        (env_tests / "test_graph.py").write_text(
            "from algolab.graph import shortest_path\n\n"
            "def test_shortest_path():\n"
            "    graph = {'a': [('b', 1), ('c', 5)], 'b': [('c', 2)], 'c': []}\n"
            "    assert shortest_path(graph, 'a', 'c') == (3, ['a', 'b', 'c'])\n",
            encoding="utf-8",
        )
        (helper_ws / "_env" / "README.md").write_text(
            "# AlgoLab\n\nA small graph algorithm package with Dijkstra-style shortest path utilities and tests.\n",
            encoding="utf-8",
        )
        return "Edited project files under _env and verified them.", []

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))
    try:
        result = await delegate._run_one_helper(
            task_id="dijkstra_impl",
            prompt="Edit _env/src/algolab/graph.py, _env/tests/test_graph.py, and _env/README.md.",
            main_workspace=str(main_ws),
            helper_workspace=str(helper_ws),
            archive_id="archive",
            group_id="group",
            user_id="user",
            resume=False,
            local_abort=asyncio.Event(),
            wait_for_register=None,
            user_lang="zh",
            kind="code",
            mode="easy",
            helper_think=False,
            expected_outputs=["src/algolab/graph.py", "tests/test_graph.py", "README.md"],
        )
    finally:
        monkeypatch.undo()

    assert result["ok"] is True
    assert result["terminal_reason"] == "completed"
    assert result["files"] == [
        "_env/README.md",
        "_env/src/algolab/graph.py",
        "_env/tests/test_graph.py",
    ]
    assert result["workspace_files"] == [
        "_env/README.md",
        "_env/src/algolab/graph.py",
        "_env/tests/test_graph.py",
    ]
    assert result["user_visible_files"] == []
    assert result["main_available_files"] == [
        "_env/README.md",
        "_env/src/algolab/graph.py",
        "_env/tests/test_graph.py",
    ]
    assert result["outputs_check"]["outputs_missing"] == []
    assert result["outputs_check"]["outputs_complete"] is True
    assert result["outputs_check"]["delivered_count"] == 3
    assert not any(
        warning.get("issue") == "suspicious_short_completion"
        for warning in result["outputs_check"]["quality_warnings"]
    )


async def test_environment_helper_counts_unchanged_expected_env_outputs(tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    main_env = main_ws / "_env"
    main_env.mkdir(parents=True)
    (main_env / "config_loader.py").write_text("def load():\n    return {'mode': 'old'}\n", encoding="utf-8")
    (main_env / "app_config.py").write_text("APP_NAME = 'demo'\n", encoding="utf-8")

    async def fake_loop(*args, **kwargs):
        helper_env = helper_ws / "_env"
        helper_env.mkdir(parents=True, exist_ok=True)
        (helper_env / "config_loader.py").write_text("def load():\n    return {'mode': 'new'}\n", encoding="utf-8")
        return "Edited config_loader.py and verified app_config.py did not require changes.", []

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))
    try:
        result = await delegate._run_one_helper(
            task_id="config_loader_fix",
            prompt="Edit _env/config_loader.py. Check _env/app_config.py and leave it unchanged if no edit is needed.",
            main_workspace=str(main_ws),
            helper_workspace=str(helper_ws),
            archive_id="archive",
            group_id="group",
            user_id="user",
            resume=False,
            local_abort=asyncio.Event(),
            wait_for_register=None,
            user_lang="zh",
            kind="code",
            mode="easy",
            helper_think=False,
            expected_outputs=["_env/config_loader.py", "_env/app_config.py"],
        )
    finally:
        monkeypatch.undo()

    assert result["ok"] is True
    assert result["outputs_check"]["outputs_missing"] == []
    assert result["outputs_check"]["outputs_complete"] is True
    assert result["outputs_check"]["unchanged_existing_outputs"] == [
        {"expected": "_env/app_config.py", "existing_path": "_env/app_config.py"}
    ]
    assert "byte-identical" in result["outputs_check"]["unchanged_existing_outputs_fact"]
    assert (main_env / "config_loader.py").read_text(encoding="utf-8") == "def load():\n    return {'mode': 'new'}\n"
    assert (main_env / "app_config.py").read_text(encoding="utf-8") == "APP_NAME = 'demo'\n"


async def test_read_helper_exposes_internal_evidence_without_user_delivery(tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    helper_ws.mkdir()
    main_ws.mkdir()

    evidence_text = (
        "VERDICT: PASS\n"
        "source_files: a.docx, b.docx\n"
        "coverage_summary: COMPLETE - all documents fully read, no missing sections\n"
        "needs_escalation: false\n"
        "line_ranges: read_dai_evidence.txt lines 1-80\n"
    )

    async def fake_loop(*args, **kwargs):
        (helper_ws / "read_dai_evidence.txt").write_text(evidence_text, encoding="utf-8")
        return (
            '```json\n{"files": ["read_dai_evidence.txt"]}\n```\n\n'
            "## 摘要\nVERDICT: PASS，所有材料已读取。"
        ), []

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))
    try:
        result = await delegate._run_one_helper(
            task_id="read_dai",
            prompt="Read all source reports and save evidence.",
            main_workspace=str(main_ws),
            helper_workspace=str(helper_ws),
            archive_id="archive",
            group_id="group",
            user_id="user",
            resume=False,
            local_abort=asyncio.Event(),
            wait_for_register=None,
            user_lang="zh",
            kind="read",
            mode="easy",
            helper_think=False,
            expected_outputs=["read_dai_evidence.txt"],
        )
    finally:
        monkeypatch.undo()

    assert result["ok"] is True
    assert result["terminal_reason"] == "completed"
    assert result["files"] == []
    assert result["internal_evidence_files"] == ["read_dai_read_dai_evidence.txt"]
    assert "内部读取证据: read_dai_read_dai_evidence.txt" in result["report"]
    assert result["read_evidence_summary"]["has_complete_evidence"] is True
    assert result["read_evidence_summary"]["verdicts"][0]["verdict"] == "PASS"
    assert "COMPLETE" in result["read_evidence_summary"]["verdicts"][0]["coverage_summary"]
    assert (main_ws / "read_dai_read_dai_evidence.txt").is_file()


async def test_read_helper_does_not_merge_env_evidence_into_project_tree(tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    (helper_ws / "_env").mkdir(parents=True)
    main_ws.mkdir()

    async def fake_loop(*args, **kwargs):
        (helper_ws / "_env").mkdir(parents=True, exist_ok=True)
        (helper_ws / "_env" / "写作_框架范文证据.txt").write_text(
            "VERDICT: PASS\ncoverage_summary: read-only extraction evidence\n",
            encoding="utf-8",
        )
        return (
            '```json\n{"files": ["_env/写作_框架范文证据.txt"]}\n```\n\n'
            "## 摘要\nVERDICT: PASS，材料已读取。"
        ), []

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))
    try:
        result = await delegate._run_one_helper(
            task_id="read_writing",
            prompt="Read writing source files and save internal evidence.",
            main_workspace=str(main_ws),
            helper_workspace=str(helper_ws),
            archive_id="archive",
            group_id="group",
            user_id="user",
            resume=False,
            local_abort=asyncio.Event(),
            wait_for_register=None,
            user_lang="zh",
            kind="read",
            mode="easy",
            helper_think=False,
            expected_outputs=["_env/写作_框架范文证据.txt"],
        )
    finally:
        monkeypatch.undo()

    assert not (main_ws / "_env" / "写作_框架范文证据.txt").exists()
    assert result["files"] == []
    assert result.get("main_available_files") in (None, [])
    assert result["copy_stats"]["env_skipped_read_evidence"] == ["_env/写作_框架范文证据.txt"]
    assert result["internal_evidence_files"] == [".helper_read_writing_full_report.txt"]
    assert result["read_evidence_summary"]["has_complete_evidence"] is True


async def test_read_helper_write_rejects_env_evidence_path(tmp_path):
    from app.core.core_processes import (
        reset_current_helper_expected_outputs,
        reset_current_helper_kind,
        set_current_helper_expected_outputs,
        set_current_helper_kind,
    )
    from app.llm.tools.workspace_file_ops import handle_write

    ws = tmp_path / "_delegate_read"
    ws.mkdir()
    kind_token = set_current_helper_kind("read")
    outputs_token = set_current_helper_expected_outputs(["_env/写作_框架范文证据.txt"])
    try:
        result = await handle_write(
            str(ws),
            "_env/写作_框架范文证据.txt",
            "VERDICT: PASS\ncoverage_summary: ok\n",
        )
    finally:
        reset_current_helper_expected_outputs(outputs_token)
        reset_current_helper_kind(kind_token)

    assert result["ok"] is False
    assert result["blocked_reason"] == "read_helper_env_evidence_forbidden"
    assert not (ws / "_env" / "写作_框架范文证据.txt").exists()


async def test_workspace_write_over_hard_limit_does_not_truncate_to_disk(tmp_path):
    from app.llm.tools.workspace_file_ops import handle_write

    content = "x" * 500_001
    result = await handle_write(str(tmp_path), "too_large.txt", content)

    assert result["ok"] is False
    assert result["error_kind"] == "workspace_write_content_too_large"
    assert not (tmp_path / "too_large.txt").exists()


async def test_workspace_write_success_returns_concise_write_facts(tmp_path):
    from app.llm.tools.workspace_file_ops import handle_write

    content = "# Report\n\nAlpha\nBeta\n"
    result = await handle_write(str(tmp_path), "report.md", content)

    assert result["ok"] is True
    assert result["path"] == "report.md"
    assert result["content_chars"] == len(content)
    assert result["line_count"] == 4
    assert len(result["sha256"]) == 64
    assert result["head_excerpt"].startswith("# Report")
    assert "a full reread is only useful for a named main-owned gap" in result["write_fact"]
    assert (tmp_path / "report.md").read_text(encoding="utf-8") == content


async def test_workspace_run_missing_unix_inventory_command_points_to_file_tools(tmp_path):
    from app.llm.tools import workspace_run as ws_run

    if sys.platform != "win32":
        pytest.skip("Windows-specific missing Unix command guidance")
    result = await ws_run.handle_run(str(tmp_path), "ls -la _env", timeout_sec=5)

    assert result["ok"] is False
    assert "Unix inventory command" in result["error"]
    assert "workspace locate" in result["FIX_HINT"]
    assert "read_file/inspect_file" in result["FIX_HINT"]
    assert "_env/project_inventory.md" in result["FIX_HINT"]


async def test_workspace_run_failed_unix_inventory_pipeline_gets_fix_hint(tmp_path):
    from app.llm.tools import workspace_run as ws_run

    if sys.platform != "win32":
        pytest.skip("Windows-specific Unix pipeline guidance")
    result = await ws_run.handle_run(str(tmp_path), "find _env -type f | head -20", timeout_sec=5)

    assert result["ok"] is False
    assert "FIX_HINT" in result
    assert "Unix-style file inventory syntax" in result["FIX_HINT"]
    assert "workspace locate" in result["FIX_HINT"]
    assert "_scratch/*.py" in result["FIX_HINT"]


async def test_workspace_run_failed_unix_inventory_with_null_redirect_gets_fix_hint(tmp_path):
    from app.llm.tools import workspace_run as ws_run

    if sys.platform != "win32":
        pytest.skip("Windows-specific Unix redirect guidance")
    result = await ws_run.handle_run(
        str(tmp_path),
        'find _env -maxdepth 1 -ls 2>/dev/null; echo "---"; ls -la _env/ 2>/dev/null',
        timeout_sec=5,
    )

    assert result["ok"] is False
    assert result.get("blocked_reason") is None
    assert "FIX_HINT" in result
    assert "Unix-style file inventory syntax" in result["FIX_HINT"]
    assert "workspace locate" in result["FIX_HINT"]
    assert "read_file/inspect_file" in result["FIX_HINT"]


async def test_workspace_run_python_c_syntax_error_suggests_scratch_script(tmp_path):
    from app.llm.tools import workspace_run as ws_run

    result = await ws_run.handle_run(
        str(tmp_path),
        "python -c \"if True print('bad')\"",
        timeout_sec=5,
    )

    assert result["ok"] is False
    assert "FIX_HINT" in result
    assert "Python `-c` failed with a SyntaxError" in result["FIX_HINT"]
    assert "_scratch/" in result["FIX_HINT"]


def test_workspace_run_dependency_install_failure_points_to_local_routes():
    from app.llm.tools.workspace_run import _dependency_install_failure_hint

    hint = _dependency_install_failure_hint(
        "python -m pip install pytest",
        "ERROR: Could not install packages due to an OSError: [WinError 5] Access is denied",
        "",
    )

    assert hint
    assert "dependency-install command failed" in hint
    assert "local/project/bundled routes" in hint
    assert "direct script/import checks" in hint
    assert "missing dependency" in hint


async def test_environment_workspace_run_python_c_syntax_error_suggests_env_python_code(tmp_path):
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools import workspace_run as ws_run

    env = EnvironmentContext(
        root_dir=str(tmp_path / "project"),
        archive_id="arch_test",
        group_id="env_user_test",
        user_id="user",
        project_key="project",
    )
    (tmp_path / "project").mkdir()

    with runtime_context("environment", env):
        result = await ws_run.handle_run(
            str(tmp_path),
            "python -c \"if True print('bad')\"",
            timeout_sec=5,
        )

    assert result["ok"] is False
    assert "FIX_HINT" in result
    assert "env_run" in result["FIX_HINT"]
    assert "python_code" in result["FIX_HINT"]


async def test_main_thread_long_framework_contract_write_delegates(tmp_path):
    from app.llm.tools.registry import _handle_workspace

    content = (
        "# Shared Framework Contract\n\n"
        "shared goal: build a multi-part database index research project.\n"
        "interface: all modules must expose a common benchmark protocol.\n"
        "schema: evidence files must include source coverage and validation fields.\n"
        "ownership: helpers own separate algorithm, benchmark, and report sections.\n"
        "validation: compile, run tests, compare outputs, and verify the final document.\n"
        "merge order: framework, implementations, benchmark data, paper outline, final docx.\n"
    ) * 8

    raw = await _handle_workspace(str(tmp_path), {
        "action": "write",
        "path": "framework_contract.txt",
        "content": content,
    })
    result = json.loads(raw)

    assert result["ok"] is False
    assert result["error_kind"] == "main_thread_project_artifact_should_delegate"
    assert result["recommended_tools"] == ["delegate"]
    assert not (tmp_path / "framework_contract.txt").exists()


async def test_main_thread_shared_framework_write_requests_delegate(tmp_path):
    from app.llm.tools.registry import _handle_workspace

    content = (
        "framework contract\n"
        "shared goal: compare algorithms\n"
        "interface: common benchmark protocol\n"
        "schema: json evidence\n"
        "validation: compile and benchmark\n"
    ) * 30

    raw = await _handle_workspace(str(tmp_path), {
        "action": "write",
        "path": "_shared/framework_contract.json",
        "content": content,
    })
    result = json.loads(raw)

    assert result["ok"] is False
    assert result["error_kind"] == "main_thread_project_artifact_should_delegate"
    assert result["recommended_tools"] == ["delegate"]
    assert not (tmp_path / "_shared" / "framework_contract.json").exists()


async def test_main_thread_large_workspace_write_requires_delegate_or_segment(tmp_path):
    from app.llm.tools.registry import _handle_workspace

    content = ("section line\n" * 1300)
    raw = await _handle_workspace(str(tmp_path), {
        "action": "write",
        "path": "large_report.md",
        "content": content,
    })
    result = json.loads(raw)

    assert result["ok"] is False
    assert result["error_kind"] == "main_thread_large_write_should_delegate_or_segment"
    assert result["content_chars"] == len(content)
    assert not (tmp_path / "large_report.md").exists()


async def test_main_thread_workspace_append_short_allowed_and_large_blocked(tmp_path):
    from app.llm.tools.registry import _handle_workspace

    target = tmp_path / "notes.md"
    target.write_text("head\n", encoding="utf-8")

    short = await _handle_workspace(str(tmp_path), {
        "action": "append",
        "path": "notes.md",
        "content": "short section\n",
    })
    short_result = json.loads(short)

    assert short_result["ok"] is True
    assert target.read_text(encoding="utf-8").endswith("short section\n")

    before = target.read_text(encoding="utf-8")
    large_content = "large section\n" * 700
    blocked = await _handle_workspace(str(tmp_path), {
        "action": "append",
        "path": "notes.md",
        "content": large_content,
    })
    blocked_result = json.loads(blocked)

    assert blocked_result["ok"] is False
    assert blocked_result["error_kind"] == "main_thread_large_write_should_delegate_or_segment"
    assert "available_recovery_shapes" in blocked_result
    assert target.read_text(encoding="utf-8") == before


async def test_main_thread_large_unbounded_read_requires_targeting(tmp_path):
    from app.llm.tools.registry import _handle_read_file

    target = tmp_path / "large.txt"
    target.write_text("x" * 130_000, encoding="utf-8")

    raw = await _handle_read_file(str(tmp_path), {"path": "large.txt"})
    result = json.loads(raw)

    assert result["ok"] is False
    assert result["error_kind"] == "main_thread_large_read_should_delegate_or_target"
    assert result["bytes"] >= 130_000

    targeted = json.loads(await _handle_read_file(
        str(tmp_path),
        {"path": "large.txt", "start_line": 1, "end_line": 1},
    ))
    assert targeted["ok"] is True


async def test_main_thread_targeted_read_uses_small_budget_but_helper_read_does_not(tmp_path):
    from app.llm.tools.registry import _handle_read_file

    target = tmp_path / "large.txt"
    target.write_text("\n".join("x" * 1000 for _ in range(80)), encoding="utf-8")

    main_result = json.loads(await _handle_read_file(
        str(tmp_path),
        {"path": "large.txt", "start_line": 1, "end_line": 80, "max_chars": 50_000},
    ))
    helper_result = json.loads(await _handle_read_file(
        str(tmp_path),
        {"path": "large.txt", "start_line": 1, "end_line": 80, "max_chars": 50_000},
        caller_kind="helper",
    ))

    assert main_result["ok"] is True
    assert main_result["truncated"] is True
    assert "main_thread_read_fact" in main_result
    assert helper_result["ok"] is True
    assert helper_result.get("error_kind") != "main_thread_large_read_should_delegate_or_target"
    assert "main_thread_read_fact" not in helper_result
    assert len(helper_result["content"]) > len(main_result["content"])
    assert main_result["content_truncated"] is True
    assert main_result["tool_result_truncated"] is True
    assert main_result["output_truncated"] is True
    assert "only the head excerpt" in main_result["visible_excerpt_policy"]
    saved = tmp_path / main_result["content_full_saved_path"]
    assert saved.is_file()
    assert "x" * 1000 in saved.read_text(encoding="utf-8")


async def test_main_thread_trusts_ready_helper_output_without_rereading_content(tmp_path):
    from app.core.filesystem import FileKind, FileRegistry, FileStatus, Visibility, intake_workspace_file
    from app.llm.tools.registry import _handle_read_file

    target = tmp_path / "classified_report.md"
    target.write_text("helper-owned final content\n", encoding="utf-8")
    registry = FileRegistry.load(
        scope_id=f"workspace:{tmp_path.resolve()}",
        workspace_root=tmp_path,
    )
    intake_workspace_file(
        registry,
        "classified_report.md",
        kind=FileKind.HELPER_OUTPUT,
        status=FileStatus.READY,
        visibility=Visibility.PROJECT,
        owner_task_id="email_triage",
        helper_kind="read",
        metadata={
            "source": "delegate_copyback",
            "outputs_complete": True,
            "producer_self_verified": True,
        },
    )

    main_result = json.loads(await _handle_read_file(str(tmp_path), {"path": "classified_report.md"}))
    forced_result = json.loads(await _handle_read_file(
        str(tmp_path),
        {"path": "classified_report.md", "force": True},
    ))
    helper_result = json.loads(await _handle_read_file(
        str(tmp_path),
        {"path": "classified_report.md"},
        caller_kind="helper",
    ))

    assert main_result["ok"] is True
    assert main_result["content_omitted"] is True
    assert main_result["content_omitted_reason"] == "helper_owned_verified_artifact"
    assert main_result["owner_task_id"] == "email_triage"
    assert "helper-owned final content" not in main_result["content"]
    assert "trust the successful helper" in main_result["helper_owned_artifact_fact"]
    assert forced_result["ok"] is True
    assert "helper-owned final content" in forced_result["content"]
    assert helper_result["ok"] is True
    assert "helper-owned final content" in helper_result["content"]


async def test_main_thread_large_ready_helper_output_returns_provenance_fact(tmp_path):
    from app.core.filesystem import FileKind, FileRegistry, FileStatus, Visibility, intake_workspace_file
    from app.llm.tools.registry import _handle_read_file

    target = tmp_path / "long_report.md"
    target.write_text("helper line\n" * 8000, encoding="utf-8")
    registry = FileRegistry.load(
        scope_id=f"workspace:{tmp_path.resolve()}",
        workspace_root=tmp_path,
    )
    intake_workspace_file(
        registry,
        "long_report.md",
        kind=FileKind.HELPER_OUTPUT,
        status=FileStatus.READY,
        visibility=Visibility.PROJECT,
        owner_task_id="writer",
        helper_kind="markdown",
        metadata={
            "source": "delegate_copyback",
            "outputs_complete": True,
            "producer_self_verified": True,
        },
    )

    result = json.loads(await _handle_read_file(str(tmp_path), {"path": "long_report.md"}))

    assert result["ok"] is True
    assert result["content_omitted_reason"] == "helper_owned_verified_artifact"
    assert result["size_bytes"] > 48_000
    assert result.get("error_kind") != "main_thread_large_read_should_delegate_or_target"


async def test_outputs_complete_without_producer_self_verified_keeps_main_read_compact(tmp_path):
    from app.core.filesystem import FileKind, FileRegistry, FileStatus, Visibility, intake_workspace_file
    from app.llm.tools.registry import _handle_read_file

    target = tmp_path / "warning_report.md"
    target.write_text("warning-bearing helper content\n", encoding="utf-8")
    registry = FileRegistry.load(
        scope_id=f"workspace:{tmp_path.resolve()}",
        workspace_root=tmp_path,
    )
    intake_workspace_file(
        registry,
        "warning_report.md",
        kind=FileKind.HELPER_OUTPUT,
        status=FileStatus.READY,
        visibility=Visibility.PROJECT,
        owner_task_id="writer",
        helper_kind="edit",
        metadata={
            "source": "delegate_copyback",
            "outputs_complete": True,
            "producer_self_verified": False,
            "quality_warning_count": 1,
        },
    )

    result = json.loads(await _handle_read_file(str(tmp_path), {"path": "warning_report.md"}))

    assert result["ok"] is True
    assert result["content_omitted"] is True
    assert result["content_omitted_reason"] == "helper_owned_unverified_artifact"
    assert "warning-bearing helper content" not in result["content"]
    assert "producer helper" in result["_next_action_instruction"]

    forced = json.loads(await _handle_read_file(
        str(tmp_path),
        {"path": "warning_report.md", "force": True},
    ))
    assert "warning-bearing helper content" in forced["content"]


async def test_main_thread_reads_normal_workspace_file_without_helper_output_guard(tmp_path):
    from app.llm.tools.registry import _handle_read_file

    target = tmp_path / "notes.md"
    target.write_text("ordinary workspace content\n", encoding="utf-8")

    result = json.loads(await _handle_read_file(str(tmp_path), {"path": "notes.md"}))

    assert result["ok"] is True
    assert result.get("content_omitted_reason") != "helper_owned_verified_artifact"
    assert "ordinary workspace content" in result["content"]


@pytest.mark.asyncio
async def test_workspace_run_long_stdout_saves_full_output(tmp_path):
    from app.llm.tools.workspace_run import handle_run

    script = tmp_path / "emit_long.py"
    script.write_text("print('A' * 70000)\n", encoding="utf-8")

    result = await handle_run(str(tmp_path), "python emit_long.py", timeout_sec=10)

    assert result["ok"] is True
    assert result["stdout_truncated"] is True
    assert result["tool_result_truncated"] is True
    assert result["output_truncated"] is True
    assert len(result["stdout"]) <= 64 * 1024
    saved = tmp_path / result["stdout_full_saved_path"]
    assert saved.is_file()
    full = saved.read_text(encoding="utf-8")
    assert len(full) >= 70000
    assert full.rstrip("\n").endswith("A")


def test_file_preview_long_content_saves_full_text(tmp_path):
    from app.core.file_preview import preview_file

    target = tmp_path / "preview.txt"
    target.write_text("C" * 5000, encoding="utf-8")

    result = preview_file(target, max_chars=1000)

    assert result["truncated"] is True
    assert result["content_truncated"] is True
    assert result["tool_result_truncated"] is True
    assert len(result["content"]) == 1000
    saved = tmp_path / result["content_full_saved_path"]
    assert saved.is_file()
    assert saved.read_text(encoding="utf-8") == "C" * 5000


def test_main_and_helper_read_file_schemas_are_distinct():
    from app.llm.tools import registry
    from app.llm.tools.delegate import _HELPER_TOOLS

    main_read = next(tool for tool in registry.ROUND2_TOOLS if tool["function"]["name"] == "read_file")
    helper_read = next(tool for tool in _HELPER_TOOLS if tool["function"]["name"] == "read_file")

    assert main_read is registry.MAIN_READ_FILE_SCHEMA
    assert helper_read is registry.READ_FILE_SCHEMA
    assert main_read is not helper_read
    assert "Main-process text spot-check reader" in main_read["function"]["description"]
    assert "helper-owned in the normal workflow" in main_read["function"]["description"]
    assert "full-file analysis" in main_read["function"]["description"]
    assert "segment" not in main_read["function"]["description"].lower()
    assert "full-file analysis to helpers" not in helper_read["function"]["description"]


def test_delegate_prompt_describes_ultra_large_file_helper_fanout():
    from app.llm.tools import registry
    from app.llm import aux_prompts

    description = registry.DELEGATE_TOOL_SCHEMA["function"]["description"]
    input_files_desc = registry.DELEGATE_TOOL_SCHEMA["function"]["parameters"]["properties"]["tasks"]["items"]["properties"]["input_files"]["description"]

    assert "one ultra-large file, long log, or long source material" in description
    assert "fan out focused `read` or `file_summary` helpers" in description
    assert "coverage summaries" in description
    assert "line ranges, page ranges, section labels" in input_files_desc
    assert "helper can reread the path from input_files" in input_files_desc
    assert "compact facts rather than complete source text" in input_files_desc
    assert "single ultra-large concrete file" in aux_prompts.TASK_QUALITY_GUARD_SYSTEM
    assert "bounded range" in aux_prompts.TASK_QUALITY_GUARD_SYSTEM


def test_edit_tools_describe_literal_old_str_not_json_escaped():
    from app.llm.tools import tool_schemas

    edit_desc = tool_schemas.EDIT_FILE_SCHEMA["function"]["description"]
    multi_desc = tool_schemas.MULTI_EDIT_SCHEMA["function"]["description"]
    multi_old_desc = (
        tool_schemas.MULTI_EDIT_SCHEMA["function"]["parameters"]["properties"]["edits"]["items"]["properties"]["old_str"]["description"]
    )
    edit_path_desc = tool_schemas.EDIT_FILE_SCHEMA["function"]["parameters"]["properties"]["path"]["description"]
    multi_path_desc = tool_schemas.MULTI_EDIT_SCHEMA["function"]["parameters"]["properties"]["path"]["description"]

    assert "`old_str` is literal file text" in edit_desc
    assert "not the JSON-escaped representation" in edit_desc
    assert "old_str 是文件原文" in multi_desc
    assert "do not copy JSON escape backslashes" in multi_old_desc
    assert "staged `_env/...` copy" in edit_path_desc
    assert "bare project-relative path" in multi_path_desc


async def test_main_thread_workspace_run_blocks_scripts_and_tests(tmp_path):
    from app.llm.tools.registry import _handle_workspace

    raw = await _handle_workspace(str(tmp_path), {
        "action": "run",
        "command": "python -m pytest",
        "timeout_sec": 30,
    })
    result = json.loads(raw)

    assert result["ok"] is False
    assert result["error_kind"] == "main_thread_workspace_run_should_delegate"


async def test_main_thread_workspace_run_allows_file_management(tmp_path):
    from app.llm.tools.registry import _handle_workspace

    raw = await _handle_workspace(str(tmp_path), {
        "action": "run",
        "command": "cmd /c dir",
        "timeout_sec": 5,
    })
    result = json.loads(raw)

    assert result["ok"] is True


async def test_completion_ledger_preserves_read_evidence_summary():
    from app.llm.tools.delegate_state import _add_to_completion_ledger, _get_completion_ledger

    trace_id = "trace_read_evidence_ledger"
    await _add_to_completion_ledger(
        trace_id,
        "read_group_a",
        {
            "task_id": "read_group_a",
            "ok": True,
            "terminal_reason": "completed",
            "elapsed_sec": 12.3,
            "files": [],
            "kind": "read",
            "mode": "easy",
            "internal_evidence_files": ["read_group_a_evidence.txt"],
            "read_evidence_summary": {
                "evidence_files": ["read_group_a_evidence.txt"],
                "verdicts": [{
                    "file": "read_group_a_evidence.txt",
                    "verdict": "PASS",
                    "coverage_summary": "coverage_summary: COMPLETE",
                }],
                "has_complete_evidence": True,
            },
            "outputs_check": {"outputs_complete": None, "outputs_missing": []},
        },
    )

    ledger = _get_completion_ledger(trace_id, last_n=1)
    assert ledger[0]["internal_evidence_files"] == ["read_group_a_evidence.txt"]
    assert ledger[0]["read_evidence_summary"]["has_complete_evidence"] is True
    assert ledger[0]["read_evidence_summary"]["verdicts"][0]["verdict"] == "PASS"


async def test_edit_helper_resource_request_freezes_instead_of_clean_completion(tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / ".temp" / "_delegate_user_paper_edit"
    main_ws = tmp_path / "main"
    helper_ws.mkdir(parents=True)
    main_ws.mkdir()
    (helper_ws / "paper.docx").write_bytes(b"PK\x03\x04" + b"x" * 6000)

    async def fake_loop(*args, dispatcher=None, **kwargs):
        assert dispatcher is not None
        blocked = await dispatcher("workspace", {
            "action": "write",
            "path": "make_charts.py",
            "content": "import matplotlib.pyplot as plt\nplt.savefig('chart.png')\n",
        })
        data = json.loads(blocked)
        assert data["suggested_tool"] == "request_resource"
        requested = await dispatcher("request_resource", {
            "kind": "draw",
            "reason": "缺少论文图表 PNG",
            "needed_outputs": ["chart.png"],
            "resume_instruction": "图表生成后嵌入 paper.docx 并修正文案",
        })
        req_data = json.loads(requested)
        assert req_data["requires_main_resource"] is True
        return '```json\n{"files": ["paper.docx"]}\n```\n\n## 摘要\n已写入文字，但绘图需要主线程资源。', []

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))
    try:
        result = await delegate._run_one_helper(
            task_id="paper_edit",
            prompt="生成含图论文",
            main_workspace=str(main_ws),
            helper_workspace=str(helper_ws),
            archive_id="archive",
            group_id="group",
            user_id="user",
            resume=False,
            local_abort=asyncio.Event(),
            wait_for_register=asyncio.Event(),
            user_lang="zh",
            kind="edit",
            mode="easy",
            helper_think=False,
            expected_outputs=["paper.docx"],
        )
    finally:
        monkeypatch.undo()

    assert result["ok"] is False
    assert result["terminal_reason"] == "resource_required"
    assert result["frozen"] is True
    assert result["resource_required"]["matching_helper_kind"] == "draw"
    assert "resource" in result["resource_required"]["resource_resolution_facts"]
    assert result["resource_required"]["suggested_helper_kind"] == "draw"
    assert "Suggested action" not in result["report"]
    assert "Resource kind fact" not in result["report"]
    assert result["outputs_check"]["outputs_complete"] is False
    assert result.get("_post_helper_action") != "output_json_directly"


async def test_delegate_reports_resource_required_without_auto_spawning(tmp_path):
    from app.core import debug
    from app.llm.tools import delegate

    debug.set_trace_id("trace-auto-resource")
    calls: list[tuple[str, str, bool]] = []

    async def fake_loop(msgs, *args, dispatcher=None, **kwargs):
        user_text = "\n".join(
            str(m.get("content") or "")
            for m in msgs
            if isinstance(m, dict) and m.get("role") == "user"
        )
        helper_kind = kwargs.get("helper_kind") or ""
        task_id = kwargs.get("task_id") or ""
        calls.append((task_id, helper_kind, "资源 helper" in user_text))
        if task_id == "paper_edit":
            requested = await dispatcher("request_resource", {
                "kind": "draw",
                "reason": "缺少论文图表 PNG",
                "needed_outputs": ["chart.png"],
                "resume_instruction": "嵌入 chart.png 后完成 paper.docx",
            })
            data = json.loads(requested)
            assert data["action"] == "request_resource"
            return '```json\n{"files": []}\n```', []
        return '```json\n{"files": []}\n```', []

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_persona_consent_guard", lambda *args, **kwargs: asyncio.sleep(0, result=(True, "", [], [], {})))
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))
    try:
        result = json.loads(await delegate.handle_delegate(
            str(tmp_path),
            {
                "action": "spawn",
                "wait_window_sec": 0,
                "tasks": [{
                    "task_id": "paper_edit",
                    "kind": "edit",
                    "prompt": "生成含图论文 paper.docx, 缺图时请求 draw 资源。",
                    "expected_outputs": ["paper.docx"],
                }],
            },
            archive_id="archive",
            group_id="group",
            user_id="user",
        ))
    finally:
        monkeypatch.undo()

    assert result["ok"] is True
    assert "auto_resource_recovery" not in result
    assert result["error_kind"] == "helper_resource_required"
    assert result["resource_required"][0]["task_id"] == "paper_edit"
    assert result["resource_required"][0]["matching_helper_kind"] == "draw"
    assert result["resource_required"][0]["suggested_helper_kind"] == "draw"
    by_tid = {r.get("task_id"): r for r in result["results"]}
    assert by_tid["paper_edit"]["ok"] is False
    assert by_tid["paper_edit"]["terminal_reason"] == "resource_required"
    assert not (tmp_path / "paper.docx").is_file()
    assert calls == [("paper_edit", "edit", False)]


async def test_delegate_batch_success_resource_pair_still_requires_consumer_resume(tmp_path):
    from app.core import debug
    from app.llm.tools import delegate

    debug.set_trace_id("trace-resource-pair")
    helper_chart_ws = tmp_path / "_delegate_user_chart_draw"
    helper_chart_ws.mkdir(parents=True)
    Image.effect_noise((400, 300), 80).convert("RGB").save(helper_chart_ws / "chart.png")

    async def fake_loop(msgs, *args, dispatcher=None, **kwargs):
        task_id = kwargs.get("task_id") or ""
        if task_id == "paper_edit":
            requested = await dispatcher("request_resource", {
                "kind": "draw",
                "reason": "need chart image before final docx assembly",
                "needed_outputs": ["chart.png"],
                "resume_instruction": "Embed chart.png into paper.docx and verify the final document.",
            })
            data = json.loads(requested)
            assert data["requires_main_resource"] is True
            return '```json\n{"files": []}\n```', []
        if task_id == "chart_draw":
            chart_path = helper_chart_ws / "chart.png"
            Image.effect_noise((400, 300), 80).convert("RGB").save(chart_path)
            return '```json\n{"files": ["chart.png"]}\n```', []
        return '```json\n{"files": []}\n```', []

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_persona_consent_guard", lambda *args, **kwargs: asyncio.sleep(0, result=(True, "", [], [], {})))
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))
    try:
        result = json.loads(await delegate.handle_delegate(
            str(tmp_path),
            {
                "action": "spawn",
                "wait_window_sec": 0,
                "tasks": [
                    {
                        "task_id": "chart_draw",
                        "kind": "draw",
                        "prompt": "Generate chart.png for the paper.",
                        "expected_outputs": ["chart.png"],
                    },
                    {
                        "task_id": "paper_edit",
                        "kind": "edit",
                        "prompt": "Write paper.docx using chart.png. If chart.png is not available, request draw resource and freeze.",
                        "expected_outputs": ["paper.docx"],
                    },
                ],
            },
            archive_id="archive",
            group_id="group",
            user_id="user",
        ))
    finally:
        monkeypatch.undo()

    assert result["ok"] is True
    assert result["task_ok"] is False
    assert result["_task_status"] == "incomplete"
    assert "task_ok=false" in result["_ok_field_meaning"]
    assert result["success_count"] == 1
    assert result["resource_required_count"] == 1
    assert result["incomplete_count"] == 1
    assert result["error_kind"] == "helper_resource_required"
    assert result["resource_required"][0]["task_id"] == "paper_edit"
    assert result["resource_required"][0]["matching_helper_kind"] == "draw"
    assert "describe failure or blocker state" in result["_evidence_policy"]
    assert "verified artifacts" in result["_evidence_policy"]
    by_tid = {r.get("task_id"): r for r in result["results"]}
    assert by_tid["chart_draw"]["ok"] is True
    assert by_tid["paper_edit"]["ok"] is False
    assert by_tid["paper_edit"]["terminal_reason"] == "resource_required"


async def test_delegate_treats_helper_partial_verdict_as_failed_status(tmp_path):
    from app.llm.tools import delegate

    async def fake_loop(*args, **kwargs):
        return (
            "VERDICT: PARTIAL\n"
            "The assigned helper kind cannot compute exact directory statistics from evidence.\n"
            "Suggested helper kind: code.\n"
            "```json\n{\"files\": []}\n```",
            [],
        )

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_persona_consent_guard", lambda *args, **kwargs: asyncio.sleep(0, result=(True, "", [], [], {})))
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))
    try:
        result = json.loads(await delegate.handle_delegate(
            str(tmp_path),
            {
                "action": "spawn",
                "wait_window_sec": 0,
                "tasks": [{
                    "task_id": "wrong_kind_stats",
                    "kind": "code",
                    "prompt": "Count exact Python files, lines, and characters in this directory.",
                }],
            },
            archive_id="archive",
            group_id="group",
            user_id="user",
        ))
    finally:
        monkeypatch.undo()

    assert result["task_ok"] is False
    assert result["success_count"] == 0
    assert result["failed_count"] == 1
    item = result["results"][0]
    assert item["ok"] is False
    assert item["report_verdict"] == "PARTIAL"
    assert item["terminal_reason"] == "failed"


async def test_interrupted_edit_helper_reports_partial_artifact_facts(tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / ".temp" / "_delegate_user_docx_partial"
    main_ws = tmp_path / "main"
    helper_ws.mkdir(parents=True)
    main_ws.mkdir()
    abort_event = asyncio.Event()

    async def fake_loop(*args, **kwargs):
        from docx import Document

        doc = Document()
        doc.add_heading("Partial Report", level=1)
        for idx in range(6):
            doc.add_paragraph(
                f"This file exists before helper final summary. Paragraph {idx + 1} contains enough "
                "content for the basic DOCX sanity checks to treat the artifact as a real document."
            )
        doc.save(helper_ws / "partial.docx")
        abort_event.set()
        return '```json\n{"files": ["partial.docx"]}\n```', []

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))
    try:
        result = await delegate._run_one_helper(
            task_id="docx_partial",
            prompt="Create partial.docx.",
            main_workspace=str(main_ws),
            helper_workspace=str(helper_ws),
            archive_id="archive",
            group_id="group",
            user_id="user",
            resume=False,
            local_abort=abort_event,
            wait_for_register=asyncio.Event(),
            user_lang="en",
            kind="edit",
            mode="easy",
            helper_think=False,
            expected_outputs=["partial.docx"],
        )
    finally:
        monkeypatch.undo()

    assert result["ok"] is True
    assert result["terminal_reason"] == "completed"
    assert result["_terminal_converged_from"] == "interrupted"
    assert (main_ws / "partial.docx").is_file()
    assert result["partial_artifacts"]["files"] == ["partial.docx"]
    assert result["partial_artifacts"]["status"] == "artifact_copied_before_helper_final_summary"
    assert "targeted inspection or verification only when needed" in result["partial_artifacts"]["fact"]


async def test_edit_helper_warns_when_pptx_missing_requested_text(tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / ".temp" / "_delegate_user_ppt_edit"
    main_ws = tmp_path / "main"
    helper_ws.mkdir(parents=True)
    main_ws.mkdir()

    async def fake_loop(*args, **kwargs):
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "数据页"
        prs.save(helper_ws / "missing_data.pptx")
        return '```json\n{"files": ["missing_data.pptx"]}\n```\n\n已检查通过。', []

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))
    try:
        result = await delegate._run_one_helper(
            task_id="ppt_edit",
            prompt="生成 PPT missing_data.pptx，数据页列出 X=12,Y=29，结论页写“链路正常”。",
            main_workspace=str(main_ws),
            helper_workspace=str(helper_ws),
            archive_id="archive",
            group_id="group",
            user_id="user",
            resume=False,
            local_abort=asyncio.Event(),
            wait_for_register=asyncio.Event(),
            user_lang="zh",
            kind="edit",
            mode="easy",
            helper_think=False,
            expected_outputs=["missing_data.pptx"],
        )
    finally:
        monkeypatch.undo()

    warnings = result["outputs_check"]["quality_warnings"]
    assert result["outputs_check"]["outputs_complete"] is True
    assert any(w["issue"] == "document_expected_text_missing" for w in warnings)
    # P135: document_expected_text_missing is a subjective warning, not blocking.
    # Helper completes; the LLM/orchestrator weighs the warning against intent.
    assert result["outputs_check"].get("quality_blocked") is not True
    assert result.get("quality_blocked") is not True
    assert result["ok"] is True
    assert result.get("terminal_reason") != "quality_blocked"
    missing = next(w["missing"] for w in warnings if w["issue"] == "document_expected_text_missing")
    assert "X=12" in missing
    assert "Y=29" in missing


async def test_edit_helper_does_not_require_forbidden_placeholder_words(tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / ".temp" / "_delegate_user_docx_forbidden_words"
    main_ws = tmp_path / "main"
    helper_ws.mkdir(parents=True)
    main_ws.mkdir()

    async def fake_loop(*args, **kwargs):
        from docx import Document

        doc = Document()
        doc.add_heading("测试指标趋势", level=1)
        doc.add_paragraph("数据为 A=10、B=14、C=13、D=19。整体呈上行态势。")
        doc.save(helper_ws / "no_placeholder.docx")
        return '```json\n{"files": ["no_placeholder.docx"]}\n```', []

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))
    try:
        result = await delegate._run_one_helper(
            task_id="docx_forbidden_words",
            prompt="生成 no_placeholder.docx，标题为“测试指标趋势”，写入数据 A=10、B=14、C=13、D=19。全文不得出现“占位”“稍后补图”“预留图表”等临时文字。",
            main_workspace=str(main_ws),
            helper_workspace=str(helper_ws),
            archive_id="archive",
            group_id="group",
            user_id="user",
            resume=False,
            local_abort=asyncio.Event(),
            wait_for_register=asyncio.Event(),
            user_lang="zh",
            kind="edit",
            mode="easy",
            helper_think=False,
            expected_outputs=["no_placeholder.docx"],
        )
    finally:
        monkeypatch.undo()

    warnings = result["outputs_check"]["quality_warnings"]
    missing_warnings = [w for w in warnings if w["issue"] == "document_expected_text_missing"]
    forbidden_missing = {
        item
        for w in missing_warnings
        for item in (w.get("missing") or [])
    }
    assert "占位" not in forbidden_missing
    assert "稍后补图" not in forbidden_missing
    assert "预留图表" not in forbidden_missing


async def test_edit_helper_does_not_require_forbidden_internal_path_words(tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / ".temp" / "_delegate_user_docx_forbidden_internal_paths"
    main_ws = tmp_path / "main"
    helper_ws.mkdir(parents=True)
    main_ws.mkdir()

    async def fake_loop(*args, **kwargs):
        from docx import Document

        doc = Document()
        doc.add_heading("Database Index Report", level=1)
        doc.add_paragraph("The document compares B+ trees, red-black trees, and skip lists.")
        doc.add_paragraph("Source handling is summarized without exposing implementation paths.")
        doc.save(helper_ws / "clean_paths.docx")
        return '```json\n{"files": ["clean_paths.docx"]}\n```', []

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))
    try:
        result = await delegate._run_one_helper(
            task_id="docx_forbidden_internal_paths",
            prompt=(
                "Generate clean_paths.docx with title \"Database Index Report\". "
                "Do not reference internal helper paths such as \"_env/\", \"_helpers_shared\", "
                "\"_helpers_shared/\", or code fragments like \"from docx import Document; d=Document(\"."
            ),
            main_workspace=str(main_ws),
            helper_workspace=str(helper_ws),
            archive_id="archive",
            group_id="group",
            user_id="user",
            resume=False,
            local_abort=asyncio.Event(),
            wait_for_register=asyncio.Event(),
            user_lang="zh",
            kind="edit",
            mode="easy",
            helper_think=False,
            expected_outputs=["clean_paths.docx"],
        )
    finally:
        monkeypatch.undo()

    missing = {
        item
        for warning in result["outputs_check"]["quality_warnings"]
        if warning["issue"] == "document_expected_text_missing"
        for item in (warning.get("missing") or [])
    }
    assert "_env/" not in missing
    assert "_helpers_shared" not in missing
    assert "_helpers_shared/" not in missing
    assert "from docx import Document; d=Document(" not in missing


async def test_edit_helper_blocks_unverified_academic_references(tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / ".temp" / "_delegate_user_paper_refs"
    main_ws = tmp_path / "main"
    helper_ws.mkdir(parents=True)
    main_ws.mkdir()

    async def fake_loop(*args, **kwargs):
        from docx import Document

        doc = Document()
        doc.add_heading("Database Index Algorithms: A Comparative Paper", level=1)
        doc.add_paragraph("Abstract. This paper compares B-trees, B+ trees, skip lists, and red-black trees.")
        doc.add_paragraph("The proposed hybrid index is evaluated conceptually against common database workloads [1].")
        doc.add_heading("References", level=1)
        doc.add_paragraph("[1] A. Researcher. Advanced Database Indexing. Journal of Databases, 2024.")
        doc.save(helper_ws / "paper.docx")
        return '```json\n{"files": ["paper.docx"]}\n```', []

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))
    try:
        result = await delegate._run_one_helper(
            task_id="paper_refs",
            prompt="生成一篇严谨论文 paper.docx，比较红黑树、跳表、B树和B+树，并提出一种新算法。",
            main_workspace=str(main_ws),
            helper_workspace=str(helper_ws),
            archive_id="archive",
            group_id="group",
            user_id="user",
            resume=False,
            local_abort=asyncio.Event(),
            wait_for_register=asyncio.Event(),
            user_lang="zh",
            kind="edit",
            mode="easy",
            helper_think=False,
            expected_outputs=["paper.docx"],
        )
    finally:
        monkeypatch.undo()

    warnings = result["outputs_check"]["quality_warnings"]
    assert any(w["issue"] == "academic_citation_unverified" for w in warnings)
    # P134: academic_citation_unverified is a warning, not blocking — the
    # judgment is subjective and the LLM should weigh it against the request.
    assert result["outputs_check"].get("quality_blocked") is not True
    assert result.get("quality_blocked") is not True
    assert result["ok"] is True
    assert result.get("terminal_reason") != "quality_blocked"


async def test_edit_helper_blocks_malformed_docx_tables(tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / ".temp" / "_delegate_user_paper_table"
    main_ws = tmp_path / "main"
    helper_ws.mkdir(parents=True)
    main_ws.mkdir()

    async def fake_loop(*args, **kwargs):
        from docx import Document

        doc = Document()
        doc.add_heading("算法比较报告", level=1)
        table = doc.add_table(rows=2, cols=8)
        table.cell(0, 0).text = "Algorithm"
        table.cell(0, 1).text = "Insert"
        table.cell(0, 2).text = "Delete"
        table.cell(0, 3).text = "Search"
        table.cell(0, 4).text = "Range"
        table.cell(0, 5).text = "Memory"
        table.cell(0, 6).text = "Concurrency"
        table.cell(0, 7).text = "Notes"
        table.cell(1, 0).text = "{'text': 'B+Tree', 'style': 'strong'}"
        table.cell(1, 7).text = "This is a paragraph-length explanation that should be prose rather than packed into a wide Word table cell. " * 4
        doc.save(helper_ws / "bad_table.docx")
        return '```json\n{"files": ["bad_table.docx"]}\n```', []

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))
    try:
        result = await delegate._run_one_helper(
            task_id="paper_table",
            prompt="生成算法比较 Word 报告 bad_table.docx，要求表格可读。",
            main_workspace=str(main_ws),
            helper_workspace=str(helper_ws),
            archive_id="archive",
            group_id="group",
            user_id="user",
            resume=False,
            local_abort=asyncio.Event(),
            wait_for_register=asyncio.Event(),
            user_lang="zh",
            kind="edit",
            mode="easy",
            helper_think=False,
            expected_outputs=["bad_table.docx"],
        )
    finally:
        monkeypatch.undo()

    issues = {w["issue"] for w in result["outputs_check"]["quality_warnings"]}
    assert "docx_table_cell_object_literal" in issues
    assert "docx_table_too_wide" in issues
    assert result["ok"] is False
    assert result["terminal_reason"] == "quality_blocked"


async def test_edit_helper_warns_when_docx_required_table_and_figure_counts_short(tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / ".temp" / "_delegate_user_docx_quantity"
    main_ws = tmp_path / "main"
    helper_ws.mkdir(parents=True)
    main_ws.mkdir()

    async def fake_loop(*args, **kwargs):
        from docx import Document

        doc = Document()
        doc.add_heading("Database Index Comparison", level=1)
        for i in range(6):
            doc.add_paragraph(
                f"Section paragraph {i}: this document has enough real prose for a short generated report."
            )
        for i in range(2):
            table = doc.add_table(rows=2, cols=2)
            table.cell(0, 0).text = "Metric"
            table.cell(0, 1).text = f"Table {i + 1}"
            table.cell(1, 0).text = "Lookup"
            table.cell(1, 1).text = "Measured"
        doc.save(helper_ws / "paper.docx")
        return '```json\n{"files": ["paper.docx"]}\n```', []

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))
    try:
        result = await delegate._run_one_helper(
            task_id="paper_quantity",
            prompt=(
                "Assemble paper.docx. At least 4 comparative tables are required. "
                "At minimum: 3 figures or charts must be present."
            ),
            main_workspace=str(main_ws),
            helper_workspace=str(helper_ws),
            archive_id="archive",
            group_id="group",
            user_id="user",
            resume=False,
            local_abort=asyncio.Event(),
            wait_for_register=asyncio.Event(),
            user_lang="zh",
            kind="edit",
            mode="easy",
            helper_think=False,
            expected_outputs=["paper.docx"],
        )
    finally:
        monkeypatch.undo()

    warnings = result["outputs_check"]["quality_warnings"]
    issues = {w["issue"] for w in warnings}
    assert "document_required_table_count_shortfall" in issues
    assert "document_required_figure_count_shortfall" in issues
    assert result["outputs_check"]["outputs_complete"] is True
    assert result["outputs_check"].get("quality_blocked") is not True
    assert result.get("quality_blocked") is not True
    assert result["ok"] is True


async def test_edit_helper_warns_when_pptx_slide_order_wrong(tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / ".temp" / "_delegate_user_ppt_order"
    main_ws = tmp_path / "main"
    helper_ws.mkdir(parents=True)
    main_ws.mkdir()

    async def fake_loop(*args, **kwargs):
        from pptx import Presentation

        prs = Presentation()
        s1 = prs.slides.add_slide(prs.slide_layouts[5])
        s1.shapes.title.text = "压测总结"
        s2 = prs.slides.add_slide(prs.slide_layouts[5])
        s2.shapes.title.text = "结论页"
        s2.shapes.add_textbox(0, 0, 4000000, 1000000).text = "链路正常"
        s3 = prs.slides.add_slide(prs.slide_layouts[5])
        s3.shapes.title.text = "数据页"
        s3.shapes.add_textbox(0, 0, 4000000, 1000000).text = "X=12\nY=29\nZ=18"
        prs.save(helper_ws / "wrong_order.pptx")
        return '```json\n{"files": ["wrong_order.pptx"]}\n```', []

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))
    try:
        result = await delegate._run_one_helper(
            task_id="ppt_order",
            prompt="生成 PPT wrong_order.pptx，共三页：第 1 页标题页，标题为“压测总结”；第 2 页数据页，清晰列出 X=12、Y=29、Z=18；第 3 页结论页写“链路正常”。",
            main_workspace=str(main_ws),
            helper_workspace=str(helper_ws),
            archive_id="archive",
            group_id="group",
            user_id="user",
            resume=False,
            local_abort=asyncio.Event(),
            wait_for_register=asyncio.Event(),
            user_lang="zh",
            kind="edit",
            mode="easy",
            helper_think=False,
            expected_outputs=["wrong_order.pptx"],
        )
    finally:
        monkeypatch.undo()

    warnings = result["outputs_check"]["quality_warnings"]
    assert result["outputs_check"]["outputs_complete"] is True
    assert any(w["issue"] == "pptx_expected_slide_order_mismatch" for w in warnings)
    # P135: pptx_expected_slide_order_mismatch is a subjective warning. Helper
    # completes; orchestrator decides whether the order matches user intent.
    assert result["outputs_check"].get("quality_blocked") is not True
    assert result.get("quality_blocked") is not True
    assert result["ok"] is True
    assert result.get("terminal_reason") != "quality_blocked"


async def test_edit_helper_does_not_require_pptx_role_words_on_slide(tmp_path):
    from app.llm.tools import delegate

    helper_ws = tmp_path / ".temp" / "_delegate_user_ppt_role_words"
    main_ws = tmp_path / "main"
    helper_ws.mkdir(parents=True)
    main_ws.mkdir()

    async def fake_loop(*args, **kwargs):
        from pptx import Presentation

        prs = Presentation()
        s1 = prs.slides.add_slide(prs.slide_layouts[5])
        s1.shapes.title.text = "压测总结"
        s2 = prs.slides.add_slide(prs.slide_layouts[5])
        s2.shapes.title.text = "数据页"
        s2.shapes.add_textbox(0, 0, 4000000, 1000000).text = "X=12\nY=29\nZ=18"
        s3 = prs.slides.add_slide(prs.slide_layouts[5])
        s3.shapes.title.text = "结论页"
        s3.shapes.add_textbox(0, 0, 4000000, 1000000).text = "链路正常"
        prs.save(helper_ws / "role_words_ok.pptx")
        return '```json\n{"files": ["role_words_ok.pptx"]}\n```', []

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(delegate, "chat_with_tools_loop", fake_loop)
    monkeypatch.setattr(delegate, "_copy_helper_debug_artifacts_to_main", lambda *args, **kwargs: None)
    monkeypatch.setattr(delegate, "_persist_pending_result", lambda *args, **kwargs: asyncio.sleep(0))
    try:
        result = await delegate._run_one_helper(
            task_id="ppt_role_words",
            prompt="生成 PPT role_words_ok.pptx，共三页：第 1 页标题页，标题为“压测总结”；第 2 页数据页，清晰列出 X=12、Y=29、Z=18；第 3 页结论页写“链路正常”。",
            main_workspace=str(main_ws),
            helper_workspace=str(helper_ws),
            archive_id="archive",
            group_id="group",
            user_id="user",
            resume=False,
            local_abort=asyncio.Event(),
            wait_for_register=asyncio.Event(),
            user_lang="zh",
            kind="edit",
            mode="easy",
            helper_think=False,
            expected_outputs=["role_words_ok.pptx"],
        )
    finally:
        monkeypatch.undo()

    warnings = result["outputs_check"]["quality_warnings"]
    assert result["outputs_check"]["outputs_complete"] is True
    assert not any(w["issue"] == "pptx_expected_slide_order_mismatch" for w in warnings)


def test_interrupted_helper_report_mentions_unmerged_shared_artifacts():
    copied_back = ["task1_report.docx"]
    declared_set = {"report.docx"}
    file_map = [{
        "helper_name": "report.docx",
        "main_name": "task1_report.docx",
        "shared_name": None,
    }]
    copy_stats = {"shared_merge_allowed": False}
    interrupted = True

    report_lines = ["done"]
    declared_to_copied = {
        str(m.get("helper_name")): str(m.get("main_name"))
        for m in file_map
        if isinstance(m, dict) and m.get("helper_name") and m.get("main_name")
    }
    declared_but_missing = sorted(
        name for name in declared_set
        if declared_to_copied.get(name, name) not in set(copied_back)
    )
    delivered_but_not_declared = []

    if declared_but_missing:
        report_lines.append(f"声明但未生成: {', '.join(declared_but_missing[:10])}")
    if delivered_but_not_declared:
        report_lines.append(f"未声明但已交付: {', '.join(delivered_but_not_declared[:10])}")
    if interrupted and not copy_stats.get("shared_merge_allowed", True):
        report_lines.append("共享支撑产物未合并: _helpers_shared/ 只保留在 helper 工作区")

    report = "\n".join(report_lines)
    assert "共享支撑产物未合并" in report
    assert "_helpers_shared/ 只保留在 helper 工作区" in report


def test_copy_results_recognizes_declared_helpers_shared_outputs_via_file_map(tmp_path):
    copied_back = ["_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png"]
    declared_set = {"_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png"}
    file_map = [{
        "helper_name": "q5_2_block_diagram.png",
        "main_name": "_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png",
        "shared_name": "_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png",
    }]

    copied_set = set(copied_back)
    declared_to_copied = {
        str(m.get("helper_name")): str(m.get("main_name"))
        for m in file_map
        if isinstance(m, dict) and m.get("helper_name") and m.get("main_name")
    }
    declared_to_copied.update({
        str(m.get("shared_name")): str(m.get("main_name"))
        for m in file_map
        if isinstance(m, dict) and m.get("shared_name") and m.get("main_name")
    })
    declared_but_missing = sorted(
        name for name in declared_set
        if declared_to_copied.get(name, name) not in copied_set
    )

    assert declared_but_missing == []


def test_copy_results_skips_helpers_shared_internal_metadata(tmp_path):
    from app.llm.tools.delegate import _copy_results_to_main

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    nested = helper_ws / "_helpers_shared" / "task"
    nested.mkdir(parents=True)
    main_ws.mkdir()
    (helper_ws / "_helpers_shared" / "step1_extract.py").write_text("print('debug')", encoding="utf-8")
    (nested / ".session_tag").write_text("tag", encoding="utf-8")
    (nested / ".todos_call_count.json").write_text("{}", encoding="utf-8")
    (nested / "task_.todos_call_count.json").write_text("{}", encoding="utf-8")
    (nested / "chart_.rewrite_count.json").write_text("{}", encoding="utf-8")
    (nested / "result.txt").write_text("done", encoding="utf-8")

    copied, _stats, _file_map = _copy_results_to_main(
        str(helper_ws),
        str(main_ws),
        task_id="task",
    )

    assert copied == ["_helpers_shared/task/result.txt"]
    assert (main_ws / "_helpers_shared" / "task" / "result.txt").is_file()
    assert not (main_ws / "_helpers_shared" / "step1_extract.py").exists()
    assert not (main_ws / "_helpers_shared" / "task" / ".session_tag").exists()
    assert not (main_ws / "_helpers_shared" / "task" / ".todos_call_count.json").exists()
    assert not (main_ws / "_helpers_shared" / "task" / "task_.todos_call_count.json").exists()
    assert not (main_ws / "_helpers_shared" / "task" / "chart_.rewrite_count.json").exists()


def test_copy_results_skips_helper_internal_counter_files(tmp_path):
    from app.llm.tools.delegate import _copy_results_to_main

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    helper_ws.mkdir()
    main_ws.mkdir()
    (helper_ws / ".todos_call_count.json").write_text("{}", encoding="utf-8")
    (helper_ws / ".rewrite_count.json").write_text("{}", encoding="utf-8")
    (helper_ws / "answer.txt").write_text("done", encoding="utf-8")

    copied, stats, file_map = _copy_results_to_main(
        str(helper_ws),
        str(main_ws),
        task_id="task1",
    )

    assert copied == ["task1_answer.txt"]
    assert not (main_ws / "task1_.todos_call_count.json").exists()
    assert not (main_ws / "task1_.rewrite_count.json").exists()
    assert file_map[0]["helper_name"] == "answer.txt"


def test_copy_results_matches_declared_prefixed_docx_by_basename(tmp_path):
    from app.llm.tools.delegate import _copy_results_to_main

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    helper_ws.mkdir()
    main_ws.mkdir()
    (helper_ws / "supplement_521_gen_docx_通信原理作业_第五章.docx").write_bytes(b"PK\x03\x04fake-docx")

    copied, stats, file_map = _copy_results_to_main(
        str(helper_ws),
        str(main_ws),
        task_id="supplement_521",
        declared_files={"gen_docx_通信原理作业_第五章.docx"},
    )

    assert copied == ["supplement_521_supplement_521_gen_docx_通信原理作业_第五章.docx"]
    assert stats["capped"] is False
    assert file_map[0]["helper_name"] == "supplement_521_gen_docx_通信原理作业_第五章.docx"


def test_copy_results_accepts_declared_helpers_shared_outputs_by_basename(tmp_path):
    from app.llm.tools.delegate import _copy_results_to_main

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    nested = helper_ws / "_helpers_shared" / "hw_redraw_charts"
    nested.mkdir(parents=True)
    main_ws.mkdir()
    (nested / "q5_2_block_diagram.png").write_bytes(b"\x89PNG\r\n\x1a\n" + b"x" * 3000)

    copied, stats, _file_map = _copy_results_to_main(
        str(helper_ws),
        str(main_ws),
        task_id="hw_redraw_charts",
        declared_files={"q5_2_block_diagram.png"},
    )

    assert copied == ["_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png"]
    assert stats["capped"] is False
    assert (main_ws / "_helpers_shared" / "hw_redraw_charts" / "q5_2_block_diagram.png").is_file()


def test_file_map_prevents_declared_missing_for_helper_prefixed_copy():
    copied_back = ["supplement_521_gen_docx_通信原理作业_第五章.docx"]
    declared_set = {"gen_docx_通信原理作业_第五章.docx"}
    file_map = [{
        "helper_name": "gen_docx_通信原理作业_第五章.docx",
        "main_name": "supplement_521_gen_docx_通信原理作业_第五章.docx",
    }]

    copied_set = set(copied_back)
    declared_to_copied = {
        str(m.get("helper_name")): str(m.get("main_name"))
        for m in file_map
        if isinstance(m, dict) and m.get("helper_name") and m.get("main_name")
    }
    declared_delivered_set = {
        declared_to_copied.get(name, name)
        for name in declared_set
    }
    declared_but_missing = sorted(
        name for name in declared_set
        if declared_to_copied.get(name, name) not in copied_set
    )
    delivered_but_not_declared = sorted(copied_set - declared_delivered_set)

    assert declared_but_missing == []
    assert delivered_but_not_declared == []


def test_file_map_matches_env_declared_output_aliases():
    from app.llm.tools.delegate import _matches_declared_output_via_mapping

    file_map = [{
        "helper_name": "reports/compression_report.docx",
        "main_name": "_env/reports/compression_report.docx",
        "shared_name": None,
    }]

    assert _matches_declared_output_via_mapping(
        "_env/reports/compression_report.docx",
        {"_env/reports/compression_report.docx"},
        file_map,
    )
    assert _matches_declared_output_via_mapping(
        "reports/compression_report.docx",
        {"_env/reports/compression_report.docx"},
        file_map,
    )


def test_delegate_result_files_hide_internal_helpers_shared_artifacts():
    from app.llm.tools.delegate import _is_internal_helper_artifact, _is_shared_support_artifact

    copied_back = [
        "_helpers_shared/.session_tag",
        "_helpers_shared/step1_extract.py",
        "_helpers_shared/bench_data.csv",
        "supplement_521_gen_docx_通信原理作业_第五章.docx",
    ]

    visible_copied_back = [
        path for path in copied_back
        if not _is_internal_helper_artifact(path)
        and not _is_shared_support_artifact(path)
    ]

    assert visible_copied_back == [
        "supplement_521_gen_docx_通信原理作业_第五章.docx",
    ]



def test_delegate_report_ignores_internal_helpers_shared_artifacts_and_prefixed_declared_outputs():
    from app.llm.tools.delegate import _is_internal_helper_artifact, _is_shared_support_artifact, _matches_declared_output_via_mapping

    copied_back = {
        "_helpers_shared/.session_tag",
        "_helpers_shared/bench_data.csv",
        "_helpers_shared/bench_full.csv",
        "_helpers_shared/step1_extract.py",
        "_helpers_shared/step1_extract_20260513_200429.py",
        "supplement_521_gen_docx_通信原理作业_第五章.docx",
    }
    declared_set = {"gen_docx_通信原理作业_第五章.docx"}
    file_map = [{
        "helper_name": "gen_docx_通信原理作业_第五章.docx",
        "main_name": "supplement_521_gen_docx_通信原理作业_第五章.docx",
        "shared_name": None,
    }]

    internal_delivered = {path for path in copied_back if _is_internal_helper_artifact(path)}
    support_delivered = {path for path in copied_back if _is_shared_support_artifact(path)}
    user_visible_delivered = copied_back - internal_delivered - support_delivered
    delivered_but_not_declared = sorted(
        path for path in user_visible_delivered
        if not _matches_declared_output_via_mapping(path, declared_set, file_map)
    )
    declared_to_copied = {
        str(m.get("helper_name")): str(m.get("main_name"))
        for m in file_map
        if isinstance(m, dict) and m.get("helper_name") and m.get("main_name")
    }
    visible_name_set = set(user_visible_delivered)
    declared_but_missing = sorted(
        name for name in declared_set
        if declared_to_copied.get(name, name) not in visible_name_set
        and not _matches_declared_output_via_mapping(
            declared_to_copied.get(name, name),
            {name},
            file_map,
        )
    )

    assert internal_delivered == {
        "_helpers_shared/.session_tag",
        "_helpers_shared/step1_extract.py",
        "_helpers_shared/step1_extract_20260513_200429.py",
    }
    assert support_delivered >= {
        "_helpers_shared/.session_tag",
        "_helpers_shared/bench_data.csv",
        "_helpers_shared/bench_full.csv",
        "_helpers_shared/step1_extract.py",
        "_helpers_shared/step1_extract_20260513_200429.py",
    }
    assert user_visible_delivered == {"supplement_521_gen_docx_通信原理作业_第五章.docx"}
    assert delivered_but_not_declared == []
    assert declared_but_missing == []


def test_delegate_report_treats_declared_helpers_shared_output_as_delivered_without_main_copy():
    from app.llm.tools.delegate import _matches_declared_output_via_mapping

    declared_set = {"_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png"}
    file_map = [{
        "helper_name": "q5_2_block_diagram.png",
        "main_name": "hw_redraw_charts_q5_2_block_diagram.png",
        "shared_name": "_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png",
    }]
    visible_copied_set = {"hw_redraw_charts_q5_2_block_diagram.png"}
    shared_to_main = {
        str(m.get("shared_name")): str(m.get("main_name"))
        for m in file_map
        if isinstance(m, dict) and m.get("shared_name") and m.get("main_name")
    }
    visible_name_set = set(visible_copied_set)
    visible_name_set.update(shared_to_main.keys())
    declared_to_copied = {
        str(m.get("helper_name")): str(m.get("main_name"))
        for m in file_map
        if isinstance(m, dict) and m.get("helper_name") and m.get("main_name")
    }
    for shared_name, main_name in shared_to_main.items():
        declared_to_copied[shared_name] = main_name

    declared_but_missing = sorted(
        name for name in declared_set
        if declared_to_copied.get(name, name) not in visible_name_set
        and not _matches_declared_output_via_mapping(
            declared_to_copied.get(name, name),
            {name},
            file_map,
        )
    )

    assert declared_but_missing == []


def test_delegate_report_counts_shared_outputs_as_declared_when_visible_only_via_shared_name():
    from app.llm.tools.delegate import _matches_declared_output_via_mapping

    declared_set = {"_helpers_shared/q518/q51_517_answers_51_517.txt"}
    file_map = [{
        "helper_name": "q51_517_answers_51_517.txt",
        "main_name": "q518_q51_517_answers_51_517.txt",
        "shared_name": "_helpers_shared/q518/q51_517_answers_51_517.txt",
    }]
    copied_back = ["q518_q51_517_answers_51_517.txt"]
    visible_name_set = set(copied_back)
    visible_name_set.add("_helpers_shared/q518/q51_517_answers_51_517.txt")
    declared_to_copied = {
        str(m.get("helper_name")): str(m.get("main_name"))
        for m in file_map
        if isinstance(m, dict) and m.get("helper_name") and m.get("main_name")
    }
    declared_to_copied["_helpers_shared/q518/q51_517_answers_51_517.txt"] = "q518_q51_517_answers_51_517.txt"
    declared_but_missing = sorted(
        name for name in declared_set
        if declared_to_copied.get(name, name) not in visible_name_set
        and not _matches_declared_output_via_mapping(
            declared_to_copied.get(name, name),
            {name},
            file_map,
        )
    )

    assert declared_but_missing == []


def test_copy_results_skips_root_internal_counter_files_without_dot_prefix(tmp_path):
    from app.llm.tools.delegate import _copy_results_to_main

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    helper_ws.mkdir()
    main_ws.mkdir()
    (helper_ws / "task_.todos_call_count.json").write_text("{}", encoding="utf-8")
    (helper_ws / "chart_.rewrite_count.json").write_text("{}", encoding="utf-8")
    (helper_ws / "note_history.json").write_text("{}", encoding="utf-8")
    (helper_ws / "note_count.json").write_text("{}", encoding="utf-8")
    (helper_ws / "answer.txt").write_text("done", encoding="utf-8")

    copied, _stats, file_map = _copy_results_to_main(
        str(helper_ws),
        str(main_ws),
        task_id="task1",
    )

    assert copied == ["task1_answer.txt"]
    assert not (main_ws / "task1_task_.todos_call_count.json").exists()
    assert not (main_ws / "task1_chart_.rewrite_count.json").exists()
    assert not (main_ws / "task1_note_history.json").exists()
    assert not (main_ws / "task1_note_count.json").exists()
    assert file_map == [{"helper_name": "answer.txt", "main_name": "task1_answer.txt", "shared_name": "_helpers_shared/task1/answer.txt"}]


def test_copy_results_filters_root_helpers_shared_internal_artifacts(tmp_path):
    from app.llm.tools.delegate import _copy_results_to_main

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    shared_root = helper_ws / "_helpers_shared"
    shared_root.mkdir(parents=True)
    main_ws.mkdir()
    (shared_root / ".session_tag").write_text("tag", encoding="utf-8")
    (shared_root / "step1_extract.py").write_text("print('debug')", encoding="utf-8")
    (shared_root / "bench_data.csv").write_text("a,b\n1,2\n", encoding="utf-8")
    (shared_root / "result.png").write_bytes(b"\x89PNG\r\n\x1a\n" + b"x" * 3000)

    copied, _stats, _file_map = _copy_results_to_main(
        str(helper_ws),
        str(main_ws),
        task_id="task1",
    )

    assert copied == ["_helpers_shared/bench_data.csv", "_helpers_shared/result.png"]
    assert not (main_ws / "_helpers_shared" / ".session_tag").exists()
    assert not (main_ws / "_helpers_shared" / "step1_extract.py").exists()


def test_copy_results_maps_declared_root_helpers_shared_python_outputs(tmp_path):
    from app.llm.tools.delegate import _copy_results_to_main, _matches_declared_output_via_mapping

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    shared_root = helper_ws / "_helpers_shared"
    shared_root.mkdir(parents=True)
    main_ws.mkdir()
    (shared_root / "reader.py").write_text("def read():\n    return 1\n", encoding="utf-8")

    copied, _stats, file_map = _copy_results_to_main(
        str(helper_ws),
        str(main_ws),
        task_id="ml_rebuild",
        declared_files={"_helpers_shared/reader.py"},
    )

    assert copied == ["_helpers_shared/reader.py"]
    assert (main_ws / "_helpers_shared" / "reader.py").is_file()
    assert file_map == [{
        "helper_name": "reader.py",
        "main_name": "_helpers_shared/reader.py",
        "shared_name": "_helpers_shared/reader.py",
    }]
    assert _matches_declared_output_via_mapping(
        "_helpers_shared/reader.py",
        {"_helpers_shared/reader.py"},
        file_map,
    )


def test_copy_results_repairs_reversible_text_mojibake(tmp_path):
    from app.llm.tools.delegate import _copy_results_to_main

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    helper_ws.mkdir()
    main_ws.mkdir()
    original = "\u4e8c\u4e00\u4e8c\u4e00\n"
    mojibake = original.encode("utf-8").decode("gbk")
    (helper_ws / "evidence.txt").write_text(mojibake, encoding="utf-8")

    copied, stats, _file_map = _copy_results_to_main(
        str(helper_ws),
        str(main_ws),
        task_id="read_ielts",
        declared_files={"evidence.txt"},
        expected_outputs=["evidence.txt"],
        helper_kind="read",
    )

    assert copied == ["evidence.txt"]
    assert (main_ws / "evidence.txt").read_text(encoding="utf-8") == original
    assert stats["text_mojibake_repaired"][0]["file"] == "evidence.txt"


def test_delegate_report_prefixed_delivered_doc_counts_as_declared_without_exact_delivered_set_match():
    from app.llm.tools.delegate import _matches_declared_output_via_mapping

    copied_back = {"hw_docx_final_通信原理作业_第五章.docx"}
    declared_set = {"通信原理作业_第五章.docx"}

    delivered_but_not_declared = sorted(
        path for path in copied_back
        if not _matches_declared_output_via_mapping(path, declared_set, [])
    )

    assert delivered_but_not_declared == []


def test_delegate_report_prefixed_declared_file_is_not_marked_missing_when_outputs_complete():
    from app.llm.tools.delegate import _matches_declared_output_via_mapping

    copied_set = {"supplement_521_gen_docx_通信原理作业_第五章.docx"}
    declared_set = {"gen_docx_通信原理作业_第五章.docx"}
    file_map = [{
        "helper_name": "gen_docx_通信原理作业_第五章.docx",
        "main_name": "supplement_521_gen_docx_通信原理作业_第五章.docx",
        "shared_name": None,
    }]
    declared_to_copied = {
        str(m.get("helper_name")): str(m.get("main_name"))
        for m in file_map
        if isinstance(m, dict) and m.get("helper_name") and m.get("main_name")
    }

    declared_but_missing = sorted(
        name for name in declared_set
        if declared_to_copied.get(name, name) not in copied_set
        and not _matches_declared_output_via_mapping(
            declared_to_copied.get(name, name),
            {name},
            file_map,
        )
    )

    assert declared_but_missing == []


def test_delegate_report_treats_helpers_shared_outputs_as_declared_when_paths_match():
    from app.llm.tools.delegate import _matches_declared_output_via_mapping

    declared_set = {
        "_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png",
        "_helpers_shared/hw_redraw_charts/q5_17_waveforms.png",
    }
    file_map = [{
        "helper_name": "q5_2_block_diagram.png",
        "main_name": "_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png",
        "shared_name": "_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png",
    }]

    assert _matches_declared_output_via_mapping(
        "_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png",
        declared_set,
        file_map,
    ) is True
    assert _matches_declared_output_via_mapping(
        "_helpers_shared/hw_redraw_charts/q5_17_waveforms.png",
        declared_set,
        file_map,
    ) is True


def test_dependency_detection_recognizes_helpers_shared_expected_outputs():
    completed_files = set()
    will_produce_files = {
        "q5_2_block_diagram.png",
        "_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png",
    }
    referenced = "q5_2_block_diagram.png".lower()

    found = (
        referenced in completed_files or
        referenced in will_produce_files or
        any(f.endswith("_" + referenced) for f in completed_files) or
        any(f.endswith("_" + referenced) for f in will_produce_files)
    )

    assert found is True


def test_delegate_outputs_check_ignores_internal_helper_artifacts_when_matching_expected_outputs():
    visible_copied_back = [
        "_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png",
        "_helpers_shared/hw_redraw_charts/q5_17_waveforms.png",
    ]
    copied_back = visible_copied_back + [
        "_helpers_shared/.session_tag",
        "hw_redraw_charts_.todos_call_count.json",
    ]
    expected_outputs = [
        "_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png",
        "_helpers_shared/hw_redraw_charts/q5_17_waveforms.png",
    ]

    delivered_basenames = {Path(f).name for f in visible_copied_back}
    delivered_paths = set(visible_copied_back)
    missing = []
    for exp in expected_outputs:
        exp_clean = exp.strip().lstrip("./").lstrip("\\")
        exp_base = Path(exp_clean).name
        exp_norm = exp_clean.replace("\\", "/")
        found_path = None
        if exp_clean in delivered_paths:
            found_path = exp_clean
        elif exp_base in delivered_basenames:
            found_path = next((f for f in visible_copied_back if Path(f).name == exp_base), None)
        else:
            found_path = next(
                (
                    f for f in visible_copied_back
                    if f.endswith(exp_clean)
                    or f.endswith(exp_base)
                    or f.endswith(exp_norm)
                    or f.endswith("_" + exp_base)
                ),
                None,
            )
        if not found_path:
            missing.append(exp)

    assert missing == []
    assert len(copied_back) > len(visible_copied_back)


def test_delegate_task_id_sanitizer_preserves_non_ascii_identity():
    from app.llm.tools.delegate import _sanitize_task_id

    first = _sanitize_task_id("read-代文静", 0)
    second = _sanitize_task_id("read-伍绍凡", 1)
    third = _sanitize_task_id("读取王謇聃组", 2)

    assert first.startswith("read-")
    assert second.startswith("read-")
    assert first != second
    assert third.startswith("task2_")
    assert all("/" not in value and "\\" not in value for value in (first, second, third))


def test_outputs_check_accepts_declared_helpers_shared_paths_without_main_copy(tmp_path):
    expected_outputs = [
        "_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png",
        "_helpers_shared/hw_redraw_charts/q5_17_waveforms.png",
    ]
    visible_copied_back = [
        "_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png",
        "_helpers_shared/hw_redraw_charts/q5_17_waveforms.png",
    ]

    delivered_basenames = {Path(f).name for f in visible_copied_back}
    delivered_paths = set(visible_copied_back)
    missing = []
    matched = []
    for exp in expected_outputs:
        exp_clean = exp.strip().lstrip("./").lstrip("\\")
        exp_base = Path(exp_clean).name
        exp_norm = exp_clean.replace("\\", "/")
        found_path = None
        if exp_clean in delivered_paths:
            found_path = exp_clean
        elif exp_base in delivered_basenames:
            found_path = next(
                (f for f in visible_copied_back if os.path.basename(f) == exp_base),
                None,
            )
        else:
            found_path = next(
                (
                    f
                    for f in visible_copied_back
                    if f.endswith(exp_clean)
                    or f.endswith(exp_base)
                    or f.endswith(exp_norm)
                    or f.endswith("_" + exp_base)
                ),
                None,
            )
        if found_path:
            matched.append((exp, found_path))
        else:
            missing.append(exp)

    assert [actual for _, actual in matched] == visible_copied_back
    assert missing == []


def test_unsatisfied_dependency_recognizes_stripped_helpers_shared_outputs():
    completed_files = set()
    delivered_files = [
        "_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png",
        "_helpers_shared/hw_redraw_charts/q5_17_waveforms.png",
    ]
    for path in delivered_files:
        norm = path.replace("\\", "/")
        base = Path(norm).name
        completed_files.add(base.lower())
        if norm.startswith("_helpers_shared/"):
            completed_files.add(norm.lower())
            completed_files.add(norm[len("_helpers_shared/"):].lower())
        if "_" in base:
            stripped = base.split("_", 1)[1] if base.count("_") <= 2 else base.split("_", base.count("_") - 1)[-1]
            completed_files.add(stripped.lower())

    referenced = {
        "q5_2_block_diagram.png",
        "q5_17_waveforms.png",
    }
    unsatisfied = []
    for ref in sorted(referenced):
        ref_l = ref.lower()
        found = (
            ref_l in completed_files
            or any(f.endswith("_" + ref_l) for f in completed_files)
        )
        if not found:
            unsatisfied.append(ref)

    assert unsatisfied == []


def test_copy_results_uses_expected_outputs_when_report_declares_nothing(tmp_path):
    from app.llm.tools.delegate import _copy_results_to_main

    helper_ws = tmp_path / "helper"
    main_ws = tmp_path / "main"
    shared_dir = helper_ws / "_helpers_shared" / "hw_redraw_charts"
    shared_dir.mkdir(parents=True)
    main_ws.mkdir()
    (shared_dir / "q5_2_block_diagram.png").write_bytes(b"\x89PNG\r\n\x1a\n" + b"x" * 3000)
    (shared_dir / "q5_17_waveforms.png").write_bytes(b"\x89PNG\r\n\x1a\n" + b"y" * 3000)

    copied, _stats, file_map = _copy_results_to_main(
        str(helper_ws),
        str(main_ws),
        task_id="hw_redraw_charts",
        declared_files={
            "_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png",
            "_helpers_shared/hw_redraw_charts/q5_17_waveforms.png",
        },
    )

    assert "_helpers_shared/hw_redraw_charts/q5_2_block_diagram.png" in copied
    assert "_helpers_shared/hw_redraw_charts/q5_17_waveforms.png" in copied
    assert all(m.get("main_name") != "hw_redraw_charts_q5_2_block_diagram.png" for m in file_map)


def test_ocr_bridge_does_not_feed_office_files_to_legacy_ocr(tmp_path, monkeypatch):
    from app.llm.tools import ocr_bridge

    docx_path = tmp_path / "sample.docx"
    docx_path.write_bytes(b"PK\x03\x04fake-docx")

    legacy_called = False

    def fake_legacy_headless(*args, **kwargs):
        nonlocal legacy_called
        legacy_called = True
        return ocr_bridge.OcrResult(ok=True, text="should not happen")

    monkeypatch.setattr(ocr_bridge, "_ocr_with_mineru_file", lambda path, timeout=120: ocr_bridge.OcrResult(ok=False, error="mineru failed"))
    monkeypatch.setattr(ocr_bridge, "_call_legacy_headless", fake_legacy_headless)

    result = ocr_bridge.ocr_file(docx_path)

    assert result.ok is False
    assert "legacy OCR does not support .docx files" in result.error
    assert legacy_called is False


def test_tool_call_args_normalize_escapes_multiline_prompt():
    from app.llm.client import _normalize_tool_call_args_for_dispatch

    raw = '{"action":"spawn","tasks":[{"task_id":"draw","prompt":"第一行\n第二行\n第三行","expected_outputs":["out.png"]}]}'

    args, err, repaired = _normalize_tool_call_args_for_dispatch(raw)

    assert err is None
    assert repaired is True
    assert args["tasks"][0]["prompt"] == "第一行\n第二行\n第三行"
    assert args["tasks"][0]["expected_outputs"] == ["out.png"]


def test_tool_call_args_normalize_escapes_inner_chinese_quotes():
    from app.llm.client import _normalize_tool_call_args_for_dispatch

    raw = '{"action":"spawn","tasks":[{"task_id":"eye","prompt":"解释其中"1"码和单一"眼睛"模式","expected_outputs":["answer.md"]}]}'

    args, err, repaired = _normalize_tool_call_args_for_dispatch(raw)

    assert err is None
    assert repaired is True
    assert args["tasks"][0]["prompt"] == '解释其中"1"码和单一"眼睛"模式'
    assert args["tasks"][0]["expected_outputs"] == ["answer.md"]


def test_tool_call_args_truncated_multiline_json_is_repaired_for_dispatch():
    from app.llm.client import _normalize_tool_call_args_for_dispatch

    raw = '{"action":"spawn","tasks":[{"task_id":"draw","prompt":"第一行\n第二行'

    args, err, repaired = _normalize_tool_call_args_for_dispatch(raw)

    assert err is None
    assert repaired is True
    assert args["action"] == "spawn"
    assert args["tasks"] == [{"task_id": "draw", "prompt": "第一行\n第二行"}]


def test_tool_args_json_broken_hint_guides_compact_delegate_requests():
    from app.llm.client import TOOL_ARGS_JSON_BROKEN_HINT

    assert "workspace manifest" in TOOL_ARGS_JSON_BROKEN_HINT
    assert "framework" in TOOL_ARGS_JSON_BROKEN_HINT
    assert "分批派发" in TOOL_ARGS_JSON_BROKEN_HINT


async def test_workspace_run_unwraps_quoted_cmd_c(tmp_path, monkeypatch):
    from app.llm.tools import workspace as ws_tool

    monkeypatch.setattr(ws_tool.sys, "platform", "win32")
    monkeypatch.setattr(ws_tool, "_translate_windows_command", lambda cmd, ws_dir: cmd)
    monkeypatch.setattr(ws_tool, "_check_bash_rate", lambda owner: asyncio.sleep(0, result=True))
    monkeypatch.setattr(ws_tool, "_is_main_thread", lambda: False)
    monkeypatch.setattr(ws_tool, "analyze_command", lambda *a, **k: SimpleNamespace(allowed=True, reason="", category="allow"))
    monkeypatch.setattr(ws_tool.shutil, "which", lambda *a, **k: None)
    monkeypatch.setattr(ws_tool, "_ensure_matplotlibrc", lambda path: None)

    captured = {}

    class FakeProc:
        pid = 12345
        returncode = 0

        async def communicate(self):
            return b"ok", b""

    async def fake_create_subprocess_exec(*args, **kwargs):
        captured["args"] = args
        captured["kwargs"] = kwargs
        return FakeProc()

    monkeypatch.setattr(ws_tool.asyncio, "create_subprocess_exec", fake_create_subprocess_exec)

    result = await ws_tool.handle_run(str(tmp_path), 'cmd /c "dir /s /b _helpers_shared\\charts\\*.png"')

    assert result["ok"] is True
    assert captured["args"][:2] == ("cmd", "/c")
    assert captured["args"][2] == "dir /s /b _helpers_shared\\charts\\*.png"


async def test_workspace_run_subprocess_permission_error_returns_tool_error(tmp_path, monkeypatch):
    from app.llm.tools import workspace as ws_tool

    monkeypatch.setattr(ws_tool.sys, "platform", "win32")
    monkeypatch.setattr(ws_tool, "_check_bash_rate", lambda owner: asyncio.sleep(0, result=True))
    monkeypatch.setattr(ws_tool, "_is_main_thread", lambda: False)
    monkeypatch.setattr(ws_tool, "analyze_command", lambda *a, **k: SimpleNamespace(allowed=True, reason="", category="allow"))
    monkeypatch.setattr(ws_tool.shutil, "which", lambda *a, **k: "python")
    monkeypatch.setattr(ws_tool, "_ensure_matplotlibrc", lambda path: None)

    async def fake_create_subprocess_exec(*args, **kwargs):
        raise PermissionError("blocked by policy")

    monkeypatch.setattr(ws_tool.asyncio, "create_subprocess_exec", fake_create_subprocess_exec)

    result = await ws_tool.handle_run(str(tmp_path), "python --version", timeout_sec=5)

    assert result["ok"] is False
    assert result["error_type"] == "PermissionError"
    assert "cmd /c fallback start failed" in result["error"]
    assert result["fallback_from"]["error_type"] == "PermissionError"
    assert "raw exec start failed" in result["fallback_from"]["error"]


def test_stuck_detector_missing_dependency_batch_injects_recovery_hint():
    from app.llm.tools.delegate_stuck import StuckDetector

    detector = StuckDetector("helper_benchmark")
    detector.record_batch([
        ("1", "read_file", '{"ok": false, "error": "file not found: _env/a.py"}', {"path": "_env/a.py"}),
        ("2", "read_file", '{"ok": false, "error": "file not found: _env/b.py"}', {"path": "_env/b.py"}),
        ("3", "read_file", '{"ok": false, "error": "file not found: _env/c.py"}', {"path": "_env/c.py"}),
        ("4", "read_file", '{"ok": false, "error": "file not found: _env/d.py"}', {"path": "_env/d.py"}),
    ])

    assert detector.stuck is False
    hint = detector.consume_soft_hint()
    assert hint is not None
    assert "missing_dependency_paths" in hint
    assert "file_map" in hint
    assert "request the exact resource" in hint


async def test_workspace_run_permission_error_retries_with_cmd_fallback(tmp_path, monkeypatch):
    from app.llm.tools import workspace as ws_tool

    monkeypatch.setattr(ws_tool.sys, "platform", "win32")
    monkeypatch.setattr(ws_tool, "_check_bash_rate", lambda owner: asyncio.sleep(0, result=True))
    monkeypatch.setattr(ws_tool, "_is_main_thread", lambda: False)
    monkeypatch.setattr(ws_tool, "analyze_command", lambda *a, **k: SimpleNamespace(allowed=True, reason="", category="allow"))
    monkeypatch.setattr(ws_tool.shutil, "which", lambda *a, **k: "node")
    monkeypatch.setattr(ws_tool, "_ensure_matplotlibrc", lambda path: None)

    calls = []

    class FakeProc:
        pid = 23456
        returncode = 0

        async def communicate(self):
            return b"v22.0.0\r\n", b""

    async def fake_create_subprocess_exec(*args, **kwargs):
        calls.append(args)
        if len(calls) == 1:
            raise PermissionError("blocked shim")
        return FakeProc()

    monkeypatch.setattr(ws_tool.asyncio, "create_subprocess_exec", fake_create_subprocess_exec)

    result = await ws_tool.handle_run(str(tmp_path), "node --version", timeout_sec=5)

    assert result["ok"] is True
    assert result["stdout"].strip() == "v22.0.0"
    assert calls[0][0] == "node"
    assert calls[1][:2] == ("cmd", "/c")
    assert calls[1][2] == "node --version"


async def test_workspace_run_blocks_bare_rename_outside_workspace(tmp_path, monkeypatch):
    from app.llm.tools import workspace as ws_tool

    monkeypatch.setattr(ws_tool, "_check_bash_rate", lambda owner: asyncio.sleep(0, result=True))
    monkeypatch.setattr(ws_tool, "_is_main_thread", lambda: False)
    monkeypatch.setattr(ws_tool, "analyze_command", lambda *a, **k: SimpleNamespace(allowed=True, reason="", category="allow"))

    result = await ws_tool.handle_run(
        str(tmp_path),
        r"ren C:\Users\win11\AppData\Local\Programs\Python\Python312\Lib\site-packages\bad.pth bad.off",
    )

    assert result["ok"] is False
    error = str(result["error"])
    assert any(marker in error for marker in ("outside .temp", "outside workspace", "blocked"))


async def test_list_generated_files_filters_internal_metadata_and_build_artifacts(tmp_path):
    from app.llm.tools.workspace import list_generated_files

    (tmp_path / "answer.txt").write_text("ok", encoding="utf-8")
    (tmp_path / "chart.png").write_bytes(b"\x89PNG\r\n\x1a\n" + b"x" * 32)
    (tmp_path / "bench_runner.exe").write_bytes(b"MZ" + b"x" * 32)
    (tmp_path / "bench_runner.o").write_bytes(b"obj")
    (tmp_path / "build_output.txt").write_text("noise", encoding="utf-8")
    (tmp_path / ".session_tag").write_text("tag", encoding="utf-8")
    (tmp_path / "task_.todos_call_count.json").write_text("{}", encoding="utf-8")
    (tmp_path / "chart_.rewrite_count.json").write_text("{}", encoding="utf-8")
    shared = tmp_path / "_helpers_shared" / "task"
    shared.mkdir(parents=True)
    (shared / "task_.todos_call_count.json").write_text("{}", encoding="utf-8")
    (shared / "chart_.rewrite_count.json").write_text("{}", encoding="utf-8")
    (shared / "plot.png").write_bytes(b"\x89PNG\r\n\x1a\n" + b"x" * 32)
    delegate = tmp_path / "_delegate_task"
    delegate.mkdir()
    (delegate / "delegate_answer.txt").write_text("hidden", encoding="utf-8")

    files = list_generated_files(str(tmp_path))

    assert "answer.txt" in files
    assert "chart.png" in files
    assert "bench_runner.exe" in files
    assert "_helpers_shared/task/plot.png" in files
    assert "bench_runner.o" not in files
    assert "build_output.txt" not in files
    assert ".session_tag" not in files
    assert "task_.todos_call_count.json" not in files
    assert "chart_.rewrite_count.json" not in files
    assert "_helpers_shared/task/task_.todos_call_count.json" not in files
    assert "_helpers_shared/task/chart_.rewrite_count.json" not in files
    assert "_delegate_task/delegate_answer.txt" not in files


async def test_workspace_locate_filters_internal_metadata_files(tmp_path):
    from app.llm.tools.workspace import handle_locate

    (tmp_path / "answer.txt").write_text("ok", encoding="utf-8")
    (tmp_path / "build_output.txt").write_text("noise", encoding="utf-8")
    (tmp_path / "module.o").write_bytes(b"obj")
    (tmp_path / ".session_tag").write_text("tag", encoding="utf-8")
    (tmp_path / "task_.todos_call_count.json").write_text("{}", encoding="utf-8")
    (tmp_path / "chart_.rewrite_count.json").write_text("{}", encoding="utf-8")
    shared = tmp_path / "_helpers_shared" / "task"
    shared.mkdir(parents=True)
    (shared / "task_.todos_call_count.json").write_text("{}", encoding="utf-8")
    (shared / "plot.png").write_bytes(b"\x89PNG\r\n\x1a\n" + b"x" * 3000)

    result = await handle_locate(str(tmp_path), "*.*")
    paths = {m["path"] for m in result["matches"]}

    assert result["total"] == 2
    assert result["returned"] == 2
    assert result["truncated"] is False
    assert "answer.txt" in paths
    assert "_helpers_shared/task/plot.png" in paths
    assert "build_output.txt" not in paths
    assert "module.o" not in paths
    assert ".session_tag" not in paths
    assert "task_.todos_call_count.json" not in paths
    assert "chart_.rewrite_count.json" not in paths
    assert "_helpers_shared/task/task_.todos_call_count.json" not in paths


async def test_workspace_locate_total_counts_all_matches_when_truncated(tmp_path, monkeypatch):
    from app.llm.tools import workspace as ws_tool

    monkeypatch.setattr(ws_tool, "_LOCATE_MAX_MATCHES", 3)
    for i in range(5):
        (tmp_path / f"match_{i}.txt").write_text(str(i), encoding="utf-8")
    for i in range(4):
        (tmp_path / f"ignored_{i}_.todos_call_count.json").write_text("{}", encoding="utf-8")

    result = await ws_tool.handle_locate(str(tmp_path), "*.txt")

    assert result["total"] == 5
    assert result["returned"] == 3
    assert result["truncated"] is True
    assert len(result["matches"]) == 3


async def test_chat_with_tools_loop_text_summary_natural_stop_skips_json_cleanup(monkeypatch):
    from app.llm import client as llm_client
    from app.llm.client import chat_with_tools_loop

    summary = "| file | status |\n| a.png | done |"
    cleanup_called = False

    async def fake_streaming_call(**kwargs):
        message = SimpleNamespace(role="assistant", content=summary, tool_calls=None, reasoning_content=None)
        response = SimpleNamespace(choices=[SimpleNamespace(message=message)])
        return response, SimpleNamespace(), "ok"

    class FakeCompletions:
        async def create(self, **kwargs):
            nonlocal cleanup_called
            cleanup_called = True
            raise AssertionError("text_summary natural stop should not request JSON cleanup")

    fake_client = SimpleNamespace(chat=SimpleNamespace(completions=FakeCompletions()))

    monkeypatch.setattr(llm_client, "_call_llm_streaming_with_idle", fake_streaming_call)
    monkeypatch.setattr(llm_client, "client", lambda: fake_client)

    content, _msgs = await chat_with_tools_loop(
        [{"role": "user", "content": "report outputs"}],
        tools=[],
        dispatcher=lambda name, args: "{}",
        finalize_kind="text_summary",
    )

    assert content == summary
    assert cleanup_called is False


def test_markdown_algorithm_sections_stay_edit_not_code():
    from app.llm.tools.delegate import _deterministic_kind_recommendations

    recs = _deterministic_kind_recommendations([
        {
            "task_id": "slice_complexity_tables",
            "kind": "edit",
            "prompt": (
                "Write the paper section comparing Red-Black Tree, Skip List, B-Tree, "
                "and B+Tree complexity. Produce Markdown tables and prose analysis only; "
                "do not run benchmarks or generate code."
            ),
            "expected_outputs": ["_env/slice_b_complexity_tables.md"],
        },
        {
            "task_id": "slice_new_algorithm",
            "kind": "edit",
            "prompt": (
                "Draft the paper section defining a new database index data structure, "
                "including motivation, operations, complexity discussion, and limitations. "
                "This is prose-only research writing."
            ),
            "expected_outputs": ["_env/slice_c_new_algorithm.md"],
        },
    ])

    assert recs == []


def test_benchmark_outputs_still_require_code_helper():
    from app.llm.tools.delegate import _deterministic_kind_recommendations

    recs = _deterministic_kind_recommendations([
        {
            "task_id": "slice_benchmark",
            "kind": "edit",
            "prompt": (
                "Create benchmark.py and generate benchmark_results.csv for comparing "
                "database index data structures."
            ),
            "expected_outputs": ["_env/benchmark/benchmark.py", "_env/benchmark_results.csv"],
        }
    ])

    assert recs
    assert recs[0]["task_id"] == "slice_benchmark"
    assert recs[0]["observed_helper_kind_name"] == "code"


def test_strict_reported_output_files_preserve_env_paths():
    from app.llm.tools.workspace_utils import _extract_reported_output_files

    report = '## Output files\n```json\n{"files": ["_env/framework.md", "_env/criteria.md"]}\n```'

    assert _extract_reported_output_files(report) == ["_env/framework.md", "_env/criteria.md"]


def test_reported_output_files_accept_output_section_bare_json():
    from app.llm.tools.workspace_utils import _extract_reported_output_files

    report = '## Output files\n{"files": ["_env/framework.md", "_env/criteria.md"]}'

    assert _extract_reported_output_files(report) == ["_env/framework.md", "_env/criteria.md"]




# ── 2026-06-10 edit-helper restructure round: guard evidence, provenance, one-pass verify ──


def test_helper_source_field_provenance_skips_prose_only_time_words():
    """Bare schedule words without digits must not trigger the raw-field directive."""
    from app.llm.tools.delegate import _mark_unverified_helper_source_interpretations

    original = (
        "Create triage_report.md. The draft reply should promise a response within the hour "
        "and ask which endpoints are affected. Count every covered email in the report body."
    )
    prompt, observations = _mark_unverified_helper_source_interpretations(
        original,
        input_files=["classification_report.txt"],
    )

    assert prompt == original
    assert observations == []


def test_helper_source_field_provenance_still_marks_numeric_field_summaries():
    from app.llm.tools.delegate import _mark_unverified_helper_source_interpretations

    prompt, observations = _mark_unverified_helper_source_interpretations(
        "Write plan.md. The hotel summary says total cost is 450 USD for 3 nights.",
        input_files=["hotel.json"],
    )

    assert "## Source Field Provenance" in prompt
    assert observations == [{"kind": "source_field_provenance_added", "input_files": 1}]


def test_task_quality_guard_system_recognizes_completed_extraction_and_candidate_reuse():
    from app.llm.aux_prompts import TASK_QUALITY_GUARD_SYSTEM

    assert "Completed extraction does not need to be repeated" in TASK_QUALITY_GUARD_SYSTEM
    assert "do not demand a fresh read helper for materials the main workflow has already read" in TASK_QUALITY_GUARD_SYSTEM
    assert "preserved blocked-create candidate" in TASK_QUALITY_GUARD_SYSTEM
    assert "provenance and explicit input facts" in TASK_QUALITY_GUARD_SYSTEM
    assert "visibly adapts to the earlier guidance" in TASK_QUALITY_GUARD_SYSTEM


def test_edit_helper_prompt_trusts_envelope_and_verifies_in_one_pass():
    from app.llm.tools.delegate import _HELPER_SYSTEM_EDIT

    assert "Trust the task envelope" in _HELPER_SYSTEM_EDIT
    assert "read each at most once" in _HELPER_SYSTEM_EDIT
    assert "Verify your own writes in one pass" in _HELPER_SYSTEM_EDIT
    assert "per-keyword search calls against a file you authored this turn" in _HELPER_SYSTEM_EDIT
    assert "一次回读或一条检查命令完成验收" in _HELPER_SYSTEM_EDIT


def test_blocked_create_candidate_reuse_guard_fact_attached():
    from app.llm.tools.delegate_actions import (
        _attach_guard_attention_facts,
        _blocked_create_candidate_reuse_facts,
    )

    specs = [{
        "task_id": "report-assembly",
        "kind": "edit",
        "prompt": "Assemble the final report from the preserved candidate.",
        "input_files": ["_env/.blocked_creates/triage_report.md"],
        "expected_outputs": ["_env/triage_report.md"],
    }]

    facts = _blocked_create_candidate_reuse_facts(specs)
    assert len(facts) == 1
    assert facts[0]["issue"] == "blocked_create_candidate_reuse"
    assert facts[0]["needs_attention"] is False
    assert facts[0]["workspace_input_files"] == ["_env/.blocked_creates/triage_report.md"]

    attached = _attach_guard_attention_facts(specs)
    observations = attached[0].get("guard_observations") or []
    reuse = [o for o in observations if o.get("issue") == "blocked_create_candidate_reuse"]
    assert len(reuse) == 1
    assert reuse[0]["needs_attention"] is False
    assert "availability and provenance fact" in reuse[0]["details"]


def test_blocked_create_candidate_reuse_fact_absent_for_plain_inputs():
    from app.llm.tools.delegate_actions import _blocked_create_candidate_reuse_facts

    assert _blocked_create_candidate_reuse_facts([{
        "task_id": "t",
        "input_files": ["_env/inbox/01.txt"],
    }]) == []


@pytest.mark.asyncio
async def test_read_file_tool_result_spill_returns_provenance_fact(tmp_path):
    from app.llm.tools.registry import _handle_read_file

    spill_dir = tmp_path / "_tool_results"
    spill_dir.mkdir()
    (spill_dir / "1781062124_read_file_content_abc.txt").write_text(
        "recovered oversized tool output\n", encoding="utf-8"
    )

    result = json.loads(await _handle_read_file(
        str(tmp_path), {"path": "_tool_results/1781062124_read_file_content_abc.txt"}
    ))

    assert result["ok"] is True
    assert result["provenance"] == "tool_result_spill"
    assert "recovery storage" in result["provenance_fact"]
    assert "recovered oversized tool output" in result["content"]


@pytest.mark.asyncio
async def test_read_file_blocked_create_candidate_returns_provenance_fact(tmp_path):
    from app.llm.tools.registry import _handle_read_file

    candidate_dir = tmp_path / "_env" / ".blocked_creates"
    candidate_dir.mkdir(parents=True)
    (candidate_dir / "triage_report.md").write_text("# preserved candidate\n", encoding="utf-8")

    result = json.loads(await _handle_read_file(
        str(tmp_path), {"path": "_env/.blocked_creates/triage_report.md"}
    ))

    assert result["ok"] is True
    assert result["provenance"] == "blocked_create_candidate"
    assert "preserved candidate" in result["provenance_fact"]
    assert "# preserved candidate" in result["content"]


@pytest.mark.asyncio
async def test_read_file_normal_paths_have_no_recovery_provenance(tmp_path):
    from app.llm.tools.registry import _handle_read_file

    (tmp_path / "notes.md").write_text("plain note\n", encoding="utf-8")
    result = json.loads(await _handle_read_file(str(tmp_path), {"path": "notes.md"}))

    assert result["ok"] is True
    assert "provenance" not in result

# ── 2026-06-10 speed-ratio round: apply auto-registration, compact agent_state, finalize/batching prompts ──


async def test_agent_state_register_artifact_dedupes_already_ready_path():
    from app.core import agent_state, debug
    from app.llm.tools import registry

    trace_id = "trace-register-dedupe"
    debug.set_trace_id(trace_id)
    agent_state.reset_trace(trace_id)
    try:
        agent_state.register_artifact(
            trace_id=trace_id,
            path="contracts/customer_event.py",
            artifact_type="code",
            created_by="env_apply_replace",
            status=agent_state.ARTIFACT_READY,
            verified_by="env_apply_replace",
        )

        result = json.loads(await registry._handle_agent_state({
            "action": "register_artifact",
            "path": "contracts/customer_event.py",
            "status": "ready",
            "task_id": "main",
        }))
    finally:
        debug.set_trace_id("")

    assert result["ok"] is True
    assert result["already_registered"] is True
    assert result["artifact"]["path"] == "contracts/customer_event.py"
    assert "status" not in result
    # No duplicate record was appended.
    ready = agent_state.structured_status(trace_id).get("artifacts_ready") or []
    assert len([r for r in ready if r.get("path") == "contracts/customer_event.py"]) == 1


async def test_agent_state_register_artifact_new_path_returns_compact_summary():
    from app.core import agent_state, debug
    from app.llm.tools import registry

    trace_id = "trace-register-compact"
    debug.set_trace_id(trace_id)
    agent_state.reset_trace(trace_id)
    try:
        result = json.loads(await registry._handle_agent_state({
            "action": "register_artifact",
            "path": "report/summary.md",
            "status": "ready",
        }))
    finally:
        debug.set_trace_id("")

    assert result["ok"] is True
    assert result["artifact"]["path"] == "report/summary.md"
    assert "status" not in result
    assert result["status_summary"]["artifacts_ready_paths"] == ["report/summary.md"]


def test_round2_prompt_requires_direct_final_json_and_batched_bookkeeping():
    from app.core.round_prompts import ROUND2_SYSTEM_TEMPLATE

    assert "the very next assistant message is this JSON object itself" in ROUND2_SYSTEM_TEMPLATE
    assert "contract self-assessment" in ROUND2_SYSTEM_TEMPLATE
    assert "should share a turn with the next action's tool calls" in ROUND2_SYSTEM_TEMPLATE
    assert "不先写自评" in ROUND2_SYSTEM_TEMPLATE


def test_round2_prompt_prefers_helpers_for_small_edits():
    from app.core.round_prompts import ROUND2_SYSTEM_TEMPLATE

    assert "Even for small edits, prefer a helper" in ROUND2_SYSTEM_TEMPLATE
    assert "source/project content, or requires quality judgment" in ROUND2_SYSTEM_TEMPLATE
    assert "narrow mechanical transfer/apply/accounting" in ROUND2_SYSTEM_TEMPLATE


def test_round2_prompt_uses_producer_boundary_for_final_checks():
    """Regression for helper-owned artifacts: main should trust clean helper
    completion instead of duplicating producer validation after transfer/apply."""
    from app.core.round_prompts import ROUND2_SYSTEM_TEMPLATE

    assert "clean helper completion with declared outputs present and producer self-verification" in ROUND2_SYSTEM_TEMPLATE
    assert "main thread should trust it and avoid re-reading or re-running checks" in ROUND2_SYSTEM_TEMPLATE
    assert "helper self-verification evidence" in ROUND2_SYSTEM_TEMPLATE
    assert "final intended applied state" in ROUND2_SYSTEM_TEMPLATE
    assert "a check before later applies validates an earlier state" in ROUND2_SYSTEM_TEMPLATE
    assert "主进程机械应用/转移不接管内容质量" in ROUND2_SYSTEM_TEMPLATE


def test_round1_prompt_keeps_routing_output_minimal():
    """round1 gates the whole turn serially; at congested decode speeds its
    ~190-token output cost ~14s of user-visible latency per run."""
    from app.core.round_prompts import ROUND1_SYSTEM

    assert "one short sentence naming the decisive routing facts" in ROUND1_SYSTEM
    assert "every extra output token adds user-visible latency" in ROUND1_SYSTEM
    assert '"rationale"' not in ROUND1_SYSTEM
    assert "_thinking 一句话即可" in ROUND1_SYSTEM


def test_round2_prompt_skips_known_red_baseline_in_main_thread():
    """20260610_140711: main ran the known-failing pytest baseline before
    delegating; the failed env_run had no same-family success within the
    recovery window and recovery_score dropped to 0.3."""
    from app.core.round_prompts import ROUND2_SYSTEM_TEMPLATE

    assert "a known-red baseline adds a failed call without new facts" in ROUND2_SYSTEM_TEMPLATE
    assert "let the helper own baseline, diagnosis, all requested outputs, and the passing rerun" in ROUND2_SYSTEM_TEMPLATE
    assert "主线程不先跑红基线" in ROUND2_SYSTEM_TEMPLATE


def test_environment_helper_prompt_maps_bare_input_files_to_staged_paths():
    src = Path("app/llm/tools/delegate.py").read_text(encoding="utf-8")

    assert "first try the staged local path `_env/app.py`" in src
    assert "same-named bare path may be absent" in src
    assert "裸项目路径先查 `_env/<路径>`" in src


@pytest.mark.asyncio
async def test_wait_any_recovers_already_collected_task_from_ledger():
    """Interlock regression: _consume_pending_result clears the completion
    event, so a second wait_any on an already-collected task_id used to block
    on a cleared event for the full window. With no active helper, wait_any
    must recover from the completion ledger immediately."""
    import time as _t
    from app.llm.tools import delegate_state as ds
    from app.llm.tools.delegate_actions import _handle_delegate_wait_any, _sync_delegate_globals
    from app.core import debug

    _sync_delegate_globals()
    trace_id = "trace-waitany-ledger"
    tid = "already-collected"
    debug.set_trace_id(trace_id)
    try:
        await ds._store_pending_result(
            trace_id, tid,
            {"task_id": tid, "ok": True, "terminal_reason": "completed",
             "report": "done", "files": ["out.txt"]},
            None,
        )
        first = await ds._consume_pending_result(trace_id, tid)
        assert first is not None

        start = _t.monotonic()
        raw = await _handle_delegate_wait_any(
            {"task_ids": [tid], "wait_window_sec": 30},
            main_owner="owner-x", trace_id=trace_id,
        )
        elapsed = _t.monotonic() - start
        result = json.loads(raw)
    finally:
        debug.set_trace_id("")

    assert elapsed < 5, f"wait_any blocked {elapsed:.1f}s on a cleared event"
    assert result["ok"] is True
    assert result["winner_task_id"] == tid
    assert result.get("recovered_from_ledger") is True
    assert result["result"]["recovered_from"] == "completion_ledger"


def test_final_contract_snapshot_skipped_for_complete_response_plan():
    """Round 7: injecting the contract snapshot AFTER a finished ResponsePlan
    forces a rewrite that reliably degrades into self-assessment prose plus a
    cleanup LLM call (cross-repo 20260610_154444 iters 13-15)."""
    from app.llm.client_tools_loop import _content_is_complete_response_plan

    plan = json.dumps({
        "intent": "Migration complete",
        "key_points": ["a", "b"],
        "tone": "rigorous-controlled",
        "length_hint": "short",
        "deliverables": ["x.py"],
    })
    assert _content_is_complete_response_plan(plan) is True
    assert _content_is_complete_response_plan("All criteria satisfied.") is False
    assert _content_is_complete_response_plan(json.dumps({
        "final_json_status": "complete", "contract_complete": True,
    })) is False
    assert _content_is_complete_response_plan("") is False


def test_document_routing_keeps_text_deliverables_on_edit_kind():
    """inbox-triage runs keep trying kind=code first for classification/draft
    work because verify_*.py scripts exist; each attempt costs a guard block
    round-trip (~20s with main-LLM latency)."""
    from app.core.orchestrator_prompts import _build_round2_system_prompts

    msgs = _build_round2_system_prompts(
        is_coding=False, is_document=True, parallelizable=False, needs_recall=False
    )
    text = "\n".join(m["content"] for m in msgs)
    assert "Classification, triage, drafting, summarizing, and report writing" in text
    assert "running an existing verifier needs no code helper" in text
    assert "not by the presence of `.py` checkers" in text
    assert "predictably costs a guard block round-trip" in text


def test_edit_helper_run_guard_redirects_file_read_shells_to_read_file():
    """20260610_161417: an edit helper ran `cat "_env/.blocked_creates/x.txt"`
    and got a generic delegate-to-code block; the failed call dinged the
    trajectory recovery score. Plain read shells now redirect to read_file."""
    from app.llm.tools.registry import _edit_helper_workspace_run_guard

    blocked = _edit_helper_workspace_run_guard("edit", 'cat "_env/.blocked_creates/report.txt"')
    assert blocked is not None
    assert blocked["error"] == "use_read_file_for_file_reads"
    assert blocked["suggested_tool"] == "read_file"
    assert blocked["suggested_args"]["path"] == "_env/.blocked_creates/report.txt"

    # verifier commands still pass
    assert _edit_helper_workspace_run_guard("edit", "python verify_all.py") is None
    # builds still blocked with the original error
    build = _edit_helper_workspace_run_guard("edit", "npm run build")
    assert build is not None
    assert build["error"] == "edit_helper_workspace_run_limited_to_existing_checks"
    # non-edit helpers unaffected
    assert _edit_helper_workspace_run_guard("code", "cat x.txt") is None


def test_helper_tech_prompt_warns_python_tool_has_no_file_io():
    """20260610_161418: a code helper used the isolated python tool with open()
    to count occurrences across files and burned a failed call; the sandbox is
    by design, so the prompt must say it up front."""
    from app.llm.tools.delegate import _HELPER_SYSTEM_CODE

    assert "isolated in-memory calculator with no file access" in _HELPER_SYSTEM_CODE
    assert "search_in_file/search_across_files" in _HELPER_SYSTEM_CODE


def test_language_directive_covers_english_users():
    """20260610_163156: an English request got a Chinese final reply because
    lang='en' produced an empty directive while the persona defaults Chinese.
    Behavior score dropped (blocker wording invisible to English scorer)."""
    from app.core.language import language_directive

    en = language_directive("en")
    assert "should be in English" in en
    assert "## Output Language" in en
    assert language_directive("zh")  # zh stays non-empty


async def test_read_file_traversal_block_suggests_env_read_in_environment_mode(tmp_path):
    r"""20260610_165331: 11 read_file calls used backslash traversal chains
    (multiple parent hops) for project files; each got a bare traversal error
    with no recovery route. In environment mode the tail is a valid path."""
    from app.core.runtime_mode import EnvironmentContext, runtime_context
    from app.llm.tools.workspace_file_ops import handle_read_file

    ws = tmp_path / "ws"
    ws.mkdir()
    project = tmp_path / "project"
    project.mkdir()

    env = EnvironmentContext(
        root_dir=str(project),
        archive_id="arch_test",
        group_id="env_user_u",
        user_id="u",
        project_key="p",
    )
    with runtime_context("environment", env):
        result = await handle_read_file(str(ws), '..\\..\\..\\inbox\\01_urgent.txt')

    assert result["ok"] is False
    assert result["suggested_tool"] == "env_read"
    assert result["suggested_args"]["path"] == "inbox/01_urgent.txt"
    assert "env_read" in result["fact"]

    # Outside environment mode the original bare error is preserved.
    plain = await handle_read_file(str(ws), r"..\..\secret.txt")
    assert plain["ok"] is False
    assert "suggested_tool" not in plain


def test_round3_plan_text_states_reply_language_from_user_message():
    """Round 7 follow-up: the language directive lives in round2's user tail
    and never reaches round3, so a Chinese-default persona replied to English
    users in Chinese (behavior 0.75 in 20260610_171622). round3 plan text now
    carries a reply-language fact."""
    from app.core import context as ctx_build
    from app.schemas.api import ResponsePlan

    plan = ResponsePlan(
        intent="report triage done",
        key_points=["8 emails classified", "phishing flagged"],
        tone="warm",
        length_hint="short",
    )
    en_msgs = ctx_build.round3_messages(
        "你是一个中文助手", plan, "user",
        "Can you go through my inbox and flag anything fishy?",
    )
    en_text = "\n".join(str(m.get("content") or "") for m in en_msgs)
    assert "the user wrote in English; reply in English" in en_text

    zh_msgs = ctx_build.round3_messages(
        "你是一个中文助手", plan, "user", "帮我整理一下收件箱里的邮件",
    )
    zh_text = "\n".join(str(m.get("content") or "") for m in zh_msgs)
    assert "用中文回复" in zh_text


async def test_main_write_of_helper_owned_path_attaches_ownership_fact(tmp_path):
    """Round 8: 20260610_163156 showed the main thread authoring
    triage_report.md while a live helper owned the same deliverable through
    expected_outputs. The write succeeds but now carries a soft ownership fact."""
    from app.core import debug
    from app.core.core_processes import registry as proc_registry
    from app.llm.tools.delegate_state import _record_task_contracts
    from app.llm.tools.workspace_file_ops import handle_write

    trace_id = "trace-coauthor-fact"
    debug.set_trace_id(trace_id)
    try:
        await _record_task_contracts(trace_id, [{
            "task_id": "email-triage",
            "kind": "edit",
            "expected_outputs": ["triage_report.md", "draft_replies.md"],
        }])

        async def _idle():
            import asyncio
            await asyncio.sleep(30)

        import asyncio
        task = asyncio.ensure_future(_idle())
        proc_id = await proc_registry().register_helper(
            owner=trace_id,
            helper_task_id="email-triage",
            task=task,
            helper_workspace="",
            abort_event=asyncio.Event(),
            description="test helper",
        )
        try:
            result = await handle_write(str(tmp_path), "triage_report.md", "# report\n")
            assert result["ok"] is True
            assert "helper_ownership_fact" in result
            assert "email-triage" in result["helper_ownership_fact"]

            other = await handle_write(str(tmp_path), "unrelated_notes.md", "notes\n")
            assert other["ok"] is True
            assert "helper_ownership_fact" not in other
        finally:
            task.cancel()
            try:
                await proc_registry().unregister(proc_id)
            except Exception:
                pass
    finally:
        debug.set_trace_id("")


async def test_edit_file_count_mismatch_reports_match_line_numbers(tmp_path):
    """Round 8: a count-mismatch retry costs a full turn; line numbers let the
    model choose expand-context vs replace-all without rereading the file."""
    from app.llm.tools.workspace_file_ops import handle_edit_file

    target = tmp_path / "f.py"
    target.write_text("value = 100\nother = 100\nthird = 100\n", encoding="utf-8")
    result = await handle_edit_file(str(tmp_path), "f.py", "= 100", "= 200", expected_count=1)
    assert result["ok"] is False
    assert result["actual_count"] == 3
    assert result["match_line_numbers"] == [1, 2, 3]
    assert "expected_count=3" in result["fact"]


async def test_guard_anchor_includes_completed_helper_evidence():
    """Round 15: the quality guard blocked a patch helper for 'missing browser
    tools' AFTER a browser helper had already completed Playwright evidence
    (20260611_162518_p16784, 2x guard_blocked, recovery 0.2). The guard anchor
    must show completed-helper evidence so follow-up delegations are judged
    against what prior helpers produced."""
    from app.core import agent_state, debug
    from app.llm.tools.delegate import _current_task_anchor_for_guard

    trace_id = "trace-guard-helper-evidence"
    debug.set_trace_id(trace_id)
    agent_state.reset_trace(trace_id)
    try:
        agent_state.add_evidence(
            trace_id=trace_id,
            source="helper",
            status=agent_state.EVIDENCE_VERIFIED,
            summary="Used Playwright/Chromium to navigate to http://127.0.0.1:53319/ and extracted the docs page text.",
            task_id="browser_docs_evidence",
            kind="code",
            data={"terminal_reason": "completed", "ok": True},
        )
        anchor = _current_task_anchor_for_guard(
            "patch report_client.py using the confirmed docs",
            [{"task_id": "patch_report_client", "kind": "code", "prompt": "patch the client"}],
        )
    finally:
        debug.set_trace_id("")

    assert "completed_helper[code]" in anchor
    assert "Playwright/Chromium" in anchor
    assert "does not need the same tools again" in anchor


def test_environment_mode_prunes_chat_memory_tools():
    """Round 17 (#2): environment/benchmark mode has no chat memory, group
    files, or avoid-mention semantics; those schemas were ~1k tokens of dead
    weight per main-loop turn. Static pruning keeps the env-mode tool hash
    stable so prefix caching within the mode is unaffected."""
    from app.llm.tools.registry import tools_for_runtime_mode

    env_names = {t["function"]["name"] for t in tools_for_runtime_mode("environment")}
    chat_names = {t["function"]["name"] for t in tools_for_runtime_mode("chat")}

    pruned = {"expand_warm", "expand_cold", "expand_kb", "mark_avoid_mention",
              "recall_thread", "search_files", "fetch_group_file"}
    assert not (pruned & env_names), pruned & env_names
    # Chat mode keeps its memory tools.
    assert "expand_warm" in chat_names
    # Core environment tools survive.
    for required in ("delegate", "task_plan", "env_run", "env_read", "env_apply_replace", "read_file"):
        assert required in env_names, required
    # Stable identity across calls (cache works).
    assert tools_for_runtime_mode("environment") is tools_for_runtime_mode("environment")


def test_framework_split_exemption_unified_entry():
    """Round 17 (#3): three framework-exemption predicates merged into one
    entry with a reason enum. Verify each variant still triggers."""
    from app.llm.tools.delegate import _framework_split_exemption

    # Variant 2: scoped fanout — component task inside a split batch.
    fanout_task = {
        "task_id": "ui-component",
        "kind": "code",
        "framework": "shared UI contract",
        "prompt": "Build the UI slice per the shared contract.",
        "expected_outputs": ["_env/ui/app.tsx", "_env/ui/styles.css"],
        "acceptance_checks": ["render check"],
    }
    assert _framework_split_exemption(
        fanout_task, total_task_count=3,
        outputs=["_env/ui/app.tsx", "_env/ui/styles.css"],
        prompt_l="build the ui slice per the shared contract.",
        enum_signals=[], comparison_pipeline=False,
    ) == "scoped_fanout"

    # Variant 3: bounded scaffold.
    scaffold_task = {
        "task_id": "project-scaffold",
        "kind": "code",
        "framework": "interfaces",
        "prompt": "Create the project scaffold with shared interfaces and harness.",
        "expected_outputs": ["_env/core/api.py"],
        "acceptance_checks": ["import check"],
    }
    assert _framework_split_exemption(
        scaffold_task, total_task_count=1,
        outputs=["_env/core/api.py"],
        prompt_l="create the project scaffold with shared interfaces and harness.",
        enum_signals=[], comparison_pipeline=False,
    ) == "bounded_scaffold"

    # Negative: enum signals present -> no exemption.
    assert _framework_split_exemption(
        scaffold_task, total_task_count=1,
        outputs=["_env/core/api.py"],
        prompt_l="create the project scaffold",
        enum_signals=["5 comparable algorithms"], comparison_pipeline=False,
    ) is None

    # Negative: outputs outside _env -> no exemption.
    assert _framework_split_exemption(
        scaffold_task, total_task_count=1,
        outputs=["src/api.py"],
        prompt_l="create the project scaffold with shared interfaces",
        enum_signals=[], comparison_pipeline=False,
    ) is None


def test_guidance_tracker_dedup_semantics():
    """Round 17 (#1): unified guidance dedup replaces ad-hoc boolean flags in
    chat_with_tools_loop. Verify once-per-loop, keyed re-arm, max_count, and
    reset semantics."""
    from app.llm.guidance_tracker import GuidanceTracker

    g = GuidanceTracker()
    # Default: once per hint.
    assert g.should_emit("snapshot") is True
    assert g.should_emit("snapshot") is False
    assert g.has_fired("snapshot") is True

    # Keyed: new key re-arms.
    assert g.should_emit("write_block", key="a.py") is True
    assert g.should_emit("write_block", key="a.py") is False
    assert g.should_emit("write_block", key="b.py") is True

    # max_count.
    assert g.should_emit("verifier", max_count=2) is True
    assert g.should_emit("verifier", max_count=2) is True
    assert g.should_emit("verifier", max_count=2) is False
    assert g.fired_count("verifier") == 2

    # reset re-arms.
    g.reset("snapshot")
    assert g.has_fired("snapshot") is False
    assert g.should_emit("snapshot") is True

    # snapshot view aggregates.
    snap = g.snapshot()
    assert snap["verifier"] == 2
