"""Office document verification handlers."""
from __future__ import annotations


def _sync_office_globals() -> None:
    from app.llm.tools import office as _office
    globals().update({
        name: value
        for name, value in vars(_office).items()
        if not name.startswith("__") and name not in {
            "_docx_verify_numbers",
            "_docx_verify_rigor",
            "_docx_verify_against_source",
            "_pptx_verify_numbers",
        }
    })


async def _docx_verify_numbers(
    workspace_dir: str, target: str, rel_path: str, args: dict,
) -> str:
    """对照一份或多份 CSV, 找出 docx 正文里可能应由 CSV 支撑的数字论断。

    args:
      csv_paths: 必填, list[str] 工作区相对路径
      tolerance: 可选 float, 默认 0.05 (5% 相对误差)
      number_pattern: 可选, 默认匹配 ms / 毫秒 / 秒 / s

    返回 JSON 含:
      claims:    [{paragraph_idx, snippet, number, context}, ...]
      csv_index: [{file, column, sample_values}]   // 给 LLM 看的
      mismatches: [{paragraph_idx, snippet, number, best_csv_match, deviation_pct,
                    severity}]
      matches_count
      mismatches_count
      hint
    """
    _sync_office_globals()

    from docx import Document
    import csv as _csv

    if not os.path.isfile(target):
        return _err(f"file not found: {rel_path}", action="verify_numbers")

    csv_paths = args.get("csv_paths") or []
    if not isinstance(csv_paths, list) or not csv_paths:
        return _err(
            "missing csv_paths",
            action="verify_numbers",
            hint="Pass csv_paths as a list of workspace-relative CSV files, "
                 "e.g. csv_paths=['bench_results.csv']",
        )
    try:
        tolerance = float(args.get("tolerance", 0.05))
    except (TypeError, ValueError):
        tolerance = 0.05

    # 默认匹配数字；后续按上下文把章节号、硬件型号、版本号、引用页码等
    # 非 CSV 事实数字降为 skipped_context_numbers，避免把工具最近值当作硬裁决。
    pat = re.compile(
        args.get("number_pattern")
        or r"(\d+(?:\.\d+)?)\s*(?:ms|毫秒|秒|s\b|µs|us)?"
    )

    def _non_csv_number_reason(text: str, num_str: str, start: int, end: int) -> str | None:
        around = text[max(0, start - 35): min(len(text), end + 35)]
        lower = around.lower()
        before = text[max(0, start - 3): start]
        after = text[end: min(len(text), end + 8)]
        try:
            num = float(num_str)
        except ValueError:
            return "not_numeric"
        if abs(num) < 1.0 and "." not in num_str:
            return "tiny_integer"
        if num <= 10 and "." not in num_str:
            return "small_integer_likely_enumeration"
        if re.search(r"\b(section|table|figure|fig\.|chapter|appendix|reference|ref\.|algorithm)\s*$", before.lower()):
            return "document_reference_number"
        if re.match(r"^\s*(?:[.)\]-]|节|章|表|图)", after):
            return "document_reference_number"
        if re.search(r"\b\d+(?:\.\d+){1,3}\b", around):
            return "version_or_section_number"
        if re.search(r"\b(?:i[3579]|ryzen|xeon|core|gcc|clang|python|windows|linux|ddr\d?|ram)\b", lower):
            return "hardware_or_software_spec"
        if re.search(r"\b(?:vol\.|no\.|pp\.|pages?|proc\.|symp\.|conference|journal|isbn|doi)\b", lower):
            return "citation_metadata"
        if "." not in num_str and 1900 <= num <= 2050:
            return "year_or_citation_year"
        return None

    def _do():
        # ── 1. 读 CSV, 建 value index ────────────────────────
        # 收集所有 CSV 里能作为浮点数的字段(数字列)
        csv_values: list[dict] = []  # {file, column, row_key, value}
        csv_summaries: list[dict] = []
        for csv_rel in csv_paths:
            try:
                csv_abs = ws_tool._safe_resolve(workspace_dir, csv_rel)
            except ValueError as e:
                csv_summaries.append({"file": csv_rel, "error": f"invalid path: {e}"})
                continue
            if not os.path.isfile(csv_abs):
                csv_summaries.append({"file": csv_rel, "error": "not found"})
                continue
            try:
                with open(csv_abs, encoding="utf-8-sig", newline="") as f:
                    reader = _csv.DictReader(f)
                    fieldnames = list(reader.fieldnames or [])
                    n_rows = 0
                    numeric_cols: dict[str, int] = {col: 0 for col in fieldnames}
                    for row in reader:
                        n_rows += 1
                        # 构造行键: 用所有非数字列作为 row_key
                        row_key_parts = []
                        for col, val in row.items():
                            try:
                                fv = float(val)
                                numeric_cols[col] = numeric_cols.get(col, 0) + 1
                                # 取行键时只用前几个非数字列
                                csv_values.append({
                                    "file": csv_rel, "column": col,
                                    "row_key": "|".join(row_key_parts[:4]),
                                    "value": fv,
                                })
                            except (TypeError, ValueError):
                                row_key_parts.append(f"{col}={val}")
                csv_summaries.append({
                    "file": csv_rel, "rows": n_rows,
                    "numeric_columns": [c for c, n in numeric_cols.items() if n > 0],
                })
            except Exception as e:
                csv_summaries.append({"file": csv_rel, "error": f"{type(e).__name__}: {e}"})

        if not csv_values:
            return None, "no numeric values found in any CSV (check csv_paths)"

        # ── 2. 扫描 docx, 提取所有数字 + context ──────────────
        try:
            doc = Document(target)
        except Exception as e:
            return None, f"open failed: {type(e).__name__}: {e}"

        claims: list[dict] = []
        skipped_context_numbers: list[dict] = []
        for p_idx, p in enumerate(doc.paragraphs):
            text = p.text or ""
            if not text.strip():
                continue
            # 跳过 heading (heading 通常不含数字论断)
            if p.style.name and p.style.name.startswith("Heading"):
                continue
            for m in re.finditer(r"\d+(?:\.\d+)?", text):
                num_str = m.group(0)
                try:
                    num = float(num_str)
                except ValueError:
                    continue
                # context: 前 25 / 后 25 字符
                start = max(0, m.start() - 25)
                end = min(len(text), m.end() + 25)
                ctx = text[start:end]
                skip_reason = _non_csv_number_reason(text, num_str, m.start(), m.end())
                if skip_reason:
                    if len(skipped_context_numbers) < 25:
                        skipped_context_numbers.append({
                            "paragraph_idx": p_idx,
                            "number": num,
                            "number_str": num_str,
                            "context": ctx,
                            "reason": skip_reason,
                        })
                    continue
                claims.append({
                    "paragraph_idx": p_idx,
                    "number": num,
                    "number_str": num_str,
                    "context": ctx,
                })

        # ── 3. 对每个候选 claim, 找 CSV 里最接近的值 ───────────────
        # 这是证据提示，不是硬裁决；LLM 需要结合语义判断该数字是否应由 CSV 支撑。
        mismatches: list[dict] = []
        matches: list[dict] = []
        for c in claims:
            target_val = c["number"]
            if target_val == 0:
                continue
            best = None
            best_diff = float("inf")
            for entry in csv_values:
                v = entry["value"]
                if v == 0:
                    continue
                diff = abs(v - target_val) / max(abs(v), abs(target_val))
                if diff < best_diff:
                    best_diff = diff
                    best = entry
            if best is None:
                continue
            rec = {
                "paragraph_idx": c["paragraph_idx"],
                "context": c["context"],
                "claim_number": target_val,
                "best_csv_match": {
                    "file": best["file"],
                    "column": best["column"],
                    "row_key": best["row_key"],
                    "value": best["value"],
                },
                "deviation_pct": round(best_diff * 100, 1),
            }
            if best_diff <= tolerance:
                matches.append(rec)
            elif best_diff <= 0.5:
                # 0.5 以上几乎可以肯定是完全不同的数 (e.g. 10x error)
                rec["severity"] = "warn"
                mismatches.append(rec)
            else:
                rec["severity"] = "high"
                mismatches.append(rec)

        return {
            "claims_extracted": len(claims),
            "matches_count": len(matches),
            "mismatches_count": len(mismatches),
            "csv_summaries": csv_summaries,
            "mismatches": mismatches[:50],   # 最多回 50 个
            "matches_sample": matches[:5],   # 抽 5 个匹配作 sanity check
            "skipped_context_numbers_count": len(skipped_context_numbers),
            "skipped_context_numbers_sample": skipped_context_numbers[:10],
        }, None

    result, err = await asyncio.to_thread(_do)
    if err:
        return _err(err, action="verify_numbers")

    hint = (
        f"Checked {result['claims_extracted']} candidate CSV-backed numeric claims; "
        f"{result['matches_count']} matched CSV within {int(tolerance*100)}% tolerance, "
        f"{result['mismatches_count']} candidate mismatches were found. "
        f"Skipped {result.get('skipped_context_numbers_count', 0)} context numbers that look like section/version/hardware/citation metadata.\n"
        f"数字候选已与 CSV 核对；章节号、版本号、硬件规格和引用元数据不作为 CSV 不匹配项。"
    )
    if result["mismatches_count"] > 0:
        # 找最严重的 deviation 作为事实样本，不要求逐项硬修。
        worst = max(result["mismatches"], key=lambda x: x["deviation_pct"])
        hint += (
            f" 最大候选偏差 {worst['deviation_pct']}%: "
            f"正文 {worst['claim_number']} vs CSV "
            f"{worst['best_csv_match']['value']} "
            f"({worst['best_csv_match']['column']}@{worst['best_csv_match']['row_key'][:30]})。"
        )
        hint += (
            " Treat mismatches as candidate facts for semantic review: edit only numbers that the document presents as CSV/benchmark results. "
            "Do not edit hardware specs, software versions, section/table/reference numbers, or citation metadata solely because this nearest-value check flagged them.\n\n"
            "候选不匹配需按语义复核；只有正文把该数字作为 CSV/benchmark 结果时才修改。"
        )
    else:
        hint += " ✓ No candidate CSV-backed number mismatches were found within this heuristic scope."

    return json.dumps({
        "ok": True,
        "action": "verify_numbers",
        "format": "docx",
        "path": rel_path,
        "tolerance_pct": int(tolerance * 100),
        "hint": hint,
        **result,
    }, ensure_ascii=False)


async def _docx_verify_rigor(
    workspace_dir: str, target: str, rel_path: str, args: dict,
) -> str:
    """全面 rigor 校验 — 9 层数据严谨性检查 (L1-L9)。

    Args (csv_paths 必填, 其余可选):
      csv_paths: list[str]            ground truth CSV
      tolerance: float = 0.05         相对误差 (5%)
      pivot_col: str | None           CSV "分组键"列名 (例 'algorithm')
      pivot_values: list[str] | None  此次评估涉及的 pivot 值
      pivot_aliases: dict | None      {pivot: [alias, ...]} — 中英文/缩写互通
      comparison_assertions: list     [{context_keyword, csv_filter, baseline_pivot,
                                       target_pivot, claimed_ratio, value_col?, tolerance?}]
      internal_facts: list            [{name, pattern (regex with 1 capture group),
                                       expected_value, tolerance?}]
      expect_reproducibility_metadata: bool = False
      scaling_n_col / scaling_value_col / scaling_group_cols / scaling_super_linear_threshold

    Returns 9 sub-checks (all 0 = ✅ ready):
      L1 number_check        绝对数字 vs CSV
      L2 ratio_check         比率论断 (X 倍) vs CSV pairwise
      L3 cherry_pick         对比段漏了某 pivot 值
      L4 unit_consistency    单位混用 (ms/s/μs/ns)
      L5 assertion_check     结构化断言违反 ← 最强 (用 comparison_assertions)
      L6 methodology_check   方法学声明 vs CSV 列
      L7 scaling_check       跨 N 复杂度观察 (observational, not violation)
      L8 internal_consistency 跨段 regex 一致性 (用 internal_facts)
      L9 reproducibility     复现元数据 (opt-in)
    """
    _sync_office_globals()

    from docx import Document
    import csv as _csv

    if not os.path.isfile(target):
        return _err(f"file not found: {rel_path}", action="verify_rigor")

    csv_paths = args.get("csv_paths") or []
    if not isinstance(csv_paths, list) or not csv_paths:
        return _err(
            "missing csv_paths",
            action="verify_rigor",
            hint="Pass csv_paths as list of workspace-relative CSV files",
        )
    try:
        tolerance = float(args.get("tolerance", 0.05))
    except (TypeError, ValueError):
        tolerance = 0.05
    pivot_col = args.get("pivot_col")
    pivot_values = args.get("pivot_values") or []
    if not isinstance(pivot_values, list):
        pivot_values = []
    # 2026-05-17 P162.1: 别名映射 (e.g. {"BPTREE": ["B+树","B+ 树"]})
    pivot_aliases = args.get("pivot_aliases") or {}
    if not isinstance(pivot_aliases, dict):
        pivot_aliases = {}
    # 2026-05-17 P162.2: 结构化比较断言, 让 LLM 显式声明 "我说的 X 倍是从 CSV 哪两行算的"
    # 每个 assertion: {context_keyword, csv_filter, baseline_pivot, target_pivot, claimed_ratio}
    comparison_assertions = args.get("comparison_assertions") or []
    if not isinstance(comparison_assertions, list):
        comparison_assertions = []

    def _do():
        # ── Load CSV ──────────────────────────────────────
        all_rows: list[dict] = []
        csv_summaries: list[dict] = []
        for csv_rel in csv_paths:
            try:
                csv_abs = ws_tool._safe_resolve(workspace_dir, csv_rel)
            except ValueError as e:
                csv_summaries.append({"file": csv_rel, "error": f"invalid path: {e}"})
                continue
            if not os.path.isfile(csv_abs):
                csv_summaries.append({"file": csv_rel, "error": "not found"})
                continue
            try:
                with open(csv_abs, encoding="utf-8-sig", newline="") as f:
                    reader = _csv.DictReader(f)
                    for row in reader:
                        row["_csv_file"] = csv_rel
                        all_rows.append(row)
                csv_summaries.append({"file": csv_rel, "rows": sum(
                    1 for r in all_rows if r.get("_csv_file") == csv_rel
                )})
            except Exception as e:
                csv_summaries.append({"file": csv_rel, "error": f"{type(e).__name__}: {e}"})

        if not all_rows:
            return None, "no rows loaded from any CSV"

        # ── Load doc ──────────────────────────────────────
        try:
            doc = Document(target)
        except Exception as e:
            return None, f"open failed: {type(e).__name__}: {e}"

        # ── (1) NUMBER CHECK (复用 verify_numbers 逻辑) ────────
        csv_values_flat: list[dict] = []
        for row in all_rows:
            for col, val in row.items():
                if col == "_csv_file":
                    continue
                try:
                    fv = float(val)
                    csv_values_flat.append({
                        "file": row["_csv_file"], "column": col,
                        "row": {k: v for k, v in row.items() if k != "_csv_file"},
                        "value": fv,
                    })
                except (TypeError, ValueError):
                    pass

        claims: list[dict] = []
        for p_idx, p in enumerate(doc.paragraphs):
            text = p.text or ""
            if not text.strip():
                continue
            if p.style.name and p.style.name.startswith("Heading"):
                continue
            for m in re.finditer(r"\d+(?:\.\d+)?", text):
                num_str = m.group(0)
                try:
                    num = float(num_str)
                except ValueError:
                    continue
                if abs(num) < 1.0 and "." not in num_str:
                    continue
                if num <= 10 and "." not in num_str:
                    continue
                if "." not in num_str and 1900 <= num <= 2050:
                    s, e = max(0, m.start() - 5), min(len(text), m.end() + 5)
                    if any(c in text[s:e] for c in ("年", "(", ")", "(", ")")):
                        continue
                # 2026-05-17 P162.5: 跳过 CPU 型号/产品编号 (前缀含字母-数字模式)
                # 例如 "i9-13980HX", "Xeon E5-2680", "GTX 1080", "M1 Pro" 等
                # 检测: 前 12 字符里含 "i\d-" / "Xeon" / "GTX" / "RTX" / "AMD" / "M\d "
                preview = text[max(0, m.start() - 12):m.start()]
                if re.search(r"i\d\s*-\s*$|i\d-?$|Xeon\s*E?\d?-?$|GTX\s*$|RTX\s*$|"
                             r"-\s*$|Ryzen\s*\d?\s*$|EPYC\s*$|M\d\s*$|A\d{1,3}\s*$",
                             preview, re.IGNORECASE):
                    continue
                # 跳过紧跟在 "×" 或 "倍" 前的数字 — 它们是比率, 归 ratio_check
                tail = text[m.end():m.end() + 3]
                if tail.startswith("×") or tail.startswith("倍"):
                    continue
                ctx = text[max(0, m.start() - 25):min(len(text), m.end() + 25)]
                claims.append({"paragraph_idx": p_idx, "number": num, "context": ctx})

        number_mismatches = []
        number_matches = 0
        for c in claims:
            t = c["number"]
            if t == 0:
                continue
            best = None
            best_diff = float("inf")
            for entry in csv_values_flat:
                v = entry["value"]
                if v == 0:
                    continue
                diff = abs(v - t) / max(abs(v), abs(t))
                if diff < best_diff:
                    best_diff = diff
                    best = entry
            if best is None:
                continue
            if best_diff <= tolerance:
                number_matches += 1
            elif best_diff > 0.3:
                number_mismatches.append({
                    "paragraph_idx": c["paragraph_idx"],
                    "context": c["context"],
                    "claim": t,
                    "best_csv_value": best["value"],
                    "best_csv_column": best["column"],
                    "deviation_pct": round(best_diff * 100, 1),
                })

        # ── (2) RATIO CHECK ───────────────────────────────────
        # 扫文本里的 "X 倍" / "X×" / "X-Y 倍" 比率论断
        # 对每个,提取数值, 看 CSV 里是否存在两个值满足该比率
        ratio_claims: list[dict] = []
        ratio_pat = re.compile(
            r"(\d+(?:\.\d+)?)\s*(?:[-–~至]\s*(\d+(?:\.\d+)?)\s*)?[倍×]"
        )
        # 2026-05-17 P162.5: 跳过 scaling / 理论预测类比率,这些不是 pivot-pair ratio
        SCALING_THEORY_KEYWORDS = (
            "增长", "增速", "scaling", "scale", "growth",
            "理论", "预测", "expected", "theoretical", "predicted",
            "渐进", "asymptotic", "complexity",
            "理论 O(", "O(N", "O(log", "O(n", "约 O", "~",
            "N=10", "N = 10", "10⁵", "10⁶", "10^5", "10^6",
        )
        for p_idx, p in enumerate(doc.paragraphs):
            text = p.text or ""
            if not text.strip():
                continue
            for m in ratio_pat.finditer(text):
                lo = float(m.group(1))
                hi = float(m.group(2)) if m.group(2) else lo
                ctx_s = max(0, m.start() - 80)
                ctx_e = min(len(text), m.end() + 30)
                ctx = text[ctx_s:ctx_e]
                # 跳过非数字比率 (例如版本号、单位换算)
                if lo < 1.0 or lo > 1000:
                    continue
                # 跳过 scaling / 理论 context 的比率
                # (这些不能简单地与 pivot-pair ratio 比对)
                if any(kw in ctx for kw in SCALING_THEORY_KEYWORDS):
                    continue
                ratio_claims.append({
                    "paragraph_idx": p_idx,
                    "claim_lo": lo,
                    "claim_hi": hi,
                    "context": ctx,
                })

        # 对每个 ratio 比率, 在 csv values 里找是否存在 (a, b) 使 a/b ≈ ratio
        ratio_mismatches: list[dict] = []
        # Pre-compute all pairwise ratios from CSV (within same operation/N/dist group)
        # 用 pivot_col 分组: 同 row 排除 pivot_col 之外都一样的两行的数值对
        all_pairs: list[float] = []
        if pivot_col and pivot_col in (all_rows[0] if all_rows else {}):
            # Group rows by everything-except-pivot
            from collections import defaultdict as _dd
            groups = _dd(list)
            for row in all_rows:
                key = tuple(
                    (k, v) for k, v in row.items()
                    if k != pivot_col and k != "_csv_file"
                       and k not in ("time_ms",)  # 数值列单独处理
                )
                # 找数值列 (time_ms 是已知数值列, 但泛化: 任何 float-able 列)
                for col, val in row.items():
                    if col == pivot_col or col == "_csv_file":
                        continue
                    try:
                        fv = float(val)
                        groups[key + ((col,),)].append((row[pivot_col], fv))
                    except (TypeError, ValueError):
                        pass
            for key, items in groups.items():
                for a_name, a_val in items:
                    for b_name, b_val in items:
                        if a_name == b_name or a_val == 0 or b_val == 0:
                            continue
                        all_pairs.append(a_val / b_val)
        else:
            # 没有 pivot_col, 用所有 numeric values 做暴力 pairwise
            vals = [e["value"] for e in csv_values_flat if e["value"] > 0]
            # 限制规模: 最多 2000 vs 2000 = 4M 对, 太多. 抽样.
            vals = vals[:500]
            for a in vals:
                for b in vals:
                    if a == 0 or b == 0 or a == b:
                        continue
                    all_pairs.append(a / b)

        for rc in ratio_claims:
            lo = rc["claim_lo"]
            hi = rc["claim_hi"]
            # 一个 ratio 论断匹配, 当 csv 里存在 some pair 落在 [lo*(1-tol), hi*(1+tol)] 区间内
            match_lo = lo * (1 - tolerance)
            match_hi = hi * (1 + tolerance)
            best_pair_ratio = None
            best_pair_diff = float("inf")
            for r in all_pairs:
                if match_lo <= r <= match_hi:
                    # 匹配 — 取第一个
                    best_pair_ratio = r
                    best_pair_diff = 0
                    break
            if best_pair_ratio is None:
                # 没匹配 — 找最接近的
                for r in all_pairs:
                    # 与中心点 (lo+hi)/2 比
                    center = (lo + hi) / 2
                    d = abs(r - center) / max(r, center)
                    if d < best_pair_diff:
                        best_pair_diff = d
                        best_pair_ratio = r
                if best_pair_ratio is not None and best_pair_diff > tolerance * 3:
                    # 3× tolerance: 默认 5% × 3 = 15%, 意思是
                    # "找不到 ±15% 内的 CSV pair ratio → 这个倍率论断可疑"
                    ratio_mismatches.append({
                        "paragraph_idx": rc["paragraph_idx"],
                        "context": rc["context"],
                        "claim_ratio_range": [lo, hi] if lo != hi else lo,
                        "closest_csv_ratio": round(best_pair_ratio, 2),
                        "deviation_pct": round(best_pair_diff * 100, 1),
                        "hint": (
                            "The document's ratio claim has no close matching value pair in the CSV evidence.\n"
                            "正文倍率论断缺少接近匹配的 CSV 证据。"
                        ),
                    })

        # ── (3) CHERRY-PICK CHECK ──────────────────────────────
        # 当文档某段提及 pivot_values 的 (N-1)/N 个时, flag
        cherry_pick_flags: list[dict] = []
        if pivot_values and len(pivot_values) >= 3:
            for p_idx, p in enumerate(doc.paragraphs):
                text = p.text or ""
                if not text.strip() or (p.style.name and p.style.name.startswith("Heading")):
                    continue
                # 2026-05-17 P162.1: 支持别名 — 一个 pivot 可有多个文字写法
                # pivot_aliases = {"BPTREE": ["BPTREE", "B+树", "B+ 树", "B+TREE"], ...}
                def _pivot_in_text(pv: str, t: str) -> bool:
                    aliases = pivot_aliases.get(pv, [pv]) if pivot_aliases else [pv]
                    return any(a in t for a in aliases)
                mentioned = [pv for pv in pivot_values if _pivot_in_text(pv, text)]
                missing = [pv for pv in pivot_values if not _pivot_in_text(pv, text)]
                # 只 flag "比较/对比/快/慢" 类语境且漏 1 个
                comparison_words = ("快", "慢", "倍", "对比", "比较", "领先", "胜过", "优于", "vs", "VS")
                if (len(mentioned) == len(pivot_values) - 1
                        and any(w in text for w in comparison_words)
                        and len(text) >= 80):
                    cherry_pick_flags.append({
                        "paragraph_idx": p_idx,
                        "snippet": text[:200],
                        "mentioned": mentioned,
                        "missing": missing,
                        "hint": (
                            f"This comparison/performance paragraph omits pivot value(s) {missing}. "
                            f"Check whether the comparison should include them or explicitly justify the scope.\n"
                            f"对比论断漏掉部分 pivot_values，需要补全或说明范围。"
                        ),
                    })

        # ── (4) UNIT CONSISTENCY ────────────────────────────────
        unit_counts: dict[str, int] = {}
        unit_pat = re.compile(r"\d+(?:\.\d+)?\s*(ms|毫秒|秒|s\b|微秒|μs|µs|us|ns|纳秒)")
        unit_paragraphs: dict[str, list[int]] = {}
        for p_idx, p in enumerate(doc.paragraphs):
            text = p.text or ""
            if not text.strip():
                continue
            for m in unit_pat.finditer(text):
                u = m.group(1).lower().replace("毫秒", "ms").replace("秒", "s").replace("微秒", "μs").replace("纳秒", "ns")
                unit_counts[u] = unit_counts.get(u, 0) + 1
                unit_paragraphs.setdefault(u, []).append(p_idx)
        dominant_unit = max(unit_counts, key=unit_counts.get) if unit_counts else None
        mixed_unit_paragraphs: list[int] = []
        if dominant_unit and len(unit_counts) > 1:
            # 段落 mixed: 一段里出现非主流单位
            for p_idx, p in enumerate(doc.paragraphs):
                seen = set()
                for m in unit_pat.finditer(p.text or ""):
                    u = m.group(1).lower().replace("毫秒", "ms").replace("秒", "s").replace("微秒", "μs").replace("纳秒", "ns")
                    seen.add(u)
                if len(seen) > 1:
                    mixed_unit_paragraphs.append(p_idx)

        # ── (5) STRUCTURED COMPARISON ASSERTIONS ─────────────────
        # 用户/LLM 显式声明"我说 X 倍是从 CSV 哪两行算的", 比 ratio_check 强得多
        # 每个 assertion: {
        #   context_keyword: 在 docx 文本里找这一段
        #   csv_filter: {col: val, ...} 锁定 CSV 行组
        #   baseline_pivot, target_pivot: 比谁
        #   claimed_ratio: 论断的倍率
        # }
        assertion_violations: list[dict] = []
        if pivot_col:
            for ai, asn in enumerate(comparison_assertions):
                if not isinstance(asn, dict):
                    continue
                ctx_kw = asn.get("context_keyword", "")
                csv_filter = asn.get("csv_filter") or {}
                baseline = asn.get("baseline_pivot")
                tgt = asn.get("target_pivot")
                claim = asn.get("claimed_ratio")
                if not (ctx_kw and baseline and tgt and claim is not None):
                    assertion_violations.append({
                        "assertion_index": ai,
                        "issue": "incomplete assertion (missing context_keyword/baseline_pivot/target_pivot/claimed_ratio)",
                        "assertion": asn,
                    })
                    continue
                # 1. context_keyword 在 docx 里能找到吗?
                hits = [p_idx for p_idx, p in enumerate(doc.paragraphs)
                        if ctx_kw in (p.text or "")]
                if not hits:
                    assertion_violations.append({
                        "assertion_index": ai,
                        "issue": "context_keyword_not_found_in_docx",
                        "context_keyword": ctx_kw,
                    })
                    continue
                # 2. csv 里按 filter 找 baseline/target 行
                def _match_row(row, flt):
                    for k, v in flt.items():
                        if str(row.get(k)) != str(v):
                            return False
                    return True
                b_row = next(
                    (r for r in all_rows
                     if _match_row(r, csv_filter) and r.get(pivot_col) == baseline),
                    None,
                )
                t_row = next(
                    (r for r in all_rows
                     if _match_row(r, csv_filter) and r.get(pivot_col) == tgt),
                    None,
                )
                if not b_row or not t_row:
                    assertion_violations.append({
                        "assertion_index": ai,
                        "issue": "no_csv_rows_match_filter",
                        "csv_filter": csv_filter,
                        "baseline_pivot": baseline,
                        "target_pivot": tgt,
                        "baseline_found": b_row is not None,
                        "target_found": t_row is not None,
                    })
                    continue
                # 3. 取数值列计算真实比率 — 默认用 "time_ms",可由 assertion 指定
                value_col = asn.get("value_col", "time_ms")
                try:
                    b_val = float(b_row.get(value_col, "nan"))
                    t_val = float(t_row.get(value_col, "nan"))
                except (TypeError, ValueError):
                    assertion_violations.append({
                        "assertion_index": ai,
                        "issue": f"value_col {value_col!r} not numeric",
                    })
                    continue
                if t_val == 0:
                    assertion_violations.append({
                        "assertion_index": ai,
                        "issue": "target value is 0 — cannot compute ratio",
                    })
                    continue
                actual_ratio = b_val / t_val
                claim_val = float(claim)
                diff = abs(actual_ratio - claim_val) / max(actual_ratio, claim_val)
                if diff > tolerance:
                    assertion_violations.append({
                        "assertion_index": ai,
                        "issue": "claimed_ratio_mismatch",
                        "context_keyword": ctx_kw,
                        "claimed_ratio": claim_val,
                        "actual_ratio": round(actual_ratio, 3),
                        "baseline": {"pivot": baseline, "value": b_val},
                        "target": {"pivot": tgt, "value": t_val},
                        "deviation_pct": round(diff * 100, 1),
                    })

        # ── (6) METHODOLOGY CHECK ───────────────────────────────
        # 2026-05-17 P162.3: 检查 "N 次/中位数/平均/std dev" 等方法学声明是否被 CSV 数据支持
        # 论文常见问题: 声明 "运行 3 次取中位数" 但 CSV 是单值列, 无法验证
        method_findings: list[dict] = []
        # 哪些 CSV 列代表 "聚合统计" — 用启发式 column name 匹配
        agg_indicators = {
            "median": ["median", "p50", "中位数"],
            "mean":   ["mean", "avg", "average", "均值", "平均"],
            "stddev": ["std", "stddev", "stderr", "variance", "var", "方差", "标准差"],
            "min":    ["min", "minimum", "最小"],
            "max":    ["max", "maximum", "最大"],
            "p95":    ["p95", "p99", "p999", "tail", "百分位"],
            "n_runs": ["n_runs", "runs", "iterations", "repetitions", "trials", "samples"],
            "ci":     ["ci", "confidence", "lower", "upper", "置信"],
        }
        first_csv_cols = list(all_rows[0].keys()) if all_rows else []
        first_csv_cols_lc = [c.lower() for c in first_csv_cols if c != "_csv_file"]
        csv_has = {k: any(any(ind in c for ind in inds) for c in first_csv_cols_lc)
                   for k, inds in agg_indicators.items()}

        # 在 docx 文本里找方法学声明
        METHOD_PATTERNS = [
            (r"(\d+)\s*次(?:运行|测试|实验|重复)?", "n_runs",
             "正文声称运行 N 次, 但 CSV 缺乏多次运行的支持列 (无 std/var/min/max 任一)。"),
            (r"中位数|median", "median",
             "正文声称取中位数, 但 CSV 缺乏 median/p50 列, 也缺乏让读者重新计算 median 的 raw 数据列。"),
            (r"(?<![a-zA-Z])(?:平均值|均值|mean|average)(?![a-zA-Z])", "mean",
             "正文声称取平均, 但 CSV 缺乏 mean/avg 列。"),
            (r"标准差|stddev|std\s*dev|方差|variance(?![\w])", "stddev",
             "正文声称报告标准差/方差, 但 CSV 没有 std/var 列。"),
            (r"置信区间|confidence\s*interval", "ci",
             "正文声称给出置信区间, 但 CSV 没有 CI 上下界列。"),
            (r"p[\s<]*0\.\d+|p\s*[-_]?\s*value|p[\s_]value", "stddev",
             "正文出现 p-value/显著性论断, 但 CSV 没有 std/variance 列, 无法做显著性检验。"),
            (r"95%|99%", "ci",
             "正文出现百分位置信表达 (95%/99%), 但 CSV 没有 CI 列。"),
        ]
        for pat, support_key, base_msg in METHOD_PATTERNS:
            for p_idx, p in enumerate(doc.paragraphs):
                text = p.text or ""
                if not text.strip():
                    continue
                m = re.search(pat, text, re.IGNORECASE)
                if not m:
                    continue
                supported = csv_has.get(support_key, False)
                # 'n_runs' 比较特殊: 如果声明运行 N 次但 CSV 无 std/var/min/max/runs 之一, 都算不可验证
                if support_key == "n_runs":
                    supported = (csv_has.get("stddev") or csv_has.get("min")
                                 or csv_has.get("max") or csv_has.get("n_runs"))
                if supported:
                    continue
                # 跳过参考文献段 (引用列表里有时含 "1996" 等数字)
                if "(" in text and any(c.isdigit() for c in text[:60]) and "pp" in text[:120]:
                    continue
                snippet_start = max(0, m.start() - 30)
                snippet_end = min(len(text), m.end() + 60)
                method_findings.append({
                    "paragraph_idx": p_idx,
                    "matched": m.group(0),
                    "snippet": text[snippet_start:snippet_end],
                    "missing_support": support_key,
                    "csv_columns_available": first_csv_cols,
                    "hint": base_msg,
                })

        # ── (7) SCALING CONSISTENCY CHECK ───────────────────────
        # 2026-05-17 P162.4: 给定 N 列, 报告跨 N 的扩展性, 帮 LLM 自查"声明的复杂度"
        # 是否被实测数据支持。这个 check 只输出观测,不判定违规——复杂度匹配需要 LLM 判断。
        scaling_obs: dict = {}
        n_col = args.get("scaling_n_col", "N")
        value_col_for_scaling = args.get("scaling_value_col", "time_ms")
        scaling_group_cols = args.get("scaling_group_cols")  # 列出按哪些列分组 (例如 ['algorithm', 'operation', 'distribution'])
        if (pivot_col and all_rows
                and n_col in all_rows[0]
                and value_col_for_scaling in all_rows[0]):
            from collections import defaultdict as _dd2
            # 默认分组: pivot_col + 所有其他非数值列 (排除 N, value_col)
            if not scaling_group_cols:
                scaling_group_cols = [pivot_col] + [
                    c for c in first_csv_cols
                    if c != "_csv_file" and c != n_col
                       and c != value_col_for_scaling
                       and c != pivot_col
                       and not c.lower().startswith(("time", "mean", "median", "std", "min", "max"))
                ]
            grouped = _dd2(dict)  # group_key -> {N: value}
            for r in all_rows:
                try:
                    n_int = int(r[n_col])
                    v = float(r[value_col_for_scaling])
                except (TypeError, ValueError, KeyError):
                    continue
                key = tuple(r.get(c) for c in scaling_group_cols)
                grouped[key][n_int] = v

            import math as _math
            # 对每个 group, 输出 (N_from, N_to) 实测 ratio 和 O(N log N) 理论 ratio
            scaling_findings: list[dict] = []
            super_linear_threshold = float(args.get("scaling_super_linear_threshold", 1.5))
            for key, ns_dict in grouped.items():
                ns = sorted(ns_dict.keys())
                if len(ns) < 2:
                    continue
                row_obs = []
                for i in range(1, len(ns)):
                    a, b = ns[i - 1], ns[i]
                    v_a = ns_dict[a]; v_b = ns_dict[b]
                    if v_a <= 0:
                        continue
                    actual = v_b / v_a
                    expected_nlogn = (b * _math.log(b)) / (a * _math.log(a))
                    expected_n = b / a
                    expected_logn = _math.log(b) / _math.log(a)
                    dev_vs_nlogn = abs(actual - expected_nlogn) / expected_nlogn
                    row_obs.append({
                        "N_from": a, "N_to": b,
                        "actual_ratio": round(actual, 2),
                        "expected_O_NlogN": round(expected_nlogn, 2),
                        "expected_O_N":    round(expected_n, 2),
                        "expected_O_logN": round(expected_logn, 2),
                        "deviation_vs_NlogN_pct": round(dev_vs_nlogn * 100, 1),
                        "super_linear": actual > expected_nlogn * super_linear_threshold,
                    })
                scaling_findings.append({
                    "group_keys": dict(zip(scaling_group_cols, key)),
                    "observations": row_obs,
                    "max_super_linear_at_largest_N": (
                        row_obs[-1]["super_linear"] if row_obs else False
                    ),
                })
            scaling_obs = {
                "group_cols": scaling_group_cols,
                "n_groups": len(scaling_findings),
                "super_linear_count": sum(
                    1 for f in scaling_findings
                    if f["max_super_linear_at_largest_N"]
                ),
                "findings": scaling_findings[:30],
            }

        # ── (8) INTERNAL CROSS-SECTION CONSISTENCY ────────────────
        # 2026-05-17 P162.8: 同一事实在多段中应该保持一致 (abstract / body / conclusion)
        # 设计教训:
        #   v1 (keyword overlap): 不同事实共享上下文词被误聚 → 大量 false positives
        #   v2 (required_keywords in same paragraph): 同段落多个不同数字被全部拉入 → 仍有 FP
        #   v3 (正则 pattern): LLM 精准指定 "fact statement 的形态", 工具捕获数字 → 几乎 0 FP
        # 每个 fact: {name, pattern (含一个捕获组定位数字), expected_value, tolerance}
        # 例: pattern=r"比红黑树(?:[^。]{0,20})快\s*约?\s*(\d+(?:\.\d+)?)\s*倍"
        #     expected_value=9.2 → 在所有 paragraph 上跑 regex, 抓数字, 对比期望值。
        internal_findings: list[dict] = []
        internal_facts = args.get("internal_facts") or []
        if not isinstance(internal_facts, list):
            internal_facts = []
        for fi, fact in enumerate(internal_facts):
            if not isinstance(fact, dict):
                continue
            f_name = fact.get("name", f"fact_{fi}")
            f_pattern = fact.get("pattern")
            f_expected = fact.get("expected_value")
            f_tol = float(fact.get("tolerance", tolerance))
            # 向下兼容旧版 required_keywords (但 pattern 优先)
            f_required_kw = fact.get("required_keywords") or []
            if f_expected is None:
                continue
            if not f_pattern and not f_required_kw:
                continue
            occurrences = []
            try:
                f_re = re.compile(f_pattern) if f_pattern else None
            except re.error:
                internal_findings.append({
                    "fact_name": f_name,
                    "error": f"invalid regex pattern: {f_pattern!r}",
                })
                continue
            for p_idx, p in enumerate(doc.paragraphs):
                text = p.text or ""
                if not text.strip():
                    continue
                if f_re:
                    # regex 模式: 每次匹配抓 group(1) 作为数字
                    for m in f_re.finditer(text):
                        try:
                            n = float(m.group(1))
                        except (ValueError, IndexError):
                            continue
                        d = abs(n - f_expected) / max(abs(n), abs(f_expected)) if max(abs(n), abs(f_expected)) > 0 else 0
                        occurrences.append({
                            "paragraph_idx": p_idx,
                            "found_number": n,
                            "deviation_pct": round(d * 100, 1),
                            "matched_text": m.group(0)[:80],
                        })
                else:
                    # 旧版 required_keywords 路径 (less precise — 仅向下兼容)
                    if not all(kw in text for kw in f_required_kw):
                        continue
                    local_nums = []
                    for m in re.finditer(r"\d+(?:\.\d+)?", text):
                        try:
                            n = float(m.group(0))
                        except ValueError:
                            continue
                        if abs(n) < 0.001:
                            continue
                        if max(n, f_expected) > 0:
                            d = abs(n - f_expected) / max(n, f_expected)
                            if d <= 0.5:
                                local_nums.append((n, d))
                    if local_nums:
                        local_nums.sort(key=lambda x: x[1])
                        picked_num, picked_diff = local_nums[0]
                        occurrences.append({
                            "paragraph_idx": p_idx,
                            "found_number": picked_num,
                            "deviation_pct": round(picked_diff * 100, 1),
                        })
            if len(occurrences) >= 1:
                vals = [o["found_number"] for o in occurrences]
                spread = (max(vals) - min(vals)) / max(vals) if max(vals) > 0 else 0
                any_too_far = any(o["deviation_pct"] > f_tol * 100 for o in occurrences)
                if spread > f_tol or any_too_far:
                    internal_findings.append({
                        "fact_name": f_name,
                        "expected_value": f_expected,
                        "occurrences_count": len(occurrences),
                        "occurrences": occurrences[:20],
                        "max_spread_pct": round(spread * 100, 1),
                        "hint": (
                            f"Fact {f_name!r} expected value {f_expected}; {len(occurrences)} occurrence(s) were found "
                            f"with spread {spread*100:.1f}%. Use consistent precision across abstract, body, and conclusion.\n"
                            f"同一事实值在文档中精度应保持一致。"
                        ),
                    })

        # ── (9) REPRODUCIBILITY METADATA CHECK ────────────────────
        # 2026-05-17 P162.9: 基准论文应该有可复现的元数据 — 随机种子、重复次数、环境配置
        # 缺这些的 paper 严格意义上无法被复现验证。这是元层面的严谨性问题。
        reproducibility_check = None
        if args.get("expect_reproducibility_metadata", False):
            # 标志位触发 — 不是所有 doc 都需要这个,只有"benchmark paper" 类
            full_text = "\n".join(p.text or "" for p in doc.paragraphs)
            required_signals = {
                "random_seed": [
                    "种子", "seed", "随机化", "random state", "random_state",
                ],
                "repeat_count": [
                    "重复", "n_runs", "n_repeats", "trials", "次实验", "次运行",
                    "中位数", "median", "mean", "平均", "标准差", "std",
                    "误差", "error bar", "置信区间", "confidence",
                ],
                "environment": [
                    "CPU", "memory", "RAM", "GHz", "GB", "操作系统", "OS",
                    "kernel", "Python", "g++", "gcc", "clang", "MSVC",
                    "version", "版本",
                ],
                "input_distribution": [
                    "随机", "random", "uniform", "skewed", "Zipf", "Gaussian",
                    "分布", "distribution",
                ],
            }
            missing_signals = {}
            for signal_name, keywords in required_signals.items():
                hit = next((kw for kw in keywords if kw.lower() in full_text.lower()), None)
                if not hit:
                    missing_signals[signal_name] = keywords[:3]
            reproducibility_check = {
                "missing_signals": missing_signals,
                "missing_count": len(missing_signals),
                "hint": (
                    "Reproducibility metadata 缺失。基准论文应该至少提及: "
                    "随机种子(seed) / 重复次数(n_runs) / 环境配置(CPU/RAM/OS) / 输入分布。"
                    "无此元数据的论文在严格意义上不可复现。"
                ) if missing_signals else "✓ 复现元数据齐全",
            }

        return {
            "number_check": {
                "claims_extracted": len(claims),
                "matches_count": number_matches,
                "mismatches_count": len(number_mismatches),
                "mismatches": number_mismatches[:30],
            },
            "ratio_check": {
                "claims_extracted": len(ratio_claims),
                "mismatches_count": len(ratio_mismatches),
                "mismatches": ratio_mismatches[:20],
            },
            "cherry_pick": {
                "flags_count": len(cherry_pick_flags),
                "flags": cherry_pick_flags[:10],
            },
            "unit_consistency": {
                "unit_counts": unit_counts,
                "dominant_unit": dominant_unit,
                "mixed_unit_paragraphs": mixed_unit_paragraphs[:10],
            },
            "assertion_check": {
                "assertions_count": len(comparison_assertions),
                "violations_count": len(assertion_violations),
                "violations": assertion_violations,
            },
            "methodology_check": {
                "csv_supports": csv_has,
                "findings_count": len(method_findings),
                "findings": method_findings[:20],
            },
            "scaling_check": scaling_obs,
            "internal_consistency": {
                "findings_count": len(internal_findings),
                "findings": internal_findings[:20],
            },
            "reproducibility_check": reproducibility_check,
            "csv_summaries": csv_summaries,
        }, None

    result, err = await asyncio.to_thread(_do)
    if err:
        return _err(err, action="verify_rigor")

    n_num = result["number_check"]["mismatches_count"]
    n_rat = result["ratio_check"]["mismatches_count"]
    n_chr = result["cherry_pick"]["flags_count"]
    n_uni = len(result["unit_consistency"]["mixed_unit_paragraphs"])
    n_ast = result["assertion_check"]["violations_count"]
    n_met = result["methodology_check"]["findings_count"]
    n_scl = result.get("scaling_check", {}).get("super_linear_count", 0)
    n_int = result["internal_consistency"]["findings_count"]
    n_rep = (result["reproducibility_check"] or {}).get("missing_count", 0) if result.get("reproducibility_check") else 0

    if (n_num == 0 and n_rat == 0 and n_chr == 0 and n_uni == 0 and n_ast == 0 
            and n_met == 0 and n_int == 0 and n_rep == 0):
        hint = (
            "✅ 全部 rigor 检查通过 (数字 / 比率 / 樱桃挑选 / 单位 / 断言 / 方法学 / "
            "内部一致性"
            + (" / 可复现性元数据" if result.get("reproducibility_check") else "")
            + ")。"
            + (f" 注: scaling_check 在 {n_scl} 个 group 检测到 N=Nmax 处的 super-linear 增长 "
               f"(可能是 cache 效应), 不算 violation 但建议正文披露。" if n_scl > 0 else "")
            + " 可以交付。"
        )
    else:
        hint_parts = []
        if n_num > 0:
            hint_parts.append(f"{n_num} 处数字与 CSV 不符")
        if n_rat > 0:
            hint_parts.append(f"{n_rat} 处比率论断 (X 倍) 与 CSV pairwise ratio 不符")
        if n_chr > 0:
            hint_parts.append(f"{n_chr} 处对比段疑似遗漏了某个 pivot 值 (cherry-picking)")
        if n_uni > 0:
            hint_parts.append(f"{n_uni} 段含混用单位")
        if n_ast > 0:
            hint_parts.append(f"{n_ast} 处结构化断言被违反")
        if n_met > 0:
            hint_parts.append(
                f"{n_met} 处方法学声明 (N 次/中位数/std/CI) 无 CSV 列支持 — "
                f"读者无法独立验证, 建议补 CSV 列或软化措辞"
            )
        if n_int > 0:
            hint_parts.append(
                f"{n_int} 处内部一致性问题 (同事实点在多段中数字不一致, "
                f"abstract/body/conclusion 应统一)"
            )
        if n_rep > 0:
            hint_parts.append(
                f"{n_rep} 类可复现性元数据缺失 (种子/重复次数/环境/分布)"
            )
        hint = (
            "⚠️ rigor 检查发现问题: " + "; ".join(hint_parts) + "。"
            "建议逐项修正后再调一次 verify_rigor 确认全绿。"
        )
        if n_scl > 0:
            hint += (
                f" 此外 scaling_check 在 {n_scl} 组数据上检测到 N=Nmax 处的 super-linear 增长 "
                f"(常见于 cache 效应; 不算 violation 但建议正文披露)。"
            )

    return json.dumps({
        "ok": True,
        "action": "verify_rigor",
        "format": "docx",
        "path": rel_path,
        "tolerance_pct": int(tolerance * 100),
        "hint": hint,
        **result,
    }, ensure_ascii=False)


async def _docx_verify_against_source(
    workspace_dir: str, target: str, rel_path: str, args: dict,
) -> str:
    """计算 docx 文本 vs source 文件的 token 覆盖率. 信息性工具, 不强制语义.

    返回 coverage_pct (0-100) + severity (LOW/MID/HIGH) + missing/invented token 列表.
    调用方根据任务意图自己判断结果含义:
    - 复现型任务 (OCR→word, 文本搬运): 高 coverage 说明 docx 忠实, 低 coverage 提示偏离
    - 解答/分析/扩展型: docx 不复述 source 时 coverage 必然低, 此工具结果与正确性无关

    算法 (token-bucket):
      1. 从 source + docx 各提取 4 类 token: LaTeX 命令, 数字 (≥2位), 数学符号, 中文 3-6 字
      2. 检测 docx 是否含图 (LaTeX → PNG 渲染):
         - 有图: coverage = 85% × chinese_cov + 15% × numsym_cov (LaTeX 不计)
         - 没图: coverage = 65% chinese + 10% numsym + 25% latex
      3. severity: ≥60% HIGH / 30-60% MID / <30% LOW

    args:
      source_path: str, source 文件 (相对工作区路径)
      source_text: str, 可选 — 直接传 source 内容 (代替 source_path)
      threshold: float, 默认 0.30, LOW/MID 边界
      warn_threshold: float, 默认 0.60, MID/HIGH 边界

    返回:
      {ok, action, coverage_pct, severity, chinese_coverage_pct, numsym_coverage_pct,
       drawings_in_docx, algorithm_note, missing_from_docx[:30], invented_in_docx[:30], hint}
    """
    _sync_office_globals()

    from docx import Document

    if not os.path.isfile(target):
        return _err(f"file not found: {rel_path}", action="verify_against_source")

    # Load source text
    source_path = str(args.get("source_path", "")).strip()
    source_text = str(args.get("source_text", ""))
    if source_path:
        try:
            source_resolved = ws_tool._safe_resolve(workspace_dir, source_path)
        except ValueError as e:
            return _err(f"invalid source_path: {e}", action="verify_against_source")
        if not os.path.isfile(source_resolved):
            return _err(f"source file not found: {source_path}", action="verify_against_source")
        try:
            with open(source_resolved, encoding="utf-8") as f:
                source_text = f.read()
        except OSError as e:
            return _err(f"failed to read source: {e}", action="verify_against_source")
    if not source_text.strip():
        return _err(
            "no source provided",
            action="verify_against_source",
            hint="Pass source_path='ocr_X.md' OR source_text='...'"
        )

    try:
        threshold = float(args.get("threshold", 0.30))
        warn_threshold = float(args.get("warn_threshold", 0.60))
    except (TypeError, ValueError):
        threshold, warn_threshold = 0.30, 0.60

    def _extract_tokens(text: str) -> dict:
        """提取 token 分类: LaTeX cmds vs Chinese phrases vs numbers vs symbols.

        分类返回是因为不同 token 类型在 docx 中可达性不同:
        - LaTeX 命令在 docx 中**渲染为 PNG**, docx 文本看不到 → 不应作为"missing"惩罚
        - 中文短语**保留为文字**, 在 docx 文本里能看到 → 是高信号
        - 数字/符号介于两者之间
        """
        tokens = {"latex": set(), "chinese": set(), "numbers": set(), "symbols": set()}
        # LaTeX commands (in docx 走 PNG 渲染, 不会出现在 text)
        tokens["latex"].update(re.findall(r"\\[a-zA-Z]+", text))
        # Numbers ≥2 digits
        for m in re.finditer(r"(?:^|[\s\(\[=,])(-?\d{2,})(?:[\s\)\],.=]|$)", text):
            tokens["numbers"].add(m.group(1))
        # Math symbols Unicode
        tokens["symbols"].update(re.findall(r"[≤≥≠≈±∞∫∑∏√→↔∂∇∈∉⊂⊃∪∩]", text))
        # Chinese phrases 3-6 chars (drop stopwords)
        STOP_PHRASES = {"分计", "本试卷", "本大题", "共22题", "共10题", "共6题",
                       "选择题", "填空题", "解答题", "原图模糊", "原图残文",
                       "高等数学", "试卷"}
        for phrase in re.findall(r"[\u4e00-\u9fff]{3,6}", text):
            if phrase not in STOP_PHRASES:
                tokens["chinese"].add(phrase)
        return tokens

    def _do():
        try:
            doc = Document(target)
        except Exception as e:
            return None, f"open failed: {type(e).__name__}: {e}"
        # 收集 docx 全部文本 (paragraphs + table cells)
        docx_parts: list[str] = []
        for p in doc.paragraphs:
            if p.text.strip():
                docx_parts.append(p.text)
        for t in doc.tables:
            for row in t.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        docx_parts.append(cell.text)
        docx_text = "\n".join(docx_parts)
        # Also count embedded images — proxy for "did helper render formulas as PNG?"
        n_drawings = 0
        for p in doc.paragraphs:
            for r in p.runs:
                n_drawings += len(r.element.findall(
                    ".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing"
                ))
        for t in doc.tables:
            for row in t.rows:
                for cell in row.cells:
                    for p in cell.paragraphs:
                        for r in p.runs:
                            n_drawings += len(r.element.findall(
                                ".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing"
                            ))
        return (source_text, docx_text, n_drawings), None

    text_pair, err = await asyncio.to_thread(_do)
    if err:
        return _err(err, action="verify_against_source")
    src_text, docx_text, n_drawings = text_pair

    ocr_buckets = _extract_tokens(src_text)
    docx_buckets = _extract_tokens(docx_text)

    # 2026-05-18 P203 算法:
    # - 中文短语**必须 overlap** (是高信号: 题目主题对不对): 60% weight
    # - 数字+符号: 30% weight
    # - LaTeX 命令: 若 docx 没渲染图 (n_drawings 极低), 走传统 token 比较 (10%);
    #   若有大量图 (说明 LaTeX 渲染成 PNG 了), LaTeX missing 不算 (0%)
    n_ocr_latex = len(ocr_buckets["latex"])
    has_pngs = n_drawings >= max(1, n_ocr_latex // 4)  # docx 有图说明确实渲染了

    chinese_overlap = len(ocr_buckets["chinese"] & docx_buckets["chinese"])
    chinese_total = len(ocr_buckets["chinese"])
    chinese_cov = chinese_overlap / chinese_total if chinese_total else 1.0

    numsym_ocr = ocr_buckets["numbers"] | ocr_buckets["symbols"]
    numsym_docx = docx_buckets["numbers"] | docx_buckets["symbols"]
    numsym_overlap = len(numsym_ocr & numsym_docx)
    numsym_total = len(numsym_ocr)
    numsym_cov = numsym_overlap / numsym_total if numsym_total else 1.0

    if has_pngs:
        # LaTeX 用 PNG 渲染了, latex token 不参与 coverage
        # P203 v3: 中文短语是最强信号 (数字常匹配靠运气, e.g. "5分" 在所有高数题都出现)
        coverage = 0.85 * chinese_cov + 0.15 * numsym_cov
        algorithm_note = (f"docx contains {n_drawings} drawings → LaTeX assumed rendered to PNG. "
                          f"Coverage = 85%×chinese + 15%×numsym")
    else:
        latex_overlap = len(ocr_buckets["latex"] & docx_buckets["latex"])
        latex_total = len(ocr_buckets["latex"])
        latex_cov = latex_overlap / latex_total if latex_total else 1.0
        coverage = 0.65 * chinese_cov + 0.10 * numsym_cov + 0.25 * latex_cov
        algorithm_note = (f"docx has only {n_drawings} drawings, LaTeX expected as text. "
                          f"Coverage = 65%×chinese + 10%×numsym + 25%×latex")

    # 总 token 视角 (用于报告 missing/invented)
    all_ocr = ocr_buckets["chinese"] | ocr_buckets["numbers"] | ocr_buckets["symbols"]
    all_docx = docx_buckets["chinese"] | docx_buckets["numbers"] | docx_buckets["symbols"]
    if not has_pngs:
        all_ocr |= ocr_buckets["latex"]
        all_docx |= docx_buckets["latex"]
    only_ocr = all_ocr - all_docx
    only_docx = all_docx - all_ocr

    if coverage < threshold:
        severity = "LOW"
        hint = (
            f"coverage {coverage*100:.0f}% < {threshold*100:.0f}% — the DOCX overlaps little with source tokens. "
            f"For source-restatement tasks, inspect missing_from_docx and invented_in_docx, then repair with replace_blocks. "
            f"For answer/analysis/expansion tasks, low coverage may be normal.\n"
            f"DOCX 与 source 覆盖较低，按任务类型决定是否修复。"
        )
    elif coverage < warn_threshold:
        severity = "MID"
        hint = (
            f"coverage {coverage*100:.0f}% is between {threshold*100:.0f}% and {warn_threshold*100:.0f}%. "
            f"Some source content is absent from the DOCX; inspect missing_from_docx and decide from the task intent.\n"
            f"部分 source 未覆盖，是否修复取决于任务意图。"
        )
    else:
        severity = "HIGH"
        hint = (
            f"coverage {coverage*100:.0f}% >= {warn_threshold*100:.0f}% — the DOCX substantially covers source tokens. "
            f"For restatement tasks, this is ready to deliver.\n"
            f"DOCX 对 source 覆盖充分。"
        )

    def _sample(s: set, n: int) -> list:
        return sorted(s)[:n]

    return json.dumps({
        "ok": True,
        "action": "verify_against_source",
        "format": "docx",
        "path": rel_path,
        "source_path": source_path if source_path else "(inline)",
        "coverage_pct": int(coverage * 100),
        "severity": severity,
        "threshold_pct": int(threshold * 100),
        "warn_threshold_pct": int(warn_threshold * 100),
        "chinese_coverage_pct": int(chinese_cov * 100),
        "numsym_coverage_pct": int(numsym_cov * 100),
        "drawings_in_docx": n_drawings,
        "algorithm_note": algorithm_note,
        "missing_from_docx_count": len(only_ocr),
        "invented_in_docx_count": len(only_docx),
        "missing_from_docx": _sample(only_ocr, 30),
        "invented_in_docx": _sample(only_docx, 30),
        "hint": hint,
    }, ensure_ascii=False)


async def _pptx_verify_numbers(
    workspace_dir: str, target: str, rel_path: str, args: dict,
) -> str:
    """对 .pptx 做与 docx 平行的数字校验。

    把每张幻灯片(含 title + body + table cells)合成一段文字, 然后复用
    与 docx 相同的数字提取/CSV 匹配逻辑。

    args 与 _docx_verify_numbers 完全一致 (csv_paths, tolerance, number_pattern)。
    """
    _sync_office_globals()

    from pptx import Presentation
    import csv as _csv

    if not os.path.isfile(target):
        return _err(f"file not found: {rel_path}", action="verify_numbers")
    csv_paths = args.get("csv_paths") or []
    if not isinstance(csv_paths, list) or not csv_paths:
        return _err(
            "missing csv_paths",
            action="verify_numbers",
            hint="Pass csv_paths as list of workspace-relative CSV files",
        )
    try:
        tolerance = float(args.get("tolerance", 0.05))
    except (TypeError, ValueError):
        tolerance = 0.05

    def _do():
        # CSV
        csv_values_flat: list[dict] = []
        csv_summaries: list[dict] = []
        for csv_rel in csv_paths:
            try:
                csv_abs = ws_tool._safe_resolve(workspace_dir, csv_rel)
            except ValueError as e:
                csv_summaries.append({"file": csv_rel, "error": f"invalid path: {e}"})
                continue
            if not os.path.isfile(csv_abs):
                csv_summaries.append({"file": csv_rel, "error": "not found"})
                continue
            try:
                with open(csv_abs, encoding="utf-8-sig", newline="") as f:
                    rows = list(_csv.DictReader(f))
                csv_summaries.append({"file": csv_rel, "rows": len(rows)})
                for row in rows:
                    row_key_parts = []
                    for col, val in row.items():
                        try:
                            fv = float(val)
                            csv_values_flat.append({
                                "file": csv_rel, "column": col,
                                "row_key": "|".join(row_key_parts[:4]),
                                "value": fv,
                            })
                        except (TypeError, ValueError):
                            row_key_parts.append(f"{col}={val}")
            except Exception as e:
                csv_summaries.append({"file": csv_rel, "error": f"{type(e).__name__}: {e}"})

        if not csv_values_flat:
            return None, "no numeric values found in any CSV"

        # Open pptx, extract slide texts
        try:
            prs = Presentation(target)
        except Exception as e:
            return None, f"open failed: {type(e).__name__}: {e}"

        # Build claims by iterating each slide as if it were a paragraph
        claims: list[dict] = []
        for s_idx, slide in enumerate(prs.slides):
            slide_text_parts = []
            # Title
            try:
                if slide.shapes.title and slide.shapes.title.has_text_frame:
                    t = (slide.shapes.title.text_frame.text or "").strip()
                    if t:
                        slide_text_parts.append(t)
            except Exception:
                pass
            # Body shapes
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                        pf = getattr(shape, "placeholder_format", None)
                        if pf and pf.idx == 0:
                            continue  # title already added
                    if shape.has_text_frame and shape.text_frame.text:
                        slide_text_parts.append(shape.text_frame.text.strip())
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                t = (cell.text or "").strip()
                                if t:
                                    slide_text_parts.append(t)
                except Exception:
                    continue
            slide_text = "\n".join(slide_text_parts)
            if not slide_text:
                continue
            # Extract numbers (same heuristics as docx)
            for m in re.finditer(r"\d+(?:\.\d+)?", slide_text):
                num_str = m.group(0)
                try:
                    num = float(num_str)
                except ValueError:
                    continue
                if abs(num) < 1.0 and "." not in num_str:
                    continue
                if num <= 10 and "." not in num_str:
                    continue
                if "." not in num_str and 1900 <= num <= 2050:
                    s, e = max(0, m.start() - 5), min(len(slide_text), m.end() + 5)
                    around = slide_text[s:e]
                    if any(c in around for c in ("年", "(", ")", "(", ")")):
                        continue
                ctx = slide_text[max(0, m.start() - 25):min(len(slide_text), m.end() + 25)]
                claims.append({
                    "slide_idx": s_idx,
                    "number": num,
                    "number_str": num_str,
                    "context": ctx,
                })

        # Match claims to CSV
        matches: list[dict] = []
        mismatches: list[dict] = []
        for c in claims:
            target_val = c["number"]
            if target_val == 0:
                continue
            best = None
            best_diff = float("inf")
            for entry in csv_values_flat:
                v = entry["value"]
                if v == 0:
                    continue
                diff = abs(v - target_val) / max(abs(v), abs(target_val))
                if diff < best_diff:
                    best_diff = diff
                    best = entry
            if best is None:
                continue
            rec = {
                "slide_idx": c["slide_idx"],
                "context": c["context"],
                "claim_number": target_val,
                "best_csv_match": {
                    "file": best["file"],
                    "column": best["column"],
                    "row_key": best["row_key"],
                    "value": best["value"],
                },
                "deviation_pct": round(best_diff * 100, 1),
            }
            if best_diff <= tolerance:
                matches.append(rec)
            elif best_diff > 0.3:
                rec["severity"] = "high" if best_diff > 0.5 else "warn"
                mismatches.append(rec)

        return {
            "claims_extracted": len(claims),
            "matches_count": len(matches),
            "mismatches_count": len(mismatches),
            "csv_summaries": csv_summaries,
            "mismatches": mismatches[:50],
            "matches_sample": matches[:5],
        }, None

    result, err = await asyncio.to_thread(_do)
    if err:
        return _err(err, action="verify_numbers")

    hint = (
        f"Checked {result['claims_extracted']} numeric claims; "
        f"{result['matches_count']} matched CSV within {int(tolerance*100)}% tolerance, "
        f"{result['mismatches_count']} did not match.\n"
        f"数字论断已与 CSV 核对。"
    )
    if result["mismatches_count"] > 0:
        worst = max(result["mismatches"], key=lambda x: x["deviation_pct"])
        hint += (
            f" 最严重一项偏差 {worst['deviation_pct']}%: "
            f"slide {worst['slide_idx']} 论断 {worst['claim_number']} vs CSV "
            f"{worst['best_csv_match']['value']}。"
        )
    else:
        hint += " ✓ 全部数字论断在容差内。"

    return json.dumps({
        "ok": True,
        "action": "verify_numbers",
        "format": "pptx",
        "path": rel_path,
        "tolerance_pct": int(tolerance * 100),
        "hint": hint,
        **result,
    }, ensure_ascii=False)
