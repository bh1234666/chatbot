"""
Office 文档工具:.docx / .pptx / .xlsx 一体化。

设计原则:
  - 单一工具入口 `office`,根据 path 后缀自动判定格式
  - 5 种 action 跨格式语义统一: read / write / append / extract_images / insert_image
  - 写文档用结构化 blocks/slides/sheets schema,模型不需要写 python-docx/pptx 代码
  - 读文档返回精简 JSON(段落 / 表格 / 图片清单),自动截断防爆 prompt

为什么不让模型直接写 python-docx 代码:
  - log 实测:模型为生成一个含图的 .docx 报告,要写 ~150 行 python(图片插入 emu vs inches、
    表格样式、cell.text 替换会丢样式),经常翻车后再补 edit_file。
  - .pptx 更糟:布局和占位符 placeholders 接口陌生,模型容易把 title 写到错误占位符里。
  - 给一个声明式 schema 把这些坑挡住,模型只管"我要这些内容、那张图、这个标题"。

依赖:
  - python-docx (.docx)    - pip install python-docx
  - python-pptx (.pptx)    - pip install python-pptx
  - openpyxl    (.xlsx)    - pip install openpyxl

任一未安装时只影响对应格式;返回明确错误信息让模型 fallback 到 workspace.run。
"""
from __future__ import annotations

import asyncio
import ast
import json
import logging
import re
import os
import hashlib
import shutil
import tempfile
import zipfile
from typing import Any

from app.llm.tools import workspace as ws_tool
from app.llm.tools.office_locking import docx_write_lock, replace_file_with_retries

log = logging.getLogger(__name__)

# docx 处理族已抽离到 office_docx.py(2026-05-20 重构);re-export 兼容。
from app.llm.tools.office_docx import (  # noqa: E402,F401
    _MAX_PARAGRAPHS_RETURNED,
    _MAX_TABLE_ROWS_RETURNED,
    _MAX_TABLE_COLS_RETURNED,
    _docx_enumerate_body,
    _drop_orphan_media,
    _clean_docx_text,
    _add_run_with_breaks,
    _setup_docx_defaults,
    _find_body_placeholder,
    _extract_tables_generic,
    _list_images_in_part,
    _docx_read,
    _docx_delete_block,
)
# pptx 处理族已抽离到 office_pptx.py(2026-05-20 重构);re-export 兼容。
from app.llm.tools.office_pptx import (  # noqa: E402,F401
    _MAX_SLIDES_RETURNED,
    _UNICODE_SUBSCRIPT,
    _UNICODE_SUPERSCRIPT,
    _script_text,
    _dump_images_from_part,
    _pptx_remove_slide,
    _pptx_delete_slide,
    _pptx_read,
)
# xlsx 操作族已抽离到 office_xlsx.py(2026-05-20 重构);re-export 兼容。
from app.llm.tools.office_xlsx import (  # noqa: E402,F401
    _MAX_XLSX_ROWS_RETURNED,
    _MAX_XLSX_COLS_RETURNED,
    _MAX_SHEETS_RETURNED,
    _openpyxl_col_letter,
    _apply_sheets_xlsx,
    _xlsx_update_cells,
    _xlsx_verify_integrity,
    _xlsx_write,
    _xlsx_append,
    _xlsx_read,
)
_DEFAULT_IMAGE_WIDTH_INCHES = 6.0
_DEFAULT_PPTX_IMAGE_WIDTH_INCHES = 8.0
_MAX_IMAGE_BYTES = 20 * 1024 * 1024
# ── 2026-05-17 P150: 自适应粒度 (Adaptive Granularity) ────────────────
# 起始大粒度,LLM json_broken 时折半,12 封底。
# 上限即 hard reject 阈值,不是建议——LLM 可自由发挥到上限,大文档一次写完。
# 之前的 (8 / 12 / 1200 / 4000) 设计太保守: log 实测 office 调用 36 次,
# 每次都要一次 LLM round-trip (~20s),累计 12 分钟,工具本身只占 5.7s。
_BLOCK_LIMIT_INITIAL    = 96
_BLOCK_LIMIT_FLOOR      = 12
_BLOCK_TEXT_INITIAL     = 5000
_BLOCK_TEXT_FLOOR       = 1200
_TOTAL_TEXT_INITIAL     = 30000
_TOTAL_TEXT_FLOOR       = 4000

_DOCX_USABLE_WIDTH_DXA = 9360

import contextvars as _ctxvars
from threading import Lock as _Lock

_current_office_key: _ctxvars.ContextVar[str] = _ctxvars.ContextVar(
    "office_adaptive_key", default="__main__"
)
_adaptive_state: dict[str, dict] = {}
_adaptive_lock = _Lock()


def set_office_adaptive_key(key: str | None) -> None:
    """供 chat_with_tools_loop 入口调用,绑定本会话(helper/main)的隔离 key.

    在 asyncio task 内 ContextVar 是隔离的,所以并行 helper 互不影响。
    """
    _current_office_key.set(key or "__main__")


def _get_limits(key: str) -> dict:
    """返回当前 key 的自适应上限快照。首次访问时初始化为 INITIAL 值。"""
    with _adaptive_lock:
        st = _adaptive_state.get(key)
        if st is None:
            st = {
                "blocks":      _BLOCK_LIMIT_INITIAL,
                "block_text":  _BLOCK_TEXT_INITIAL,
                "total_text":  _TOTAL_TEXT_INITIAL,
                "shrinks":     0,
                "last_reason": None,
                # 2026-05-18 P176: 跟踪最近 append/write 的 block 数量, 用来检测粒度过小
                "recent_block_counts": [],
            }
            _adaptive_state[key] = st
        return dict(st)


def _record_block_count_and_get_nudge(key: str, action: str, n_blocks: int) -> dict | None:
    """P176: 记录粒度并在过度细化时返回正向建议。

    检测: 最近 4 次 write/append 中, ≥3 次单块 → helper 在做"逐题写"反模式。
    返回 nudge 字典 (None 表示不 nudge), 作为 office 工具结果的附加字段。
    """
    if action not in ("write", "append"):
        return None
    with _adaptive_lock:
        st = _adaptive_state.get(key)
        if st is None:
            return None
        recent = st.get("recent_block_counts") or []
        recent = (recent + [n_blocks])[-6:]  # 保留最近 6 次
        st["recent_block_counts"] = recent

        # nudge 条件: 最近 4 次至少有 3 次单块, 且无自适应缩小过
        if len(recent) >= 4 and st["shrinks"] == 0:
            recent_window = recent[-4:]
            tiny_count = sum(1 for n in recent_window if n <= 1)
            if tiny_count >= 3:
                return {
                    "nudge_kind": "granularity_too_fine",
                    "recent_block_counts": recent,
                    "current_blocks_ceiling": st["blocks"],
                    "advice": (
                        f"Recent office {action} granularity fact: {tiny_count} of the last 4 calls wrote "
                        f"1 block while the current adaptive ceiling is {st['blocks']} blocks and no shrink "
                        "has occurred. A single call can carry multiple coherent paragraphs, tables, or "
                        "section blocks within the current ceiling; choose the next granularity from the "
                        "remaining acceptance gaps and JSON reliability.\n"
                        f"中文摘要：最近 4 次 {action} 中 {tiny_count} 次只写 1 个 block，当前上限 "
                        f"{st['blocks']}；可按剩余验收缺口决定是否合并多个相关 block。"
                    ),
                }
        return None


def report_office_failure(key: str | None, reason: str) -> dict:
    """外部接口: llm/client.py 在检测到 office 调用的 json_broken 时调用。

    每次调用: blocks/block_text/total_text 各折半,但不低于对应 FLOOR。
    返回新的 limits 快照(供调用方日志)。
    """
    k = key or "__main__"
    with _adaptive_lock:
        st = _adaptive_state.setdefault(k, {
            "blocks":      _BLOCK_LIMIT_INITIAL,
            "block_text":  _BLOCK_TEXT_INITIAL,
            "total_text":  _TOTAL_TEXT_INITIAL,
            "shrinks":     0,
            "last_reason": None,
            "recent_block_counts": [],
        })
        if st["blocks"] <= _BLOCK_LIMIT_FLOOR:
            # 已经在 floor,只记录原因
            st["last_reason"] = reason
            snap = dict(st)
            already_at_floor = True
        else:
            st["blocks"]     = max(st["blocks"] // 2,      _BLOCK_LIMIT_FLOOR)
            st["block_text"] = max(st["block_text"] // 2,  _BLOCK_TEXT_FLOOR)
            st["total_text"] = max(st["total_text"] // 2,  _TOTAL_TEXT_FLOOR)
            st["shrinks"]    += 1
            st["last_reason"] = reason
            snap = dict(st)
            already_at_floor = False
    try:
        from app.core import debug
        if already_at_floor:
            debug.log(
                "office.adaptive.at_floor",
                f"key={k}: already at floor(blocks={snap['blocks']}), "
                f"shrinks={snap['shrinks']}, reason={reason[:80]}",
            )
        else:
            debug.log(
                "office.adaptive.shrunk",
                f"key={k}: blocks→{snap['blocks']} block_text→{snap['block_text']} "
                f"total_text→{snap['total_text']} (shrinks={snap['shrinks']}) "
                f"reason={reason[:80]}",
            )
    except Exception:
        pass
    return snap


def reset_office_adaptive_state(key: str | None = None) -> None:
    """重置某个 key 的自适应状态。用于 helper 完成后由 delegate 调用清理。

    若 key=None 则不做任何事(避免误清主线程状态)。
    """
    if not key:
        return
    with _adaptive_lock:
        _adaptive_state.pop(key, None)


# ─────────────────────────────────────────────────────────────────
# 入口分发
# ─────────────────────────────────────────────────────────────────

async def handle_office(workspace_dir: str, args: dict) -> str:
    """office 工具主分发。返回 JSON 字符串。"""
    if not workspace_dir:
        return _err("workspace not available (easy path)")

    action = str(args.get("action", "")).strip().lower()
    path = str(args.get("path", "")).strip()
    if not action:
        return _office_arg_recovery_error("", ["action", "path"])
    if not path:
        return _office_arg_recovery_error(action, ["path"])

    # 2026-05-18 P197: action=ocr_images 也接受**裸图片路径** (jpg/png/...)。
    # 实测病例 (trace 20:48-21:05): helper 收到主线程图片 (`_downloaded_media/img_X.jpg`)
    # 后调 `office ocr_images <jpg>` → 之前直接报 `unsupported extension '.jpg'`
    # → helper 走灾难工作流: 建空 .docx → 嵌图 → ocr_images 再读 → 用了 3 倍时间
    # → 上下文塞满 3 份 OCR 重复输出 → 上下文 truncate → 丢失原文 → 凭训练集脑补不同题目。
    # 修复: 裸图片直接调 ocr_file, 无需 zip 解包。
    _IMAGE_EXTS_RAW = {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".tiff", ".gif"}
    ext_check = os.path.splitext(path)[1].lower()
    if action == "ocr_images" and ext_check in _IMAGE_EXTS_RAW:
        try:
            target = ws_tool._safe_resolve(workspace_dir, path)
        except ValueError as e:
            return _err(f"invalid path: {e}", action=action,
                        hint="Use a relative path within the workspace")
        return await _raw_image_ocr(workspace_dir, target, path, args)

    # 后缀决定格式
    ext = os.path.splitext(path)[1].lower()
    if ext == ".docx":
        fmt = "docx"
    elif ext == ".pptx":
        fmt = "pptx"
    elif ext in (".xlsx", ".xlsm"):
        fmt = "xlsx"
    else:
        # 2026-05-18 P198: 若是图片后缀但不是 ocr_images action, 给定向 hint.
        if ext_check in _IMAGE_EXTS_RAW:
            return _err(
                f"unsupported extension {ext!r} for action={action!r}",
                hint=(
                    f"Image file {path!r} is handled by office action='ocr_images' directly. "
                    f"The current action is {action!r}; retry with action='ocr_images' for this image, "
                    "or use an Office document path for other office actions.\n\n"
                    "图片文件直接使用 action='ocr_images'；其它 office action 需要 Office 文档路径。"
                )
            )
        return _err(
            f"unsupported extension {ext!r}",
            hint="Rename the file to use a supported extension: .docx (Word), .pptx (PowerPoint), .xlsx/.xlsm (Excel)"
        )

    # 路径解析
    try:
        target = ws_tool._safe_resolve(workspace_dir, path)
    except ValueError as e:
        return _err(f"invalid path: {e}", action=action, hint="Use a relative path within the workspace, e.g. 'output/report.docx'")

    # 库可用性检查
    err = _check_lib(fmt)
    if err:
        return _err(err, action=action, hint="Install the missing library via workspace.run first, then retry the office call")

    # 分发
    valid_actions = (
        "read", "write", "append",
        # 2026-05-04 v19.1: 增量编辑能力
        # 大文档(论文/PPT)一次性 write 容易触发 LLM thinking timeout,
        # 提供"按章节/按 block 增量改"的能力,让模型分多次小调用完成。
        "replace_section", "replace_block", "replace_blocks", "fill_empty_headings", "delete_block", "insert_block",
        # pptx 专用
        "replace_slide", "insert_slide", "delete_slide",
        # xlsx 专用
        "update_cells",
        "extract_images", "insert_image",
        # 2026-05-16: 一键 OCR 所有嵌入图片 (docx/pptx/xlsx)
        # 替代 extract_images + ocr(每张) 多步骤的工具组合
        "ocr_images",
        # 2026-05-17 P161: 对照 CSV 验证 docx/pptx 里的数字, 防 LLM 凭印象编数
        "verify_numbers",
        # 2026-05-17 P162 系列: docx 9 项数据严谨性检查 (L1-L9)
        "verify_rigor",
        # 2026-05-17 P162.7: xlsx 跨 sheet 一致性 + 公式 cached_value 健全性
        "verify_integrity",
        # 2026-05-18 P203: docx 内容 vs source 文件一致性检查 (防 OCR hallucination)
        "verify_against_source",
    )
    if action not in valid_actions:
        return _err(f"unknown action: {action!r}", hint=f"Valid actions: {', '.join(valid_actions)}")

    handler = _DISPATCH.get((fmt, action))
    if handler is None:
        supported = sorted(a for f, a in _DISPATCH if f == fmt)
        if fmt == "docx":
            hint = (
                "For DOCX, use read for structure/body inspection, write/append/replace/insert/delete for edits, "
                "verify_numbers or verify_rigor for data claims, and ocr_images for embedded image text. "
                "verify_integrity is XLSX-only."
                "\n\nDOCX 使用 read/编辑/数字或严谨性验证；嵌入图片文字用 ocr_images，verify_integrity 只适用于 XLSX。"
            )
            next_call_fact = (
                "This is a DOCX file. Structural acceptance facts come from office(action='read'); "
                "numeric/data checks use verify_numbers or verify_rigor with csv_paths; verify_integrity is only for XLSX."
            )
        elif fmt == "xlsx":
            hint = (
                "For XLSX, use read/write/append/update_cells for sheets, ocr_images for embedded image text, "
                "and verify_integrity for cross-sheet/formula integrity."
                "\n\nXLSX 使用表格读写、单元格更新、嵌入图片 OCR 和跨表/公式完整性验证。"
            )
            next_call_fact = "This is an XLSX file; verify_integrity can be used for cross-sheet/formula integrity checks."
        elif fmt == "pptx":
            hint = (
                "For PPTX, use read/write/append/slide edit actions, extract_images or ocr_images for media, "
                "and verify_numbers for numeric claims."
                "\n\nPPTX 使用幻灯片读写编辑、图片提取/OCR 和数字声明验证。"
            )
            next_call_fact = "This is a PPTX file; structural checks use read, slide edits use slide actions, and numeric checks use verify_numbers."
        else:
            hint = f"For .{fmt} files, valid actions are: {', '.join(supported)}."
            next_call_fact = f"For .{fmt} files, valid Office actions are: {', '.join(supported)}."
        return _err(
            f"action {action!r} not supported for .{fmt}",
            hint=f"{hint}\n\n可用 action: {', '.join(supported)}",
            extra={"supported_actions": supported, "next_call_fact": next_call_fact},
        )

    try:
        return await handler(workspace_dir, target, path, args)
    except Exception as e:
        log.exception("office tool failed: fmt=%s action=%s", fmt, action)
        return _err(
            f"{type(e).__name__}: {e}",
            action=action,
            hint="This is an unexpected internal error. Check the error message for details before retrying."
        )


def _office_arg_recovery_error(action: str, required_fields: list[str]) -> str:
    # 2026-05-18 P174: 区分两类错误根因
    # 旧版本: 不论 path 丢失还是 blocks 太大, 都笼统说"JSON 太大请用更小 blocks 重试"
    # → helper 看到错误 (即使根因是 path 字段丢失) 也会缩小 blocks 数量, 误判降粒度。
    # 修复: 根据 required_fields 区分提示。
    is_missing_path = "path" in required_fields
    is_missing_blocks = "blocks" in required_fields

    if is_missing_path and not is_missing_blocks:
        hint = (
            "The call is missing the `path` field. This is not caused by too many blocks; the tool-call JSON lost a required field, often because quotes or newlines inside strings were not escaped. "
            "Retry with the intended target path. The previous block count may still be usable; adjust blocks only if a separate payload-size or JSON-stability warning appears.\n"
            "缺少 path 字段；补齐 path；只有出现大小或 JSON 稳定性问题时才需要改 blocks 粒度。"
        )
        recovery_advice = "Add the missing path field; keep or adjust blocks according to any separate size/JSON warning."
        next_call_fact = "The previous Office call was missing `path`; retry the same action with the target path and intended blocks."
    elif is_missing_blocks:
        hint = (
            "The `blocks` field is empty or missing. It must be a non-empty array such as [{type:..., text:...}, ...]. "
            "Current Office limits are adaptive; the tool result reports the active block and text ceilings when size matters.\n"
            "blocks 需要是非空数组；大小相关上限以工具结果中的当前限制为准。"
        )
        recovery_advice = "Provide a non-empty blocks array; use the current adaptive limits from tool facts when relevant."
        next_call_fact = "The previous Office call had no usable `blocks`; provide a non-empty blocks array."
    else:
        # 多字段缺失 / action 错误
        hint = (
            "The tool arguments are empty or missing required fields such as action/path/blocks. "
            "Check that JSON serialization is complete and that quotes/backslashes inside strings are escaped correctly.\n"
            "核心字段缺失；检查 JSON 与字符串转义。"
        )
        recovery_advice = "Verify all required fields are present and JSON is valid."
        next_call_fact = "The previous Office call was missing required fields; include action/path and content fields required by that action."

    return _err(
        "invalid_or_empty_args",
        action=action,
        hint=hint,
        extra={
            "recovery": {
                "action": "retry_with_complete_args",  # 旧的 "retry_with_smaller_blocks" 是误导
                "required_fields": required_fields,
                "advice": recovery_advice,
            },
            "next_call_fact": next_call_fact,
        },
    )


def _office_blocks_sizing_warning(action: str, blocks: list) -> dict | None:
    # 2026-05-17 P150: 改用自适应 limits (按 task_id 隔离, 失败折半 12 封底)
    key = _current_office_key.get()
    limits = _get_limits(key)
    blocks_ceiling = limits["blocks"]
    block_text_max = limits["block_text"]
    total_max      = limits["total_text"]

    oversize_texts: list[dict] = []
    total_text_chars = 0
    for idx, block in enumerate(blocks):
        if not isinstance(block, dict):
            continue
        text = block.get("text")
        if isinstance(text, str):
            total_text_chars += len(text)
            if len(text) > block_text_max:
                oversize_texts.append({"index": idx, "chars": len(text)})

    # 自适应模式: ceiling 即 hard limit; 低于 ceiling 不警告 (让 LLM 自由发挥)
    too_many = len(blocks) > blocks_ceiling
    too_many_hard = too_many
    total_text_too_large = total_text_chars > total_max
    if not too_many and not oversize_texts and not total_text_too_large:
        return None

    shrunk_hint = ""
    if limits["shrinks"] > 0:
        shrunk_hint = (
            f" Adaptive limits have already shrunk {limits['shrinks']} time(s) in this task; "
            f"last reason: {(limits['last_reason'] or '?')[:60]}."
        )
    floor_hint = ""
    if blocks_ceiling <= _BLOCK_LIMIT_FLOOR:
        floor_hint = (
            f" Reached the minimum block limit {_BLOCK_LIMIT_FLOOR}; keep this block size. "
            "If it still fails, inspect whether the JSON output contains unescaped newlines or quotes.\n"
            "已到最小块数，后续应检查 JSON 转义。"
        )

    return {
        "blocks_count": len(blocks),
        "current_blocks_limit": blocks_ceiling,
        "current_block_text_chars_limit": block_text_max,
        "current_total_text_chars_limit": total_max,
        "adaptive_shrinks_so_far": limits["shrinks"],
        "adaptive_floor_reached": blocks_ceiling <= _BLOCK_LIMIT_FLOOR,
        "total_text_chars": total_text_chars,
        "too_many_blocks": too_many,
        "too_many_blocks_hard": too_many_hard,
        "total_text_too_large": total_text_too_large,
        "oversize_text_blocks": oversize_texts[:8],
        "next_action_instruction": (
            f"The {action} call exceeds the current adaptive office block limits: "
            f"blocks <= {blocks_ceiling}, block.text <= {block_text_max} characters, "
            f"total text <= {total_max} characters (compact: blocks≤{blocks_ceiling}, "
            f"block.text≤{block_text_max}, total_text≤{total_max}).{shrunk_hint}{floor_hint}\n"
            "拆小 blocks 或正文后重试；达到封底时检查 JSON 转义。"
        ),
    }


def _office_blocks_hard_limit_error(action: str, blocks: list) -> str | None:
    warning = _office_blocks_sizing_warning(action, blocks)
    if not warning:
        return None
    if not warning.get("too_many_blocks_hard"):
        return None
    floor_extra = (
        " The minimum block limit has been reached; if this still fails, inspect JSON escaping."
        if warning.get("adaptive_floor_reached")
        else f" Adaptive limits have already shrunk {warning['adaptive_shrinks_so_far']} time(s) in this task."
    )
    return _err(
        "office_args_too_large",
        action=action,
        hint=(
            f"The {action} call has {warning['blocks_count']} blocks, which exceeds the current adaptive limit "
            f"{warning['current_blocks_limit']}. Split the document operation into smaller calls and keep each block "
            f"within the current text limits.{floor_extra}\n"
            "Office blocks 超过当前上限；拆成更小调用并保持当前粒度。"
        ),
        extra={
            "arg_size_warning": warning,
            "recovery": {
                "action": "retry_with_smaller_blocks",
                "required_fields": ["action", "path", "blocks"],
            },
        },
    )


def _office_blocks_format_warning(action: str, blocks: list) -> dict | None:
    """检测 block.text 里残留的 LaTeX/Markdown 数学标记 (^{...} / _{...} / \\frac /
    \\sqrt 等), 但不硬改。给 LLM 可执行的提示让它自己决定 (公式残留就重写, 故意
    保留 raw 风格则忽略)。

    实测 trace fee099 (2026-06-05): docx 文末出现 `e^{-r/2}` `e^{-r/4}` 等 6 处
    plain-text 形式没有渲染成上下标也没用 LaTeX `$...$`, 看起来像残骸。这里只
    报告位置, 由 LLM 决定是否补救。
    """
    findings: list[dict] = []
    seen: set[tuple[int, str]] = set()
    # 检测 ^{...} / _{...} 风格 (raw markdown/LaTeX 风格的上下标 + 大括号包裹)
    pat_brace = re.compile(r"([\^_])\{([^{}\n]+)\}")
    # 检测裸 LaTeX 命令字符串(被剥掉 $ 后的命令残骸)
    pat_latex_cmd = re.compile(r"\\(?:frac|sqrt|sum|int|prod|alpha|beta|gamma|delta|"
                                r"theta|sigma|tau|phi|psi|omega|partial|nabla|cdot|"
                                r"text|mathrm|mathbf|boxed|begin|end|left|right|"
                                r"infty|leq|geq|neq|approx|equiv)\b")
    # 检测未配对的 $ 符号 (奇数个 $ 通常意味着 LaTeX 包裹失败)
    for idx, block in enumerate(blocks or []):
        if not isinstance(block, dict):
            continue
        text = block.get("text") or block.get("caption") or ""
        if not isinstance(text, str) or not text:
            # 也扫表格 cells
            cells = block.get("cells")
            if isinstance(cells, list):
                text = "\n".join(
                    str(c) for row in cells if isinstance(row, list)
                    for c in row if isinstance(c, str)
                )
            if not text:
                continue
        # ^{...} / _{...}
        for m in pat_brace.finditer(text):
            kind = "superscript" if m.group(1) == "^" else "subscript"
            sample = m.group(0)
            key = (idx, sample)
            if key in seen:
                continue
            seen.add(key)
            findings.append({"block_index": idx, "kind": kind, "sample": sample})
            if len(findings) >= 12:
                break
        if len(findings) >= 12:
            break
        # \frac \sqrt 等命令
        for m in pat_latex_cmd.finditer(text):
            sample = m.group(0)
            key = (idx, "cmd:" + sample)
            if key in seen:
                continue
            seen.add(key)
            findings.append({"block_index": idx, "kind": "raw_latex_command", "sample": sample})
            if len(findings) >= 12:
                break
        if len(findings) >= 12:
            break
        # 奇数 $ 数量 (一段里 $ 出现奇数次 → 包裹未闭合)
        n_dollar = text.count("$")
        if n_dollar % 2 == 1 and n_dollar > 0:
            findings.append({"block_index": idx, "kind": "unbalanced_dollar", "count": n_dollar})

    if not findings:
        return None
    return {
        "issue": "raw_math_markup_in_plain_text",
        "findings_count": len(findings),
        "findings": findings,
        "advisory": (
            f"The {action} call has {len(findings)} block-text occurrence(s) of LaTeX/Markdown style "
            f"math markup (like `^{{-r/4}}` or `\\frac{{a}}{{b}}` or unbalanced `$`) sitting as plain text. "
            f"In docx these will render literally as `^{{-r/4}}` rather than as a real superscript or formula. "
            f"If these are intentional code samples / notation explanations, ignore this. Otherwise consider "
            f"either: (a) wrapping the math in `$...$` so the office tool renders it via LaTeX→Unicode/OMML, "
            f"or (b) rewriting with Unicode super/subscript characters directly (e.g. `e⁻ʳ⁄⁴` or `e^(-r/4)`).\n"
            f"docx 中残留的 ^{{...}} / \\frac 等 raw 数学标记不会自动渲染；按需用 $...$ 包或改 Unicode 上下标，"
            f"如果是代码示例可忽略。"
        ),
    }


def _attach_block_sizing_warning(payload: dict, action: str, blocks: list) -> dict:
    warning = _office_blocks_sizing_warning(action, blocks)
    if warning:
        payload["arg_size_warning"] = warning
    # 2026-06-05: 公式/数学残留检测 (软提示, 不改文)
    fmt_warning = _office_blocks_format_warning(action, blocks)
    if fmt_warning:
        payload["format_warnings"] = fmt_warning
    # 2026-05-18 P176: 粒度过细 nudge — 记录最近 block 数, 若过度细化给 helper 正向引导
    try:
        key = _current_office_key.get()
        nudge = _record_block_count_and_get_nudge(key, action, len(blocks) if blocks else 0)
        if nudge:
            payload["granularity_nudge"] = nudge
    except Exception:
        pass
    normalization = _office_blocks_shape_normalization_note(blocks)
    if normalization:
        payload["input_normalizations"] = normalization
    return payload


def _office_blocks_shape_normalization_note(blocks: list) -> dict | None:
    """Return compact facts for tolerated DOCX block-shape aliases.

    These are not acceptance decisions. They tell the model which input shapes were
    interpreted as the nearest supported Office schema shape so future calls can use
    the canonical form directly.
    """
    if not isinstance(blocks, list):
        return None
    subtitle_blocks = 0
    table_row_objects = 0
    for b in blocks:
        if not isinstance(b, dict):
            continue
        if b.get("_office_normalized_type_from") == "subtitle":
            subtitle_blocks += 1
        table_row_objects += int(b.get("_office_normalized_table_row_objects") or 0)
    if not subtitle_blocks and not table_row_objects:
        return None
    facts: dict[str, Any] = {}
    if subtitle_blocks:
        facts["subtitle_blocks_as_paragraph"] = subtitle_blocks
    if table_row_objects:
        facts["table_rows_cells_objects_as_arrays"] = table_row_objects
    facts["canonical_docx_fact"] = (
        "DOCX blocks use type=paragraph for subtitles/plain prose; table rows are non-empty 2D arrays, "
        "for example rows:[[\"A\",\"B\"],[\"C\",\"D\"]]."
    )
    return facts


# office 通用 helper 已抽离到 office_common.py(2026-05-20 重构);re-export 兼容。
from app.llm.tools.office_common import (  # noqa: E402,F401
    _err,
    _jsonable,
    _check_lib,
    _resolve_out_dir,
)




# ═══════════════════════════════════════════════════════════════
# Word (.docx)
# ═══════════════════════════════════════════════════════════════


def _save_docx_atomic(doc, target: str) -> tuple[int, str | None]:
    """Save a DOCX through a verified temp file before replacing target."""
    directory = os.path.dirname(target) or "."
    os.makedirs(directory, exist_ok=True)
    fd, tmp_path = tempfile.mkstemp(prefix=".docx_tmp_", suffix=".docx", dir=directory)
    os.close(fd)
    bak_path = target + ".bak"
    try:
        doc.save(tmp_path)
        with zipfile.ZipFile(tmp_path, "r") as zf:
            bad_member = zf.testzip()
        if bad_member:
            return 0, f"save verification failed: bad zip member {bad_member}"
        if os.path.exists(target):
            try:
                shutil.copy2(target, bak_path)
            except OSError:
                pass
        replace_err = replace_file_with_retries(tmp_path, target)
        if replace_err is not None:
            return 0, (
                f"save failed: {replace_err} "
                "The target DOCX is likely still locked by another Office/read/write operation. "
                "Wait briefly, avoid parallel edits to the same file, then retry the same operation from the current evidence.\n"
                "保存失败:目标 DOCX 仍被占用；请避免并行编辑同一文件，稍后基于当前证据重试。"
            )
        size = os.path.getsize(target) if os.path.exists(target) else 0
        return size, None
    except Exception as e:
        return 0, f"save failed: {type(e).__name__}: {e}"
    finally:
        try:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
        except OSError:
            pass


def _docx_open_error(e: Exception, rel_path: str) -> str:
    return (
        f"open failed: {type(e).__name__}: {e}. The DOCX package may be damaged or half-written: {rel_path}. "
        "If a sibling .bak file exists, copy it to a new clean filename and continue from that backup; otherwise "
        "rebuild the document from verified source/evidence instead of repeatedly editing this damaged file.\n"
        "DOCX 可能已损坏或半写入；优先用 .bak 复制成干净文件继续，没有备份则从证据重建。"
    )


def _open_docx_detached(target: str):
    """Open a DOCX from a private temp copy so Windows can replace the target later."""
    from docx import Document

    directory = os.path.dirname(target) or "."
    fd, tmp_path = tempfile.mkstemp(prefix=".docx_read_", suffix=".docx", dir=directory)
    os.close(fd)
    try:
        shutil.copy2(target, tmp_path)
        return Document(tmp_path)
    finally:
        try:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
        except OSError:
            pass


async def _docx_write(workspace_dir: str, target: str, rel_path: str, args: dict) -> str:
    from docx import Document

    blocks = args.get("blocks") or []
    if not isinstance(blocks, list):
        return _err("blocks must be an array", action="write", hint="Pass blocks as a JSON array: [{\"type\":\"heading\",\"level\":1,\"text\":\"Title\"}, ...]")

    title = args.get("title")

    def _do_write():
        doc = Document()
        # 2026-05-18 P196: A4 + 中文字体 + 含图段落自适应行高
        _setup_docx_defaults(doc)
        if isinstance(title, str) and title.strip():
            doc.add_heading(_clean_docx_text(title), level=0)
        err = _apply_blocks_docx(doc, workspace_dir, blocks)
        if err:
            return None, err
        size, save_err = _save_docx_atomic(doc, target)
        if save_err:
            return None, save_err
        return size, None

    async with docx_write_lock(target):
        size, err = await asyncio.to_thread(_do_write)
    if err:
        return _err(err)
    payload = _attach_block_sizing_warning({
        "ok": True, "action": "write", "format": "docx", "path": rel_path,
        "blocks_written": len(blocks), "size_bytes": size,
    }, "write", blocks)
    return json.dumps(payload, ensure_ascii=False)


async def _docx_append(workspace_dir: str, target: str, rel_path: str, args: dict) -> str:
    if not os.path.isfile(target):
        return _err(f"file not found: {rel_path}")

    blocks = args.get("blocks") or []
    if not isinstance(blocks, list) or not blocks:
        return _office_arg_recovery_error("append", ["action", "path", "blocks"])

    hard_limit_error = _office_blocks_hard_limit_error("append", blocks)
    if hard_limit_error:
        return hard_limit_error

    def _do_append():
        try:
            doc = _open_docx_detached(target)
        except Exception as e:
            return _docx_open_error(e, rel_path)
        err = _apply_blocks_docx(doc, workspace_dir, blocks)
        if err:
            return err
        _, save_err = _save_docx_atomic(doc, target)
        if save_err:
            return save_err
        return None

    async with docx_write_lock(target):
        err = await asyncio.to_thread(_do_append)
    if err:
        return _err(err)

    payload = _attach_block_sizing_warning({
        "ok": True, "action": "append", "format": "docx", "path": rel_path,
        "blocks_appended": len(blocks),
    }, "append", blocks)
    return json.dumps(payload, ensure_ascii=False)


# ───────────────────────────────────────────────────────────────
# 2026-05-04 v19.1: docx 增量编辑
# 大文档(论文/报告)一次性 write 整篇容易触发 LLM thinking timeout。
# 增量编辑允许模型分多次小调用完成:
#   1. write(只写第 1 章)         → 几百字,几秒
#   2. append(第 2 章)             → 几百字,几秒
#   3. replace_section(找"摘要"重写) → 几百字,几秒
#   4. replace_block(idx=5 那段重写) → 一段,几秒
#
# 设计要点:
# - 用 _docx_iter_blocks 同步枚举所有 block 的 (index, type, text/heading-level/...)
# - replace_section: 按 heading_text 找,heading 之后的内容直到下一个 同级或更高级
#   heading 之间的所有 block 都被替换
# - replace_block: 直接按 paragraph index 替换(0-based,数到第 N 个 paragraph/heading/table/image)
# - delete_block: 同上,删除指定 index
# - insert_block: 在指定 index 之前插入 blocks
#
# 注:python-docx 的 doc.paragraphs / doc.tables 是分别的列表,
# 真实文档顺序需要从 doc.element.body 的子元素遍历。
# ───────────────────────────────────────────────────────────────






async def _docx_replace_section(
    workspace_dir: str, target: str, rel_path: str, args: dict,
) -> str:
    """按 heading 文字找章节并重写其 body(到下一个 同级或更高级 heading 之前)。

    args:
      heading_text: 要替换的章节标题(精确匹配 strip 后)
      blocks: 新内容(数组,可不含 heading 自己)
      keep_heading: 是否保留原 heading(默认 True)
    """
    from docx import Document

    if not os.path.isfile(target):
        return _err(f"file not found: {rel_path}")
    heading_text = str(args.get("heading_text", "")).strip()
    if not heading_text:
        return _err(
            "missing heading_text",
            action="replace_section",
            hint=(
                "Provide the exact heading text to find. Use the text copied from office(action='read') "
                "or another verified document inspection result; the match is exact after trimming whitespace.\n\n"
                "提供从文档读取结果中确认过的精确章节标题。"
            ),
        )
    blocks = args.get("blocks") or []
    if not isinstance(blocks, list):
        return _err("blocks must be array", action="replace_section", hint="Pass blocks as a JSON array with the new content to replace the section")
    keep_heading = bool(args.get("keep_heading", True))

    def _do():
        try:
            doc = _open_docx_detached(target)
        except Exception as e:
            return None, _docx_open_error(e, rel_path)

        items = _docx_enumerate_body(doc)
        # 找匹配的 heading
        target_idx = None
        target_level = None
        for it in items:
            if it["kind"] == "heading" and it["text"].strip() == heading_text:
                target_idx = it["index"]
                target_level = it["level"]
                break
        if target_idx is None:
            available = [
                f"[{x['index']}] L{x['level']} {x['text'][:30]!r}"
                for x in items if x["kind"] == "heading"
            ][:20]
            return None, (
                f"heading not found: {heading_text!r}; "
                f"available headings (first 20): {available}"
            )

        # 找区间结束:下一个同级或更高级 heading 之前
        end_idx = len(items)  # 默认到文档结尾
        for it in items[target_idx + 1:]:
            if it["kind"] == "heading" and it["level"] <= target_level:
                end_idx = it["index"]
                break

        # 删除区间元素(保留或不保留 heading 自己)
        start_delete = target_idx if not keep_heading else target_idx + 1
        elements_to_delete = [items[i]["_element"] for i in range(start_delete, end_idx)]
        # 倒序删除避免 index 失效
        for el in reversed(elements_to_delete):
            parent = el.getparent()
            if parent is not None:
                parent.remove(el)

        # 在 heading 之后(或被删的位置)插入新 blocks
        # 用 _apply_blocks_docx 在末尾添加,然后 diff 出新 element 移到目标位置。
        # 关键:python-docx 的 add_paragraph 会复用 lxml element 的内存地址(C 层),
        # 所以 Python id() 不可靠。改用 element 引用列表(lxml Element 的 == 是
        # 引用比较,与 is 等价,可以放进 list 里 in-check)。
        body = doc.element.body
        before_list = list(body.iterchildren())

        err = _apply_blocks_docx(doc, workspace_dir, blocks) if blocks else None
        if err:
            return None, err

        new_elements = [
            c for c in body.iterchildren()
            if not any(c is x for x in before_list)
        ]

        # 找插入位置:keep_heading 时是 heading 之后;否则是 heading 之前(已删)
        if keep_heading:
            heading_el = items[target_idx]["_element"]
            anchor = heading_el
            for ne in new_elements:
                parent = ne.getparent()
                if parent is not None:
                    parent.remove(ne)
                anchor.addnext(ne)
                anchor = ne
        else:
            if target_idx > 0:
                anchor = items[target_idx - 1]["_element"]
                for ne in new_elements:
                    parent = ne.getparent()
                    if parent is not None:
                        parent.remove(ne)
                    anchor.addnext(ne)
                    anchor = ne
            # 否则就让它们留在末尾(罕见情况:第一个 heading 就是 target,
            # 不 keep heading + 没有前置元素 = 文档变成纯新 blocks)

        # 2026-05-17 P152: 清理孤儿 image (replace 后旧图可能不再引用)
        orphan_info = _drop_orphan_media(doc)

        _, save_err = _save_docx_atomic(doc, target)
        if save_err:
            return None, save_err
        return {
            "deleted_blocks": end_idx - start_delete,
            "added_blocks": len(blocks),
            "section_level": target_level,
            "orphan_cleanup": orphan_info,
        }, None

    async with docx_write_lock(target):
        info, err = await asyncio.to_thread(_do)
    if err:
        return _err(err)
    return json.dumps({
        "ok": True, "action": "replace_section", "format": "docx",
        "path": rel_path,
        "heading_text": heading_text,
        **info,
    }, ensure_ascii=False)


async def _docx_replace_blocks(
    workspace_dir: str, target: str, rel_path: str, args: dict,
) -> str:
    """批量替换多个 block — 一次 tool 调用改多处.

    2026-05-18 P201: 实测灾难 (trace 0518 21:04-21:09): helper 写完 22 题 docx 后
    发现 ~15 个 block 内容错误, 进入 read → replace_block → read → replace_block
    死循环, **20 次 replace_block + 8 次 read = 28 次工具调用, 累计 5-6 分钟**.
    每次 replace_block 改 1 block, LLM 还要在 round-trip 之间思考"下一个改哪". 
    一次性传 edits[] 数组, 让 helper 一次拿出全部错位并修复.

    args:
      edits: [{"index": 3, "blocks": [...]}, {"index": 7, "blocks": [...]}, ...]
        - 每个 edit 的 index 取自 read 输出
        - blocks 是新内容 (可多 block 替换原 1 个)
        - 多个 edit 的 index 不能重叠 (P201 内部按 index 倒序处理保证 index 不漂移)

    Returns:
      {ok, action, replaced_count, edits_applied: [{index, added_blocks, old_kind}], 
       orphan_cleanup}
    """
    from docx import Document

    if not os.path.isfile(target):
        return _err(f"file not found: {rel_path}", action="replace_blocks")

    edits = args.get("edits") or []
    if not isinstance(edits, list) or not edits:
        return _err(
            "edits must be a non-empty array",
            action="replace_blocks",
            hint='Pass edits:[{"index":3,"blocks":[{"type":"paragraph","text":"..."}]}, ...]'
        )

    # Validate each edit
    cleaned_edits = []
    seen_indices = set()
    for i, e in enumerate(edits):
        if not isinstance(e, dict):
            return _err(f"edits[{i}] must be an object", action="replace_blocks")
        try:
            idx = int(e.get("index", -1))
        except (TypeError, ValueError):
            return _err(f"edits[{i}].index must be int", action="replace_blocks")
        if idx < 0:
            return _err(f"edits[{i}].index must be >= 0", action="replace_blocks")
        if idx in seen_indices:
            return _err(
                f"edits[{i}].index={idx} is duplicate",
                action="replace_blocks",
                hint="Each index can be edited once per call. Merge edits to the same block, or call replace_blocks twice."
            )
        seen_indices.add(idx)
        blocks = e.get("blocks") or []
        if not isinstance(blocks, list) or not blocks:
            return _err(
                f"edits[{i}].blocks must be a non-empty array",
                action="replace_blocks"
            )
        cleaned_edits.append({"index": idx, "blocks": blocks})

    # Sort by index DESCENDING — process tail first so earlier indices stay valid
    cleaned_edits.sort(key=lambda e: e["index"], reverse=True)

    def _do():
        try:
            doc = _open_docx_detached(target)
        except Exception as e:
            return None, _docx_open_error(e, rel_path)
        items = _docx_enumerate_body(doc)
        max_idx = len(items)
        for e in cleaned_edits:
            if e["index"] >= max_idx:
                return None, (
                    f"edits index {e['index']} out of range "
                    f"(doc has {max_idx} blocks). Use 'read' to refresh indices."
                )

        body = doc.element.body
        applied: list[dict] = []
        for e in cleaned_edits:
            idx = e["index"]
            new_blocks = e["blocks"]
            target_el = items[idx]["_element"]
            before_list = list(body.iterchildren())
            err = _apply_blocks_docx(doc, workspace_dir, new_blocks)
            if err:
                return None, f"edit at index {idx} failed: {err}"
            new_elements = [
                c for c in body.iterchildren()
                if not any(c is x for x in before_list)
            ]
            prev_anchor = target_el.getprevious()
            for ne in new_elements:
                parent = ne.getparent()
                if parent is not None:
                    parent.remove(ne)
                if prev_anchor is None:
                    body.insert(0, ne)
                else:
                    prev_anchor.addnext(ne)
                prev_anchor = ne
            parent = target_el.getparent()
            if parent is not None:
                parent.remove(target_el)
            applied.append({
                "index": idx,
                "added_blocks": len(new_blocks),
                "old_kind": items[idx]["kind"],
                "old_text_preview": items[idx]["text"][:60],
            })

        # P152: orphan media cleanup once at end
        orphan_info = _drop_orphan_media(doc)

        _, save_err = _save_docx_atomic(doc, target)
        if save_err:
            return None, save_err
        return {
            "replaced_count": len(applied),
            "edits_applied": applied,
            "orphan_cleanup": orphan_info,
        }, None

    async with docx_write_lock(target):
        info, err = await asyncio.to_thread(_do)
    if err:
        return _err(err)
    return json.dumps({
        "ok": True, "action": "replace_blocks", "format": "docx",
        "path": rel_path,
        **info,
    }, ensure_ascii=False)


async def _docx_replace_block(
    workspace_dir: str, target: str, rel_path: str, args: dict,
) -> str:
    """替换指定 index 的单个 block。

    args:
      index: 要替换的 block 序号(从 0 开始;用 read 看到的 index)
      blocks: 新内容(数组,可多个 block 替换原 1 个)

    2026-05-18 P201: 若要改多处, 用 `replace_blocks` (复数) 一次传 edits[] 数组
    避免反复 read→replace 循环. 实测改 15 个 block 用 replace_blocks 1 次调用
    替代 replace_block × 15 次, 省 ~3 分钟.
    """
    from docx import Document

    if not os.path.isfile(target):
        return _err(f"file not found: {rel_path}", action="replace_block")
    try:
        index = int(args.get("index", -1))
    except (TypeError, ValueError):
        return _err("index must be integer", action="replace_block", hint="Provide a 0-based block_index from the read output, e.g. index:3")
    if index < 0:
        return _err("index must be >= 0", action="replace_block", hint="The block_index starts at 0; check the read output for valid indices")
    blocks = args.get("blocks") or []
    if not isinstance(blocks, list) or not blocks:
        return _err("blocks must be a non-empty array", action="replace_block", hint="Pass new content as blocks:[{type:'paragraph', text:'...'}, ...]")

    def _do():
        try:
            doc = _open_docx_detached(target)
        except Exception as e:
            return None, _docx_open_error(e, rel_path)
        items = _docx_enumerate_body(doc)
        if index >= len(items):
            return None, f"index {index} out of range (have {len(items)} blocks)"

        target_el = items[index]["_element"]
        body = doc.element.body
        before_list = list(body.iterchildren())
        err = _apply_blocks_docx(doc, workspace_dir, blocks)
        if err:
            return None, err
        new_elements = [
            c for c in body.iterchildren()
            if not any(c is x for x in before_list)
        ]

        # 把新元素移到目标位置之前
        prev_anchor = target_el.getprevious()
        for ne in new_elements:
            parent = ne.getparent()
            if parent is not None:
                parent.remove(ne)
            if prev_anchor is None:
                # 插到 body 最前面
                body.insert(0, ne)
            else:
                prev_anchor.addnext(ne)
            prev_anchor = ne

        # 删除原 block
        parent = target_el.getparent()
        if parent is not None:
            parent.remove(target_el)

        # 2026-05-17 P152: 清理孤儿 image
        orphan_info = _drop_orphan_media(doc)

        _, save_err = _save_docx_atomic(doc, target)
        if save_err:
            return None, save_err
        return {
            "replaced_index": index,
            "added_blocks": len(blocks),
            "old_kind": items[index]["kind"],
            "old_text_preview": items[index]["text"][:60],
            "orphan_cleanup": orphan_info,
        }, None

    async with docx_write_lock(target):
        info, err = await asyncio.to_thread(_do)
    if err:
        return _err(err)
    return json.dumps({
        "ok": True, "action": "replace_block", "format": "docx",
        "path": rel_path,
        **info,
    }, ensure_ascii=False)


async def _docx_fill_empty_headings(
    workspace_dir: str, target: str, rel_path: str, args: dict,
) -> str:
    """Fill empty DOCX heading paragraphs in one call, optionally inserting body blocks.

    The model still supplies all titles and body content. This tool only applies
    that decided content efficiently to existing empty Heading paragraphs.
    """
    if not os.path.isfile(target):
        return _err(f"file not found: {rel_path}", action="fill_empty_headings")
    headings = args.get("headings") or []
    if not isinstance(headings, list) or not headings:
        return _err(
            "headings must be a non-empty array",
            action="fill_empty_headings",
            hint=(
                "Pass headings:[{text:'...', level:1, after_blocks:[{type:'paragraph', text:'...'}]}, ...]. "
                "Entries are applied to empty heading paragraphs in document order; level is optional and used as a guard."
            ),
        )

    cleaned: list[dict] = []
    for i, item in enumerate(headings):
        if not isinstance(item, dict):
            return _err(f"headings[{i}] must be an object", action="fill_empty_headings")
        text = str(item.get("text") or "").strip()
        if not text:
            return _err(f"headings[{i}].text is required", action="fill_empty_headings")
        level = item.get("level")
        if level is not None:
            try:
                level = max(1, min(9, int(level)))
            except (TypeError, ValueError):
                return _err(f"headings[{i}].level must be an integer", action="fill_empty_headings")
        after_blocks = item.get("after_blocks") or item.get("blocks_after") or []
        if after_blocks and not isinstance(after_blocks, list):
            return _err(f"headings[{i}].after_blocks must be an array", action="fill_empty_headings")
        cleaned.append({"text": text, "level": level, "after_blocks": after_blocks})

    def _style_level(paragraph) -> int | None:
        name = str(getattr(getattr(paragraph, "style", None), "name", "") or "")
        if "heading" not in name.lower():
            return None
        m = re.search(r"(\d+)", name)
        return int(m.group(1)) if m else 1

    def _set_paragraph_text(paragraph, text: str) -> None:
        while len(paragraph.runs) > 1:
            paragraph._p.remove(paragraph.runs[-1]._r)
        if paragraph.runs:
            paragraph.runs[0].text = text
        else:
            paragraph.add_run(text)

    def _do():
        try:
            doc = _open_docx_detached(target)
        except Exception as e:
            return None, _docx_open_error(e, rel_path)

        targets = []
        for idx, paragraph in enumerate(doc.paragraphs):
            level = _style_level(paragraph)
            if level is None:
                continue
            if str(paragraph.text or "").strip():
                continue
            targets.append((idx, paragraph, level))
        if len(cleaned) > len(targets):
            return None, (
                f"received {len(cleaned)} heading entries, but only {len(targets)} empty heading paragraph(s) exist. "
                "Read the document structure and retry with the exact remaining empty headings."
            )

        filled = []
        body = doc.element.body
        for item, (idx, paragraph, actual_level) in zip(cleaned, targets):
            expected_level = item.get("level")
            if expected_level is not None and expected_level != actual_level:
                return None, (
                    f"heading level mismatch at paragraph {idx}: expected level {expected_level}, actual level {actual_level}. "
                    "Read the document structure and align the heading list before retrying."
                )
            _set_paragraph_text(paragraph, item["text"])
            anchor = paragraph._p
            after_blocks = item.get("after_blocks") or []
            if after_blocks:
                before_list = list(body.iterchildren())
                err = _apply_blocks_docx(doc, workspace_dir, after_blocks)
                if err:
                    return None, f"after_blocks for heading {idx} failed: {err}"
                new_elements = [
                    c for c in body.iterchildren()
                    if not any(c is x for x in before_list)
                ]
                for ne in new_elements:
                    parent = ne.getparent()
                    if parent is not None:
                        parent.remove(ne)
                    anchor.addnext(ne)
                    anchor = ne
            filled.append({
                "paragraph_index": idx,
                "level": actual_level,
                "text": item["text"],
                "after_blocks": len(after_blocks),
            })

        _, save_err = _save_docx_atomic(doc, target)
        if save_err:
            return None, save_err
        remaining = 0
        for paragraph in doc.paragraphs:
            if _style_level(paragraph) is not None and not str(paragraph.text or "").strip():
                remaining += 1
        return {
            "filled_count": len(filled),
            "remaining_empty_headings": remaining,
            "filled": filled[:80],
        }, None

    async with docx_write_lock(target):
        info, err = await asyncio.to_thread(_do)
    if err:
        return _err(err)
    return json.dumps({
        "ok": True, "action": "fill_empty_headings", "format": "docx",
        "path": rel_path,
        **info,
    }, ensure_ascii=False)




async def _docx_insert_block(
    workspace_dir: str, target: str, rel_path: str, args: dict,
) -> str:
    """在指定 index 之前插入 blocks。

    args:
      index: 在哪个 block 之前插入(0 表示插到最前)
      blocks: 要插入的内容
    """
    from docx import Document

    if not os.path.isfile(target):
        return _err(f"file not found: {rel_path}", action="insert_block")
    try:
        index = int(args.get("index", -1))
    except (TypeError, ValueError):
        return _err("index must be integer", action="insert_block", hint="Provide a 0-based position, e.g. index:0 to insert at the beginning")
    if index < 0:
        return _err("index must be >= 0", action="insert_block")
    blocks = args.get("blocks") or []
    if not isinstance(blocks, list) or not blocks:
        return _err("blocks must be a non-empty array", action="insert_block", hint="Pass the content to insert as blocks:[{type:'paragraph', text:'...'}, ...]")

    def _do():
        try:
            doc = _open_docx_detached(target)
        except Exception as e:
            return None, _docx_open_error(e, rel_path)
        items = _docx_enumerate_body(doc)
        # 允许 index == len (= 等价 append)
        if index > len(items):
            return None, f"index {index} out of range (max {len(items)})"

        body = doc.element.body
        before_list = list(body.iterchildren())
        err = _apply_blocks_docx(doc, workspace_dir, blocks)
        if err:
            return None, err
        new_elements = [
            c for c in body.iterchildren()
            if not any(c is x for x in before_list)
        ]

        if index == len(items):
            # append 末尾,不动
            pass
        else:
            target_el = items[index]["_element"]
            prev = target_el.getprevious()
            for ne in new_elements:
                parent = ne.getparent()
                if parent is not None:
                    parent.remove(ne)
                if prev is None:
                    body.insert(0, ne)
                else:
                    prev.addnext(ne)
                prev = ne

        _, save_err = _save_docx_atomic(doc, target)
        if save_err:
            return None, save_err
        return {
            "inserted_at": index,
            "added_blocks": len(blocks),
        }, None

    async with docx_write_lock(target):
        info, err = await asyncio.to_thread(_do)
    if err:
        return _err(err)
    return json.dumps({
        "ok": True, "action": "insert_block", "format": "docx",
        "path": rel_path,
        **info,
    }, ensure_ascii=False)



# ── Opt 4+5: LaTeX 公式渲染 + 富文本自动展开 ──────────────────────
# 2026-05-18 P194/P195: 渲染参数调优
_LATEX_PNG_DPI = 200                  # 150 → 200, 1" 宽公式从 ~150px 提升到 200px
# P195: inline 用 13pt (与 Word 12pt 正文视觉接近, 不撑高行高); display 用 16pt
_LATEX_PNG_FONTSIZE_INLINE = 13       # 行内公式: 字号略大于 12pt 正文, 视觉协调
_LATEX_PNG_FONTSIZE_DISPLAY = 16      # 居中独立段落: 大些更清晰
# 兼容旧 caller (没传 fontsize 参数): 用 inline 默认值
_LATEX_PNG_FONTSIZE = _LATEX_PNG_FONTSIZE_INLINE
_LATEX_PNG_PAD = 0.02                 # 之前 0.05" + 内层 0.15" = 0.20"; 现在 0.02"

# 匹配 display math $$...$$ 和 inline math $...$
_LATEX_DISPLAY_RE = re.compile(r"\$\$(.+?)\$\$", re.DOTALL)
_LATEX_INLINE_RE = re.compile(r"\$(.+?)\$")
# 匹配 markdown 图片语法 ![alt](path)
_MD_IMAGE_RE = re.compile(r"!\[.*?\]\(([^)\s]+)\)")


_MATRIX_ENV_RE = re.compile(
    r"\\begin\{(vmatrix|bmatrix|pmatrix|Bmatrix|Vmatrix|matrix|cases|array)\}"
    r"(.*?)\\end\{\1\}", re.DOTALL
)


# LaTeX/数学公式 → OMML 转换族已抽离到 office_latex.py(2026-05-20 重构);re-export 兼容。
from app.llm.tools.office_latex import (  # noqa: E402,F401
    _latex_contains_cjk,
    _strip_cjk_text_commands,
    _rewrite_choose_to_binom,
    _is_broken_after_strip,
    _looks_like_bare_math_identifier,
    _read_brace_group,
    _latex_text_to_omml_runs,
    _classify_latex_complexity,
)




def _text_to_safe_formula_image(text: str, workspace_dir: str, stem: str,
                                fontsize: int = _LATEX_PNG_FONTSIZE_DISPLAY) -> str | None:
    plain = _latex_plain_fallback(text)
    if not plain:
        return None
    out_name = f"_latex_text_{stem}_{hashlib.md5(plain.encode('utf-8')).hexdigest()[:8]}v1.png"
    out_path = os.path.join(workspace_dir, out_name)
    if os.path.isfile(out_path) and os.path.getsize(out_path) > 200:
        return out_name
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        from matplotlib import font_manager
    except ImportError:
        return None
    try:
        font_prop = None
        for font_name in ("Microsoft YaHei", "SimHei", "Noto Sans CJK SC", "Arial Unicode MS"):
            try:
                font_path = font_manager.findfont(font_name, fallback_to_default=False)
            except Exception:
                font_path = ""
            if font_path and os.path.isfile(font_path):
                font_prop = font_manager.FontProperties(fname=font_path)
                break
        fig, ax = plt.subplots(figsize=(0.01, 0.01))
        ax.text(0.5, 0.5, plain, fontsize=fontsize, fontproperties=font_prop,
                ha="center", va="center", transform=ax.transAxes)
        ax.axis("off")
        fig.patch.set_alpha(0.0)
        fig.savefig(out_path, dpi=_LATEX_PNG_DPI,
                    bbox_inches="tight", pad_inches=_LATEX_PNG_PAD,
                    transparent=True, edgecolor="none")
        plt.close(fig)
        return out_name if os.path.isfile(out_path) and os.path.getsize(out_path) > 200 else None
    except Exception:
        try:
            plt.close(fig)
        except Exception:
            pass
        return None


def _render_matrix_env_to_png(equation: str, workspace_dir: str, stem: str,
                              fontsize: int = 16) -> str | None:
    """渲染 LaTeX 矩阵/cases 环境为 PNG (mathtext 不支持这些).

    2026-05-18 P204: 实测灾难 (trace 0518 v5 docx): Q10 `\\begin{vmatrix}` 4×4 行列式
    + Q8 `\\begin{cases}` 分段函数 + Q16 矩阵 `A=\\begin{pmatrix}` 全部触发
    `[⚠️ 公式渲染失败 ...]` 占位符, helper 看到后整个 docx **重写一遍**
    (22:26:37 单次 37 items write) 浪费 ~90s + 上下文.

    支持的环境: vmatrix (行列式), bmatrix, pmatrix, Bmatrix, Vmatrix, matrix,
                cases (分段函数), array.

    实现: 解析 \\begin{X}...\\end{X} 内的行/列 (\\\\ 分行, & 分列), 用 matplotlib
    自动测量每个 cell 的真实宽度建自适应列宽, 然后用 ax.plot 画 delimiter (vmatrix 直线
    bmatrix 直角括号 pmatrix 弧线 cases Bezier 大括号), 单元格用 ax.text 各自渲染.
    """
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import numpy as np
        from matplotlib.patches import PathPatch
        from matplotlib.path import Path as MplPath
    except ImportError:
        return None

    out_name = f"_latex_{stem}.png"
    out_path = os.path.join(workspace_dir, out_name)
    if os.path.isfile(out_path) and os.path.getsize(out_path) > 200:
        return out_name  # cache hit

    m = _MATRIX_ENV_RE.search(equation)
    if not m:
        return None
    env, body = m.group(1), m.group(2).strip()

    # array 可能有列对齐说明 {lcr} — 跳过
    if env == "array":
        body = re.sub(r"^\s*\{[^}]+\}", "", body).strip()

    # 切行 & 切列
    rows_raw = [r.strip() for r in re.split(r"\\\\", body) if r.strip()]
    if not rows_raw:
        return None
    matrix = [[c.strip() for c in re.split(r"&", row)] for row in rows_raw]
    n_rows = len(matrix)
    n_cols = max(len(r) for r in matrix) if matrix else 0
    if n_cols == 0:
        return None

    # 用 matplotlib 实测每 cell 宽度 — 共享一个 fig 测全部 (避免 N×M 次 fig 创建)
    try:
        mfig, max_ = plt.subplots()
        renderer = mfig.canvas.get_renderer()
        def _measure(s: str) -> float:
            if not s:
                return 0.4
            try:
                t = max_.text(0, 0, f"${s}$", fontsize=fontsize)
                bbox = t.get_window_extent(renderer=renderer)
                t.remove()
                return bbox.width / 72.0 + 0.15
            except Exception:
                # 字符数估计回退
                visible = re.sub(r"\\[a-zA-Z]+|\{|\}", "X", s)
                return max(0.4, len(visible) * (fontsize / 72.0) * 0.5)
        col_widths = []
        for j in range(n_cols):
            max_w = 0.4
            for row in matrix:
                if j < len(row):
                    max_w = max(max_w, _measure(row[j]))
            col_widths.append(max_w)
        plt.close(mfig)
    except Exception:
        # 测量整体失败 — 退化到字符数估计
        col_widths = []
        for j in range(n_cols):
            max_w = 0.4
            for row in matrix:
                if j < len(row):
                    visible = re.sub(r"\\[a-zA-Z]+|\{|\}", "X", row[j])
                    max_w = max(max_w, max(0.4, len(visible) * (fontsize / 72.0) * 0.5))
            col_widths.append(max_w)

    cell_h = fontsize / 72.0 * 1.7  # 行距 (足够给 \frac 留位置)
    inner_w = sum(col_widths)
    inner_h = n_rows * cell_h
    delim_w = 0.20
    fig_w = inner_w + 2 * delim_w
    fig_h = inner_h + 0.10

    try:
        fig, ax = plt.subplots(figsize=(fig_w, fig_h))
        ax.set_xlim(0, fig_w)
        ax.set_ylim(0, fig_h)
        ax.set_aspect("equal", adjustable="box")
        ax.axis("off")

        x0 = delim_w
        for i, row in enumerate(matrix):
            full_row = row + [""] * (n_cols - len(row))
            x_cursor = x0
            for j, cell in enumerate(full_row):
                cx = x_cursor + col_widths[j] / 2
                cy = fig_h - 0.05 - (i + 0.5) * cell_h
                if cell.strip():
                    try:
                        ax.text(cx, cy, f"${cell}$", fontsize=fontsize,
                                ha="center", va="center")
                    except Exception:
                        ax.text(cx, cy, cell, fontsize=fontsize,
                                ha="center", va="center")
                x_cursor += col_widths[j]

        top, bottom = fig_h - 0.03, 0.03
        # Delimiter 绘制
        if env in ("vmatrix", "Vmatrix"):
            # 行列式: 双竖线 (Vmatrix 是双线, 简化为单线)
            ax.plot([x0 - 0.05] * 2, [bottom, top], "k-", lw=1.2)
            ax.plot([x0 + inner_w + 0.05] * 2, [bottom, top], "k-", lw=1.2)
        elif env in ("bmatrix", "Bmatrix"):
            bx = 0.08
            ax.plot([x0 - 0.08, x0 - 0.08 + bx, x0 - 0.08 + bx, x0 - 0.08],
                    [bottom, bottom, top, top], "k-", lw=1.2)
            ax.plot([x0 + inner_w + 0.08, x0 + inner_w + 0.08 - bx,
                     x0 + inner_w + 0.08 - bx, x0 + inner_w + 0.08],
                    [bottom, bottom, top, top], "k-", lw=1.2)
        elif env == "pmatrix":
            tt = np.linspace(0, 1, 30)
            ys = bottom + (top - bottom) * tt
            bulge = 0.08
            xs_left = (x0 - 0.05) - bulge * (1 - 4 * (tt - 0.5) ** 2)
            ax.plot(xs_left, ys, "k-", lw=1.2)
            xs_right = (x0 + inner_w + 0.05) + bulge * (1 - 4 * (tt - 0.5) ** 2)
            ax.plot(xs_right, ys, "k-", lw=1.2)
        elif env == "cases":
            # 左大括号 — Bezier
            h = top - bottom
            mid = (top + bottom) / 2
            bx = x0 - 0.05
            bulge = 0.10
            verts = [
                (bx, top),
                (bx - bulge, top - h / 4), (bx - bulge, mid + h / 8),
                (bx - bulge * 1.5, mid),
                (bx - bulge, mid - h / 8), (bx - bulge, bottom + h / 4),
                (bx, bottom),
            ]
            codes = [MplPath.MOVETO,
                     MplPath.CURVE4, MplPath.CURVE4, MplPath.CURVE4,
                     MplPath.CURVE4, MplPath.CURVE4, MplPath.CURVE4]
            p = PathPatch(MplPath(verts, codes), fc="none", ec="black", lw=1.2)
            ax.add_patch(p)
        # matrix / array: 无 delimiter

        fig.patch.set_alpha(0.0)
        fig.savefig(out_path, dpi=_LATEX_PNG_DPI,
                    bbox_inches="tight", pad_inches=_LATEX_PNG_PAD,
                    transparent=True, edgecolor="none")
        plt.close(fig)
        return out_name
    except Exception:
        try:
            plt.close(fig)
        except Exception:
            pass
        return None


def _render_latex_to_png(equation: str, workspace_dir: str, stem: str,
                         fontsize: int = _LATEX_PNG_FONTSIZE) -> str | None:
    """用 matplotlib mathtext 渲染 LaTeX 公式为 PNG,返回工作区相对路径。

    不依赖系统 LaTeX 安装,只用 matplotlib 内置数学渲染器(mathtext)。
    支持常见 LaTeX 数学语法:\\frac, \\sqrt, \\sum, \\int, \\alpha, \\beta 等。
    不支持的宏(\\text, \\begin{aligned} 等)会静默回退到纯文本渲染。

    2026-05-18 P186: stem 已含 content hash (P179), 同公式落同文件。这里加文件存在
    检查, 二次调用直接复用 ~250ms/张, 22 题 docx 写完省 3-5 秒。

    2026-05-18 P194: 视觉优化 (实测 trace 0518 19:29 docx 33 张公式 PNG, 全部是白底
    RGBA 255/255/255/255, 1.05" 宽 0.6" 高, font 14pt @ 150dpi):
      - **透明背景** (`transparent=True`): 公式行内嵌入与正文文字流自然衔接, 不再有
        刺眼白矩形挤断中文段落。文件名 corner_pixel 实测: (255,255,255,255) → (255,255,255,0).
      - **DPI 150 → 200**: 1.05" 宽公式在 150dpi 下 157px 显得模糊, 200dpi 提升到
        ~210px 边缘锐利无白晕 (transparent BG 也消除了原 antialias 白边).
      - **font 14 → 16pt**: 与 Word 默认正文 12pt × 1.3 行距 ≈ 15.6pt 匹配, 行内公式
        不再比周围中文小一圈。
      - **padding 0.05" + 内 0.15" → 0.02"**: 之前的 fig2 重渲染留了双层 padding,
        渲染时间 126ms → 31ms (4× 提速, 22 题省 ~2 秒).
      - **单次渲染**: 原代码先建 fig1 测 bbox 再建 fig2 输出, bbox_inches='tight'
        自带尺寸计算, 没必要双层 figure.
    """
    out_name = f"_latex_{stem}.png"
    out_path = os.path.join(workspace_dir, out_name)
    # 2026-05-18 P186: cache hit
    if os.path.isfile(out_path) and os.path.getsize(out_path) > 200:
        return out_name

    # 2026-05-18 P204: 矩阵 / cases / array 环境走专用渲染器 (mathtext 不支持)
    if _MATRIX_ENV_RE.search(equation):
        result = _render_matrix_env_to_png(equation, workspace_dir, stem, fontsize)
        if result:
            return result
        # 若专用渲染也失败, 继续往下走 mathtext fallback (大概率也失败,
        # 但会触发统一的失败路径生成清晰 placeholder)

    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except ImportError:
        return None  # matplotlib 未安装,静默回退

    try:
        # P194: 单次渲染 — bbox_inches='tight' 自动算尺寸, 不需要先 fig1 测 bbox
        fig, ax = plt.subplots(figsize=(0.01, 0.01))
        ax.text(0.5, 0.5, f"${equation}$",
                fontsize=fontsize,
                ha="center", va="center",
                transform=ax.transAxes)
        ax.axis("off")
        fig.patch.set_alpha(0.0)  # transparent figure background
        fig.savefig(out_path, dpi=_LATEX_PNG_DPI,
                    bbox_inches="tight", pad_inches=_LATEX_PNG_PAD,
                    transparent=True, edgecolor="none")
        plt.close(fig)
        return out_name
    except Exception:
        return None




# LaTeX 文本渲染/预处理族已抽离到 office_render.py(2026-05-20 重构);re-export 兼容。
from app.llm.tools.office_render import (  # noqa: E402,F401
    _GREEK_LETTERS,
    _NARY_OPS,
    _SIMPLE_OPS,
    _normalize_latex_for_render,
    _preprocess_simple_latex_macros,
    _apply_latex_scripts_to_text,
    _read_sub_or_sup_arg,
    _read_subsup,
    _latex_plain_fallback,
    _preprocess_simple_latex_macros_omml_safe,
    _try_unicode_render,
    _latex_to_omml_simple,
    _render_inline_latex_to_text_only,
)




def _render_text_formula_to_png(text: str, workspace_dir: str, stem: str,
                                fontsize: int = _LATEX_PNG_FONTSIZE) -> str | None:
    # 2026-05-18 P186: cache hit (stem 含 content hash from P179)
    out_name = f"_latex_text_{stem}.png"
    out_path = os.path.join(workspace_dir, out_name)
    if os.path.isfile(out_path) and os.path.getsize(out_path) > 200:
        return out_name

    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except ImportError:
        return None

    plain = _latex_plain_fallback(text)
    if not plain:
        return None
    # 2026-05-18 P178: 不要把含 \\begin{vmatrix}/\\displaystyle 等剥离残骸的 plain 文本
    # 假装渲染成"公式 PNG" — 这会让用户看到 "beginvmatrix 0 & a & b..." 这种
    # 灾难输出而不报错。直接返回 None, caller 改走文字 placeholder。
    if _is_broken_after_strip(plain, text):
        return None
    try:
        # 2026-05-18 P194: 同 _render_latex_to_png — 单次渲染 + 透明背景 + 优化 DPI/字号
        fig, ax = plt.subplots(figsize=(0.01, 0.01))
        ax.text(0.5, 0.5, plain,
                fontsize=fontsize,
                ha="center", va="center",
                transform=ax.transAxes)
        ax.axis("off")
        fig.patch.set_alpha(0.0)
        fig.savefig(out_path, dpi=_LATEX_PNG_DPI,
                    bbox_inches="tight", pad_inches=_LATEX_PNG_PAD,
                    transparent=True, edgecolor="none")
        plt.close(fig)
        return out_name
    except Exception:
        return None


def _render_formula_to_png(equation: str, workspace_dir: str, stem: str,
                            is_display: bool = False) -> str | None:
    # 2026-05-18 P179: 用 equation 内容 hash 作为 stem 后缀, 避免同名 PNG 文件
    # 在多次 office API 调用间被覆盖, 然后被 python-docx 按路径 dedup 误用同一 rId。
    # 实测病例 (trace 0518 06:46): 22 道题里 _latex_text_i0_fallback.png 被 3 道
    # 不同题轮流覆写, python-docx 看到"同名 path 出现 3 次"就发同一 rId, 导致 Q1/Q2/Q3
    # 都嵌入了最后一次写入的图片内容。
    #
    # 2026-05-18 P194: 加 render-version 标记 (`v2`) — 当渲染参数 (DPI/font/transparent)
    # 变化时, 老 cache 文件不会被复用 (P186 cache 会因为文件名不同而 miss → 重新渲染),
    # 避免旧 docx 残留的白底 150DPI 文件冒充新参数生成的透明 200DPI 公式。
    #
    # 2026-05-18 P195: 区分 inline / display 字号. inline 用 13pt (与 Word 12pt 正文
    # 视觉接近, 行间公式不撑高行高过多); display 用 16pt (居中独立段落, 大些更清晰).
    # 之前一律 fontsize=16, inline 公式比正文大一圈让段落行高被撑到 0.46", 而正文 12pt
    # 行高才 0.19" — 公式行明显比纯文本行高 ~2.4×。
    content_hash = hashlib.md5(equation.encode("utf-8")).hexdigest()[:8]
    # v3 标记 (P195 字号分支再次失效旧 cache)
    suffix = "v3d" if is_display else "v3i"
    hashed_stem = f"{stem}_{content_hash}{suffix}"

    fontsize = _LATEX_PNG_FONTSIZE_DISPLAY if is_display else _LATEX_PNG_FONTSIZE_INLINE
    png = _render_latex_to_png(equation, workspace_dir, hashed_stem, fontsize=fontsize)
    if png:
        return png
    normalized = _normalize_latex_for_render(equation)
    if normalized and normalized != equation:
        png = _render_latex_to_png(normalized, workspace_dir, f"{hashed_stem}_retry", fontsize=fontsize)
        if png:
            return png
    return _render_text_formula_to_png(equation, workspace_dir, f"{hashed_stem}_fallback", fontsize=fontsize)


def _png_natural_width_inches(workspace_dir: str, png_name: str,
                                source_dpi: int = _LATEX_PNG_DPI) -> float | None:
    """读 PNG 实际像素, 按生成 DPI 反算自然显示宽度 (inch)。

    用途: 公式 PNG 渲染时按 16pt 字号生成, 实际尺寸约 0.2-2 inch 宽。
    但嵌入 docx 时若不指定 width 会走默认 6 inch → 拉伸 3-5x 模糊。
    用此函数取实际尺寸, 让公式按"原比例"嵌入。

    2026-05-18 P194: min 0.3 → 0.15 inch. 之前的 0.3 floor 是为了"避免单字符太小看不见",
    但 P194 把 font 14→16 + DPI 150→200 后, 短公式 (`\\partial f / \\partial x`) 自然宽
    ~0.22", 强行放大到 0.3" 会引入 36% 模糊拉伸. 0.15 floor 保留极端 case 保护 (单符号)
    同时让短公式按原比例显示, 视觉与周围文字大小匹配。

    返回 None 表示无法读取 (PIL 缺失/文件不存在), 调用方应 fallback 到 inline 默认。
    """
    try:
        from PIL import Image
        fp = os.path.join(workspace_dir, png_name)
        if not os.path.exists(fp):
            return None
        with Image.open(fp) as img:
            px_w = img.size[0]
        # PIL 默认 96 DPI, 但我们的 PNG 是 200 DPI 渲染的, 用 source_dpi 反算
        natural = px_w / float(source_dpi)
        # P194: 约束 0.15" - 6.5" (页面 90%)
        return max(0.15, min(6.5, natural))
    except (ImportError, OSError, AttributeError):
        return None





def _expand_rich_text_blocks(workspace_dir: str, blocks: list) -> list:
    """预处理 blocks:展开 LaTeX 公式和 markdown 图片语法为结构化 image blocks。

    - `![alt](path)` → {type: "image", path: "path", ...}
    - `$$...$$`     → {type: "image", path: "_latex_N.png", caption: "display formula"}
    - `$...$`       → {type: "image", path: "_latex_N.png", width_inches: 2.0}
    """
    _latex_counter = 0
    expanded: list = []

    for b in blocks:
        if not isinstance(b, dict):
            expanded.append(b)
            continue

        btype = str(b.get("type", "")).strip().lower()

        if btype == "paragraph":
            text = str(b.get("text", ""))
            _sub_blocks = _parse_rich_paragraph(workspace_dir, text, b, _latex_counter)
            # count how many latex blocks generated
            for sb in _sub_blocks:
                if isinstance(sb, dict) and sb.get("_latex"):
                    _latex_counter += 1
            expanded.extend(_sub_blocks)

        elif btype == "heading":
            # 2026-05-18 P184: heading 也允许 inline `$...$` (走 Unicode-only 路径).
            # 之前 heading 不被 _expand_rich_text_blocks 处理, helper 写
            # `{type:'heading', text:'求 $\\int f$'}` 会原样写入 Word 标题。
            # heading 不适合嵌入 image (会让标题挤变形), 所以 PNG fallback 时
            # 退回 _latex_plain_fallback 文字。
            text = str(b.get("text", ""))
            if "$" in text:
                expanded.append({**b, "text": _render_inline_latex_to_text_only(text)})
            else:
                expanded.append(b)

        elif btype == "list":
            # 2026-05-18 P184: list items 用 Unicode-only LaTeX 渲染. 之前的实现把
            # PNG fallback 退化成 `[见图: path]` 占位, 现在退化为 plain_fallback 文字,
            # 或被 P178 detector 标 `[⚠️ ...]`。list 视觉上不适合嵌图。
            items = b.get("items") or []
            if isinstance(items, list):
                new_items = []
                for item in items:
                    item_text = str(item)
                    if "$" in item_text:
                        new_items.append(_render_inline_latex_to_text_only(item_text))
                    else:
                        new_items.append(item_text)
                expanded.append({**b, "items": new_items})
            else:
                expanded.append(b)

        else:
            expanded.append(b)

    return expanded


def _parse_rich_paragraph(workspace_dir: str, text: str, orig_block: dict, latex_seed: int) -> list:
    """将段落文本解析为段落块列表。

    2026-05-17 P163: 内联公式 (单 $...$) 和文字应保留在**同一段落**内, 表示为
    `{type:'paragraph', runs:[{kind:'text',...}, {kind:'image',...}, ...]}` 结构。
    早期实现把每个 $...$ 拆成独立段落 → 1 段变 5+ 段, docx 95% 段落是 <20 字碎片。

    处理:
      - 显示公式 $$...$$ → 独立段落 (居中, 较大)
      - 内联公式 $...$    → 在原段落中作 inline image run
      - markdown 图片     → 独立段落 image
    """
    # 先按 display math $$...$$ 切段, 每段是独立 paragraph
    segments = _LATEX_DISPLAY_RE.split(text)
    results: list = []
    _ltx_idx = latex_seed

    def _process_inline_segment(seg_text: str) -> dict | None:
        """处理一段没有 $$ 的文本: 含 inline $...$ 时返回 runs-paragraph block; 否则返回普通 paragraph。"""
        nonlocal _ltx_idx
        if not seg_text.strip():
            return None
        # 拆分 inline math $...$ 和 markdown image ![](path)
        combined_re = re.compile(r"(\$[^$\n]+\$)|(!\[.*?\]\([^)]+\))")
        parts = combined_re.split(seg_text)
        runs: list = []
        has_inline_element = False
        for part in parts:
            if part is None or part == "":
                continue
            if _LATEX_INLINE_RE.match(part):
                eq = part[1:-1].strip()
                # Unicode fast path
                uni = _try_unicode_render(eq)
                if uni is not None:
                    runs.append({"kind": "text", "text": uni})
                    continue
                # 2026-05-20 P205: 公式里混中文时, 先把中文说明移出 math,
                # 公式本体优先 Word 原生 OMML；仍不支持时才走图片兜底。
                if _latex_contains_cjk(eq):
                    math_eq, notes = _strip_cjk_text_commands(eq)
                    math_eq = math_eq or eq
                    omml_xml = _latex_to_omml_simple(math_eq)
                    if omml_xml:
                        runs.append({"kind": "equation", "latex": math_eq})
                        has_inline_element = True
                    else:
                        png = _render_formula_to_png(math_eq, workspace_dir, f"i{_ltx_idx}", is_display=False)
                        _ltx_idx += 1
                        if png:
                            natural_w = _png_natural_width_inches(workspace_dir, png) or 0.8
                            runs.append({"kind": "image", "path": png,
                                         "width_inches": natural_w, "_latex": True})
                            has_inline_element = True
                        else:
                            plain_png = _text_to_safe_formula_image(eq, workspace_dir, f"i{_ltx_idx}_cjk",
                                                                    fontsize=_LATEX_PNG_FONTSIZE_INLINE)
                            if plain_png:
                                natural_w = _png_natural_width_inches(workspace_dir, plain_png) or 0.8
                                runs.append({"kind": "image", "path": plain_png,
                                             "width_inches": natural_w, "_latex": True})
                                has_inline_element = True
                            else:
                                runs.append({"kind": "text", "text": _latex_plain_fallback(eq)})
                    if notes:
                        runs.append({"kind": "text", "text": "（" + "；".join(notes) + "）"})
                    continue
                # 2026-05-20 P206: 优先 Word 原生 OMML, 尽量避免图片公式。
                omml_xml = _latex_to_omml_simple(eq)
                if omml_xml:
                    runs.append({"kind": "equation", "latex": eq})
                    has_inline_element = True
                    continue
                # PNG fallback
                png = _render_formula_to_png(eq, workspace_dir, f"i{_ltx_idx}", is_display=False)
                _ltx_idx += 1
                if png:
                    natural_w = _png_natural_width_inches(workspace_dir, png) or 0.8
                    runs.append({"kind": "image", "path": png,
                                 "width_inches": natural_w, "_latex": True})
                    has_inline_element = True
                else:
                    # 2026-05-18 P178: 渲染失败 — 不要静默把"被剥离的命令残骸"当文本嵌入。
                    # 用 fallback_text 时也要透出原 LaTeX, 让用户清楚是渲染失败。
                    fallback_plain = _latex_plain_fallback(eq)
                    if _is_broken_after_strip(fallback_plain, eq):
                        # 明确标记: 公式渲染失败, 给出原 LaTeX 让人识别
                        runs.append({"kind": "text",
                                     "text": f"[⚠️ 公式渲染失败 ${eq}$]"})
                    else:
                        runs.append({"kind": "text", "text": fallback_plain})
            elif _MD_IMAGE_RE.match(part):
                m = _MD_IMAGE_RE.match(part)
                path = m.group(1).strip()
                runs.append({"kind": "image", "path": path})
                has_inline_element = True
            else:
                # plain text — keep as-is (paragraph break preserved by segment split)
                runs.append({"kind": "text", "text": part})
        # 合并相邻 text runs (减少不必要的 run 数量)
        merged_runs = []
        for r in runs:
            if (merged_runs and merged_runs[-1].get("kind") == "text"
                    and r.get("kind") == "text"):
                merged_runs[-1]["text"] = merged_runs[-1]["text"] + r["text"]
            else:
                merged_runs.append(r)
        if not merged_runs:
            return None
        # 没有 image (纯文字), 退回成普通 paragraph
        if not has_inline_element:
            joined = "".join(r["text"] for r in merged_runs if r.get("kind") == "text")
            if not joined.strip():
                return None
            return {**orig_block, "type": "paragraph", "text": joined}
        # 含 inline image — 输出 runs 结构
        out = {**orig_block, "type": "paragraph", "runs": merged_runs}
        out.pop("text", None)
        return out

    for i, segment in enumerate(segments):
        if i % 2 == 0:
            # plain text segment (between $$...$$ or at boundaries)
            block = _process_inline_segment(segment)
            if block is not None:
                results.append(block)
        else:
            # display math segment
            display_latex = segment.strip()
            if not display_latex:
                continue
            uni = _try_unicode_render(display_latex)
            if uni is not None:
                results.append({**orig_block, "type": "paragraph", "text": uni})
                continue
            # 2026-05-20 P205: 公式里混中文时, 把中文说明移到公式外,
            # 公式本体优先 Word 原生 OMML；仍不支持时才走图片兜底。
            if _latex_contains_cjk(display_latex):
                math_latex, notes = _strip_cjk_text_commands(display_latex)
                math_latex = math_latex or display_latex
                if _latex_to_omml_simple(math_latex):
                    results.append({**orig_block, "type": "equation", "latex": math_latex, "display": True})
                    if notes:
                        results.append({**orig_block, "type": "paragraph", "text": "；".join(notes)})
                else:
                    png = _render_formula_to_png(math_latex, workspace_dir, f"d{_ltx_idx}", is_display=True)
                    _ltx_idx += 1
                    if not png:
                        png = _text_to_safe_formula_image(display_latex, workspace_dir, f"d{_ltx_idx}_cjk")
                    if png:
                        natural_w = _png_natural_width_inches(workspace_dir, png) or 3.0
                        results.append({
                            **orig_block,
                            "type": "image", "path": png,
                            "width_inches": natural_w,
                            "caption": "",
                            "_latex": True,
                        })
                        if notes:
                            results.append({**orig_block, "type": "paragraph", "text": "；".join(notes)})
                    else:
                        text = _latex_plain_fallback(display_latex)
                        if notes:
                            text = f"{text}（{'；'.join(notes)}）"
                        results.append({**orig_block, "type": "paragraph", "text": text})
                continue
            if _latex_to_omml_simple(display_latex):
                results.append({**orig_block, "type": "equation", "latex": display_latex, "display": True})
                continue
            png = _render_formula_to_png(display_latex, workspace_dir, f"d{_ltx_idx}", is_display=True)
            _ltx_idx += 1
            if png:
                natural_w = _png_natural_width_inches(workspace_dir, png) or 3.0
                # 2026-05-18 P187: caption 用 _latex_plain_fallback 输出, 若残骸
                # (`begin...end...`) 不要作 caption — 直接省略 caption.
                caption_plain = _latex_plain_fallback(display_latex)
                if _is_broken_after_strip(caption_plain, display_latex):
                    caption_plain = ""
                img_block: dict = {
                    "type": "image", "path": png,
                    "width_inches": natural_w,
                    "caption": caption_plain,
                    "_latex": True,
                }
                results.append(img_block)
            else:
                # 2026-05-18 P178: 同样, 显示公式渲染失败时给清晰标记
                fallback_plain = _latex_plain_fallback(display_latex)
                if _is_broken_after_strip(fallback_plain, display_latex):
                    results.append({**orig_block, "type": "paragraph",
                                    "text": f"[⚠️ 公式渲染失败 $${display_latex}$$]"})
                else:
                    results.append({**orig_block, "type": "paragraph",
                                    "text": fallback_plain})

    return results





















def _apply_blocks_docx(doc, workspace_dir: str, blocks: list) -> str | None:
    from docx.shared import Pt

    # ── Opt 4+5: 预处理 LaTeX 公式 & markdown 图片 → 结构化 image blocks ──
    blocks = _expand_rich_text_blocks(workspace_dir, blocks)
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    align_map = {
        "left": WD_ALIGN_PARAGRAPH.LEFT,
        "center": WD_ALIGN_PARAGRAPH.CENTER,
        "right": WD_ALIGN_PARAGRAPH.RIGHT,
        "justify": WD_ALIGN_PARAGRAPH.JUSTIFY,
    }

    for i, b in enumerate(blocks):
        if not isinstance(b, dict):
            return f"block[{i}] must be an object"
        btype = str(b.get("type", "")).strip().lower()
        if not btype:
            if "rows" in b:
                btype = "table"
            elif "items" in b:
                btype = "list"
            elif "path" in b:
                btype = "image"
            elif "latex" in b:
                btype = "equation"
            elif "text" in b or "runs" in b:
                btype = "paragraph"
        if btype in {"numbered_list", "ordered_list"}:
            btype = "list"
            b.setdefault("ordered", True)
        elif btype in {"bullet_list", "bulleted_list", "unordered_list"}:
            btype = "list"
        elif btype == "subtitle":
            btype = "paragraph"
            b["_office_normalized_type_from"] = "subtitle"
        else:
            heading_match = re.fullmatch(r"h(?:eading)?([1-6])", btype)
            if heading_match:
                btype = "heading"
                b.setdefault("level", int(heading_match.group(1)))

        if btype == "heading":
            level = max(1, min(6, int(b.get("level", 1))))
            doc.add_heading(_clean_docx_text(b.get("text", "")), level=level)

        elif btype == "paragraph":
            p = doc.add_paragraph()
            # 2026-05-17 P163: 含 inline runs (文字 + 公式图片混合) 走 runs 路径,
            # 保持同一段落不被拆分
            runs_data = b.get("runs")
            if isinstance(runs_data, list) and runs_data:
                for r_item in runs_data:
                    if not isinstance(r_item, dict):
                        continue
                    kind = r_item.get("kind", "text")
                    if kind == "text":
                        # 2026-05-18 P169: text run 也走 break-aware 路径
                        _add_run_with_breaks(
                            p, _clean_docx_text(r_item.get("text", "")),
                            bold=bool(b.get("bold")),
                            italic=bool(b.get("italic")),
                            font_size_pt=b.get("font_size_pt"),
                        )
                    elif kind == "equation":
                        latex = str(r_item.get("latex") or r_item.get("text") or "").strip()
                        omml_xml = _latex_to_omml_simple(latex) if latex else None
                        if omml_xml:
                            try:
                                from lxml import etree
                                p._p.append(etree.fromstring(omml_xml))
                            except Exception:
                                p.add_run(_latex_plain_fallback(latex))
                        elif latex:
                            p.add_run(_latex_plain_fallback(latex))
                    elif kind == "image":
                        # 在同一段落里 add picture - 作为 inline image run
                        img_path_rel = r_item.get("path", "")
                        if not img_path_rel:
                            continue
                        try:
                            img_abs = ws_tool._safe_resolve(workspace_dir, img_path_rel)
                        except ValueError:
                            continue
                        if not os.path.isfile(img_abs):
                            continue
                        try:
                            from docx.shared import Inches
                            w = r_item.get("width_inches")
                            run = p.add_run()
                            if isinstance(w, (int, float)) and w > 0:
                                run.add_picture(img_abs, width=Inches(w))
                            else:
                                run.add_picture(img_abs)
                        except Exception:
                            # 图片插入失败,降级为 placeholder 文本
                            p.add_run(f"[图: {os.path.basename(img_path_rel)}]")
            else:
                # 普通文字段落 (无 inline 公式/图)
                # 2026-05-18 P169: \n 转软换行, 不再被 Word 折叠成空白
                _add_run_with_breaks(
                    p, _clean_docx_text(b.get("text", "")),
                    bold=bool(b.get("bold")),
                    italic=bool(b.get("italic")),
                    font_size_pt=b.get("font_size_pt"),
                )
            align = str(b.get("align", "")).lower()
            if align in align_map:
                p.alignment = align_map[align]

        elif btype == "list":
            items = b.get("items") or []
            if not isinstance(items, list):
                return f"block[{i}].items must be array"
            style = "List Number" if b.get("ordered") else "List Bullet"
            for item in items:
                doc.add_paragraph(_clean_docx_text(item), style=style)

        elif btype == "table":
            err = _apply_table_docx(doc, workspace_dir, i, b)
            if err:
                return err

        elif btype == "image":
            err = _apply_image_docx(doc, workspace_dir, i, b)
            if err:
                return err

        elif btype == "equation":
            # 2026-05-11 P13d: Word 原生公式 (OMML) 块
            # 主进程 LLM 可显式发: {type:"equation", latex:"10^4", display:false}
            # display=true 时居中放大显示, false 时 inline 大小
            err = _apply_equation_docx(doc, workspace_dir, i, b)
            if err:
                return err

        elif btype == "page_break":
            doc.add_page_break()

        else:
            return (f"block[{i}].type unknown: {btype!r}; "
                    f"valid: heading/paragraph/list/table/image/equation/page_break. "
                    f"Use paragraph for plain prose, list for bullets or numbered items, "
                    f"and image only with an existing workspace path.\n"
                    f"普通正文用 paragraph，项目符号用 list，图片块必须提供已存在路径。")

    return None


def _apply_table_docx(doc, workspace_dir: str, idx: int, b: dict) -> str | None:
    """Render a table block. Each cell may contain `$...$` inline LaTeX.

    2026-05-18 P182: 之前 cell.text = raw_str 让 $\\frac{...}{...}$ 这种 LaTeX
    原样保留为文本(没渲染), 实测病例 (trace 1779100524438): 选择题选项表里
    `(A) $\\lambda=\\frac{1}{2}$` 等公式全部以原始 LaTeX 字符串显示, 用户看到
    一堆 $ 和反斜杠。
    现在: 把每个 cell 文本送进 _parse_rich_paragraph 渲染公式 (Unicode 优先,
    PNG 兜底), 再把结果 runs 应用到 cell 的第 1 个 paragraph。
    """
    rows = _normalize_docx_table_rows_shape(b)
    if not isinstance(rows, list) or not rows:
        return f"block[{idx}].rows must be a non-empty 2D array"
    bad_rows: list[dict] = []
    for r_idx, row in enumerate(rows):
        if not isinstance(row, list):
            bad_rows.append({"row_index": r_idx, "issue": "row_not_array", "value_type": type(row).__name__})
            continue
        non_empty_cells = []
        for c_idx, cell_val in enumerate(row):
            cell_text, _ = _normalize_docx_table_cell(cell_val)
            if str(cell_text or "").strip():
                non_empty_cells.append(c_idx)
        if not non_empty_cells:
            bad_rows.append({
                "row_index": r_idx,
                "issue": "row_has_no_non_empty_cells",
                "cell_count": len(row),
            })
    if bad_rows:
        return (
            f"block[{idx}].rows contains invalid table row(s): {json.dumps(bad_rows[:8], ensure_ascii=False)}. "
            "Each row must be an array and contain at least one non-empty cell after text/value/content normalization. "
            "Remove empty rows or represent explanatory prose as paragraph/list blocks; keep real table rows as a 2D array.\n"
            "表格行需为数组，且每行至少有一个归一化后的非空单元格；空行请删除，说明性文字用 paragraph/list。"
        )
    n_cols = max((len(r) if isinstance(r, list) else 0) for r in rows)
    if n_cols == 0:
        return (
            f"block[{idx}].rows: every row is empty; remove empty rows or use a paragraph/list block "
            f"when there is no real table data.\n表格必须有非空二维数据；没有表格数据时改用 paragraph/list。"
        )
    tbl = doc.add_table(rows=len(rows), cols=n_cols)
    try:
        tbl.style = b.get("style", "Light Grid Accent 1") or "Light Grid Accent 1"
    except KeyError:
        pass  # 样式名平台可能不一致,保持默认
    _configure_docx_table_geometry(tbl, b, n_cols)

    for r_idx, r in enumerate(rows):
        if not isinstance(r, list):
            continue
        for c_idx, cell_val in enumerate(r):
            if c_idx >= n_cols:
                break
            cell_text, cell_style = _normalize_docx_table_cell(cell_val)
            _render_cell_with_latex(tbl.cell(r_idx, c_idx),
                                    workspace_dir,
                                    _clean_docx_text(cell_text))
            if cell_style in {"strong", "bold"}:
                for p in tbl.cell(r_idx, c_idx).paragraphs:
                    for run in p.runs:
                        run.bold = True

    if b.get("header") and tbl.rows:
        for cell in tbl.rows[0].cells:
            for p in cell.paragraphs:
                for run in p.runs:
                    run.bold = True
    return None


def _normalize_docx_table_rows_shape(block: dict) -> list:
    """Normalize tolerated table row aliases into the canonical 2D array.

    The canonical model-visible schema is rows:[[...], ...]. Logs show models may
    naturally emit rows:[{"cells":[...]}]. That carries the same table data, so
    the tool accepts it and reports the normalization in the success payload.
    """
    rows = block.get("rows") or []
    if not isinstance(rows, list):
        return rows
    normalized: list = []
    changed = 0
    for row in rows:
        if isinstance(row, dict) and isinstance(row.get("cells"), list):
            normalized.append(row.get("cells") or [])
            changed += 1
        else:
            normalized.append(row)
    if changed:
        block["rows"] = normalized
        block["_office_normalized_table_row_objects"] = changed
    return normalized


def _configure_docx_table_geometry(tbl, block: dict, n_cols: int) -> None:
    """Apply stable table geometry so generated Word tables do not rely on autofit."""
    if n_cols <= 0:
        return
    try:
        from docx.shared import Twips
        from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
        tbl.autofit = False
        tbl.alignment = WD_TABLE_ALIGNMENT.LEFT
        margin = 120
        usable = _DOCX_USABLE_WIDTH_DXA - margin
        raw_widths = block.get("widths") or block.get("col_widths") or []
        widths: list[int] = []
        if isinstance(raw_widths, list) and len(raw_widths) == n_cols:
            for value in raw_widths:
                try:
                    widths.append(max(360, int(float(value))))
                except (TypeError, ValueError):
                    widths = []
                    break
        if not widths:
            widths = [max(720, usable // n_cols)] * n_cols
        total = sum(widths) or 1
        if total > usable:
            widths = [max(360, int(w * usable / total)) for w in widths]
        try:
            tbl._tbl.tblPr.tblW.type = "dxa"
            tbl._tbl.tblPr.tblW.w = str(sum(widths))
        except Exception:
            pass
        try:
            grid = tbl._tbl.tblGrid
            for idx, col in enumerate(grid.gridCol_lst):
                if idx < len(widths):
                    col.w = str(widths[idx])
        except Exception:
            pass
        for row in tbl.rows:
            for c_idx, cell in enumerate(row.cells):
                if c_idx < len(widths):
                    cell.width = Twips(widths[c_idx])
                try:
                    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
                except Exception:
                    pass
                _set_docx_cell_margins(cell, top=80, start=120, bottom=80, end=120)
    except Exception:
        return


def _set_docx_cell_margins(cell, *, top: int, start: int, bottom: int, end: int) -> None:
    try:
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn
        tc_pr = cell._tc.get_or_add_tcPr()
        tc_mar = tc_pr.first_child_found_in("w:tcMar")
        if tc_mar is None:
            tc_mar = OxmlElement("w:tcMar")
            tc_pr.append(tc_mar)
        for key, value in (("top", top), ("start", start), ("bottom", bottom), ("end", end)):
            node = tc_mar.find(qn(f"w:{key}"))
            if node is None:
                node = OxmlElement(f"w:{key}")
                tc_mar.append(node)
            node.set(qn("w:w"), str(value))
            node.set(qn("w:type"), "dxa")
    except Exception:
        return


def _normalize_docx_table_cell(value) -> tuple[str, str]:
    """Normalize structured table-cell values for DOCX output.

    Cells may be passed as plain scalars or as objects like
    {"text": "...", "style": "strong"}. The latter must become rich cell text,
    not a Python/JSON object literal in the document.

    表格单元格支持 text/style 对象；写入文档前转成文本和样式。
    """
    if isinstance(value, str):
        stripped = value.strip()
        if re.match(r"^\{['\"]?(?:text|value|content|style)['\"]?\s*:", stripped):
            try:
                parsed = ast.literal_eval(stripped)
            except (ValueError, SyntaxError):
                parsed = None
            if isinstance(parsed, dict):
                value = parsed
    if isinstance(value, dict):
        text = value.get("text")
        if text is None:
            text = value.get("value")
        if text is None:
            text = value.get("content")
        style = str(value.get("style") or "").strip().lower()
        return str(text or ""), style
    return str(value or ""), ""


def _render_cell_with_latex(cell, workspace_dir: str, text: str) -> None:
    """渲染单个 cell 内容, 处理 inline `$...$` 公式。

    2026-05-18 P182: 用 _parse_rich_paragraph 拿到 sub_blocks (paragraph + runs
    或 paragraph + text 或 image), 然后清空 cell 默认段落, 按 sub_blocks 重建。
    """
    # 快速路径: 没有 $ → 直接 text
    if "$" not in text:
        cell.text = text
        return

    # 用 _parse_rich_paragraph 渲染
    try:
        sub_blocks = _parse_rich_paragraph(
            workspace_dir, text,
            {"type": "paragraph", "text": text},
            latex_seed=0,
        )
    except Exception:
        # 兜底: 失败就当普通文本
        cell.text = text
        return

    if not sub_blocks:
        cell.text = text
        return

    # 清空默认 paragraph 后逐个写入
    # python-docx 创建 cell 时已有 1 空 paragraph, 用第一块填它, 其余 add
    cell.text = ""  # 重置
    first_p = cell.paragraphs[0] if cell.paragraphs else cell.add_paragraph()
    paragraphs_to_use = [first_p]
    for sb_idx, sb in enumerate(sub_blocks):
        if not isinstance(sb, dict):
            continue
        if sb.get("type") == "image":
            # cell 里 image — 用 add_paragraph 加一段独立段
            img_p = cell.add_paragraph() if sb_idx > 0 else first_p
            img_path_rel = sb.get("path", "")
            if img_path_rel:
                try:
                    img_abs = ws_tool._safe_resolve(workspace_dir, img_path_rel)
                    if os.path.isfile(img_abs):
                        from docx.shared import Inches
                        w = sb.get("width_inches")
                        run = img_p.add_run()
                        if isinstance(w, (int, float)) and w > 0:
                            run.add_picture(img_abs, width=Inches(w))
                        else:
                            run.add_picture(img_abs)
                except (ValueError, OSError, Exception):
                    img_p.add_run(f"[图: {os.path.basename(img_path_rel)}]")
        elif sb.get("type") == "paragraph":
            # 段落: 可能含 runs[] (inline 公式图混合) 或普通 text
            target_p = first_p if sb_idx == 0 else cell.add_paragraph()
            runs_data = sb.get("runs")
            if isinstance(runs_data, list) and runs_data:
                for r_item in runs_data:
                    if not isinstance(r_item, dict):
                        continue
                    kind = r_item.get("kind", "text")
                    if kind == "text":
                        target_p.add_run(r_item.get("text", ""))
                    elif kind == "equation":
                        latex = str(r_item.get("latex") or r_item.get("text") or "").strip()
                        omml_xml = _latex_to_omml_simple(latex) if latex else None
                        if omml_xml:
                            try:
                                from lxml import etree
                                target_p._p.append(etree.fromstring(omml_xml))
                            except Exception:
                                target_p.add_run(_latex_plain_fallback(latex))
                        elif latex:
                            target_p.add_run(_latex_plain_fallback(latex))
                    elif kind == "image":
                        img_rel = r_item.get("path", "")
                        if not img_rel:
                            continue
                        try:
                            img_abs = ws_tool._safe_resolve(workspace_dir, img_rel)
                            if os.path.isfile(img_abs):
                                from docx.shared import Inches
                                w = r_item.get("width_inches")
                                run = target_p.add_run()
                                if isinstance(w, (int, float)) and w > 0:
                                    run.add_picture(img_abs, width=Inches(w))
                                else:
                                    run.add_picture(img_abs)
                        except (ValueError, OSError, Exception):
                            target_p.add_run(f"[图: {os.path.basename(img_rel)}]")
            else:
                target_p.add_run(sb.get("text", ""))


def _resolve_main_named_to_sandbox(workspace_dir: str, img_path: str) -> str | None:
    """主区命名 → helper sandbox 路径反推 (P47 辅助)。

    helper sandbox 内文件命名约定:
      主区命名: `<task_id>_<sub_path>` (例 `gen_charts_chart1_overview.png`)
      sandbox: `_helpers_shared/<task_id>/<sub_path>`

    LLM 写 docx/pptx 时容易拿主区命名引用图,但 helper sandbox 内只有 sandbox 命名,
    导致 office 找不到 → STUCK. 此函数枚举 _helpers_shared/ 下的 task_id 目录,
    若 input 以某 task_id 开头, 构造 sandbox 路径并校验存在性。

    Returns: 真实可读的 sandbox 路径 (绝对/相对), 或 None.
    """
    if not img_path or not workspace_dir:
        return None
    helpers_dir = os.path.join(workspace_dir, "_helpers_shared")
    if not os.path.isdir(helpers_dir):
        return None
    try:
        # 优先匹配长 task_id (避免 "gen" 错配 "gen_charts")
        for task_id in sorted(os.listdir(helpers_dir), key=len, reverse=True):
            task_dir = os.path.join(helpers_dir, task_id)
            if not os.path.isdir(task_dir):
                continue
            prefix = task_id + "_"
            if img_path.startswith(prefix):
                sub_path = img_path[len(prefix):]
                candidate = os.path.join(task_dir, sub_path)
                if os.path.isfile(candidate):
                    return candidate
    except OSError:
        pass
    return None




def _caption_text_to_runs(workspace_dir: str, text: str, latex_seed: int) -> tuple[list[dict], int]:
    """Split figure captions so bare identifiers like f_c / B = 2R_B render as formula runs."""
    pattern = re.compile(r"([A-Za-z][A-Za-z0-9_]*(?:\s*=\s*[-+]?\d+(?:\.\d+)?(?:\s*[A-Za-z]+)?)?)")
    runs: list[dict] = []
    pos = 0
    for m in pattern.finditer(text or ""):
        token = m.group(1)
        if m.start() > 0 and text[m.start() - 1] in "^_":
            continue
        if m.end() < len(text) and text[m.end()] in "^_":
            continue
        if not _looks_like_bare_math_identifier(token):
            continue
        if m.start() > pos:
            runs.append({"kind": "text", "text": text[pos:m.start()]})
        omml_xml = _latex_to_omml_simple(token)
        if omml_xml:
            runs.append({"kind": "equation", "latex": token})
        else:
            uni = _try_unicode_render(token)
            runs.append({"kind": "text", "text": uni if uni is not None else token})
        latex_seed += 1
        pos = m.end()
    if pos < len(text or ""):
        runs.append({"kind": "text", "text": text[pos:]})
    if not runs:
        runs.append({"kind": "text", "text": text or ""})
    return runs, latex_seed


def _render_docx_caption_paragraph(doc, workspace_dir: str, caption: str, italic: bool = True,
                                   font_size_pt: int = 10):
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches, Pt

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cap = _clean_docx_text(caption)
    if not cap:
        return p
    sub_blocks = _parse_rich_paragraph(
        workspace_dir,
        cap,
        {"type": "paragraph", "text": cap, "align": "center", "font_size_pt": font_size_pt, "italic": italic},
        latex_seed=0,
    ) if "$" in cap else [{"type": "paragraph", "runs": _caption_text_to_runs(workspace_dir, cap, 0)[0]}]

    for sb_idx, sb in enumerate(sub_blocks or []):
        if sb_idx > 0:
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        if not isinstance(sb, dict):
            continue
        if sb.get("type") == "image":
            try:
                img_abs = ws_tool._safe_resolve(workspace_dir, sb.get("path", ""))
                if os.path.isfile(img_abs):
                    w = sb.get("width_inches")
                    run = p.add_run()
                    if isinstance(w, (int, float)) and w > 0:
                        run.add_picture(img_abs, width=Inches(w))
                    else:
                        run.add_picture(img_abs)
            except Exception:
                pass
            continue
        runs_data = sb.get("runs")
        if isinstance(runs_data, list) and runs_data:
            for r_item in runs_data:
                if not isinstance(r_item, dict):
                    continue
                if r_item.get("kind") == "text":
                    run = p.add_run(_clean_docx_text(r_item.get("text", "")))
                    run.italic = italic
                    run.font.size = Pt(font_size_pt)
                elif r_item.get("kind") == "equation":
                    latex = str(r_item.get("latex") or r_item.get("text") or "").strip()
                    omml_xml = _latex_to_omml_simple(latex) if latex else None
                    if omml_xml:
                        try:
                            from lxml import etree
                            p._p.append(etree.fromstring(omml_xml))
                        except Exception:
                            run = p.add_run(_latex_plain_fallback(latex))
                            run.italic = italic
                            run.font.size = Pt(font_size_pt)
                    elif latex:
                        run = p.add_run(_latex_plain_fallback(latex))
                        run.italic = italic
                        run.font.size = Pt(font_size_pt)
                elif r_item.get("kind") == "image":
                    try:
                        img_abs = ws_tool._safe_resolve(workspace_dir, r_item.get("path", ""))
                        if os.path.isfile(img_abs):
                            w = r_item.get("width_inches")
                            run = p.add_run()
                            if isinstance(w, (int, float)) and w > 0:
                                run.add_picture(img_abs, width=Inches(w))
                            else:
                                run.add_picture(img_abs)
                    except Exception:
                        pass
        else:
            run = p.add_run(_clean_docx_text(sb.get("text", "")))
            run.italic = italic
            run.font.size = Pt(font_size_pt)
    return p


def _apply_image_docx(doc, workspace_dir: str, idx: int, b: dict) -> str | None:
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    img_path = str(b.get("path", "")).strip()
    if not img_path:
        return (
            f"block[{idx}].path required for image; generate, fetch, or copy the image into the workspace first, "
            f"then pass that relative path.\n图片块必须引用已在工作区存在的相对路径。"
        )
    try:
        img_target = ws_tool._safe_resolve(workspace_dir, img_path)
    except ValueError as e:
        return f"block[{idx}] image path invalid: {e}"
    if not os.path.isfile(img_target):
        # 2026-05-12 P47: office 找图自动重定向 (类似 P42 read_file)
        # 实施 2 层:
        #   1. 规则反推: 主区命名 `<task_id>_<sub>` → sandbox `_helpers_shared/<task_id>/<sub>`
        #      (helper sandbox 内文件命名约定固定, 反推 100% 确定)
        #   2. fuzzy 兜底: score>=95 自动重定向; 否则给 LLM 候选 hint
        # 病因(实测 23:55 trace): gen_paper helper 用主区命名 `gen_charts_chart1_overview.png`
        # 但 helper sandbox 内只有 `_helpers_shared/gen_charts/chart1_overview.png`,
        # 找不到 → office 失败 → 连续 4 次失败 → helper STUCK.
        _sandbox_redirect = _resolve_main_named_to_sandbox(workspace_dir, img_path)
        if _sandbox_redirect and os.path.isfile(_sandbox_redirect):
            img_target = _sandbox_redirect
            # 规则反推命中, 用真实路径继续 (fall through)
        else:
            try:
                _suggestions = ws_tool._suggest_similar_files(workspace_dir, img_path)
                if _suggestions and _suggestions[0]["score"] >= 95:
                    _redirect = _suggestions[0]["path"]
                    _redirect_target = ws_tool._safe_resolve(workspace_dir, _redirect)
                    if os.path.isfile(_redirect_target):
                        img_target = _redirect_target
                    else:
                        return (
                            f"block[{idx}] image not found: {img_path} "
                            f"(fuzzy match found {_redirect}, but that file is still not readable). "
                            "Generate or stage the image in the workspace first.\n"
                            "图片需先在工作区生成或暂存。"
                        )
                else:
                    _hint = ""
                    if _suggestions:
                        _cands = ", ".join(s["path"] for s in _suggestions[:3])
                        _hint = f" Similar workspace files: {_cands}"
                    return (
                        f"block[{idx}] image not found: {img_path} "
                        f"(the image must exist in the workspace before insertion).{_hint}\n"
                        "图片需先在工作区生成或暂存。"
                    )
            except (ValueError, OSError, AttributeError):
                return (
                    f"block[{idx}] image not found: {img_path} "
                    "(the image must exist in the workspace before insertion).\n"
                    "图片需先在工作区生成或暂存。"
                )
    try:
        if os.path.getsize(img_target) > _MAX_IMAGE_BYTES:
            return f"block[{idx}] image too large (>{_MAX_IMAGE_BYTES // (1024*1024)}MB)"
    except OSError:
        pass
    try:
        w = float(b.get("width_inches") or _DEFAULT_IMAGE_WIDTH_INCHES)
        w = max(0.5, min(7.5, w))
    except (TypeError, ValueError):
        w = _DEFAULT_IMAGE_WIDTH_INCHES
    try:
        doc.add_picture(img_target, width=Inches(w))
    except Exception as e:
        return f"block[{idx}] add_picture failed for {img_path}: {type(e).__name__}: {e}"
    cap = _clean_docx_text(b.get("caption", ""))
    if cap:
        _render_docx_caption_paragraph(doc, workspace_dir, cap, italic=True, font_size_pt=10)
    return None


# 2026-05-11 P13d: Word 原生公式 (OMML) 支持
# OMML = Office Math Markup Language, Word 2007+ 原生格式, 可编辑可复制可缩放
# 用法: blocks: [{"type":"equation", "latex":"10^4", "display":false}]
# 实现: 把简单 LaTeX 转 OMML XML, lxml 直接插入 paragraph
# 复杂公式无法转 OMML 时降级到 PNG (走原 image block 路径)
#
# 2026-05-17 P151: 扩展支持 \sum/\int/\prod/\lim/\binom + 嵌套 \frac/\sqrt
# 之前这些直接 return None 走 PNG,损失"在 Word 里可编辑"的能力。
# 工程论文常见的 \sum_{i=1}^n x_i 这类写法现在原生 OMML 渲染。

















def _apply_equation_docx(doc, workspace_dir: str, idx: int, b: dict) -> str | None:
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    """处理 equation block — 优先 OMML 原生, 失败降级 PNG.

    args:
      latex: 公式源码 (LaTeX 风格)
      display: true 时居中放大, false (默认) 时 inline 大小
    """
    latex = str(b.get("latex") or b.get("text") or "").strip()
    if not latex:
        return f"block[{idx}] equation: missing 'latex' field"
    is_display = bool(b.get("display"))

    # 2026-05-17 P160: complexity preflight — 不支持的语法直接报错, 不假装成功
    level, hint = _classify_latex_complexity(latex)
    if level == "unsupported":
        return (
            f"block[{idx}] equation: unsupported LaTeX syntax. {hint} "
            f"LaTeX source preview: {latex[:120]!r}"
        )

    # 优先尝试 OMML
    omml_xml = _latex_to_omml_simple(latex)
    if omml_xml:
        try:
            from lxml import etree
            from docx.oxml.ns import qn
            # 解析 OMML XML
            omml_elem = etree.fromstring(omml_xml)
            # 加到新 paragraph
            p = doc.add_paragraph()
            if is_display:
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p._p.append(omml_elem)
            return None
        except Exception as _e:
            # OMML 嵌入失败 → 降级 PNG
            pass

    # 降级到 PNG 路径 — 2026-05-18 P193: 用 _render_formula_to_png 统一路径,
    # 自带 P179 内容 hash 文件名 + P186 缓存 + P180/P181/P183/P191 normalize +
    # P178 broken-detector. 之前直接调 _render_latex_to_png 会绕开这套保护。
    # 2026-05-18 P195: 把 is_display 传过去, 公式 block 默认 display 字号 16pt.
    png = _render_formula_to_png(latex, workspace_dir, f"eq{idx}", is_display=is_display)
    if png is None:
        # 最后降级: 公式无法渲染时, 用 P178 一致的明确标记
        p = doc.add_paragraph()
        if is_display:
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        wrap_left = "$$" if is_display else "$"
        wrap_right = "$$" if is_display else "$"
        # P193: 不再写 italic 原始 LaTeX (这是 P178 修过的灾难) — 改用 broken 标记
        fallback_plain = _latex_plain_fallback(latex)
        if _is_broken_after_strip(fallback_plain, latex):
            p.add_run(f"[⚠️ 公式渲染失败 {wrap_left}{latex}{wrap_right}]")
        else:
            p.add_run(fallback_plain)
        return None

    natural_w = _png_natural_width_inches(workspace_dir, png)
    width = natural_w or (3.0 if is_display else 0.8)
    # 走 image 嵌入
    img_block = {"path": png, "width_inches": width}
    if is_display:
        img_block["caption"] = ""  # 可加 latex 源作 caption
    return _apply_image_docx(doc, workspace_dir, idx, img_block)



# 2026-05-17 P161: docx 数字校验 — 对照 CSV 找潜在编造的数字
from app.llm.tools.office_verify import (  # noqa: E402,F401
    _docx_verify_numbers,
    _docx_verify_rigor,
    _docx_verify_against_source,
    _pptx_verify_numbers,
)



# 2026-05-17 P162: 综合 rigor 检测器 — 9 层数据严谨性检查 (L1-L9)
# 详见函数 docstring。




async def _docx_extract_images(workspace_dir: str, target: str, rel_path: str, args: dict) -> str:
    from docx import Document

    if not os.path.isfile(target):
        return _err(f"file not found: {rel_path}")

    out_dir = _resolve_out_dir(args, rel_path)
    try:
        out_target = ws_tool._safe_resolve(workspace_dir, out_dir)
    except ValueError as e:
        return _err(f"out_dir invalid: {e}")
    os.makedirs(out_target, exist_ok=True)

    def _do_extract():
        try:
            doc = Document(target)
        except Exception as e:
            return None, f"open failed: {type(e).__name__}: {e}"
        extracted = _dump_images_from_part(doc.part, out_target, out_dir)
        return extracted, None

    extracted, err = await asyncio.to_thread(_do_extract)
    if err:
        return _err(err)
    return json.dumps({
        "ok": True, "action": "extract_images", "format": "docx", "path": rel_path,
        "out_dir": out_dir, "extracted": extracted, "count": len(extracted),
    }, ensure_ascii=False)


async def _raw_image_ocr(workspace_dir: str, target: str, rel_path: str,
                          args: dict) -> str:
    """OCR a single raw image file (jpg/png/...) — bypass office zip extraction.

    2026-05-18 P197: 让 helper 能直接对 jpg/png OCR 不必走 "嵌入临时 docx 再 OCR"
    的灾难工作流. 复用底层 ocr_file primitive, 输出与 _office_ocr_images 同结构
    (items[] + merged_text), 让 helper 用同一个 schema 处理。

    2026-05-18 P200: 增加 `save_to` 参数. 当 OCR 输出 ≥ 2KB (典型多题试卷) 时
    helper 应传 `save_to='ocr_X.md'`, 系统**只返回 metadata + preview**而非全文,
    helper 再用 read_file(start_line=N, end_line=M) 局部读取写 docx, 避免单次
    30KB OCR 全文塞满上下文 → truncate → 脑补的灾难循环。
    """
    if not os.path.isfile(target):
        return _err(f"file not found: {rel_path}")

    try:
        from app.llm.tools.ocr_bridge import ocr_file
    except ImportError as e:
        return _err(f"ocr_bridge not available: {e}")

    per_image_timeout = int(args.get("per_image_timeout", 60))
    max_size_mb = float(args.get("max_size_mb", 50.0))
    save_to = str(args.get("save_to", "")).strip()  # P200: optional output file
    try:
        size_bytes = os.path.getsize(target)
    except OSError:
        size_bytes = 0
    if size_bytes > max_size_mb * 1024 * 1024:
        return _err(
            f"image too large ({size_bytes/1024/1024:.1f}MB > {max_size_mb}MB)",
            action="ocr_images",
            hint=f"Pass max_size_mb={int(size_bytes/1024/1024)+1} to override",
        )

    def _do_ocr():
        try:
            from pathlib import Path
            r = ocr_file(Path(target), timeout=per_image_timeout)
            return {
                "ok": r.ok, "text": r.text, "error": r.error,
                "engine": r.engine, "elapsed_ms": r.elapsed_ms,
            }
        except Exception as e:
            return {"ok": False, "error": f"ocr_file failed: {type(e).__name__}: {e}"}

    result = await asyncio.to_thread(_do_ocr)

    if not result or not result.get("ok"):
        return _err(result.get("error", "ocr failed") if result else "ocr returned None")

    text = (result.get("text") or "").strip()
    name = os.path.basename(rel_path)
    merged_text = f"[图 {name}]\n{text}" if text else ""

    # P200: save_to 模式 — 写文件, 返回 metadata + 短 preview (前 600 字符 + 行数)
    if save_to:
        try:
            save_path = ws_tool._safe_resolve(workspace_dir, save_to)
        except ValueError as e:
            return _err(f"invalid save_to path: {e}", action="ocr_images",
                        hint="save_to must be a relative path within workspace")
        try:
            os.makedirs(os.path.dirname(save_path) or ".", exist_ok=True)
            with open(save_path, "w", encoding="utf-8") as f:
                f.write(merged_text)
        except OSError as e:
            return _err(f"failed to write {save_to}: {e}")
        lines = merged_text.count("\n") + (1 if merged_text and not merged_text.endswith("\n") else 0)
        preview = merged_text[:600]
        if len(merged_text) > 600:
            preview += f"\n... ({len(merged_text) - 600} more chars, use read_file to load by range)"
        payload = {
            "ok": True,
            "action": "ocr_images",
            "format": "image",
            "path": rel_path,
            "saved_to": save_to,
            "size_bytes": len(merged_text.encode("utf-8")),
            "lines": lines,
            "preview": preview,
            "engine": result.get("engine", ""),
            "elapsed_ms": result.get("elapsed_ms", 0),
            "hint": (
                f"OCR text was saved to {save_to} ({lines} lines, {len(merged_text.encode('utf-8'))} bytes). "
                f"Use read_file('{save_to}', start_line=N, end_line=M) to read it in chunks before writing DOCX content. "
                f"Do not load the entire OCR text back into one prompt.\n"
                f"OCR 已保存到文件，应分段读取。"
            ),
        }
        return json.dumps(payload, ensure_ascii=False)

    # 兼容老行为: 直接返回全文 (P197 schema)
    payload = {
        "ok": True,
        "action": "ocr_images",
        "format": "image",
        "path": rel_path,
        "total_images": 1,
        "ocr_count": 1 if text else 0,
        "skipped_no_text": 0 if text else 1,
        "items": [{"name": name, "size_bytes": size_bytes, "text": text}],
        "merged_text": merged_text,
        "engine": result.get("engine", ""),
        "elapsed_ms": result.get("elapsed_ms", 0),
    }
    # P200: 如果输出大 (≥2KB), 在 hint 里建议用 save_to
    if len(merged_text) >= 2048:
        payload["hint"] = (
            f"OCR text of {len(merged_text)} characters was returned directly into context. "
            "For similar large OCR tasks, set `save_to='ocr_X.md'`, then read the saved file by line range to avoid repeated large-context truncation.\n\n"
            "大段 OCR 文本建议保存为文件后分段读取。"
        )
    return json.dumps(payload, ensure_ascii=False)


async def _office_ocr_images(workspace_dir: str, target: str, rel_path: str,
                             args: dict, fmt: str) -> str:
    """一键 OCR office 文件 (docx/pptx/xlsx) 内所有嵌入图片.
    
    比 extract_images + ocr(每张) 省 N-1 个工具调用.
    
    args:
        max_images: int (默认 30), 单文件 OCR 上限
        max_size_mb: float (默认 50.0), 单图最大 MB
    
    返回:
        {ok: true, total_images, ocr_count, items: [{name, size_bytes, text}, ...], merged_text}
    """
    if not os.path.isfile(target):
        return _err(f"file not found: {rel_path}")
    
    try:
        from app.llm.tools.ocr_bridge import ocr_office_images
    except ImportError as e:
        return _err(f"ocr_bridge not available: {e}")
    
    max_images = int(args.get("max_images", 30))
    image_offset = int(args.get("image_offset", 0) or 0)
    max_size_mb = float(args.get("max_size_mb", 50.0))
    # 2026-05-18 P172: 之前 per_image_timeout=90, 16 张图 = 24min worst case →
    # API stall (90s) 必触发。改 30s 默认 + 4 并发 worker, 单图卡死最多 30s,
    # 整体 16 张 ≤ 2 分钟 (4 并发) 而非 24 分钟。
    per_image_timeout = int(args.get("per_image_timeout", 30))
    max_workers = int(args.get("max_workers", 4))
    
    def _do_ocr():
        return ocr_office_images(
            target,
            max_images=max_images,
            image_offset=image_offset,
            max_size_mb=max_size_mb,
            per_image_timeout=per_image_timeout,
            max_workers=max_workers,
        )
    
    result = await asyncio.to_thread(_do_ocr)
    
    if not result.get("ok"):
        return _err(result.get("error", "ocr failed"))
    
    # 添加 action / format / path 字段
    result["action"] = "ocr_images"
    result["format"] = fmt
    result["path"] = rel_path

    # 2026-05-18 P200: save_to 模式 — 写文件而非塞全文到上下文
    save_to = str(args.get("save_to", "")).strip()
    if save_to:
        merged_text = result.get("merged_text", "")
        try:
            save_path = ws_tool._safe_resolve(workspace_dir, save_to)
        except ValueError as e:
            return _err(f"invalid save_to path: {e}", action="ocr_images",
                        hint="save_to must be a relative path within workspace")
        try:
            os.makedirs(os.path.dirname(save_path) or ".", exist_ok=True)
            with open(save_path, "w", encoding="utf-8") as f:
                f.write(merged_text)
        except OSError as e:
            return _err(f"failed to write {save_to}: {e}")
        lines = merged_text.count("\n") + (1 if merged_text and not merged_text.endswith("\n") else 0)
        preview = merged_text[:600]
        if len(merged_text) > 600:
            preview += f"\n... ({len(merged_text) - 600} more chars, use read_file to load by range)"
        # 不带 items + merged_text 全文, 只返 metadata + preview
        item_summaries = [
            {"name": it.get("name", ""), "size_bytes": it.get("size_bytes", 0),
             "text_chars": len(it.get("text", ""))}
            for it in result.get("items", [])
        ]
        return json.dumps({
            "ok": True, "action": "ocr_images", "format": fmt, "path": rel_path,
            "saved_to": save_to,
            "size_bytes": len(merged_text.encode("utf-8")),
            "lines": lines,
            "preview": preview,
            "total_images": result.get("total_images", 0),
            "image_offset": result.get("image_offset", image_offset),
            "processed_images": result.get("processed_images", 0),
            "has_more_images": result.get("has_more_images", False),
            **({"next_image_offset": result.get("next_image_offset")} if result.get("has_more_images") else {}),
            "ocr_count": result.get("ocr_count", 0),
            "skipped_too_large": result.get("skipped_too_large", 0),
            "skipped_no_text": result.get("skipped_no_text", 0),
            "items_summary": item_summaries,
            "hint": (
                f"OCR text was saved to {save_to} ({lines} lines, {len(merged_text.encode('utf-8'))} bytes). "
                f"Use read_file('{save_to}', start_line=N, end_line=M) to read it in chunks before writing DOCX content.\n"
                f"OCR 已保存到文件，应分段读取。"
            ),
        }, ensure_ascii=False)

    # 老行为: 全文返回
    merged_text = result.get("merged_text", "")
    if len(merged_text) >= 2048:
        result["hint"] = (
            f"OCR text of {len(merged_text)} characters was returned directly into context. "
            "For similar large OCR tasks, set `save_to='ocr_X.md'`, then read the saved file by line range to avoid repeated large-context truncation.\n\n"
            "大段 OCR 文本建议保存为文件后分段读取。"
        )
    if result.get("has_more_images"):
        result["batch_hint"] = (
            "More embedded images remain. Call office(action='ocr_images') again with "
            f"image_offset={result.get('next_image_offset')} and save_to set to another .txt/.md file, "
            "then read OCR output by line range.\n\n"
            "仍有图片未 OCR；继续分批保存并分段读取。"
        )
    return json.dumps(result, ensure_ascii=False)


async def _docx_ocr_images(workspace_dir: str, target: str, rel_path: str, args: dict) -> str:
    return await _office_ocr_images(workspace_dir, target, rel_path, args, "docx")


async def _pptx_ocr_images(workspace_dir: str, target: str, rel_path: str, args: dict) -> str:
    return await _office_ocr_images(workspace_dir, target, rel_path, args, "pptx")


async def _xlsx_ocr_images(workspace_dir: str, target: str, rel_path: str, args: dict) -> str:
    return await _office_ocr_images(workspace_dir, target, rel_path, args, "xlsx")


async def _docx_insert_image(workspace_dir: str, target: str, rel_path: str, args: dict) -> str:
    img_path = str(args.get("image_path", "")).strip()
    if not img_path:
        return _err("image_path required", action="insert_image", hint="Provide the workspace-relative path to an existing image file, e.g. image_path:'chart.png'")
    block: dict = {"type": "image", "path": img_path}
    if args.get("width_inches") is not None:
        block["width_inches"] = args["width_inches"]
    if args.get("caption"):
        block["caption"] = args["caption"]
    return await _docx_append(workspace_dir, target, rel_path, {"blocks": [block]})


# ═══════════════════════════════════════════════════════════════
# PowerPoint (.pptx)
# ═══════════════════════════════════════════════════════════════



async def _pptx_write(workspace_dir: str, target: str, rel_path: str, args: dict) -> str:
    from pptx import Presentation

    slides = args.get("slides") or []
    if not isinstance(slides, list) or not slides:
        return _err("slides must be a non-empty array", action="write", hint="Pass slides as [{layout:'title', title:'Title'}, {layout:'title_content', title:'...', body:[...]}, ...]")

    def _do_write():
        prs = Presentation()
        # 2026-05-18 P196: pptx 默认 4:3 (10"×7.5") 已过时, 改 16:9 (13.33"×7.5")
        # 匹配现代显示器/投影仪. helper 写 layout='title'/'title_content' 等会更宽.
        from pptx.util import Inches as PptxInches
        prs.slide_width = PptxInches(13.33)
        prs.slide_height = PptxInches(7.5)
        err = _apply_slides_pptx(prs, workspace_dir, slides)
        if err:
            return None, err
        os.makedirs(os.path.dirname(target) or ".", exist_ok=True)
        try:
            prs.save(target)
        except Exception as e:
            return None, f"save failed: {type(e).__name__}: {e}"
        size = os.path.getsize(target) if os.path.exists(target) else 0
        return size, None

    size, err = await asyncio.to_thread(_do_write)
    if err:
        return _err(err)
    return json.dumps({
        "ok": True, "action": "write", "format": "pptx", "path": rel_path,
        "slides_written": len(slides), "size_bytes": size,
    }, ensure_ascii=False)


async def _pptx_append(workspace_dir: str, target: str, rel_path: str, args: dict) -> str:
    from pptx import Presentation

    if not os.path.isfile(target):
        return _err(f"file not found: {rel_path}")

    slides = args.get("slides") or []
    if not isinstance(slides, list) or not slides:
        return _err("slides must be a non-empty array")

    def _do_append():
        try:
            prs = Presentation(target)
        except Exception as e:
            return f"open failed: {type(e).__name__}: {e}"
        err = _apply_slides_pptx(prs, workspace_dir, slides)
        if err:
            return err
        try:
            prs.save(target)
        except Exception as e:
            return f"save failed: {type(e).__name__}: {e}"
        return None

    err = await asyncio.to_thread(_do_append)
    if err:
        return _err(err)
    return json.dumps({
        "ok": True, "action": "append", "format": "pptx", "path": rel_path,
        "slides_appended": len(slides),
    }, ensure_ascii=False)


# ───────────────────────────────────────────────────────────────
# 2026-05-04 v19.1: pptx 增量编辑
# 长 PPT(15+ slides)一次性 write 容易触发 LLM thinking timeout。
# 模型可以分:1) write 前 5 页骨架 → 2) append 中间 5 页 → 3) replace_slide 改某页
# ───────────────────────────────────────────────────────────────



async def _pptx_replace_slide(
    workspace_dir: str, target: str, rel_path: str, args: dict,
) -> str:
    """替换指定 index 的一页。

    args:
      index: 页码(0-based)
      slide: 单页字典(同 write 的 slides[i] 结构)
        或者 slides: 多页数组(替换 1 页变成多页)
    """
    from pptx import Presentation

    if not os.path.isfile(target):
        return _err(f"file not found: {rel_path}", action="replace_slide")
    try:
        index = int(args.get("index", -1))
    except (TypeError, ValueError):
        return _err("index must be integer", action="replace_slide", hint="Provide a 0-based slide index from read output, e.g. index:2")
    if index < 0:
        return _err("index must be >= 0", action="replace_slide")

    # 单页 or 多页
    new_slides = args.get("slides")
    if new_slides is None:
        single = args.get("slide")
        if single is None:
            return _err("missing 'slide' (single dict) or 'slides' (array)", action="replace_slide", hint="Pass either slide:{{layout:'title_content', title:'...', body:[...]}} or slides:[...]")
        new_slides = [single]
    if not isinstance(new_slides, list) or not new_slides:
        return _err("slides must be non-empty array", action="replace_slide", hint="Pass the replacement slide content in slides:[{...}] or slide:{...}")

    def _do():
        try:
            prs = Presentation(target)
        except Exception as e:
            return None, f"open failed: {type(e).__name__}: {e}"
        total = len(prs.slides)
        if index >= total:
            return None, f"index {index} out of range (have {total} slides)"

        # 策略:先在末尾追加新页,然后用 XML 操作把它们移到 index 位置,再删除原 index 页
        # 简单实现:删除目标页,然后 append 新页,再用 XML 把 append 出的新页移到目标位置
        sldIdLst = prs.slides._sldIdLst
        original_ids = list(sldIdLst)
        old_count = len(original_ids)

        # 1. 先 append 新页
        err = _apply_slides_pptx(prs, workspace_dir, new_slides)
        if err:
            return None, err

        # 2. 找到刚 append 出来的新页 ids(在 sldIdLst 末尾)
        new_ids = list(sldIdLst)[old_count:]

        # 3. 删除原 index 页
        try:
            _pptx_remove_slide(prs, index)
        except IndexError as e:
            return None, str(e)

        # 4. 把新加的 ids 移到 index 位置
        # 现在 sldIdLst 是: original (去掉 index) + new_ids
        # 移除 new_ids 然后 insert 到 index 位置
        for nid in new_ids:
            sldIdLst.remove(nid)
        # insert 时倒序,因为 insert(index, x) 之后 index 位置是 x
        for nid in reversed(new_ids):
            sldIdLst.insert(index, nid)

        try:
            prs.save(target)
        except Exception as e:
            return None, f"save failed: {type(e).__name__}: {e}"
        return {
            "replaced_index": index,
            "new_slide_count": len(new_slides),
        }, None

    info, err = await asyncio.to_thread(_do)
    if err:
        return _err(err)
    return json.dumps({
        "ok": True, "action": "replace_slide", "format": "pptx",
        "path": rel_path,
        **info,
    }, ensure_ascii=False)


async def _pptx_insert_slide(
    workspace_dir: str, target: str, rel_path: str, args: dict,
) -> str:
    """在指定 index 之前插入一页或多页。

    args:
      index: 页码(0-based;index=N 表示插到第 N 页之前)
      slides: 数组
    """
    from pptx import Presentation

    if not os.path.isfile(target):
        return _err(f"file not found: {rel_path}")
    try:
        index = int(args.get("index", -1))
    except (TypeError, ValueError):
        return _err("index must be integer")
    if index < 0:
        return _err("index must be >= 0")
    new_slides = args.get("slides") or []
    if not isinstance(new_slides, list) or not new_slides:
        return _err("slides must be non-empty array")

    def _do():
        try:
            prs = Presentation(target)
        except Exception as e:
            return None, f"open failed: {type(e).__name__}: {e}"
        total = len(prs.slides)
        if index > total:
            return None, f"index {index} out of range (max {total})"

        sldIdLst = prs.slides._sldIdLst
        old_count = len(list(sldIdLst))
        err = _apply_slides_pptx(prs, workspace_dir, new_slides)
        if err:
            return None, err
        new_ids = list(sldIdLst)[old_count:]

        if index < total:
            # 移到 index 位置
            for nid in new_ids:
                sldIdLst.remove(nid)
            for nid in reversed(new_ids):
                sldIdLst.insert(index, nid)
        # else: append 到尾部,不动

        try:
            prs.save(target)
        except Exception as e:
            return None, f"save failed: {type(e).__name__}: {e}"
        return {
            "inserted_at": index,
            "new_slide_count": len(new_slides),
        }, None

    info, err = await asyncio.to_thread(_do)
    if err:
        return _err(err)
    return json.dumps({
        "ok": True, "action": "insert_slide", "format": "pptx",
        "path": rel_path,
        **info,
    }, ensure_ascii=False)





def _apply_slides_pptx(prs, workspace_dir: str, slides: list) -> str | None:
    """根据 layout 字段创建 slide。完整重写(2026-05 修):

    支持的 layout(每种用对应的 python-pptx 内置 slide_layout 索引):
      - title         (idx 0): 封面页, 用 title + subtitle
      - section       (idx 2): 章节分隔, 用 title + subtitle
      - title_content (idx 1): 标题 + 项目符号正文, 用 title + body[str list 或 dict list]
      - two_column    (idx 3): 双栏, 用 title + left[] + right[]
      - image         (idx 5): 标题 + 居中大图, 用 title + image{path,width_inches,caption}
      - table         (idx 5): 标题 + 表格, 用 title + table{rows[[]], header}
      - blank         (idx 6): 空白, 用 title?(可选) + text(单段长文本)
                              或 body[](自由组合 text/bullets/image/table)

    旧版 bug(trace da389409 实测):
      - 模型按 schema 写 subtitle/text/顶层 image — 全被忽略
      - "image" layout 把图放顶层 image 字段 — 旧 handler 只看 body[].type=='image' → 图片永远不入 slide
      - 模型在 blank/title_content 层级混用各种格式 — handler 报 'has non-object element'
    重写后所有 7 种 layout 完整闭环, 字符串/dict 混排都宽容处理。
    """
    from pptx.util import Inches, Pt

    # python-pptx 默认模板 layout 索引(标准 Office Theme):
    #   0=Title Slide   1=Title and Content   2=Section Header   3=Two Content
    #   4=Comparison    5=Title Only          6=Blank             7=Content w/ Caption
    #   8=Picture w/ Caption
    n_layouts = len(prs.slide_layouts)

    def _layout(idx: int):
        if idx < n_layouts:
            return prs.slide_layouts[idx]
        return prs.slide_layouts[0]

    layout_map = {
        "title":         _layout(0),
        "section":       _layout(2 if n_layouts > 2 else 0),
        "title_content": _layout(1 if n_layouts > 1 else 0),
        "two_column":    _layout(3 if n_layouts > 3 else 1),
        "image":         _layout(5 if n_layouts > 5 else 6 if n_layouts > 6 else 0),
        "table":         _layout(5 if n_layouts > 5 else 6 if n_layouts > 6 else 0),
        "blank":         _layout(6 if n_layouts > 6 else 5 if n_layouts > 5 else 0),
    }

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for s_idx, s in enumerate(slides):
        if not isinstance(s, dict):
            return f"slide[{s_idx}] must be an object"

        layout_name = str(s.get("layout", "")).strip().lower()
        title = str(s.get("title", "")).strip()
        # 不指定 layout: 有 title → title_content; 否则 blank
        if not layout_name:
            layout_name = "title_content" if title else "blank"
        layout = layout_map.get(layout_name)
        if layout is None:
            # 未知 layout — 宽容降级而不是报错
            debug_log_unknown = True
            layout = layout_map["title_content"] if title else layout_map["blank"]
        else:
            debug_log_unknown = False
        slide = prs.slides.add_slide(layout)

        # 设标题
        if title and slide.shapes.title is not None:
            # 2026-05-18 P190: 标题里允许 inline `$...$` LaTeX (Unicode-only path)
            slide.shapes.title.text = (
                _render_inline_latex_to_text_only(title) if "$" in title else title
            )

        # ── 按 layout 类型分发 ──
        try:
            if layout_name == "title":
                err = _pptx_render_title(slide, s)
            elif layout_name == "section":
                err = _pptx_render_section(slide, s, slide_w)
            elif layout_name == "title_content":
                err = _pptx_render_title_content(slide, s, slide_w, slide_h, title)
            elif layout_name == "two_column":
                err = _pptx_render_two_column(slide, s, slide_w, slide_h, title)
            elif layout_name == "image":
                err = _pptx_render_image_layout(slide, workspace_dir, s_idx, s, slide_w, slide_h, title)
            elif layout_name == "table":
                err = _pptx_render_table_layout(slide, s_idx, s, slide_w, slide_h, title)
            else:  # blank or unknown
                err = _pptx_render_blank(slide, workspace_dir, s_idx, s, slide_w, slide_h, title)
            if err:
                return err
        except Exception as e:
            return f"slide[{s_idx}] render failed: {type(e).__name__}: {e}"

    return None


# ─── Per-layout renderers ────────────────────────────────────────────

def _pptx_render_title(slide, s: dict) -> str | None:
    """title layout: title + subtitle (placeholder idx 1 in Title Slide).

    2026-05-18 P190: subtitle 也走 LaTeX 渲染 (Unicode/plain only, 同 _fill_text_frame).
    """
    subtitle = str(s.get("subtitle", "")).strip()
    if not subtitle:
        return None
    # P190: 渲染 LaTeX
    if "$" in subtitle:
        subtitle = _render_inline_latex_to_text_only(subtitle)
    for shp in slide.placeholders:
        pf = getattr(shp, "placeholder_format", None)
        if pf and pf.idx == 1 and shp.has_text_frame:
            shp.text_frame.text = subtitle
            return None
    # Section Header layout 没有 idx=1 的 subtitle placeholder,fallback 到任意非 title 文本框
    for shp in slide.placeholders:
        try:
            if shp != slide.shapes.title and shp.has_text_frame:
                shp.text_frame.text = subtitle
                return None
        except Exception:
            continue
    return None


def _pptx_render_section(slide, s: dict, slide_w) -> str | None:
    """section layout: title + subtitle. Section Header 模板 placeholder 通常是 idx 1。"""
    return _pptx_render_title(slide, s)


def _pptx_render_title_content(slide, s: dict, slide_w, slide_h, title: str) -> str | None:
    """title_content layout: title + bullet 正文。

    支持三种 body 形式:
      - body=["要点1", "要点2"]            纯字符串数组 → 每条作为 bullet
      - body=[{"type":"bullets","items":[...]}, {"type":"text","text":"..."}]   旧格式
      - body=[{"type":"image",...}]        会触发 fallback 到 _pptx_render_blank
    """
    body = s.get("body")
    text_field = str(s.get("text", "")).strip()
    if body is None and not text_field:
        return None

    # body 含图片/表格时降级到 blank 渲染(图片/表混排)
    if isinstance(body, list) and any(
        isinstance(b, dict) and str(b.get("type", "")).lower() in ("image", "table")
        for b in body
    ):
        return _pptx_render_blank(slide, "", -1, s, slide_w, slide_h, title)

    bullets, paragraphs = _normalize_text_items(body, text_field)
    if not bullets and not paragraphs:
        return None

    body_ph = _find_body_placeholder(slide)
    if body_ph is None:
        # title_content layout 应该有 idx=1 placeholder, 没有就新建文本框
        from pptx.util import Inches
        body_ph = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(1.5) if title else Inches(0.5),
            slide_w - Inches(1.0),
            slide_h - (Inches(2.0) if title else Inches(1.0)),
        )
    _fill_text_frame(body_ph.text_frame, bullets, paragraphs)
    return None


def _pptx_render_two_column(slide, s: dict, slide_w, slide_h, title: str) -> str | None:
    """two_column layout: title + left[] + right[].

    Two Content 标准布局有 placeholder idx 1 (left) + idx 2 (right)。
    """
    left = s.get("left")
    right = s.get("right")
    if left is None and right is None:
        return None
    # 找两个 body placeholder
    ph_left = None
    ph_right = None
    for shp in slide.placeholders:
        pf = getattr(shp, "placeholder_format", None)
        if pf is None or not shp.has_text_frame:
            continue
        if pf.idx == 1 and ph_left is None:
            ph_left = shp
        elif pf.idx == 2 and ph_right is None:
            ph_right = shp
    from pptx.util import Inches
    half_w = (slide_w - Inches(1.5)) // 2
    top = Inches(1.5) if title else Inches(0.5)
    h = slide_h - top - Inches(0.5)
    if ph_left is None:
        ph_left = slide.shapes.add_textbox(Inches(0.5), top, half_w, h)
    if ph_right is None:
        ph_right = slide.shapes.add_textbox(Inches(0.5) + half_w + Inches(0.5), top, half_w, h)

    if left is not None:
        b, p = _normalize_text_items(left, "")
        _fill_text_frame(ph_left.text_frame, b, p)
    if right is not None:
        b, p = _normalize_text_items(right, "")
        _fill_text_frame(ph_right.text_frame, b, p)
    return None


def _pptx_render_image_layout(slide, workspace_dir: str, s_idx: int, s: dict,
                              slide_w, slide_h, title: str) -> str | None:
    """image layout: title + 大图(顶层 image 字段, 不在 body 里)。

    关键修复: 模型按 schema 直接传顶层 image:{path,width_inches,caption}。
    旧 handler 只看 body[].type=='image' → 图片完全不出现。
    """
    from pptx.util import Inches, Pt

    img = s.get("image")
    if not isinstance(img, dict):
        # 模型可能误塞在 body 里, 兼容一下
        body = s.get("body") or []
        for b in body:
            if isinstance(b, dict) and str(b.get("type", "")).lower() == "image":
                img = b
                break
    if not isinstance(img, dict):
        return None  # 没图就只剩 title, 也不是错

    img_path = str(img.get("path", "")).strip()
    if not img_path:
        return f"slide[{s_idx}].image.path required"
    try:
        img_target = ws_tool._safe_resolve(workspace_dir, img_path)
    except ValueError as e:
        return f"slide[{s_idx}] image path invalid: {e}"
    if not os.path.isfile(img_target):
        # 2026-05-12 P47: pptx 找图自动重定向 (与 docx 一致, 2 层)
        _sandbox_redirect = _resolve_main_named_to_sandbox(workspace_dir, img_path)
        if _sandbox_redirect and os.path.isfile(_sandbox_redirect):
            img_target = _sandbox_redirect
        else:
            try:
                _sug = ws_tool._suggest_similar_files(workspace_dir, img_path)
                if _sug and _sug[0]["score"] >= 95:
                    _r = ws_tool._safe_resolve(workspace_dir, _sug[0]["path"])
                    if os.path.isfile(_r):
                        img_target = _r
                    else:
                        return f"slide[{s_idx}] image not found: {img_path}"
                else:
                    _hint = ""
                    if _sug:
                        _cands = ", ".join(s["path"] for s in _sug[:3])
                        _hint = f" 工作区内相似文件: {_cands}"
                    return f"slide[{s_idx}] image not found: {img_path}{_hint}"
            except (ValueError, OSError, AttributeError):
                return f"slide[{s_idx}] image not found: {img_path}"

    try:
        w_in = float(img.get("width_inches") or _DEFAULT_PPTX_IMAGE_WIDTH_INCHES)
        w_in = max(1.0, min(13.0, w_in))
    except (TypeError, ValueError):
        w_in = _DEFAULT_PPTX_IMAGE_WIDTH_INCHES
    width = Inches(w_in)

    caption = str(img.get("caption", "")).strip()
    img_top = Inches(1.5) if title else Inches(0.5)
    available_h = slide_h - img_top - (Inches(0.6) if caption else Inches(0.3))
    # 居中, 限制最大高度避免溢出
    left = (slide_w - width) // 2
    try:
        pic = slide.shapes.add_picture(img_target, left, img_top, width=width)
        # 若按宽度撑出 slide, 等比缩
        if pic.height > available_h:
            ratio = available_h / pic.height
            new_w = int(pic.width * ratio)
            pic.width = new_w
            pic.height = available_h
            pic.left = (slide_w - new_w) // 2
    except Exception as e:
        return f"slide[{s_idx}] add_picture failed: {type(e).__name__}: {e}"

    # 加 caption(图下小字, 居中)
    if caption:
        # 2026-05-18 P190: caption 也支持 inline `$...$` LaTeX
        if "$" in caption:
            caption = _render_inline_latex_to_text_only(caption)
        cap_top = pic.top + pic.height + Inches(0.1)
        cap_h = Inches(0.4)
        cap_box = slide.shapes.add_textbox(Inches(0.5), cap_top, slide_w - Inches(1.0), cap_h)
        tf = cap_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = caption
        p.alignment = 2  # center (PP_ALIGN.CENTER = 2)
        for run in p.runs:
            run.font.size = Pt(11)
            run.font.italic = True
    return None


def _pptx_render_table_layout(slide, s_idx: int, s: dict,
                              slide_w, slide_h, title: str) -> str | None:
    """table layout: title + 表格(顶层 table 字段)。"""
    from pptx.util import Inches

    tbl_def = s.get("table")
    if not isinstance(tbl_def, dict):
        body = s.get("body") or []
        for b in body:
            if isinstance(b, dict) and str(b.get("type", "")).lower() == "table":
                tbl_def = b
                break
    if not isinstance(tbl_def, dict):
        return None

    rows = tbl_def.get("rows") or []
    if not isinstance(rows, list) or not rows:
        return f"slide[{s_idx}].table.rows required"
    n_cols = max((len(r) if isinstance(r, list) else 0) for r in rows)
    if n_cols == 0:
        return f"slide[{s_idx}].table.rows all empty"

    top = Inches(1.5) if title else Inches(0.5)
    width = slide_w - Inches(1.0)
    available_h = slide_h - top - Inches(0.5)
    height = Inches(min(0.45 * len(rows) + 0.2, available_h / 914400))  # EMU→inches
    left = Inches(0.5)

    shape = slide.shapes.add_table(len(rows), n_cols, left, top, width, height)
    tbl = shape.table
    has_header = bool(tbl_def.get("header"))
    for r_idx, r in enumerate(rows):
        if not isinstance(r, list):
            continue
        for c_idx, cell_val in enumerate(r):
            if c_idx >= n_cols:
                break
            cell = tbl.cell(r_idx, c_idx)
            # 2026-05-18 P190: cell 也支持 inline `$...$` LaTeX
            cell_text = str(cell_val) if cell_val is not None else ""
            if "$" in cell_text:
                cell_text = _render_inline_latex_to_text_only(cell_text)
            cell.text = cell_text
            if has_header and r_idx == 0:
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.bold = True
    return None


def _pptx_render_blank(slide, workspace_dir: str, s_idx: int, s: dict,
                       slide_w, slide_h, title: str) -> str | None:
    """blank layout: 自由组合 — 优先用 text 字段(单段), 否则用 body[] 数组。

    body 元素类型: text / bullets / image / table。
    """
    from pptx.util import Inches

    # 优先简单 text 字段
    text_field = str(s.get("text", "")).strip()
    body = s.get("body") or []
    if not isinstance(body, list):
        return f"slide[{s_idx}].body must be array"

    # 收集文本/非文本元素
    bullets, paragraphs = _normalize_text_items(body, text_field)
    non_text: list[dict] = []
    for b in body:
        if isinstance(b, dict):
            t = str(b.get("type", "")).lower()
            if t in ("image", "table"):
                non_text.append(b)

    # 文本框区域
    text_top = Inches(1.5) if title else Inches(0.5)
    if (bullets or paragraphs) and not non_text:
        # 全是文字, 占满下半区
        h = slide_h - text_top - Inches(0.5)
        tb = slide.shapes.add_textbox(Inches(0.5), text_top, slide_w - Inches(1.0), h)
        _fill_text_frame(tb.text_frame, bullets, paragraphs)
        used_top = text_top + h
    elif bullets or paragraphs:
        # 文字 + 后面有非文本, 文字占上 1/3
        h = (slide_h - text_top - Inches(0.5)) // 3
        tb = slide.shapes.add_textbox(Inches(0.5), text_top, slide_w - Inches(1.0), h)
        _fill_text_frame(tb.text_frame, bullets, paragraphs)
        used_top = text_top + h + Inches(0.2)
    else:
        used_top = text_top

    # 处理非文本(image/table 顺次堆叠)
    for nt_idx, b in enumerate(non_text):
        btype = b.get("type", "").lower()
        if btype == "image":
            err = _add_image_to_slide(slide, workspace_dir, s_idx, nt_idx, b, used_top, slide_w)
            if err:
                return err
            used_top = used_top + Inches(2.5)
        elif btype == "table":
            err = _add_table_to_slide(slide, s_idx, nt_idx, b, used_top, slide_w)
            if err:
                return err
            used_top = used_top + Inches(2.0)
    return None


# ─── 辅助函数 ────────────────────────────────────────────────────────

def _normalize_text_items(body, text_field: str) -> tuple[list[str], list[str]]:
    """把各种格式的 body 归一化为 (bullets, paragraphs)。

    body 可以是:
      - None 或空
      - 字符串数组(每条作为 bullet)
      - dict 数组,元素 type=text/bullets
    text_field 是 s.get("text") — 当 body 为空时作为单段落使用。
    """
    bullets: list[str] = []
    paragraphs: list[str] = []
    if isinstance(body, list):
        for b in body:
            if isinstance(b, str):
                # 裸字符串当 bullet 处理(title_content 默认场景)
                if b.strip():
                    bullets.append(b.strip())
            elif isinstance(b, dict):
                btype = str(b.get("type", "")).lower()
                if btype == "text":
                    txt = str(b.get("text", "")).strip()
                    if txt:
                        paragraphs.append(txt)
                elif btype == "bullets":
                    items = b.get("items") or []
                    if isinstance(items, list):
                        for it in items:
                            if str(it).strip():
                                bullets.append(str(it).strip())
                # image/table 由调用方单独处理
    if text_field and not bullets and not paragraphs:
        # text 字段支持 \n 分多段
        for line in text_field.split("\n"):
            line = line.strip()
            if line:
                paragraphs.append(line)
    return bullets, paragraphs




def _fill_text_frame(tf, bullets: list[str], paragraphs: list[str]) -> None:
    """把文字塞进 text_frame, paragraphs 在前(普通段落), bullets 在后。

    2026-05-18 P190: 在写入前对每条文字调 `_render_inline_latex_to_text_only`,
    把 `$...$` 内联公式渲染成 Unicode (10^4 → 10⁴) 或 plain fallback (10^4 → 10^4).
    pptx text frame 不像 docx paragraph 那样能插 inline image, 所以只走 Unicode/plain
    路径 (不渲染 PNG). 复杂公式建议用 helper 单独写成 image slide 或 paragraph 描述。
    """
    tf.word_wrap = True
    first = True
    for para in paragraphs:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        # 2026-05-18 P190: LaTeX 渲染
        p.text = _render_inline_latex_to_text_only(para) if "$" in para else para
    for bullet in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        # 2026-05-18 P190: LaTeX 渲染
        rendered = _render_inline_latex_to_text_only(bullet) if "$" in bullet else bullet
        if rendered.startswith(("•", "·", "-", "*", "►")):
            p.text = rendered
        else:
            p.text = "• " + rendered


def _add_image_to_slide(slide, workspace_dir: str, s_idx: int, nt_idx: int,
                       b: dict, top, slide_w) -> str | None:
    from pptx.util import Inches

    img_path = str(b.get("path", "")).strip()
    if not img_path:
        return f"slide[{s_idx}].body[image#{nt_idx}].path required"
    try:
        img_target = ws_tool._safe_resolve(workspace_dir, img_path)
    except ValueError as e:
        return f"slide[{s_idx}] image path invalid: {e}"
    if not os.path.isfile(img_target):
        # 2026-05-12 P47: pptx body image 重定向 (与 docx 一致)
        _sandbox_redirect = _resolve_main_named_to_sandbox(workspace_dir, img_path)
        if _sandbox_redirect and os.path.isfile(_sandbox_redirect):
            img_target = _sandbox_redirect
        else:
            try:
                _sug = ws_tool._suggest_similar_files(workspace_dir, img_path)
                if _sug and _sug[0]["score"] >= 95:
                    _r = ws_tool._safe_resolve(workspace_dir, _sug[0]["path"])
                    if os.path.isfile(_r):
                        img_target = _r
                    else:
                        return f"slide[{s_idx}] image not found: {img_path}"
                else:
                    _hint = ""
                    if _sug:
                        _cands = ", ".join(s["path"] for s in _sug[:3])
                        _hint = f" 工作区内相似文件: {_cands}"
                    return f"slide[{s_idx}] image not found: {img_path}{_hint}"
            except (ValueError, OSError, AttributeError):
                return f"slide[{s_idx}] image not found: {img_path}"
    try:
        if os.path.getsize(img_target) > _MAX_IMAGE_BYTES:
            return f"slide[{s_idx}] image too large"
    except OSError:
        pass

    try:
        w_in = float(b.get("width_inches") or _DEFAULT_PPTX_IMAGE_WIDTH_INCHES)
        w_in = max(1.0, min(13.0, w_in))
    except (TypeError, ValueError):
        w_in = _DEFAULT_PPTX_IMAGE_WIDTH_INCHES
    width = Inches(w_in)
    left = (slide_w - width) // 2  # 居中
    try:
        slide.shapes.add_picture(img_target, left, top, width=width)
    except Exception as e:
        return f"slide[{s_idx}] add_picture failed: {type(e).__name__}: {e}"
    return None


def _add_table_to_slide(slide, s_idx: int, nt_idx: int, b: dict, top, slide_w) -> str | None:
    from pptx.util import Inches

    rows = b.get("rows") or []
    if not isinstance(rows, list) or not rows:
        return f"slide[{s_idx}].table[{nt_idx}].rows required"
    n_cols = max((len(r) if isinstance(r, list) else 0) for r in rows)
    if n_cols == 0:
        return f"slide[{s_idx}].table[{nt_idx}].rows all empty"
    width = slide_w - Inches(1.0)
    height = Inches(min(0.4 * len(rows) + 0.2, 5.0))
    left = Inches(0.5)
    shape = slide.shapes.add_table(len(rows), n_cols, left, top, width, height)
    tbl = shape.table
    for r_idx, r in enumerate(rows):
        if not isinstance(r, list):
            continue
        for c_idx, cell_val in enumerate(r):
            if c_idx >= n_cols:
                break
            # 2026-05-18 P190: pptx table cell 也支持 inline `$...$` LaTeX
            cell_text = str(cell_val) if cell_val is not None else ""
            if "$" in cell_text:
                cell_text = _render_inline_latex_to_text_only(cell_text)
            tbl.cell(r_idx, c_idx).text = cell_text
    return None


# 2026-05-17 P162.6: pptx 平行 rigor — 主要捕捉幻灯片里的数字编造


async def _pptx_extract_images(workspace_dir: str, target: str, rel_path: str, args: dict) -> str:
    from pptx import Presentation

    if not os.path.isfile(target):
        return _err(f"file not found: {rel_path}")

    out_dir = _resolve_out_dir(args, rel_path)
    try:
        out_target = ws_tool._safe_resolve(workspace_dir, out_dir)
    except ValueError as e:
        return _err(f"out_dir invalid: {e}")
    os.makedirs(out_target, exist_ok=True)

    def _do_extract():
        try:
            prs = Presentation(target)
        except Exception as e:
            return None, f"open failed: {type(e).__name__}: {e}"
        extracted = _dump_images_from_part(prs.part, out_target, out_dir)
        return extracted, None

    extracted, err = await asyncio.to_thread(_do_extract)
    if err:
        return _err(err)
    return json.dumps({
        "ok": True, "action": "extract_images", "format": "pptx", "path": rel_path,
        "out_dir": out_dir, "extracted": extracted, "count": len(extracted),
    }, ensure_ascii=False)


async def _pptx_insert_image(workspace_dir: str, target: str, rel_path: str, args: dict) -> str:
    img_path = str(args.get("image_path", "")).strip()
    if not img_path:
        return _err("image_path required", action="insert_image", hint="Provide the workspace-relative path to an existing image file, e.g. image_path:'chart.png'")
    body_block: dict = {"type": "image", "path": img_path}
    if args.get("width_inches") is not None:
        body_block["width_inches"] = args["width_inches"]
    title = str(args.get("caption", "")).strip()
    return await _pptx_append(workspace_dir, target, rel_path, {
        "slides": [{"title": title, "body": [body_block]}],
    })


# ═══════════════════════════════════════════════════════════════
# Excel (.xlsx)
# ═══════════════════════════════════════════════════════════════















# ───────────────────────────────────────────────────────────────
# 2026-05-04 v19.1: xlsx 增量更新单元格
# 大型数据表(900+ 行 benchmark)write 整张耗时长,且改少量数据没必要重写。
# update_cells 接受 [{sheet, ref, value}] 数组,精准更新若干格子。
# ───────────────────────────────────────────────────────────────




async def _xlsx_extract_images(workspace_dir: str, target: str, rel_path: str, args: dict) -> str:
    """xlsx 内嵌图片:openpyxl 把它们存在 ws._images,直接读 image bytes。"""
    import openpyxl

    if not os.path.isfile(target):
        return _err(f"file not found: {rel_path}")

    out_dir = _resolve_out_dir(args, rel_path)
    try:
        out_target = ws_tool._safe_resolve(workspace_dir, out_dir)
    except ValueError as e:
        return _err(f"out_dir invalid: {e}")
    os.makedirs(out_target, exist_ok=True)

    def _do_extract():
        try:
            wb = openpyxl.load_workbook(target)
        except Exception as e:
            return None, f"open failed: {type(e).__name__}: {e}"

        extracted: list[dict] = []
        counter = 0
        for sheet_name in wb.sheetnames:
            ws_sheet = wb[sheet_name]
            for img in getattr(ws_sheet, "_images", []) or []:
                counter += 1
                try:
                    blob = img._data() if callable(getattr(img, "_data", None)) else getattr(img, "ref", b"")
                    if isinstance(blob, str):
                        if os.path.isfile(blob):
                            with open(blob, "rb") as f:
                                blob = f.read()
                        else:
                            continue
                    if not blob:
                        continue
                except Exception:
                    continue
                ext = ".png"
                try:
                    fmt_attr = str(getattr(img, "format", "") or "")
                    if fmt_attr.lower() in ("jpeg", "jpg", "png", "gif", "bmp"):
                        ext = "." + fmt_attr.lower().replace("jpeg", "jpg")
                except Exception:
                    pass
                out_name = f"image{counter}{ext}"
                out_full = os.path.join(out_target, out_name)
                try:
                    with open(out_full, "wb") as f:
                        f.write(blob)
                except OSError:
                    continue
                extracted.append({
                    "filename": os.path.join(out_dir, out_name).replace("\\", "/"),
                    "size_bytes": len(blob),
                    "sheet": sheet_name,
                })
        return extracted, None

    extracted, err = await asyncio.to_thread(_do_extract)
    if err:
        return _err(err)

    return json.dumps({
        "ok": True, "action": "extract_images", "format": "xlsx", "path": rel_path,
        "out_dir": out_dir, "extracted": extracted, "count": len(extracted),
    }, ensure_ascii=False)


# ═══════════════════════════════════════════════════════════════
# 通用 helpers
# ═══════════════════════════════════════════════════════════════









# ─────────────────────────────────────────────────────────────────
# 分发表
# ─────────────────────────────────────────────────────────────────

_DISPATCH = {
    ("docx", "read"): _docx_read,
    ("docx", "write"): _docx_write,
    ("docx", "append"): _docx_append,
    ("docx", "replace_section"): _docx_replace_section,
    ("docx", "replace_block"): _docx_replace_block,
    ("docx", "replace_blocks"): _docx_replace_blocks,  # 2026-05-18 P201
    ("docx", "fill_empty_headings"): _docx_fill_empty_headings,
    ("docx", "delete_block"): _docx_delete_block,
    ("docx", "insert_block"): _docx_insert_block,
    ("docx", "extract_images"): _docx_extract_images,
    ("docx", "ocr_images"): _docx_ocr_images,
    ("docx", "insert_image"): _docx_insert_image,
    ("docx", "verify_numbers"): _docx_verify_numbers,  # 2026-05-17 P161
    ("docx", "verify_rigor"): _docx_verify_rigor,      # 2026-05-17 P162
    ("docx", "verify_against_source"): _docx_verify_against_source,  # 2026-05-18 P203

    ("pptx", "read"): _pptx_read,
    ("pptx", "write"): _pptx_write,
    ("pptx", "append"): _pptx_append,
    ("pptx", "replace_slide"): _pptx_replace_slide,
    ("pptx", "insert_slide"): _pptx_insert_slide,
    ("pptx", "delete_slide"): _pptx_delete_slide,
    ("pptx", "extract_images"): _pptx_extract_images,
    ("pptx", "ocr_images"): _pptx_ocr_images,
    ("pptx", "insert_image"): _pptx_insert_image,
    ("pptx", "verify_numbers"): _pptx_verify_numbers,  # 2026-05-17 P162.6

    ("xlsx", "read"): _xlsx_read,
    ("xlsx", "write"): _xlsx_write,
    ("xlsx", "append"): _xlsx_append,
    ("xlsx", "update_cells"): _xlsx_update_cells,
    ("xlsx", "extract_images"): _xlsx_extract_images,
    ("xlsx", "ocr_images"): _xlsx_ocr_images,
    ("xlsx", "verify_integrity"): _xlsx_verify_integrity,  # 2026-05-17 P162.7
    # xlsx insert_image 较罕见,模型需要时可用 extract_images + 手写
}
