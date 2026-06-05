"""pptx 演示文稿处理:脚本文本、part 内图片导出、删除/移除幻灯片、读取 + 上下标常量、上限。

2026-05-20 重构: 从 llm/tools/office.py 原样抽出。closure 自包含(8 符号, 0 unsafe);
复用 office_common(_err)与 office_docx(_list_images_in_part / 表格上限常量),依赖链
common←docx←pptx 无环;模块自建 logger;python-pptx 函数内惰性 import。office.py re-export 兼容。
"""
from __future__ import annotations

import asyncio
import json
import logging
import os

from app.llm.tools.office_common import _err
from app.llm.tools.office_docx import _list_images_in_part, _MAX_TABLE_ROWS_RETURNED, _MAX_TABLE_COLS_RETURNED

log = logging.getLogger(__name__)


_MAX_SLIDES_RETURNED = 200


# 2026-05-11 P13: 简单 inline 公式 Unicode fast path
# 病因(实测 trace 21:11 gen_paper docx 嵌入 51 张图, 其中 45 张是
# "10^4"/"4.716×10^6" 这种简单数字, 每张 2-3 KB PNG, 占文档 100+ KB
# 且显示效果差). 修复: 简单上标/下标/数字用 Unicode 字符直接表达,
# 不走 matplotlib 渲染. 复杂公式(\frac/\sum/\sqrt 等)仍走 PNG.
_UNICODE_SUPERSCRIPT = {
    "0": "⁰", "1": "¹", "2": "²", "3": "³", "4": "⁴", "5": "⁵",
    "6": "⁶", "7": "⁷", "8": "⁸", "9": "⁹", "+": "⁺", "-": "⁻",
    "=": "⁼", "(": "⁽", ")": "⁾", "a": "ᵃ", "b": "ᵇ", "c": "ᶜ",
    "d": "ᵈ", "e": "ᵉ", "f": "ᶠ", "g": "ᵍ", "h": "ʰ", "i": "ⁱ",
    "j": "ʲ", "k": "ᵏ", "l": "ˡ", "m": "ᵐ", "n": "ⁿ", "o": "ᵒ",
    "p": "ᵖ", "r": "ʳ", "s": "ˢ", "t": "ᵗ", "u": "ᵘ", "v": "ᵛ",
    "w": "ʷ", "x": "ˣ", "y": "ʸ", "z": "ᶻ",
}


_UNICODE_SUBSCRIPT = {
    "0": "₀", "1": "₁", "2": "₂", "3": "₃", "4": "₄", "5": "₅",
    "6": "₆", "7": "₇", "8": "₈", "9": "₉", "+": "₊", "-": "₋",
    "=": "₌", "(": "₍", ")": "₎", "a": "ₐ", "e": "ₑ", "h": "ₕ",
    "i": "ᵢ", "j": "ⱼ", "k": "ₖ", "l": "ₗ", "m": "ₘ", "n": "ₙ",
    "o": "ₒ", "p": "ₚ", "r": "ᵣ", "s": "ₛ", "t": "ₜ", "u": "ᵤ",
    "v": "ᵥ", "x": "ₓ",
}


def _script_text(script: str, raw: str) -> str:
    mapping = _UNICODE_SUPERSCRIPT if script == "^" else _UNICODE_SUBSCRIPT
    return "".join(mapping.get(ch, ch) for ch in raw)


async def _pptx_read(workspace_dir: str, target: str, rel_path: str, args: dict) -> str:
    if not os.path.isfile(target):
        return _err(f"file not found: {rel_path}")

    from pptx import Presentation

    def _do_read():
        try:
            prs = Presentation(target)
        except Exception as e:
            return None, f"cannot open as pptx ({type(e).__name__}: {e})"

        slides_out: list[dict] = []
        truncated_slides = False
        for s_idx, slide in enumerate(prs.slides):
            if s_idx >= _MAX_SLIDES_RETURNED:
                truncated_slides = True
                break
            title_text = ""
            body_texts: list[str] = []
            tables: list[dict] = []
            try:
                if slide.shapes.title and slide.shapes.title.has_text_frame:
                    title_text = (slide.shapes.title.text_frame.text or "").strip()
            except Exception:
                pass
            for shape in slide.shapes:
                try:
                    # 跳过 title 占位符,避免重复
                    if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                        pf = getattr(shape, "placeholder_format", None)
                        if pf and pf.idx == 0:
                            continue
                    if shape.has_text_frame and shape.text_frame.text:
                        txt = shape.text_frame.text.strip()
                        if txt:
                            body_texts.append(txt[:1500])
                    if shape.has_table:
                        tbl = shape.table
                        rows = []
                        for r_idx, row in enumerate(tbl.rows):
                            if r_idx >= _MAX_TABLE_ROWS_RETURNED:
                                break
                            cells = []
                            for c_idx, cell in enumerate(row.cells):
                                if c_idx >= _MAX_TABLE_COLS_RETURNED:
                                    break
                                t = (cell.text or "").strip()
                                if len(t) > 500:
                                    t = t[:500] + "…"
                                cells.append(t)
                            rows.append(cells)
                        tables.append({"rows": rows})
                except Exception:
                    continue
            slides_out.append({
                "index": s_idx,
                "title": title_text,
                "body_texts": body_texts,
                "tables": tables,
            })

        images = _list_images_in_part(prs.part)
        return {
            "slides_out": slides_out,
            "truncated_slides": truncated_slides,
            "slide_count": len(prs.slides),
            "images": images,
        }, None

    result, err = await asyncio.to_thread(_do_read)
    if err:
        return _err(err)

    return json.dumps({
        "ok": True, "action": "read", "format": "pptx", "path": rel_path,
        "slide_count": result["slide_count"],
        "image_count": len(result["images"]),
        "slides": result["slides_out"],
        "images": result["images"],
        **({"truncated_slides": True} if result["truncated_slides"] else {}),
        "hint": (
            "Image metadata is listed. Use office(action='ocr_images') only when text inside images is needed; "
            "use action='extract_images' to release embedded images as files.\n"
            "图片元数据已列出；需要图中文字时再 OCR。"
        ),
    }, ensure_ascii=False)


def _pptx_remove_slide(prs, slide_index: int):
    """从 prs.slides 中移除第 slide_index 页(python-pptx 没暴露这个 API,要手工搞 XML)。"""
    sldIdLst = prs.slides._sldIdLst
    slide_ids = list(sldIdLst)
    if slide_index < 0 or slide_index >= len(slide_ids):
        raise IndexError(f"slide index {slide_index} out of range (have {len(slide_ids)})")
    rId = slide_ids[slide_index].rId
    sldIdLst.remove(slide_ids[slide_index])
    # 同时清理 part 关系
    prs.part.drop_rel(rId)


async def _pptx_delete_slide(
    workspace_dir: str, target: str, rel_path: str, args: dict,
) -> str:
    """删除一或多页。

    args:
      index: 起始页(0-based)
      count: 删除页数(默认 1)
    """
    from pptx import Presentation

    if not os.path.isfile(target):
        return _err(f"file not found: {rel_path}")
    try:
        index = int(args.get("index", -1))
        count = int(args.get("count", 1))
    except (TypeError, ValueError):
        return _err("index/count must be integers")
    if index < 0 or count < 1:
        return _err("index >= 0 and count >= 1 required")

    def _do():
        try:
            prs = Presentation(target)
        except Exception as e:
            return None, f"open failed: {type(e).__name__}: {e}"
        total = len(prs.slides)
        if index >= total:
            return None, f"index {index} out of range (have {total} slides)"
        end = min(index + count, total)
        # 倒序删
        for i in range(end - 1, index - 1, -1):
            try:
                _pptx_remove_slide(prs, i)
            except IndexError as e:
                return None, str(e)

        try:
            prs.save(target)
        except Exception as e:
            return None, f"save failed: {type(e).__name__}: {e}"
        return {
            "deleted_count": end - index,
            "start_index": index,
        }, None

    info, err = await asyncio.to_thread(_do)
    if err:
        return _err(err)
    return json.dumps({
        "ok": True, "action": "delete_slide", "format": "pptx",
        "path": rel_path,
        **info,
    }, ensure_ascii=False)


def _dump_images_from_part(part, out_target: str, out_dir_rel: str) -> list[dict]:
    """把 docx/pptx Part 中的所有图片落盘到 out_target 目录。"""
    extracted: list[dict] = []
    seen = set()
    try:
        rels = part.rels
        for rel_id, rel in rels.items():
            if "image" not in rel.reltype:
                continue
            partname = str(getattr(rel.target_part, "partname", ""))
            if partname in seen:
                continue
            seen.add(partname)
            try:
                blob = rel.target_part.blob
            except Exception:
                continue
            ext = os.path.splitext(partname)[1].lower() or ".png"
            out_name = f"image{len(extracted) + 1}{ext}"
            out_full = os.path.join(out_target, out_name)
            try:
                with open(out_full, "wb") as f:
                    f.write(blob)
            except OSError:
                continue
            extracted.append({
                "filename": os.path.join(out_dir_rel, out_name).replace("\\", "/"),
                "size_bytes": len(blob),
            })
    except Exception:
        log.exception("dump_images_from_part failed")
    return extracted
