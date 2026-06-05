"""
D17: Office (docx/xlsx/pptx) 深度总结

问题: 现有 extract_docx_text 只取前 max_chars, 长文档后面找不到
目标: dream 后台跑完整 + LLM 总结 + 章节大纲 + 关键词

工具链 (3 步, ~15s):
1. Extract 完整文本
2. LLM 总结
3. UPDATE cold_nodes

阈值: 1 (event-driven)
"""
from __future__ import annotations

from app.core.dream.prompt_catalog import (
    D17_OFFICE_INDEX_PROMPT,
)
_LLM_PROMPT = D17_OFFICE_INDEX_PROMPT


import asyncio
import json
import os
import time
from typing import Any

from app.core.dream.dream_log import dream_log
from app.core.dream.event_bus import event_bus
from app.core.dream.registry import register_dream_task
from app.core.dream.task_base import InfoDrivenTask
from app.core.dream.tasks.file_searchability.file_meta import (
    emit_file_indexed,
    is_active_file_metadata,
)


D17_THRESHOLD = 1
D17_MAX_FILES_PER_RUN = 3
_D17_RAW_SCAN_LIMIT = 5000
_D17_MAX_CANDIDATES_PER_RUN = 200
_OFFICE_EXTS = {".docx", ".xlsx", ".pptx"}




def _validate(raw):
    if not isinstance(raw, dict):
        return False
    s = raw.get("summary", "")
    return 30 <= len(s) <= 800


def _extract_docx_full(path: str, do_ocr: bool = True) -> tuple[str, list[str]]:
    """完整提取 docx 文本 + 章节标题 (段落 + 表格 + 嵌入图片 OCR).
    
    Args:
        do_ocr: 是否对嵌入图片做 OCR. demote 时设 False 加快.
    """
    try:
        import docx
        doc = docx.Document(path)
    except Exception:
        return ("", [])
    
    parts = []
    headings = []
    
    # 1. 段落
    for p in doc.paragraphs:
        t = (p.text or "").strip()
        if not t:
            continue
        parts.append(t)
        try:
            sn = (p.style.name or "").lower() if p.style else ""
            if "heading" in sn or "标题" in sn:
                headings.append(t)
        except Exception:
            pass
    
    # 2. 表格
    try:
        for tbl in doc.tables:
            for row in tbl.rows:
                cells = []
                for cell in row.cells:
                    ct = (cell.text or "").strip()
                    if ct:
                        cells.append(ct)
                if cells:
                    parts.append(" | ".join(cells))
    except Exception:
        pass
    
    # 3. 嵌入图片 OCR (慢, demote 时跳过)
    if do_ocr:
        image_ocr_text = _extract_office_zip_images_ocr(
            path, media_prefix="word/media/",
            max_images=30, max_size_mb=50.0,
        )
        if image_ocr_text:
            parts.append("\n[嵌入图片 OCR 内容]\n" + image_ocr_text)
    
    return ("\n".join(parts), headings)


def _extract_office_zip_images_ocr(
    file_path: str,
    media_prefix: str,
    max_images: int = 50,
    max_size_mb: float = 50.0,
) -> str:
    """D17 调用 ocr_bridge.ocr_office_images, 返回合并文本 (兼容旧接口).
    
    2026-05-16: 抽到 ocr_bridge.py 作公共函数, 主线程 office(action='ocr_images')
    + D17 共用. 大图默认 max_size_mb=50MB (基本不限), 让 LLM/dream 主动调时不被卡.
    """
    try:
        from app.llm.tools.ocr_bridge import ocr_office_images_scheduled
        result = asyncio.run(ocr_office_images_scheduled(
            file_path,
            media_prefix=media_prefix,
            max_images=max_images,
            max_size_mb=max_size_mb,
            per_image_timeout=90,
        ))
    except Exception as e:
        dream_log.warn(
            "dream.task.d17_office_index.image_ocr_call_failed",
            f"err={e!r}"[:200],
        )
        return ""
    
    if not result.get("ok"):
        return ""
    
    ocr_count = result.get("ocr_count", 0)
    total = result.get("total_images", 0)
    skipped = result.get("skipped_too_large", 0)
    
    if ocr_count > 0:
        dream_log.log(
            "dream.task.d17_office_index.office_images_ocred",
            f"file={os.path.basename(file_path)} "
            f"ocred={ocr_count}/{total} images (skipped_too_large={skipped})",
        )
    
    return result.get("merged_text", "")


def _extract_xlsx_summary(path: str) -> tuple[str, list[str]]:
    """xlsx: sheet 名 + 头部数据样本."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    except Exception:
        return ("", [])
    
    parts = []
    sheets = []
    for sheet_name in wb.sheetnames[:10]:
        sheets.append(sheet_name)
        ws = wb[sheet_name]
        rows = []
        for i, row in enumerate(ws.iter_rows(values_only=True, max_row=30)):
            cells = [str(c) if c is not None else "" for c in row[:10]]
            rows.append("\t".join(cells))
            if i > 30:
                break
        parts.append(f"## {sheet_name}\n" + "\n".join(rows))
    
    try:
        wb.close()
    except Exception:
        pass
    
    return ("\n\n".join(parts), sheets)


def _extract_pptx_outline(path: str, do_ocr: bool = True) -> tuple[str, list[str]]:
    """pptx: 每页文字 + 标题 + 嵌入图片 OCR.
    
    Args:
        do_ocr: 是否对嵌入图片做 OCR. demote 时设 False 加快.
    """
    try:
        from pptx import Presentation
        prs = Presentation(path)
    except Exception:
        return ("", [])
    
    parts = []
    titles = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    slide_text.append(t)
                    if not titles or len(titles) < i:
                        first_para = t.split("\n")[0][:50]
                        if first_para:
                            titles.append(first_para)
        if slide_text:
            parts.append(f"--- Slide {i} ---\n" + "\n".join(slide_text))
    
    # 嵌入图片 OCR (慢, demote 时跳过)
    if do_ocr:
        image_ocr_text = _extract_office_zip_images_ocr(
            path, media_prefix="ppt/media/",
            max_images=50, max_size_mb=50.0,
        )
        if image_ocr_text:
            parts.append("\n[嵌入图片 OCR 内容]\n" + image_ocr_text)
    
    return ("\n\n".join(parts), titles)


async def _extract_docx_full_async(path: str, do_ocr: bool = True) -> tuple[str, list[str]]:
    text, headings = await asyncio.to_thread(_extract_docx_full, path, do_ocr=False)
    if do_ocr:
        image_ocr_text = await _extract_office_zip_images_ocr_async(
            path, media_prefix="word/media/", max_images=30, max_size_mb=50.0
        )
        if image_ocr_text:
            text = text + "\n\n[嵌入图片 OCR 内容]\n" + image_ocr_text if text else image_ocr_text
    return text, headings


async def _extract_pptx_outline_async(path: str, do_ocr: bool = True) -> tuple[str, list[str]]:
    text, titles = await asyncio.to_thread(_extract_pptx_outline, path, do_ocr=False)
    if do_ocr:
        image_ocr_text = await _extract_office_zip_images_ocr_async(
            path, media_prefix="ppt/media/", max_images=50, max_size_mb=50.0
        )
        if image_ocr_text:
            text = text + "\n\n[嵌入图片 OCR 内容]\n" + image_ocr_text if text else image_ocr_text
    return text, titles


async def _extract_office_zip_images_ocr_async(
    file_path: str,
    media_prefix: str,
    max_images: int = 50,
    max_size_mb: float = 50.0,
) -> str:
    try:
        from app.llm.tools.ocr_bridge import ocr_office_images_scheduled
        result = await ocr_office_images_scheduled(
            file_path,
            media_prefix=media_prefix,
            max_images=max_images,
            max_size_mb=max_size_mb,
            per_image_timeout=90,
        )
    except Exception as e:
        dream_log.warn(
            "dream.task.d17_office_index.image_ocr_call_failed",
            f"err={e!r}"[:200],
        )
        return ""

    if not result.get("ok"):
        return ""

    ocr_count = result.get("ocr_count", 0)
    total = result.get("total_images", 0)
    skipped = result.get("skipped_too_large", 0)

    if ocr_count > 0:
        dream_log.log(
            "dream.task.d17_office_index.office_images_ocred",
            f"file={os.path.basename(file_path)} "
            f"ocred={ocr_count}/{total} images (skipped_too_large={skipped})",
        )

    return result.get("merged_text", "")


async def _llm_summarize(text: str, filename: str, headings: list[str],
                         lite_first: bool = False) -> dict | None:
    from app.llm import client as llm
    
    user_text = (
        f"## Office File Metadata\n"
        f"filename: {filename}\n"
        f"detected_headings: {headings[:20] if headings else '(none)'}\n\n"
        f"## Extracted Document Text ({len(text)} chars)\n{text[:20000]}\n\n"
        "Return the indexing JSON.\n\n输出索引 JSON。"
    )
    messages = [
        {"role": "system", "content": _LLM_PROMPT},
        {"role": "user", "content": user_text},
    ]
    return await llm.chat_json_with_upgrade(
        messages,
        validate=_validate,
        label="dream_d17_office",
        lite_first=lite_first,
    )


async def _get_unprocessed_office(limit: int) -> list[dict]:
    from app.db.pool import pool
    
    async with pool().acquire() as conn:
        rows = await conn.fetch("""
            SELECT id, archive_id, group_id, headline, content, file_metadata
            FROM cold_nodes
            WHERE scope = 'kb' AND node_type = 'file'
              AND (file_metadata->>'dream_processed' IS NULL
                   OR file_metadata->>'dream_processed' = 'false')
              AND file_metadata->>'workspace_path' IS NOT NULL
              AND file_metadata->>'workspace_path' != ''
              AND (file_metadata->>'deleted' IS NULL
                   OR file_metadata->>'deleted' != 'true')
              AND (file_metadata->>'download_status' IS NULL
                   OR file_metadata->>'download_status' = ''
                   OR file_metadata->>'download_status' = 'done')
            ORDER BY created_at DESC
            LIMIT $1
        """, _D17_RAW_SCAN_LIMIT)
    
    out = []
    for r in rows:
        meta = r["file_metadata"]
        if isinstance(meta, str):
            try:
                meta = json.loads(meta)
            except Exception:
                continue
        if not isinstance(meta, dict):
            continue
        if not is_active_file_metadata(meta):
            continue
        filename = meta.get("filename", "") or meta.get("file_name", "")
        ext = os.path.splitext(filename)[1].lower()
        if ext not in _OFFICE_EXTS:
            continue
        out.append({
            "id": r["id"], "archive_id": r["archive_id"], "group_id": r["group_id"],
            "headline": r["headline"], "content": r["content"] or "",
            "file_metadata": meta, "filename": filename, "ext": ext,
        })
        if len(out) >= limit:
            break
    return out


def _get_office_path(file_meta: dict) -> str | None:
    ws_path = file_meta.get("workspace_path") or ""
    if not ws_path:
        return None
    from app.core.dream.cache import _get_workspace_root
    ws_root = _get_workspace_root()
    archive_id = file_meta.get("archive_id", "")
    group_id = file_meta.get("group_id", "")
    if not ws_root:
        return None
    full = os.path.join(ws_root, archive_id, group_id, ws_path)
    return full if os.path.isfile(full) else None


@register_dream_task
class D17OfficeIndex(InfoDrivenTask):
    """D17: Office 深度索引."""
    
    name = "d17_office_index"
    threshold = D17_THRESHOLD
    uses_llm = True
    startup_sweep = True
    
    async def info_fn(self) -> float:
        """信息水位 = file_uploaded 事件累计.
        
        2026-05-16 Round 14 fallback:
        event_bus.total_count 进程启动归零 → 第一次启动 D17 永不触发
        (Round 9 已为 D27 修过, 但 D17/D15/D16/D18 同样问题).
        若 event_bus 为 0 但数据库有未处理 office 文件 → 用候选数作信息水位,
        让 D17 至少跑一次 backfill.
        """
        try:
            return float(len(await _get_unprocessed_office(limit=1000)))
        except Exception as e:
            dream_log.warn(
                f"dream.task.{self.name}.info_fn_backfill_failed",
                repr(e)[:200],
            )
            return float(event_bus.total_count("file_uploaded"))

    async def should_run(self) -> bool:
        import time as _t
        if self.suspended_until > _t.time():
            return False
        try:
            return await self.info_fn() >= self.threshold
        except Exception:
            return False
    
    async def _do_work(self) -> None:
        candidates = await _get_unprocessed_office(_D17_MAX_CANDIDATES_PER_RUN)
        if not candidates:
            return
        
        success = 0
        attempted = 0
        skipped = 0
        for node in candidates:
            path = _get_office_path(node["file_metadata"])
            if not path:
                skipped += 1
                await self._mark_skipped(node, "path_invalid")
                continue
            if attempted >= D17_MAX_FILES_PER_RUN:
                break
            attempted += 1
            
            # 按 ext 提取
            ext = node["ext"]
            try:
                if ext == ".docx":
                    text, headings = await _extract_docx_full_async(path, do_ocr=not self.demoted)
                elif ext == ".xlsx":
                    text, headings = await asyncio.to_thread(_extract_xlsx_summary, path)
                elif ext == ".pptx":
                    text, headings = await _extract_pptx_outline_async(path, do_ocr=not self.demoted)
                else:
                    continue
            except asyncio.CancelledError:
                raise
            except Exception as e:
                dream_log.warn(
                    "dream.task.d17_office_index.extract_failed",
                    f"id={node['id']} err={e!r}"[:200],
                )
                continue
            
            if not text or len(text) < 30:
                # 文件可能损坏 - 标记 skipped
                await self._mark_skipped(node, "no_extractable_text")
                continue
            
            # LLM 总结
            try:
                info = await _llm_summarize(
                    text, node["filename"], headings,
                    lite_first=self.demoted,  # cancel ≥5 次后强制 lite
                )
            except asyncio.CancelledError:
                raise
            except Exception as e:
                dream_log.warn(
                    "dream.task.d17_office_index.llm_failed",
                    f"id={node['id']} err={e!r}"[:200],
                )
                info = None
            
            # UPDATE (即使 LLM 失败, 也写完整文本进 content)
            await self._apply_update(node, text, headings, info)
            success += 1
        
        if success or skipped:
            dream_log.log(
                "dream.task.d17_office_index.cycle_done",
                f"indexed={success} skipped={skipped}",
            )
    
    async def _apply_update(self, node, full_text, headings, info):
        from app.db.pool import pool
        from app.memory.kb import sanitize_summary
        
        parts = []
        if info and info.get("summary"):
            parts.append(f"[摘要]\n{info['summary']}")
        if info and info.get("outline"):
            parts.append("[大纲]\n" + "\n".join(info["outline"][:30]))
        elif headings:
            parts.append("[章节]\n" + "\n".join(headings[:20]))
        if info and info.get("keywords"):
            parts.append("[关键词] " + ", ".join(info["keywords"]))
        # 加部分原文
        if full_text:
            parts.append(f"[内容片段]\n{full_text[:3000]}")
        
        new_content = sanitize_summary("\n\n".join(parts))
        
        meta = dict(node["file_metadata"])
        meta.update({
            "dream_processed": True,
            "dream_processed_at": time.time(),
            "dream_task": "d17_office_index",
            "topics": info.get("topics", []) if info else [],
            "keywords": info.get("keywords", []) if info else [],
            "doc_type": info.get("doc_type", "") if info else "",
            "extracted_chars": len(full_text),
        })
        
        try:
            async with pool().acquire() as conn:
                await conn.execute("""
                    UPDATE cold_nodes
                    SET content = $1, file_metadata = $2, updated_at = NOW()
                    WHERE id = $3
                """, new_content, json.dumps(meta, ensure_ascii=False), node["id"])
            await emit_file_indexed(node, self.name)
            dream_log.log(
                "dream.task.d17_office_index.indexed",
                f"id={node['id']} ext={node['ext']} chars={len(full_text)}",
            )
        except Exception as e:
            dream_log.error("dream.task.d17_office_index.update_failed", repr(e)[:200])
    
    async def _mark_skipped(self, node, reason):
        from app.db.pool import pool
        meta = dict(node["file_metadata"])
        meta.update({
            "dream_processed": True,
            "dream_processed_at": time.time(),
            "dream_skipped": True,
            "dream_skip_reason": reason,
        })
        try:
            async with pool().acquire() as conn:
                await conn.execute(
                    "UPDATE cold_nodes SET file_metadata = $1 WHERE id = $2",
                    json.dumps(meta, ensure_ascii=False), node["id"],
                )
        except Exception:
            pass
